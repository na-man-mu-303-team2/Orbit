from io import BytesIO

from fastapi.testclient import TestClient
from pptx import Presentation

import app.main as api_module
from tests.test_config import VALID_ENV


def client() -> TestClient:
    api_module.app.state.config = api_module.load_config(VALID_ENV)
    return TestClient(api_module.app)


def test_pptx_import_endpoint_returns_editable_deck() -> None:
    response = client().post(
        "/pptx/import",
        files={
            "file": (
                "team-update.pptx",
                create_sample_pptx(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
        data={
            "project_id": "project_demo_1",
            "file_id": "file_1",
        },
    )

    assert response.status_code == 200
    payload = response.json()
    deck = payload["deck"]

    assert deck["projectId"] == "project_demo_1"
    assert deck["metadata"]["sourceType"] == "import"
    assert deck["title"] == "team-update"
    assert len(deck["slides"]) == 1
    assert deck["slides"][0]["title"] == "팀 업데이트"
    assert deck["slides"][0]["elements"][0]["props"]["text"] == "팀 업데이트"
    assert payload["warnings"] == []


def create_sample_pptx() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "팀 업데이트"
    slide.placeholders[1].text = "진행 현황\n다음 단계"

    output = BytesIO()
    presentation.save(output)
    return output.getvalue()
