import base64
import copy
import hashlib
import importlib
import shutil
import subprocess
import zipfile
from io import BytesIO
from pathlib import Path
from xml.etree import ElementTree as ET

import pytest
from fastapi.testclient import TestClient
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from app.ai.pptx_ooxml_generation import (
    CanvasSpec,
    PptxRenderUnavailableError,
    generate_pptx_ooxml,
    render_pdf_to_png_assets,
    render_pptx_to_png_assets,
    shape_fallback_assets,
    strip_text_from_pptx_package,
    sync_pptx_ooxml,
)
from app.ai.pptx_design_importer import ImportedDesignAsset
from app.ai import pptx_ooxml_generation, pptx_rendering
from app.ai.pptx_rendering import (
    PptxNotesRenderError,
    PptxNotesRenderErrorCode,
    render_pptx_notes_to_png_assets,
)
from app.ai.pptx_package_security import (
    PPTX_ACTIVE_CONTENT_BLOCKED,
    PPTX_EXTERNAL_RELATIONSHIP_BLOCKED,
)
from app.ai.pptx_render_resource_limits import PptxRenderResourceLimitError
import app.main as api_module
from app.main import app

IMPORT_FIDELITY_NOTES_FIXTURE = (
    Path(__file__).parent / "fixtures" / "pptx" / "import-fidelity-notes.pptx"
)
PRESENTATION_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
RELATIONSHIP_NS = "http://schemas.openxmlformats.org/package/2006/relationships"


@pytest.mark.parametrize(
    ("form_data", "expected_preference"),
    [
        ({}, "editability-first"),
        ({"import_preference": "appearance-first"}, "appearance-first"),
    ],
)
def test_generation_endpoint_validates_and_forwards_import_preference(
    monkeypatch: pytest.MonkeyPatch,
    form_data: dict[str, str],
    expected_preference: str,
) -> None:
    captured: dict[str, str] = {}

    def fake_generate(
        _path: Path,
        file_id: str,
        *,
        import_preference: str,
    ) -> pptx_ooxml_generation.PptxOoxmlGenerationResult:
        captured.update(
            {"file_id": file_id, "import_preference": import_preference}
        )
        return pptx_ooxml_generation.PptxOoxmlGenerationResult(
            canvas={},
            blueprint={},
            templateBlueprint={},
            qualityReport={},
        )

    monkeypatch.setattr(api_module, "generate_pptx_ooxml", fake_generate)
    response = TestClient(app).post(
        "/ai/pptx-ooxml-generation",
        files={"file": ("template.pptx", b"synthetic", "application/octet-stream")},
        data={"file_id": "file_template", **form_data},
    )

    assert response.status_code == 200
    assert captured == {
        "file_id": "file_template",
        "import_preference": expected_preference,
    }


def test_generation_endpoint_rejects_unknown_import_preference(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    called = False

    def fake_generate(
        _path: Path,
        _file_id: str,
        *,
        import_preference: str,
    ) -> pptx_ooxml_generation.PptxOoxmlGenerationResult:
        nonlocal called
        called = True
        del import_preference
        return pptx_ooxml_generation.PptxOoxmlGenerationResult(
            canvas={},
            blueprint={},
            templateBlueprint={},
            qualityReport={},
        )

    monkeypatch.setattr(api_module, "generate_pptx_ooxml", fake_generate)

    response = TestClient(app).post(
        "/ai/pptx-ooxml-generation",
        files={"file": ("template.pptx", b"synthetic", "application/octet-stream")},
        data={"file_id": "file_template", "import_preference": "balanced"},
    )

    assert response.status_code == 422
    assert called is False


def test_generation_endpoint_returns_bounded_package_security_code() -> None:
    unsafe_package = BytesIO()
    with zipfile.ZipFile(unsafe_package, "w") as package:
        package.writestr("../payload.xml", b"unsafe")

    response = TestClient(app).post(
        "/ai/pptx-ooxml-generation",
        files={
            "file": (
                "unsafe.pptx",
                unsafe_package.getvalue(),
                "application/octet-stream",
            )
        },
        data={"file_id": "file_unsafe"},
    )

    assert response.status_code == 400
    assert response.json() == {"detail": "PPTX_PACKAGE_PATH_TRAVERSAL"}


def template_slide_id(generated: object, slide_index: int = 0) -> str:
    return generated.template_blueprint["slides"][slide_index]["slideId"]


def test_import_fidelity_fixture_contains_baseline_risks() -> None:
    with zipfile.ZipFile(IMPORT_FIDELITY_NOTES_FIXTURE, "r") as package:
        slide_xml = package.read("ppt/slides/slide1.xml")
        slide_rels = ET.fromstring(package.read("ppt/slides/_rels/slide1.xml.rels"))
        notes_root = ET.fromstring(package.read("ppt/notesSlides/notesSlide1.xml"))
        notes_master_xml = package.read("ppt/notesMasters/notesMaster1.xml")
        master_root = ET.fromstring(package.read("ppt/slideMasters/slideMaster1.xml"))

    assert b'wedgeRoundRectCallout' in slide_xml
    assert b'spc="120"' in slide_xml
    assert b"<p:grpSp>" in slide_xml
    assert b"<p:pic>" in slide_xml
    assert any(
        relationship.get("Type", "").endswith("/notesSlide")
        for relationship in slide_rels.findall(f"{{{RELATIONSHIP_NS}}}Relationship")
    )

    namespaces = {"a": DRAWING_NS, "p": PRESENTATION_NS}
    body_shapes = []
    for shape in notes_root.findall("./p:cSld/p:spTree/p:sp", namespaces):
        placeholder = shape.find("./p:nvSpPr/p:nvPr/p:ph", namespaces)
        if placeholder is not None and placeholder.get("type") == "body":
            body_shapes.append(shape)
    assert len(body_shapes) == 1
    paragraphs = body_shapes[0].findall("./p:txBody/a:p", namespaces)
    assert len(paragraphs) == 3
    assert paragraphs[1].find(".//a:t", namespaces) is None
    assert paragraphs[2].find("./a:br", namespaces) is not None
    assert b"NOTES_NON_BODY_DO_NOT_IMPORT" in ET.tostring(notes_root)
    assert b"NOTES_MASTER_DECORATION_DO_NOT_IMPORT" in notes_master_xml

    title_style = master_root.find(
        "./p:txStyles/p:titleStyle/a:lvl1pPr/a:defRPr", namespaces
    )
    assert title_style is not None
    assert title_style.get("sz") == "4400"
    slide_root = ET.fromstring(slide_xml)
    title_run_properties = next(
        run_properties
        for run_properties in slide_root.findall(".//a:rPr", namespaces)
        if run_properties.get("spc") == "120"
    )
    assert title_run_properties.get("sz") is None


def test_import_fidelity_fixture_imports_notes_and_effective_title_style() -> None:
    result = generate_pptx_ooxml(
        IMPORT_FIDELITY_NOTES_FIXTURE,
        "file_import_fidelity_notes",
        render=False,
    )
    imported_slide = result.blueprint["slides"][0]
    template_slide = result.template_blueprint["slides"][0]
    title = next(
        element
        for element in imported_slide["elements"]
        if element.get("props", {}).get("text") == "Inherited title style"
    )

    assert imported_slide["speakerNotes"] == "첫 번째 문단\n\n수동\n줄바꿈"
    assert template_slide["notesPage"] == {
        "status": "preserved",
        "sourceNotesPart": "ppt/notesSlides/notesSlide1.xml",
        "sourceNotesMasterPart": "ppt/notesMasters/notesMaster1.xml",
        "bodyShapeId": "3",
        "bodyWritable": True,
        "notesWidthEmu": 6_858_000,
        "notesHeightEmu": 9_144_000,
        "hasNonBodyContent": True,
    }
    assert result.quality_report["notesDiagnostics"] == {
        "total": 1,
        "imported": 1,
        "rendered": 0,
        "writable": 1,
        "warnings": [],
    }
    assert title["props"]["fontFamily"] == "Pretendard"
    assert title["props"]["fontWeight"] == 600
    assert title["props"]["fontSize"] == 88
    assert title["props"]["color"] == "#000000"
    assert title["props"]["letterSpacing"] == 2.4
    assert title["props"]["runs"][0]["letterSpacing"] == 2.4
    assert not any(
        warning.startswith("PPTX_RICH_TEXT_UNSUPPORTED_LETTER_SPACING")
        for warning in result.warnings
    )
    assert any(
        "unsupported preset wedgeRoundRectCallout" in warning
        for warning in result.warnings
    )


def test_notes_renderer_uses_notes_only_filter_and_notes_size_ratio(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    commands: list[list[str]] = []
    temporary_paths: list[Path] = []

    monkeypatch.setattr(
        "app.ai.pptx_rendering.find_libreoffice_executable",
        lambda: "/test/soffice",
    )

    def fake_run(command: list[str], **_kwargs: object) -> subprocess.CompletedProcess[str]:
        commands.append(command)
        output = Path(command[command.index("--outdir") + 1])
        source = Path(command[-1])
        temporary_paths.extend((output.parent, output, source))
        write_test_pdf(output / "source.pdf", page_count=1)
        return subprocess.CompletedProcess(command, 0, "", "")

    monkeypatch.setattr("app.ai.pptx_rendering.subprocess.run", fake_run)

    assets = render_pptx_notes_to_png_assets(
        IMPORT_FIDELITY_NOTES_FIXTURE.read_bytes(),
        notes_width_emu=6_858_000,
        notes_height_emu=9_144_000,
        expected_page_count=1,
    )

    assert len(assets) == 1
    assert assets[0].asset_id == "notes_render_1"
    assert assets[0].mime_type == "image/png"
    with Image.open(BytesIO(base64.b64decode(assets[0].content_base64))) as image:
        assert image.size == (960, 1280)
    assert "ExportNotesPages" in commands[0][commands[0].index("--convert-to") + 1]
    assert "ExportOnlyNotesPages" in commands[0][
        commands[0].index("--convert-to") + 1
    ]
    assert any(argument.startswith("-env:UserInstallation=file:") for argument in commands[0])
    assert all(not path.exists() for path in temporary_paths)


def test_notes_renderer_times_out_and_cleans_temporary_files(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    temporary_paths: list[Path] = []
    monkeypatch.setattr(
        "app.ai.pptx_rendering.find_libreoffice_executable",
        lambda: "/test/soffice",
    )

    def timeout(command: list[str], **kwargs: object) -> subprocess.CompletedProcess[str]:
        output = Path(command[command.index("--outdir") + 1])
        temporary_paths.append(output.parent)
        raise subprocess.TimeoutExpired(command, kwargs["timeout"])

    monkeypatch.setattr("app.ai.pptx_rendering.subprocess.run", timeout)

    with pytest.raises(PptxNotesRenderError) as error:
        render_pptx_notes_to_png_assets(
            IMPORT_FIDELITY_NOTES_FIXTURE.read_bytes(),
            notes_width_emu=6_858_000,
            notes_height_emu=9_144_000,
            expected_page_count=1,
            timeout_seconds=0.01,
        )

    assert error.value.code == "PPTX_NOTES_RENDER_TIMEOUT"
    assert all(not path.exists() for path in temporary_paths)


def test_notes_renderer_rejects_page_count_mismatch(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setattr(
        "app.ai.pptx_rendering.find_libreoffice_executable",
        lambda: "/test/soffice",
    )

    def fake_run(command: list[str], **_kwargs: object) -> subprocess.CompletedProcess[str]:
        output = Path(command[command.index("--outdir") + 1])
        write_test_pdf(output / "source.pdf", page_count=2)
        return subprocess.CompletedProcess(command, 0, "", "")

    monkeypatch.setattr("app.ai.pptx_rendering.subprocess.run", fake_run)

    with pytest.raises(PptxNotesRenderError) as error:
        render_pptx_notes_to_png_assets(
            IMPORT_FIDELITY_NOTES_FIXTURE.read_bytes(),
            notes_width_emu=6_858_000,
            notes_height_emu=9_144_000,
            expected_page_count=1,
        )

    assert error.value.code == "PPTX_NOTES_PAGE_COUNT_MISMATCH"


def test_notes_and_source_bitmap_decode_timeouts_use_bounded_codes(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    monkeypatch.setattr(
        "app.ai.pptx_rendering.find_libreoffice_executable",
        lambda: "/test/soffice",
    )

    def fake_run(command: list[str], **_kwargs: object) -> subprocess.CompletedProcess[str]:
        output = Path(command[command.index("--outdir") + 1])
        write_test_pdf(output / "source.pdf", page_count=1)
        return subprocess.CompletedProcess(command, 0, "", "")

    def decode_timeout(*_args: object, **_kwargs: object) -> object:
        timeout_code = str(_kwargs["timeout_code"])
        raise PptxRenderResourceLimitError(timeout_code)

    monkeypatch.setattr("app.ai.pptx_rendering.subprocess.run", fake_run)
    monkeypatch.setattr(
        pptx_rendering,
        "run_bitmap_decode_with_timeout",
        decode_timeout,
    )
    with pytest.raises(PptxNotesRenderError) as notes_error:
        render_pptx_notes_to_png_assets(
            IMPORT_FIDELITY_NOTES_FIXTURE.read_bytes(),
            notes_width_emu=6_858_000,
            notes_height_emu=9_144_000,
            expected_page_count=1,
        )
    assert notes_error.value.code == "PPTX_NOTES_PREVIEW_DECODE_TIMEOUT"

    pdf_path = tmp_path / "source.pdf"
    write_test_pdf(pdf_path, page_count=1)
    monkeypatch.setattr(
        pptx_ooxml_generation,
        "run_bitmap_decode_with_timeout",
        decode_timeout,
    )
    with pytest.raises(PptxRenderUnavailableError) as source_error:
        render_pdf_to_png_assets(
            pdf_path,
            CanvasSpec("wide-16-9", 1920, 1080, "16:9"),
        )
    assert str(source_error.value) == "PPTX_SOURCE_RENDER_DECODE_TIMEOUT"


def test_generate_pptx_ooxml_maps_notes_preview_only_after_proven_render(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    notes_asset = make_test_png_asset("notes_render_1", "notes-01.png")
    monkeypatch.setattr(
        pptx_ooxml_generation,
        "render_pptx_to_png_assets",
        lambda _package, _canvas: [make_test_png_asset("slide_render_1", "slide-01.png")],
    )
    monkeypatch.setattr(
        pptx_ooxml_generation,
        "render_pptx_notes_to_png_assets",
        lambda *_args, **_kwargs: [notes_asset],
    )

    result = generate_pptx_ooxml(
        IMPORT_FIDELITY_NOTES_FIXTURE,
        "file_notes",
        render=True,
    )

    notes_page = result.template_blueprint["slides"][0]["notesPage"]
    assert notes_page["status"] == "rendered"
    assert notes_page["renderAssetFileId"] == "asset:notes_render_1"
    assert result.quality_report["notesDiagnostics"]["rendered"] == 1
    assert any(asset.asset_id == "notes_render_1" for asset in result.assets)


def test_generation_preserves_unsafe_source_but_sanitizes_renderer_input(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    source_bytes = package_with_external_relationship_and_macro(
        IMPORT_FIDELITY_NOTES_FIXTURE.read_bytes()
    )
    source_path = tmp_path / "unsafe-content.pptx"
    source_path.write_bytes(source_bytes)
    rendered_packages: list[bytes] = []

    def capture_render(package_bytes: bytes, _canvas: object) -> list[ImportedDesignAsset]:
        rendered_packages.append(package_bytes)
        return [make_test_png_asset("slide_render_1", "slide-01.png")]

    monkeypatch.setattr(
        pptx_ooxml_generation,
        "render_pptx_to_png_assets",
        capture_render,
    )
    monkeypatch.setattr(
        pptx_ooxml_generation,
        "render_pptx_notes_to_png_assets",
        lambda *_args, **_kwargs: [
            make_test_png_asset("notes_render_1", "notes-01.png")
        ],
    )

    result = generate_pptx_ooxml(source_path, "file_unsafe", render=True)

    assert result.warnings[:2] == [
        PPTX_EXTERNAL_RELATIONSHIP_BLOCKED,
        PPTX_ACTIVE_CONTENT_BLOCKED,
    ]
    assert "example.invalid" not in " ".join(result.warnings)
    assert current_package_bytes(result.assets) == source_bytes
    assert rendered_packages
    with zipfile.ZipFile(BytesIO(rendered_packages[0]), "r") as package:
        assert "ppt/vbaProject.bin" not in package.namelist()
        relationships = package.read("ppt/slides/_rels/slide1.xml.rels")
    assert b"rIdExternalBlocked" not in relationships


@pytest.mark.parametrize(
    "code",
    [
        "PPTX_NOTES_RENDERER_UNAVAILABLE",
        "PPTX_NOTES_RENDER_TIMEOUT",
        "PPTX_NOTES_PAGE_COUNT_MISMATCH",
    ],
)
def test_generate_pptx_ooxml_preserves_package_when_notes_render_is_unavailable(
    monkeypatch: pytest.MonkeyPatch,
    code: PptxNotesRenderErrorCode,
) -> None:
    monkeypatch.setattr(
        pptx_ooxml_generation,
        "render_pptx_to_png_assets",
        lambda _package, _canvas: [make_test_png_asset("slide_render_1", "slide-01.png")],
    )

    def unavailable(*_args: object, **_kwargs: object) -> list[ImportedDesignAsset]:
        raise PptxNotesRenderError(code)

    monkeypatch.setattr(
        pptx_ooxml_generation,
        "render_pptx_notes_to_png_assets",
        unavailable,
    )

    result = generate_pptx_ooxml(
        IMPORT_FIDELITY_NOTES_FIXTURE,
        "file_notes",
        render=True,
    )

    assert current_package_bytes(result.assets) == IMPORT_FIDELITY_NOTES_FIXTURE.read_bytes()
    assert result.blueprint["slides"][0]["speakerNotes"] == (
        "첫 번째 문단\n\n수동\n줄바꿈"
    )
    assert result.template_blueprint["slides"][0]["notesPage"]["status"] == (
        "render-unavailable"
    )
    assert result.quality_report["notesDiagnostics"]["warnings"] == [
        {"code": code, "count": 1}
    ]


def test_pure_generation_preserves_package_entries_and_source_text(
    tmp_path: Path,
) -> None:
    pptx_path = sample_pptx(tmp_path)

    result = generate_pptx_ooxml(pptx_path, "file_template", render=False)
    package_asset = next(
        asset for asset in result.assets if asset.asset_id == "current_package"
    )
    package_bytes = base64.b64decode(package_asset.content_base64)

    assert package_bytes == pptx_path.read_bytes()
    assert zip_entry_hashes(package_bytes) == zip_entry_hashes(pptx_path.read_bytes())
    assert len(result.template_blueprint["slides"]) == 1
    assert result.template_blueprint["currentPackageFileId"] == "asset:current_package"
    assert result.blueprint["slides"][0]["elements"]
    assert any(
        element["type"] == "text"
        and element["props"]["text"] == "Placeholder Title"
        for element in result.blueprint["slides"][0]["elements"]
    )


def test_apply_slot_texts_route_is_not_registered() -> None:
    assert "/ai/pptx-ooxml-apply-slot-texts" not in app.openapi()["paths"]


def test_generation_blueprint_uses_detected_canvas(tmp_path: Path) -> None:
    pptx_path = sample_pptx(tmp_path, wide=False)

    result = generate_pptx_ooxml(pptx_path, "file_template", render=False)
    background = next(
        element
        for element in result.blueprint["slides"][0]["elements"]
        if element["role"] == "background"
    )
    title_slot = next(
        slot
        for slot in result.template_blueprint["slides"][0]["slots"]
        if slot["usage"] == "content-slot"
    )

    assert result.canvas["preset"] == "standard-4-3"
    assert result.blueprint["canvas"] == {"width": 1024, "height": 768}
    assert background["width"] == 1024
    assert background["height"] == 768
    assert title_slot["bounds"]["width"] <= 1024
    assert title_slot["bounds"]["height"] <= 768


def test_extracts_slot_mapping(tmp_path: Path) -> None:
    pptx_path = sample_pptx(tmp_path)
    result = generate_pptx_ooxml(pptx_path, "file_template", render=False)
    slots = result.template_blueprint["slides"][0]["slots"]

    assert any(slot["usage"] == "content-slot" for slot in slots)
    assert any(
        slot["usage"] == "media-slot" and slot["replaceMode"] == "replace"
        for slot in slots
    )
    assert all(
        slot["source"].get("slidePart") == "ppt/slides/slide1.xml" for slot in slots
    )


def test_generation_includes_imported_image_assets(tmp_path: Path) -> None:
    pptx_path = sample_pptx(tmp_path)

    result = generate_pptx_ooxml(pptx_path, "file_template", render=False)
    image_asset = next(asset for asset in result.assets if asset.asset_id == "image_1")

    assert image_asset.mime_type == "image/png"
    assert base64.b64decode(image_asset.content_base64).startswith(b"\x89PNG")


def test_generation_deduplicates_repeated_media_bytes_by_content_hash(
    tmp_path: Path,
) -> None:
    pptx_path = sample_pptx(tmp_path)
    presentation = Presentation(pptx_path)
    presentation.slides[0].shapes.add_picture(
        str(tmp_path / "image.png"),
        Inches(9),
        Inches(2),
        Inches(2),
        Inches(2),
    )
    presentation.save(pptx_path)

    result = generate_pptx_ooxml(pptx_path, "file_template", render=False)
    image_assets = [
        asset for asset in result.assets if asset.asset_id.startswith("image_")
    ]
    image_sources = {
        element["props"]["src"]
        for element in result.blueprint["slides"][0]["elements"]
        if element.get("type") == "image"
    }

    assert len(image_assets) == 1
    assert image_sources == {f"asset:{image_assets[0].asset_id}"}


def test_sync_pptx_ooxml_applies_text_and_frame_patch(tmp_path: Path) -> None:
    pptx_path = sample_pptx(tmp_path)
    generated = generate_pptx_ooxml(pptx_path, "file_template", render=False)
    title_slot = next(
        slot
        for slot in generated.template_blueprint["slides"][0]["slots"]
        if slot["usage"] == "content-slot"
    )
    title_element = next(
        element
        for element in generated.blueprint["slides"][0]["elements"]
        if element["elementId"] == title_slot["elementId"]
    )

    result = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        deck_canvas={
            "preset": "wide-16-9",
            "width": 1920,
            "height": 1080,
            "aspectRatio": "16:9",
        },
        synced_deck_version=2,
        render=False,
        operations=[
            {
                "type": "update_element_props",
                "slideId": template_slide_id(generated),
                "elementId": title_slot["elementId"],
                "props": {"text": "Synced Title"},
            },
            {
                "type": "update_element_frame",
                "slideId": template_slide_id(generated),
                "elementId": title_slot["elementId"],
                "frame": {
                    "role": title_element.get("role"),
                    "x": 96,
                    "y": 48,
                    "width": 640,
                    "height": 120,
                    "rotation": title_element["rotation"],
                    "opacity": title_element["opacity"],
                    "zIndex": title_element["zIndex"],
                    "locked": title_element["locked"],
                    "visible": title_element["visible"],
                },
            },
        ],
    )
    package_asset = next(
        asset for asset in result.assets if asset.asset_id == "current_package"
    )
    package_bytes = base64.b64decode(package_asset.content_base64)

    with zipfile.ZipFile(BytesIO(package_bytes), "r") as package:
        slide_xml = package.read("ppt/slides/slide1.xml")

    assert b"Synced Title" in slide_xml
    assert b"Placeholder Title" not in slide_xml
    assert b'x="609600"' in slide_xml
    assert b'cx="4064000"' in slide_xml
    slide_root = ET.fromstring(slide_xml)
    assert any(
        body_pr.get("horzOverflow") == "clip"
        and body_pr.get("vertOverflow") == "clip"
        and body_pr.get("wrap") == "square"
        for body_pr in slide_root.findall(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr"
        )
    )


def test_sync_pptx_ooxml_updates_only_notes_body_and_reimports_text(
    tmp_path: Path,
) -> None:
    source_bytes = IMPORT_FIDELITY_NOTES_FIXTURE.read_bytes()
    generated = generate_pptx_ooxml(
        IMPORT_FIDELITY_NOTES_FIXTURE,
        "file_notes_sync",
        render=False,
    )
    speaker_notes = "첫 번째 문단\n\n수정된 발표 지시문\n마지막 줄"

    result = sync_pptx_ooxml(
        IMPORT_FIDELITY_NOTES_FIXTURE,
        template_blueprint=generated.template_blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
        operations=[
            {
                "type": "update_speaker_notes",
                "slideId": template_slide_id(generated),
                "speakerNotes": speaker_notes,
            }
        ],
    )
    synced_bytes = current_package_bytes(result.assets)

    assert [item.model_dump(by_alias=True, exclude_none=True) for item in result.applied_operations] == [
        {
            "operationType": "update_speaker_notes",
            "slideId": template_slide_id(generated),
        }
    ]
    assert result.unsupported_operations == []
    assert package_entry(source_bytes, "ppt/notesSlides/_rels/notesSlide1.xml.rels") == package_entry(
        synced_bytes,
        "ppt/notesSlides/_rels/notesSlide1.xml.rels",
    )
    assert package_entry(source_bytes, "ppt/notesMasters/notesMaster1.xml") == package_entry(
        synced_bytes,
        "ppt/notesMasters/notesMaster1.xml",
    )
    assert notes_non_body_semantic_hash(source_bytes) == notes_non_body_semantic_hash(
        synced_bytes
    )
    assert notes_body_structure_hash(source_bytes) == notes_body_structure_hash(
        synced_bytes
    )
    source_paragraphs = notes_body_paragraphs(source_bytes)
    synced_paragraphs = notes_body_paragraphs(synced_bytes)
    assert ET.tostring(source_paragraphs[0]) == ET.tostring(synced_paragraphs[0])
    assert {
        name: digest
        for name, digest in zip_entry_hashes(source_bytes).items()
        if name != "ppt/notesSlides/notesSlide1.xml"
    } == {
        name: digest
        for name, digest in zip_entry_hashes(synced_bytes).items()
        if name != "ppt/notesSlides/notesSlide1.xml"
    }

    synced_path = tmp_path / "notes-body-synced.pptx"
    synced_path.write_bytes(synced_bytes)
    reimported = generate_pptx_ooxml(
        synced_path,
        "file_notes_reimported",
        render=False,
    )
    assert reimported.blueprint["slides"][0]["speakerNotes"] == speaker_notes


def test_sync_pptx_ooxml_inherits_notes_paragraph_and_run_style(
    tmp_path: Path,
) -> None:
    source_path = tmp_path / "styled-notes.pptx"
    source_path.write_bytes(notes_package_with_body_style())
    generated = generate_pptx_ooxml(source_path, "file_styled_notes", render=False)

    result = sync_pptx_ooxml(
        source_path,
        template_blueprint=generated.template_blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
        operations=[
            {
                "type": "update_speaker_notes",
                "slideId": template_slide_id(generated),
                "speakerNotes": "첫 번째 문단\n\n새 발표 지시문",
            }
        ],
    )

    source_paragraph = notes_body_paragraphs(source_path.read_bytes())[2]
    synced_paragraph = notes_body_paragraphs(current_package_bytes(result.assets))[2]
    assert ET.tostring(source_paragraph.find(f"./{{{DRAWING_NS}}}pPr")) == ET.tostring(
        synced_paragraph.find(f"./{{{DRAWING_NS}}}pPr")
    )
    assert ET.tostring(
        source_paragraph.find(f"./{{{DRAWING_NS}}}r/{{{DRAWING_NS}}}rPr")
    ) == ET.tostring(
        synced_paragraph.find(f"./{{{DRAWING_NS}}}r/{{{DRAWING_NS}}}rPr")
    )


@pytest.mark.parametrize("locator_case", ["missing", "ambiguous"])
def test_sync_pptx_ooxml_rejects_unsafe_notes_body_locator_without_changes(
    locator_case: str,
) -> None:
    generated = generate_pptx_ooxml(
        IMPORT_FIDELITY_NOTES_FIXTURE,
        "file_notes_locator",
        render=False,
    )
    template_blueprint = copy.deepcopy(generated.template_blueprint)
    slide = template_blueprint["slides"][0]
    if locator_case == "missing":
        slide["notesPage"].pop("bodyShapeId")
    else:
        template_blueprint["slides"].append(copy.deepcopy(slide))

    result = sync_pptx_ooxml(
        IMPORT_FIDELITY_NOTES_FIXTURE,
        template_blueprint=template_blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
        operations=[
            {
                "type": "update_speaker_notes",
                "slideId": template_slide_id(generated),
                "speakerNotes": "안전하지 않은 locator에서는 저장하지 않음",
            }
        ],
    )

    assert current_package_bytes(result.assets) == IMPORT_FIDELITY_NOTES_FIXTURE.read_bytes()
    assert [
        item.model_dump(by_alias=True, exclude_none=True)
        for item in result.unsupported_operations
    ] == [
        {
            "operationType": "update_speaker_notes",
            "slideId": template_slide_id(generated),
            "reasonCode": "NOTES_BODY_LOCATOR_UNSAFE",
        }
    ]


def test_sync_pptx_ooxml_regenerates_notes_preview_after_body_edit(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    generated = generate_pptx_ooxml(
        IMPORT_FIDELITY_NOTES_FIXTURE,
        "file_notes_preview_sync",
        render=False,
    )
    monkeypatch.setattr(
        pptx_ooxml_generation,
        "render_pptx_to_png_assets",
        lambda *_args, **_kwargs: [
            make_test_png_asset("slide_render_1", "slide-01.png")
        ],
    )
    monkeypatch.setattr(
        pptx_ooxml_generation,
        "render_pptx_notes_to_png_assets",
        lambda *_args, **_kwargs: [
            make_test_png_asset("notes_render_1", "notes-01.png")
        ],
    )

    result = sync_pptx_ooxml(
        IMPORT_FIDELITY_NOTES_FIXTURE,
        template_blueprint=generated.template_blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=True,
        operations=[
            {
                "type": "update_speaker_notes",
                "slideId": template_slide_id(generated),
                "speakerNotes": "preview 갱신",
            }
        ],
    )

    assert {asset.asset_id for asset in result.assets} >= {
        "current_package",
        "slide_render_1",
        "notes_render_1",
    }


def test_sync_pptx_ooxml_creates_notes_page_without_existing_master(
    tmp_path: Path,
) -> None:
    source_path = sample_pptx(tmp_path)
    source_bytes = source_path.read_bytes()
    generated = generate_pptx_ooxml(source_path, "file_notes_create", render=False)
    assert generated.template_blueprint["slides"][0]["notesPage"]["status"] == (
        "absent"
    )
    speaker_notes = "새 발표 메모\n\n마지막 지시문"

    result = sync_pptx_ooxml(
        source_path,
        template_blueprint=generated.template_blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
        operations=[
            {
                "type": "update_speaker_notes",
                "slideId": template_slide_id(generated),
                "speakerNotes": speaker_notes,
            }
        ],
    )
    synced_bytes = current_package_bytes(result.assets)

    assert result.unsupported_operations == []
    assert [item.model_dump(by_alias=True) for item in result.notes_pages] == [
        {
            "slideId": template_slide_id(generated),
            "notesPage": {
                "status": "preserved",
                "sourceNotesPart": "ppt/notesSlides/notesSlide1.xml",
                "sourceNotesMasterPart": "ppt/notesMasters/notesMaster1.xml",
                "bodyShapeId": "3",
                "bodyWritable": True,
                "notesWidthEmu": 6_858_000,
                "notesHeightEmu": 9_144_000,
                "hasNonBodyContent": False,
            },
        }
    ]
    with zipfile.ZipFile(BytesIO(synced_bytes), "r") as package:
        names = set(package.namelist())
        assert "ppt/notesSlides/notesSlide1.xml" in names
        assert "ppt/notesSlides/_rels/notesSlide1.xml.rels" in names
        assert "ppt/notesMasters/notesMaster1.xml" in names
        assert "ppt/notesMasters/_rels/notesMaster1.xml.rels" in names
        assert notes_relationship_targets(
            package.read("ppt/slides/_rels/slide1.xml.rels")
        ) == [("notesSlide", "../notesSlides/notesSlide1.xml")]
        assert notes_relationship_targets(
            package.read("ppt/notesSlides/_rels/notesSlide1.xml.rels")
        ) == [
            ("notesMaster", "../notesMasters/notesMaster1.xml"),
            ("slide", "../slides/slide1.xml"),
        ]
        assert [
            item
            for item in notes_relationship_targets(
                package.read("ppt/_rels/presentation.xml.rels")
            )
            if item[0] == "notesMaster"
        ] == [("notesMaster", "notesMasters/notesMaster1.xml")]
        assert notes_content_type_parts(package.read("[Content_Types].xml")) == {
            "/ppt/notesMasters/notesMaster1.xml",
            "/ppt/notesSlides/notesSlide1.xml",
        }
    assert {
        name: digest
        for name, digest in zip_entry_hashes(source_bytes).items()
        if name
        not in {
            "[Content_Types].xml",
            "ppt/_rels/presentation.xml.rels",
            "ppt/slides/_rels/slide1.xml.rels",
        }
    } == {
        name: digest
        for name, digest in zip_entry_hashes(synced_bytes).items()
        if name in zip_entry_hashes(source_bytes)
        and name
        not in {
            "[Content_Types].xml",
            "ppt/_rels/presentation.xml.rels",
            "ppt/slides/_rels/slide1.xml.rels",
        }
    }

    synced_path = tmp_path / "created-notes.pptx"
    synced_path.write_bytes(synced_bytes)
    assert Presentation(synced_path).slides[0].notes_slide.notes_text_frame.text == (
        speaker_notes
    )
    reimported = generate_pptx_ooxml(
        synced_path,
        "file_notes_created_reimport",
        render=False,
    )
    assert reimported.blueprint["slides"][0]["speakerNotes"] == speaker_notes


def test_sync_pptx_ooxml_reuses_existing_notes_master_for_new_page(
    tmp_path: Path,
) -> None:
    source_path = sample_pptx_with_one_of_two_notes_pages(tmp_path)
    source_bytes = source_path.read_bytes()
    generated = generate_pptx_ooxml(
        source_path,
        "file_notes_existing_master",
        render=False,
    )
    target_slide_id = template_slide_id(generated, 1)
    assert generated.template_blueprint["slides"][1]["notesPage"]["status"] == (
        "absent"
    )

    result = sync_pptx_ooxml(
        source_path,
        template_blueprint=generated.template_blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
        operations=[
            {
                "type": "update_speaker_notes",
                "slideId": target_slide_id,
                "speakerNotes": "두 번째 슬라이드 발표 메모",
            }
        ],
    )
    synced_bytes = current_package_bytes(result.assets)

    assert result.unsupported_operations == []
    assert package_entry(source_bytes, "ppt/notesSlides/notesSlide1.xml") == (
        package_entry(synced_bytes, "ppt/notesSlides/notesSlide1.xml")
    )
    assert package_entry(source_bytes, "ppt/notesMasters/notesMaster1.xml") == (
        package_entry(synced_bytes, "ppt/notesMasters/notesMaster1.xml")
    )
    with zipfile.ZipFile(BytesIO(synced_bytes), "r") as package:
        assert sorted(
            name
            for name in package.namelist()
            if name.startswith("ppt/notesMasters/notesMaster")
            and name.endswith(".xml")
        ) == ["ppt/notesMasters/notesMaster1.xml"]
        assert notes_relationship_targets(
            package.read("ppt/notesSlides/_rels/notesSlide2.xml.rels")
        ) == [
            ("notesMaster", "../notesMasters/notesMaster1.xml"),
            ("slide", "../slides/slide2.xml"),
        ]

    synced_path = tmp_path / "existing-master-notes.pptx"
    synced_path.write_bytes(synced_bytes)
    reimported = generate_pptx_ooxml(
        synced_path,
        "file_notes_existing_master_reimport",
        render=False,
    )
    assert reimported.blueprint["slides"][1]["speakerNotes"] == (
        "두 번째 슬라이드 발표 메모"
    )


def test_sync_pptx_ooxml_creates_multiple_notes_pages_sharing_new_master(
    tmp_path: Path,
) -> None:
    pptx_path = tmp_path / "two-notesless-slides.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(pptx_path)

    generated = generate_pptx_ooxml(pptx_path, "file_two_notes", render=False)
    assert generated.template_blueprint["slides"][0]["notesPage"]["status"] == (
        "absent"
    )
    assert generated.template_blueprint["slides"][1]["notesPage"]["status"] == (
        "absent"
    )

    result = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
        operations=[
            {
                "type": "update_speaker_notes",
                "slideId": template_slide_id(generated, 0),
                "speakerNotes": "첫 번째 슬라이드 발표 메모",
            },
            {
                "type": "update_speaker_notes",
                "slideId": template_slide_id(generated, 1),
                "speakerNotes": "두 번째 슬라이드 발표 메모",
            },
        ],
    )
    synced_bytes = current_package_bytes(result.assets)

    # The first operation creates the notes master (in added_entries) and the
    # second must reuse it instead of failing NOTES_MASTER_CAPABILITY_UNSAFE.
    assert result.unsupported_operations == []
    notes_pages = [item.model_dump(by_alias=True) for item in result.notes_pages]
    assert len(notes_pages) == 2
    assert {
        page["notesPage"]["sourceNotesMasterPart"] for page in notes_pages
    } == {"ppt/notesMasters/notesMaster1.xml"}

    with zipfile.ZipFile(BytesIO(synced_bytes), "r") as package:
        names = set(package.namelist())
        assert "ppt/notesSlides/notesSlide1.xml" in names
        assert "ppt/notesSlides/notesSlide2.xml" in names
        assert sorted(
            name
            for name in names
            if name.startswith("ppt/notesMasters/notesMaster")
            and name.endswith(".xml")
        ) == ["ppt/notesMasters/notesMaster1.xml"]
        assert notes_relationship_targets(
            package.read("ppt/notesSlides/_rels/notesSlide2.xml.rels")
        ) == [
            ("notesMaster", "../notesMasters/notesMaster1.xml"),
            ("slide", "../slides/slide2.xml"),
        ]

    synced_path = tmp_path / "two-created-notes.pptx"
    synced_path.write_bytes(synced_bytes)
    reopened = Presentation(synced_path)
    assert reopened.slides[0].notes_slide.notes_text_frame.text == (
        "첫 번째 슬라이드 발표 메모"
    )
    assert reopened.slides[1].notes_slide.notes_text_frame.text == (
        "두 번째 슬라이드 발표 메모"
    )
    reimported = generate_pptx_ooxml(
        synced_path,
        "file_two_notes_reimport",
        render=False,
    )
    assert reimported.blueprint["slides"][0]["speakerNotes"] == (
        "첫 번째 슬라이드 발표 메모"
    )
    assert reimported.blueprint["slides"][1]["speakerNotes"] == (
        "두 번째 슬라이드 발표 메모"
    )


def test_sync_pptx_ooxml_rejects_unsafe_minimal_notes_master_atomically(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    source_path = sample_pptx(tmp_path)
    generated = generate_pptx_ooxml(
        source_path,
        "file_notes_unsafe_master",
        render=False,
    )
    monkeypatch.setattr(
        pptx_ooxml_generation,
        "minimal_notes_package_template",
        lambda: None,
    )

    result = sync_pptx_ooxml(
        source_path,
        template_blueprint=generated.template_blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
        operations=[
            {
                "type": "update_speaker_notes",
                "slideId": template_slide_id(generated),
                "speakerNotes": "생성 불가 시 저장하지 않음",
            }
        ],
    )

    assert current_package_bytes(result.assets) == source_path.read_bytes()
    assert [
        item.model_dump(by_alias=True, exclude_none=True)
        for item in result.unsupported_operations
    ] == [
        {
            "operationType": "update_speaker_notes",
            "slideId": template_slide_id(generated),
            "reasonCode": "NOTES_MASTER_CAPABILITY_UNSAFE",
        }
    ]


def test_sync_pptx_ooxml_rejects_unsafe_existing_notes_master_atomically(
    tmp_path: Path,
) -> None:
    source_path = sample_pptx_with_one_of_two_notes_pages(tmp_path)
    notes_master_rels_part = (
        "ppt/notesMasters/_rels/notesMaster1.xml.rels"
    )
    root = ET.fromstring(
        package_entry(source_path.read_bytes(), notes_master_rels_part)
    )
    theme_relationship = next(
        relationship
        for relationship in root
        if str(relationship.get("Type", "")).endswith("/theme")
    )
    theme_relationship.set("TargetMode", "External")
    theme_relationship.set("Target", "https://invalid.example/theme.xml")
    source_path.write_bytes(
        replace_package_entry(
            source_path.read_bytes(),
            notes_master_rels_part,
            ET.tostring(root, encoding="utf-8", xml_declaration=True),
        )
    )
    source_bytes = source_path.read_bytes()
    generated = generate_pptx_ooxml(
        source_path,
        "file_notes_unsafe_existing_master",
        render=False,
    )

    result = sync_pptx_ooxml(
        source_path,
        template_blueprint=generated.template_blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
        operations=[
            {
                "type": "update_speaker_notes",
                "slideId": template_slide_id(generated, 1),
                "speakerNotes": "안전하지 않은 master에는 생성하지 않음",
            }
        ],
    )

    assert current_package_bytes(result.assets) == source_bytes
    assert [
        item.model_dump(by_alias=True, exclude_none=True)
        for item in result.unsupported_operations
    ] == [
        {
            "operationType": "update_speaker_notes",
            "slideId": template_slide_id(generated, 1),
            "reasonCode": "NOTES_MASTER_CAPABILITY_UNSAFE",
        }
    ]


def test_created_notes_page_opens_in_libreoffice(tmp_path: Path) -> None:
    if not (shutil.which("libreoffice") or shutil.which("soffice")):
        pytest.skip("LibreOffice is not installed.")
    source_path = sample_pptx(tmp_path)
    generated = generate_pptx_ooxml(source_path, "file_notes_open", render=False)
    result = sync_pptx_ooxml(
        source_path,
        template_blueprint=generated.template_blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
        operations=[
            {
                "type": "update_speaker_notes",
                "slideId": template_slide_id(generated),
                "speakerNotes": "LibreOffice 재개방 검증",
            }
        ],
    )

    assert result.unsupported_operations == []
    assets = render_pptx_notes_to_png_assets(
        current_package_bytes(result.assets),
        notes_width_emu=6_858_000,
        notes_height_emu=9_144_000,
        expected_page_count=1,
    )
    assert [asset.asset_id for asset in assets] == ["notes_render_1"]


def test_sync_pptx_ooxml_skips_grouped_child_frame_patch(tmp_path: Path) -> None:
    pptx_path = sample_scaled_group_pptx(tmp_path)
    original_bytes = pptx_path.read_bytes()
    generated = generate_pptx_ooxml(pptx_path, "file_template", render=False)
    target = next(
        element
        for element in generated.blueprint["slides"][0]["elements"]
        if element.get("props", {}).get("text") == "Grouped frame target"
    )
    original_frame = {
        key: target[key] for key in ("x", "y", "width", "height", "rotation")
    }

    result = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
        operations=[
            {
                "type": "update_element_frame",
                "slideId": template_slide_id(generated),
                "elementId": target["elementId"],
                "frame": {
                    "x": target["x"] + 100,
                    "y": target["y"] + 50,
                    "width": target["width"],
                    "height": target["height"],
                },
            }
        ],
    )
    package_bytes = current_package_bytes(result.assets)

    assert package_bytes == original_bytes
    assert result.warnings == [
        f"OOXML grouped frame sync skipped for {target['elementId']}."
    ]
    assert result.applied_operations == []
    assert [
        operation.reason_code for operation in result.unsupported_operations
    ] == ["GROUPED_FRAME_UNSUPPORTED"]

    synced_path = tmp_path / "grouped-frame-synced.pptx"
    synced_path.write_bytes(package_bytes)
    reimported = generate_pptx_ooxml(synced_path, "file_template", render=False)
    reimported_target = next(
        element
        for element in reimported.blueprint["slides"][0]["elements"]
        if element.get("props", {}).get("text") == "Grouped frame target"
    )
    assert {
        key: reimported_target[key]
        for key in ("x", "y", "width", "height", "rotation")
    } == original_frame


def test_sync_pptx_ooxml_round_trips_text_and_target_image(
    tmp_path: Path,
) -> None:
    pptx_path = sample_round_trip_pptx(tmp_path)
    original_bytes = pptx_path.read_bytes()
    generated = generate_pptx_ooxml(pptx_path, "file_template", render=False)
    elements = generated.blueprint["slides"][0]["elements"]
    sources = generated.template_blueprint["slides"][0]["elementSources"]
    title = next(element for element in elements if element["type"] == "text")
    images = [element for element in elements if element["type"] == "image"]
    target_image = next(
        element for element in images if element["props"]["src"] == "asset:image_2"
    )
    untouched_image = next(
        element for element in images if element["props"]["src"] == "asset:image_1"
    )
    fallback_source = next(source for source in sources if source.get("fallbackReason"))
    target_source = source_for_element(sources, target_image["elementId"])
    untouched_source = source_for_element(sources, untouched_image["elementId"])
    replacement_bytes = png_bytes("#00ff00")
    replacement_data_url = (
        "data:image/png;base64," + base64.b64encode(replacement_bytes).decode("ascii")
    )

    result = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
        operations=[
            {
                "type": "update_element_props",
                "slideId": template_slide_id(generated),
                "elementId": title["elementId"],
                "props": {"text": "Synced round-trip title"},
            },
            {
                "type": "update_element_props",
                "slideId": template_slide_id(generated),
                "elementId": target_image["elementId"],
                "props": {"src": replacement_data_url},
            },
        ],
    )
    package_bytes = current_package_bytes(result.assets)
    target_relationship_id = picture_relationship_id(
        package_bytes, target_source["shapeId"]
    )

    assert result.warnings == []
    assert len(result.applied_operations) == 2
    assert result.unsupported_operations == []
    assert target_relationship_id != target_source["relationshipId"]
    assert source_for_element(
        result.element_sources,
        target_image["elementId"],
    )["relationshipId"] == target_relationship_id
    assert (
        picture_relationship_id(package_bytes, untouched_source["shapeId"])
        == untouched_source["relationshipId"]
    )
    assert relationship_blob(
        package_bytes,
        "ppt/slides/slide1.xml",
        target_relationship_id,
    ) == replacement_bytes
    assert relationship_blob(
        package_bytes,
        "ppt/slides/slide1.xml",
        untouched_source["relationshipId"],
    ) == relationship_blob(
        original_bytes,
        "ppt/slides/slide1.xml",
        untouched_source["relationshipId"],
    )
    assert shape_xml(package_bytes, fallback_source["shapeId"]) == shape_xml(
        original_bytes,
        fallback_source["shapeId"],
    )

    round_trip_path = tmp_path / "round-trip.pptx"
    round_trip_path.write_bytes(package_bytes)
    round_trip = generate_pptx_ooxml(
        round_trip_path,
        "file_round_trip",
        render=False,
    )

    assert any(
        element["type"] == "text"
        and element["props"]["text"] == "Synced round-trip title"
        for element in round_trip.blueprint["slides"][0]["elements"]
    )
    assert replacement_bytes in {
        base64.b64decode(asset.content_base64)
        for asset in round_trip.assets
        if asset.mime_type == "image/png"
    }
    assert any(
        source.get("fallbackReason") == fallback_source["fallbackReason"]
        for source in round_trip.template_blueprint["slides"][0]["elementSources"]
    )


def test_sync_pptx_ooxml_round_trips_image_crop_and_rejects_unsafe_capability(
    tmp_path: Path,
) -> None:
    pptx_path = sample_round_trip_pptx(tmp_path)
    original_bytes = pptx_path.read_bytes()
    generated = generate_pptx_ooxml(pptx_path, "file_template", render=False)
    target = next(
        element
        for element in generated.blueprint["slides"][0]["elements"]
        if element["type"] == "image" and element["props"]["src"] == "asset:image_2"
    )
    source = source_for_element(
        generated.template_blueprint["slides"][0]["elementSources"],
        target["elementId"],
    )
    crop = {"left": 0.2, "top": 0.1, "right": 0.15, "bottom": 0.05}
    operation = {
        "type": "update_element_props",
        "slideId": template_slide_id(generated),
        "elementId": target["elementId"],
        "props": {"crop": crop},
    }

    assert source["ooxmlEditCapabilities"]["crop"] == "picture"

    synced = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
        operations=[operation],
    )
    synced_bytes = current_package_bytes(synced.assets)

    assert synced.warnings == []
    assert len(synced.applied_operations) == 1
    assert synced.unsupported_operations == []
    assert picture_crop_rect(synced_bytes, source["shapeId"]) == {
        "l": "20000",
        "t": "10000",
        "r": "15000",
        "b": "5000",
    }

    synced_path = tmp_path / "crop-synced.pptx"
    synced_path.write_bytes(synced_bytes)
    reimported = generate_pptx_ooxml(synced_path, "file_crop", render=False)
    reimported_target = next(
        element
        for element in reimported.blueprint["slides"][0]["elements"]
        if element["type"] == "image" and element["props"].get("crop") == crop
    )
    assert reimported_target["props"]["crop"] == crop

    unsafe_blueprint = copy.deepcopy(generated.template_blueprint)
    unsafe_source = source_for_element(
        unsafe_blueprint["slides"][0]["elementSources"], target["elementId"]
    )
    unsafe_source["ooxmlEditCapabilities"]["crop"] = "none"
    rejected = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=unsafe_blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
        operations=[operation],
    )

    assert current_package_bytes(rejected.assets) == original_bytes
    assert rejected.applied_operations == []
    assert [
        unsupported.reason_code for unsupported in rejected.unsupported_operations
    ] == ["CROP_CAPABILITY_UNSAFE"]


def test_sync_pptx_ooxml_adds_writable_text_rect_and_image(
    tmp_path: Path,
) -> None:
    pptx_path = sample_round_trip_pptx(tmp_path)
    generated = generate_pptx_ooxml(pptx_path, "file_template", render=False)
    first_image = png_bytes("#00ff00")
    second_image = png_bytes("#ffff00")
    added_elements = [
        {
            "elementId": "el_added_text",
            "type": "text",
            "x": 100,
            "y": 600,
            "width": 500,
            "height": 100,
            "props": {"text": "Added text"},
        },
        {
            "elementId": "el_added_rect",
            "type": "rect",
            "x": 700,
            "y": 600,
            "width": 300,
            "height": 100,
            "props": {
                "fill": "#336699",
                "stroke": "#0090FF",
                "strokeWidth": 3,
                "borderRadius": 18,
            },
        },
        {
            "elementId": "el_added_image",
            "type": "image",
            "x": 1100,
            "y": 550,
            "width": 240,
            "height": 180,
            "props": {
                "src": "data:image/png;base64,"
                + base64.b64encode(first_image).decode("ascii")
            },
        },
    ]

    added = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
        operations=[
            {
                "type": "add_element",
                "slideId": template_slide_id(generated),
                "element": element,
            }
            for element in added_elements
        ],
    )

    assert added.warnings == []
    added_sources = {
        source["elementId"]: source for source in added.element_sources
    }
    assert set(added_sources) == {
        "el_added_text",
        "el_added_rect",
        "el_added_image",
    }
    assert all(
        source["shapeId"] != "0" and source["writable"] is True
        for source in added_sources.values()
    )
    assert added_sources["el_added_image"]["relationshipId"].startswith("rId")
    assert (
        added_sources["el_added_image"]["ooxmlEditCapabilities"]["crop"]
        == "picture"
    )

    added_path = tmp_path / "added.pptx"
    added_path.write_bytes(current_package_bytes(added.assets))
    next_blueprint = copy.deepcopy(generated.template_blueprint)
    next_blueprint["slides"][0]["elementSources"].extend(added.element_sources)
    edited = sync_pptx_ooxml(
        added_path,
        template_blueprint=next_blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=3,
        render=False,
        operations=[
            {
                "type": "update_element_props",
                "slideId": template_slide_id(generated),
                "elementId": "el_added_text",
                "props": {"text": "Edited added text"},
            },
            {
                "type": "update_element_frame",
                "slideId": template_slide_id(generated),
                "elementId": "el_added_rect",
                "frame": {"x": 320, "y": 600, "width": 300, "height": 100},
            },
            {
                "type": "update_element_props",
                "slideId": template_slide_id(generated),
                "elementId": "el_added_rect",
                "props": {"fill": "#FEFEFE"},
            },
            {
                "type": "update_element_props",
                "slideId": template_slide_id(generated),
                "elementId": "el_added_rect",
                "props": {"stroke": "#FEFEFE"},
            },
            {
                "type": "update_element_props",
                "slideId": template_slide_id(generated),
                "elementId": "el_added_image",
                "props": {
                    "src": "data:image/png;base64,"
                    + base64.b64encode(second_image).decode("ascii"),
                    "crop": {
                        "left": 0.1,
                        "top": 0.2,
                        "right": 0.15,
                        "bottom": 0.05,
                    },
                },
            },
            {
                "type": "update_element_frame",
                "slideId": template_slide_id(generated),
                "elementId": "el_added_image",
                "frame": {
                    "opacity": 0.29,
                    "visible": False,
                    "zIndex": 17,
                },
            },
        ],
    )
    edited_bytes = current_package_bytes(edited.assets)

    assert edited.warnings == []
    assert b"Edited added text" in shape_xml(
        edited_bytes, added_sources["el_added_text"]["shapeId"]
    )
    assert b'x="2032000"' in shape_xml(
        edited_bytes, added_sources["el_added_rect"]["shapeId"]
    )
    edited_rect_xml = shape_xml(
        edited_bytes, added_sources["el_added_rect"]["shapeId"]
    )
    assert b'prst="roundRect"' in edited_rect_xml
    assert b'name="adj" fmla="val 18000"' in edited_rect_xml
    assert edited_rect_xml.count(b'<a:srgbClr val="FEFEFE"') == 2
    edited_image_source = source_for_element(
        edited.element_sources,
        "el_added_image",
    )
    assert relationship_blob(
        edited_bytes,
        "ppt/slides/slide1.xml",
        edited_image_source["relationshipId"],
    ) == second_image
    assert picture_crop_rect(
        edited_bytes, added_sources["el_added_image"]["shapeId"]
    ) == {"l": "10000", "t": "20000", "r": "15000", "b": "5000"}
    edited_image_xml = shape_xml(
        edited_bytes, added_sources["el_added_image"]["shapeId"]
    )
    assert b'alphaModFix amt="29000"' in edited_image_xml
    assert b'hidden="1"' in edited_image_xml

    edited_path = tmp_path / "edited-added.pptx"
    edited_path.write_bytes(edited_bytes)
    reshown_blueprint = copy.deepcopy(next_blueprint)
    reshown_blueprint["slides"][0]["elementSources"] = [
        edited_image_source
        if source["elementId"] == "el_added_image"
        else source
        for source in reshown_blueprint["slides"][0]["elementSources"]
    ]
    reshown = sync_pptx_ooxml(
        edited_path,
        template_blueprint=reshown_blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=4,
        render=False,
        operations=[
            {
                "type": "update_element_frame",
                "slideId": template_slide_id(generated),
                "elementId": "el_added_image",
                "frame": {"opacity": 0.29, "visible": True},
            }
        ],
    )
    reshown_bytes = current_package_bytes(reshown.assets)
    reshown_image_xml = shape_xml(
        reshown_bytes, added_sources["el_added_image"]["shapeId"]
    )
    assert b'alphaModFix amt="29000"' in reshown_image_xml
    assert b'hidden="1"' not in reshown_image_xml

    edited_path.write_bytes(reshown_bytes)
    round_trip = generate_pptx_ooxml(
        edited_path,
        "file_round_trip",
        render=False,
    )
    round_trip_shape_ids = {
        source["shapeId"]
        for source in round_trip.template_blueprint["slides"][0]["elementSources"]
    }

    assert {
        source["shapeId"] for source in added_sources.values()
    }.issubset(round_trip_shape_ids)
    assert any(
        element["type"] == "text" and element["props"]["text"] == "Edited added text"
        for element in round_trip.blueprint["slides"][0]["elements"]
    )
    assert second_image in {
        base64.b64decode(asset.content_base64)
        for asset in round_trip.assets
        if asset.mime_type == "image/png"
    }


def test_sync_pptx_ooxml_rasterizes_and_updates_authored_visual_elements(
    tmp_path: Path,
) -> None:
    pptx_path = sample_round_trip_pptx(tmp_path)
    generated = generate_pptx_ooxml(pptx_path, "file_template", render=False)
    slide_id = template_slide_id(generated)
    line = {
        "elementId": "el_raster_line",
        "type": "line",
        "x": 140,
        "y": 180,
        "width": 420,
        "height": 120,
        "rotation": 5,
        "props": {"stroke": "#2563EB", "strokeWidth": 8},
    }
    chart = {
        "elementId": "el_raster_chart",
        "type": "chart",
        "x": 620,
        "y": 180,
        "width": 520,
        "height": 320,
        "props": {
            "type": "bar",
            "title": "분기별 매출",
            "data": [
                {"label": "1Q", "value": 20},
                {"label": "2Q", "value": 42},
            ],
            "style": {"showDataLabels": True},
        },
    }
    theme = {
        "name": "Orbit",
        "fontFamily": "Inter",
        "textColor": "#111827",
        "accentColor": "#2563EB",
    }
    fallbacks = {
        "theme": theme,
        "elements": [
            {"slideId": slide_id, "element": line},
            {"slideId": slide_id, "element": chart},
        ],
    }

    added = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
        authored_element_fallbacks=fallbacks,
        operations=[
            {"type": "add_element", "slideId": slide_id, "element": line},
            {"type": "add_element", "slideId": slide_id, "element": chart},
        ],
    )

    assert added.unsupported_operations == []
    assert [item.element_id for item in added.applied_operations] == [
        "el_raster_line",
        "el_raster_chart",
    ]
    sources = {
        source["elementId"]: source for source in added.element_sources
    }
    assert set(sources) == {"el_raster_line", "el_raster_chart"}
    assert all(source["fallbackMode"] == "rasterized" for source in sources.values())
    assert all(source["sourceType"] == "image" for source in sources.values())
    added_bytes = current_package_bytes(added.assets)
    first_line_png = relationship_blob(
        added_bytes,
        "ppt/slides/slide1.xml",
        sources["el_raster_line"]["relationshipId"],
    )
    assert Image.open(BytesIO(first_line_png)).mode == "RGBA"

    updated_line = copy.deepcopy(line)
    updated_line["x"] = 260
    updated_line["props"] = {"stroke": "#EF4444", "strokeWidth": 14}
    next_blueprint = copy.deepcopy(generated.template_blueprint)
    next_blueprint["slides"][0]["elementSources"].extend(added.element_sources)
    added_path = tmp_path / "raster-added.pptx"
    added_path.write_bytes(added_bytes)
    updated = sync_pptx_ooxml(
        added_path,
        template_blueprint=next_blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=3,
        render=False,
        authored_element_fallbacks={
            "theme": theme,
            "elements": [{"slideId": slide_id, "element": updated_line}],
        },
        operations=[
            {
                "type": "update_element_props",
                "slideId": slide_id,
                "elementId": "el_raster_line",
                "props": updated_line["props"],
            },
            {
                "type": "update_element_frame",
                "slideId": slide_id,
                "elementId": "el_raster_line",
                "frame": {"x": 260},
            },
        ],
    )

    assert updated.unsupported_operations == []
    assert [item.operation_type for item in updated.applied_operations] == [
        "update_element_props",
        "update_element_frame",
    ]
    updated_source = source_for_element(updated.element_sources, "el_raster_line")
    assert updated_source["relationshipId"] == sources["el_raster_line"][
        "relationshipId"
    ]
    updated_bytes = current_package_bytes(updated.assets)
    assert relationship_blob(
        updated_bytes,
        "ppt/slides/slide1.xml",
        updated_source["relationshipId"],
    ) != first_line_png
    assert b'x="' in shape_xml(updated_bytes, updated_source["shapeId"])


def test_sync_pptx_ooxml_keeps_package_atomic_when_raster_candidate_is_missing(
    tmp_path: Path,
) -> None:
    pptx_path = sample_round_trip_pptx(tmp_path)
    original_bytes = pptx_path.read_bytes()
    generated = generate_pptx_ooxml(pptx_path, "file_template", render=False)
    result = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
        authored_element_fallbacks={"theme": {"name": "Orbit"}, "elements": []},
        operations=[
            {
                "type": "add_element",
                "slideId": template_slide_id(generated),
                "element": {
                    "elementId": "el_missing_fallback",
                    "type": "arrow",
                    "x": 10,
                    "y": 20,
                    "width": 300,
                    "height": 80,
                    "props": {"stroke": "#2563EB", "strokeWidth": 4},
                },
            }
        ],
    )

    assert current_package_bytes(result.assets) == original_bytes
    assert result.applied_operations == []
    assert [
        item.reason_code for item in result.unsupported_operations
    ] == ["AUTHORED_RASTER_FALLBACK_FAILED"]


def test_sync_pptx_ooxml_scopes_duplicate_element_ids_to_slide_part(
    tmp_path: Path,
) -> None:
    pptx_path, shape_ids = sample_duplicate_element_ids_pptx(tmp_path)
    original_bytes = pptx_path.read_bytes()
    generated = generate_pptx_ooxml(pptx_path, "file_template", render=False)
    blueprint = copy.deepcopy(generated.template_blueprint)
    mapped_sources: list[dict[str, dict]] = []
    for slide_index, ids in enumerate(shape_ids):
        sources = blueprint["slides"][slide_index]["elementSources"]
        text_source = next(source for source in sources if source["shapeId"] == ids["text"])
        image_source = next(
            source for source in sources if source["shapeId"] == ids["image"]
        )
        delete_source = next(
            source for source in sources if source["shapeId"] == ids["delete"]
        )
        text_source["elementId"] = "el_shared_text"
        image_source["elementId"] = "el_shared_image"
        delete_source["elementId"] = "el_shared_delete"
        assert image_source["ooxmlEditCapabilities"]["frame"] is True
        image_source["ooxmlEditCapabilities"]["frame"] = False
        assert delete_source["ooxmlEditCapabilities"]["delete"] is True
        delete_source["ooxmlEditCapabilities"]["delete"] = False
        mapped_sources.append(
            {"text": text_source, "image": image_source, "delete": delete_source}
        )

    replacement_image = png_bytes("#ffff00")
    synced = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
        operations=[
            {
                "type": "update_element_props",
                "slideId": blueprint["slides"][1]["slideId"],
                "elementId": "el_shared_text",
                "props": {"text": "Only slide two"},
            },
            {
                "type": "update_element_frame",
                "slideId": blueprint["slides"][1]["slideId"],
                "elementId": "el_shared_image",
                "frame": {"x": 400, "y": 300, "width": 240, "height": 180},
            },
            {
                "type": "update_element_props",
                "slideId": blueprint["slides"][1]["slideId"],
                "elementId": "el_shared_image",
                "props": {
                    "src": "data:image/png;base64,"
                    + base64.b64encode(replacement_image).decode("ascii")
                },
            },
            {
                "type": "delete_element",
                "slideId": blueprint["slides"][1]["slideId"],
                "elementId": "el_shared_delete",
            },
        ],
    )
    synced_bytes = current_package_bytes(synced.assets)

    assert synced.warnings == []
    assert package_entry(synced_bytes, "ppt/slides/slide1.xml") == package_entry(
        original_bytes, "ppt/slides/slide1.xml"
    )
    assert package_entry(
        synced_bytes, "ppt/slides/_rels/slide1.xml.rels"
    ) == package_entry(original_bytes, "ppt/slides/_rels/slide1.xml.rels")
    assert b"Only slide two" in shape_xml(
        synced_bytes,
        mapped_sources[1]["text"]["shapeId"],
        "ppt/slides/slide2.xml",
    )
    assert b'x="2540000"' in shape_xml(
        synced_bytes,
        mapped_sources[1]["image"]["shapeId"],
        "ppt/slides/slide2.xml",
    )
    synced_image_source = next(
        source
        for source in synced.element_sources
        if source["slidePart"] == "ppt/slides/slide2.xml"
        and source["elementId"] == "el_shared_image"
    )
    assert relationship_blob(
        synced_bytes,
        "ppt/slides/slide2.xml",
        synced_image_source["relationshipId"],
    ) == replacement_image
    with pytest.raises(AssertionError):
        shape_xml(
            synced_bytes,
            mapped_sources[1]["delete"]["shapeId"],
            "ppt/slides/slide2.xml",
        )

    synced_path = tmp_path / "duplicate-ids-synced.pptx"
    synced_path.write_bytes(synced_bytes)
    added = sync_pptx_ooxml(
        synced_path,
        template_blueprint=blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=3,
        render=False,
        operations=[
            {
                "type": "add_element",
                "slideId": blueprint["slides"][slide_index - 1]["slideId"],
                "element": {
                    "elementId": "el_shared_added",
                    "type": "text",
                    "x": 100,
                    "y": 600,
                    "width": 400,
                    "height": 80,
                    "props": {"text": f"Added on slide {slide_index}"},
                },
            }
            for slide_index in (1, 2)
        ],
    )
    added_sources = [
        source
        for source in added.element_sources
        if source["elementId"] == "el_shared_added"
    ]

    assert added.warnings == []
    assert {source["slidePart"] for source in added_sources} == {
        "ppt/slides/slide1.xml",
        "ppt/slides/slide2.xml",
    }


@pytest.mark.parametrize(
    "src",
    [
        "not-a-data-url",
        "data:image/svg+xml;base64,PHN2Zy8+",
        "data:image/png;base64,!!!",
        "data:image/png;base64,bm90LWFuLWltYWdl",
    ],
)
def test_sync_pptx_ooxml_rejects_invalid_image_data_without_package_changes(
    tmp_path: Path,
    src: str,
) -> None:
    pptx_path = sample_pptx(tmp_path)
    generated = generate_pptx_ooxml(pptx_path, "file_template", render=False)
    image = next(
        element
        for element in generated.blueprint["slides"][0]["elements"]
        if element["type"] == "image"
    )

    result = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
        operations=[
            {
                "type": "update_element_props",
                "slideId": template_slide_id(generated),
                "elementId": image["elementId"],
                "props": {"src": src},
            },
            {
                "type": "add_element",
                "slideId": template_slide_id(generated),
                "element": {
                    "elementId": "el_invalid_image",
                    "type": "image",
                    "x": 0,
                    "y": 0,
                    "width": 100,
                    "height": 100,
                    "props": {"src": src},
                },
            },
        ],
    )

    assert current_package_bytes(result.assets) == pptx_path.read_bytes()
    assert len(result.warnings) == 2
    assert all("OOXML" in warning and "image" in warning for warning in result.warnings)


def test_renders_slide_pngs_when_libreoffice_is_available(tmp_path: Path) -> None:
    if not (shutil.which("libreoffice") or shutil.which("soffice")):
        pytest.skip("LibreOffice is not installed.")

    pptx_path = sample_pptx(tmp_path)
    result = generate_pptx_ooxml(pptx_path, "file_template", render=False)
    canvas = type(
        "Canvas",
        (),
        {"width": 1920, "height": 1080, "preset": "wide-16-9", "aspect_ratio": "16:9"},
    )()

    try:
        assets = render_pptx_to_png_assets(
            base64.b64decode(result.assets[0].content_base64),
            canvas,
        )
    except PptxRenderUnavailableError as error:
        pytest.skip(str(error))

    assert len(assets) == 1
    assert assets[0].mime_type == "image/png"
    assert base64.b64decode(assets[0].content_base64).startswith(b"\x89PNG")


def test_shape_fallback_assets_crop_from_slide_render() -> None:
    slide_render = BytesIO()
    Image.new("RGB", (100, 80), "#336699").save(slide_render, format="PNG")
    warnings: list[str] = []

    assets = shape_fallback_assets(
        {
            "slides": [
                {
                    "sourceSlideIndex": 1,
                    "elements": [
                        {
                            "type": "image",
                            "x": 10,
                            "y": 15,
                            "width": 30,
                            "height": 25,
                            "props": {
                                "src": "asset:shape_render_1_slide_2",
                            },
                        }
                    ],
                }
            ]
        },
        [
            ImportedDesignAsset(
                assetId="slide_render_1",
                fileName="slide-01.png",
                mimeType="image/png",
                contentBase64=base64.b64encode(slide_render.getvalue()).decode(
                    "ascii"
                ),
            )
        ],
        warnings,
    )
    crop = Image.open(BytesIO(base64.b64decode(assets[0].content_base64)))

    assert warnings == []
    assert assets[0].asset_id == "shape_render_1_slide_2"
    assert assets[0].mime_type == "image/png"
    assert crop.size == (30, 25)


def test_strip_text_from_pptx_package_removes_text_bodies(tmp_path: Path) -> None:
    pptx_path = sample_pptx(tmp_path)

    stripped = strip_text_from_pptx_package(pptx_path.read_bytes())

    with zipfile.ZipFile(BytesIO(stripped), "r") as package:
        slide_xml = package.read("ppt/slides/slide1.xml")

    assert b"Placeholder Title" not in slide_xml
    assert b"<p:txBody>" not in slide_xml


def test_sync_pptx_ooxml_adds_authored_slide_and_same_batch_image(
    tmp_path: Path,
) -> None:
    pptx_path = sample_pptx(tmp_path)
    generated = generate_pptx_ooxml(pptx_path, "file_template", render=False)
    blueprint = copy.deepcopy(generated.template_blueprint)
    blueprint["slides"].append(
        {
            "slideId": "slide_authored",
            "slideIndex": 2,
            "sourceSlideIndex": 2,
            "sourceSlidePart": "ppt/slides/slide2.xml",
            "ooxmlOrigin": "authored",
            "slots": [],
            "elementSources": [],
        }
    )
    image_src = "data:image/png;base64," + base64.b64encode(
        png_bytes("#22c55e")
    ).decode("ascii")
    authored_line = {
        "elementId": "el_authored_line",
        "type": "line",
        "x": 160,
        "y": 500,
        "width": 480,
        "height": 80,
        "props": {"stroke": "#7C3AED", "strokeWidth": 6},
    }

    result = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
        authored_element_fallbacks={
            "theme": {"name": "Orbit", "accentColor": "#2563EB"},
            "elements": [
                {"slideId": "slide_authored", "element": authored_line}
            ],
        },
        operations=[
            {
                "type": "add_slide",
                "sourceSlidePart": "ppt/slides/slide2.xml",
                "slide": {
                    "slideId": "slide_authored",
                    "order": 2,
                    "title": "Authored slide",
                    "elements": [
                        {
                            "elementId": "el_authored_text",
                            "type": "text",
                            "x": 120,
                            "y": 100,
                            "width": 600,
                            "height": 100,
                            "props": {"text": "Authored title"},
                        },
                        {
                            "elementId": "el_authored_rect",
                            "type": "rect",
                            "x": 120,
                            "y": 260,
                            "width": 400,
                            "height": 180,
                            "props": {"fill": "#2563EB"},
                        },
                        authored_line,
                    ],
                },
            },
            {
                "type": "add_element",
                "slideId": "slide_authored",
                "sourceSlidePart": "ppt/slides/slide2.xml",
                "element": {
                    "elementId": "el_authored_image",
                    "type": "image",
                    "x": 800,
                    "y": 260,
                    "width": 320,
                    "height": 180,
                    "props": {"src": image_src, "fit": "contain"},
                },
            },
            {
                "type": "reorder_slides",
                "slideOrders": [
                    {
                        "slideId": "slide_authored",
                        "order": 1,
                        "sourceSlidePart": "ppt/slides/slide2.xml",
                    },
                    {
                        "slideId": template_slide_id(generated),
                        "order": 2,
                        "sourceSlidePart": "ppt/slides/slide1.xml",
                    },
                ],
            },
        ],
    )

    assert result.unsupported_operations == []
    assert [item.operation_type for item in result.applied_operations] == [
        "add_slide",
        "add_element",
        "reorder_slides",
    ]
    assert {source["elementId"] for source in result.element_sources} == {
        "el_authored_text",
        "el_authored_rect",
        "el_authored_line",
        "el_authored_image",
    }
    line_source = source_for_element(result.element_sources, "el_authored_line")
    assert line_source["fallbackMode"] == "rasterized"
    package_bytes = current_package_bytes(result.assets)
    assert relationship_blob(
        package_bytes,
        "ppt/slides/slide2.xml",
        line_source["relationshipId"],
    ).startswith(b"\x89PNG")
    round_trip = Presentation(BytesIO(package_bytes))
    assert len(round_trip.slides) == 2
    assert any(
        "Authored title" in shape.text
        for shape in round_trip.slides[0].shapes
        if hasattr(shape, "text")
    )
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as package:
        assert "ppt/slides/slide2.xml" in package.namelist()
        rels = ET.fromstring(package.read("ppt/slides/_rels/slide2.xml.rels"))
        assert any(
            str(relationship.get("Type", "")).endswith("/slideLayout")
            for relationship in rels
        )
        assert b'/ppt/slides/slide2.xml' in package.read("[Content_Types].xml")


def sample_pptx(tmp_path: Path, *, wide: bool = True) -> Path:
    pptx_path = tmp_path / "template.pptx"
    image_path = tmp_path / "image.png"
    Image.new("RGB", (32, 32), "#ff0000").save(image_path)

    presentation = Presentation()
    presentation.slide_width = Inches(13.333333 if wide else 10)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text_frame.text = "Placeholder Title"
    slide.placeholders[1].text_frame.text = "Placeholder Subtitle"
    slide.shapes.add_picture(
        str(image_path), Inches(7), Inches(2), Inches(2), Inches(2)
    )
    presentation.save(pptx_path)
    return pptx_path


def sample_pptx_with_one_of_two_notes_pages(tmp_path: Path) -> Path:
    pptx_path = tmp_path / "one-of-two-notes.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    first = presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.slides.add_slide(presentation.slide_layouts[6])
    first.notes_slide.notes_text_frame.text = "기존 발표 메모"
    presentation.save(pptx_path)
    return pptx_path


def sample_round_trip_pptx(tmp_path: Path) -> Path:
    pptx_path = tmp_path / "round-trip-source.pptx"
    first_image_path = tmp_path / "first.png"
    second_image_path = tmp_path / "second.png"
    first_image_path.write_bytes(png_bytes("#ff0000"))
    second_image_path.write_bytes(png_bytes("#0000ff"))

    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(
        Inches(1), Inches(0.5), Inches(5), Inches(0.8)
    ).text_frame.text = "Original round-trip title"
    slide.shapes.add_picture(
        str(first_image_path), Inches(1), Inches(2), Inches(2), Inches(2)
    )
    slide.shapes.add_picture(
        str(second_image_path), Inches(4), Inches(2), Inches(2), Inches(2)
    )
    slide.shapes.add_shape(
        MSO_SHAPE.CLOUD,
        Inches(8),
        Inches(2),
        Inches(2),
        Inches(1.5),
    )
    presentation.save(pptx_path)
    return pptx_path


def sample_scaled_group_pptx(tmp_path: Path) -> Path:
    pptx_path = tmp_path / "scaled-group.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1),
        Inches(1),
        Inches(2),
        Inches(1),
    )
    target = slide.shapes.add_textbox(
        Inches(1.2),
        Inches(1.2),
        Inches(1.4),
        Inches(0.5),
    )
    target.text_frame.text = "Grouped frame target"
    group = slide.shapes.add_group_shape([box, target])
    group.left = Inches(4)
    group.top = Inches(2)
    group.width = Inches(5)
    group.height = Inches(2.5)
    presentation.save(pptx_path)
    return pptx_path


def sample_duplicate_element_ids_pptx(
    tmp_path: Path,
) -> tuple[Path, list[dict[str, str]]]:
    pptx_path = tmp_path / "duplicate-element-ids.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    shape_ids: list[dict[str, str]] = []
    for slide_index, color in enumerate(("#ff0000", "#0000ff"), start=1):
        image_path = tmp_path / f"duplicate-{slide_index}.png"
        image_path.write_bytes(png_bytes(color))
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        text = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(5), Inches(0.8)
        )
        text.text_frame.text = f"Slide {slide_index} title"
        image = slide.shapes.add_picture(
            str(image_path), Inches(1), Inches(2), Inches(2), Inches(2)
        )
        delete_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(4),
            Inches(2),
            Inches(2),
            Inches(1),
        )
        shape_ids.append(
            {
                "text": str(text.shape_id),
                "image": str(image.shape_id),
                "delete": str(delete_shape.shape_id),
            }
        )
    presentation.save(pptx_path)
    return pptx_path, shape_ids


def write_test_pdf(path: Path, *, page_count: int) -> None:
    fitz = importlib.import_module("fitz")
    document = fitz.open()
    try:
        for _index in range(page_count):
            document.new_page(width=600, height=800)
        document.save(str(path))
    finally:
        document.close()


def make_test_png_asset(asset_id: str, file_name: str) -> ImportedDesignAsset:
    return ImportedDesignAsset(
        assetId=asset_id,
        fileName=file_name,
        mimeType="image/png",
        contentBase64=base64.b64encode(png_bytes("#336699")).decode("ascii"),
    )


def png_bytes(color: str) -> bytes:
    output = BytesIO()
    Image.new("RGB", (8, 8), color).save(output, format="PNG")
    return output.getvalue()


def current_package_bytes(assets: list[ImportedDesignAsset]) -> bytes:
    package = next(asset for asset in assets if asset.asset_id == "current_package")
    return base64.b64decode(package.content_base64)


def source_for_element(sources: list[dict], element_id: str) -> dict:
    return next(source for source in sources if source["elementId"] == element_id)


def picture_relationship_id(package_bytes: bytes, shape_id: str) -> str:
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as package:
        root = ET.fromstring(package.read("ppt/slides/slide1.xml"))
    for picture in root.iter():
        if not picture.tag.endswith("pic"):
            continue
        c_nv_pr = next(
            (node for node in picture.iter() if node.tag.endswith("cNvPr")),
            None,
        )
        if c_nv_pr is None or c_nv_pr.get("id") != shape_id:
            continue
        blip = next(node for node in picture.iter() if node.tag.endswith("blip"))
        return next(value for key, value in blip.attrib.items() if key.endswith("embed"))
    raise AssertionError(f"picture shape not found: {shape_id}")


def picture_crop_rect(package_bytes: bytes, shape_id: str) -> dict[str, str]:
    root = ET.fromstring(shape_xml(package_bytes, shape_id))
    source_rect = next(
        (node for node in root.iter() if node.tag.endswith("srcRect")),
        None,
    )
    if source_rect is None:
        raise AssertionError(f"picture crop not found: {shape_id}")
    return dict(source_rect.attrib)


def relationship_blob(
    package_bytes: bytes,
    slide_part: str,
    relationship_id: str,
) -> bytes:
    slide_path = Path(slide_part)
    rels_part = f"{slide_path.parent.as_posix()}/_rels/{slide_path.name}.rels"
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as package:
        root = ET.fromstring(package.read(rels_part))
        relationship = next(
            child for child in root if child.get("Id") == relationship_id
        )
        target = relationship.get("Target", "")
        media_part = str((Path(slide_part).parent / target).resolve()).replace("\\", "/")
        media_part = media_part.split("/ppt/", maxsplit=1)[-1]
        return package.read(f"ppt/{media_part}")


def shape_xml(
    package_bytes: bytes,
    shape_id: str,
    slide_part: str = "ppt/slides/slide1.xml",
) -> bytes:
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as package:
        root = ET.fromstring(package.read(slide_part))
    for shape in root.iter():
        if not shape.tag.endswith(("sp", "pic")):
            continue
        c_nv_pr = next(
            (node for node in shape.iter() if node.tag.endswith("cNvPr")),
            None,
        )
        if c_nv_pr is not None and c_nv_pr.get("id") == shape_id:
            return ET.tostring(shape)
    raise AssertionError(f"shape not found: {shape_id}")


def package_entry(package_bytes: bytes, name: str) -> bytes:
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as package:
        return package.read(name)


def replace_package_entry(package_bytes: bytes, name: str, content: bytes) -> bytes:
    output = BytesIO()
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as source:
        with zipfile.ZipFile(output, "w") as destination:
            for info in source.infolist():
                destination.writestr(
                    info,
                    content if info.filename == name else source.read(info),
                )
    return output.getvalue()


def package_with_external_relationship_and_macro(package_bytes: bytes) -> bytes:
    relationships_part = "ppt/slides/_rels/slide1.xml.rels"
    relationships_root = ET.fromstring(
        package_entry(package_bytes, relationships_part)
    )
    ET.SubElement(
        relationships_root,
        f"{{{RELATIONSHIP_NS}}}Relationship",
        {
            "Id": "rIdExternalBlocked",
            "Type": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
            "Target": "https://example.invalid/image.png",
            "TargetMode": "External",
        },
    )
    content_types_root = ET.fromstring(
        package_entry(package_bytes, "[Content_Types].xml")
    )
    ET.SubElement(
        content_types_root,
        "{http://schemas.openxmlformats.org/package/2006/content-types}Override",
        {
            "PartName": "/ppt/vbaProject.bin",
            "ContentType": "application/vnd.ms-office.vbaProject",
        },
    )
    replacements = {
        relationships_part: ET.tostring(
            relationships_root,
            encoding="utf-8",
            xml_declaration=True,
        ),
        "[Content_Types].xml": ET.tostring(
            content_types_root,
            encoding="utf-8",
            xml_declaration=True,
        ),
    }
    output = BytesIO()
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as source:
        with zipfile.ZipFile(output, "w") as destination:
            for info in source.infolist():
                destination.writestr(
                    info,
                    replacements.get(info.filename, source.read(info)),
                )
            destination.writestr("ppt/vbaProject.bin", b"synthetic macro")
    return output.getvalue()


def notes_relationship_targets(rels_xml: bytes) -> list[tuple[str, str]]:
    root = ET.fromstring(rels_xml)
    return [
        (relationship_type, str(relationship.get("Target", "")))
        for relationship in root
        if (
            relationship_type := str(relationship.get("Type", "")).rsplit(
                "/", maxsplit=1
            )[-1]
        )
        in {"notesMaster", "notesSlide", "slide"}
    ]


def notes_content_type_parts(content_types_xml: bytes) -> set[str]:
    root = ET.fromstring(content_types_xml)
    return {
        str(item.get("PartName", ""))
        for item in root
        if str(item.get("ContentType", "")).endswith(
            ("notesMaster+xml", "notesSlide+xml")
        )
    }


def notes_non_body_semantic_hash(package_bytes: bytes) -> str:
    root = ET.fromstring(
        package_entry(package_bytes, "ppt/notesSlides/notesSlide1.xml")
    )
    shape_tree = root.find(f"./{{{PRESENTATION_NS}}}cSld/{{{PRESENTATION_NS}}}spTree")
    assert shape_tree is not None
    for shape in list(shape_tree):
        placeholder = shape.find(
            f"./{{{PRESENTATION_NS}}}nvSpPr/{{{PRESENTATION_NS}}}nvPr/"
            f"{{{PRESENTATION_NS}}}ph"
        )
        if placeholder is not None and placeholder.get("type") == "body":
            shape_tree.remove(shape)
    return hashlib.sha256(ET.tostring(root)).hexdigest()


def notes_body_structure_hash(package_bytes: bytes) -> str:
    root = ET.fromstring(
        package_entry(package_bytes, "ppt/notesSlides/notesSlide1.xml")
    )
    namespaces = {"a": DRAWING_NS, "p": PRESENTATION_NS}
    body_shape = next(
        shape
        for shape in root.findall("./p:cSld/p:spTree/p:sp", namespaces)
        if (
            placeholder := shape.find("./p:nvSpPr/p:nvPr/p:ph", namespaces)
        )
        is not None
        and placeholder.get("type") == "body"
    )
    structure = copy.deepcopy(body_shape)
    text_body = structure.find("./p:txBody", namespaces)
    assert text_body is not None
    for paragraph in text_body.findall("./a:p", namespaces):
        text_body.remove(paragraph)
    return hashlib.sha256(ET.tostring(structure)).hexdigest()


def notes_body_paragraphs(package_bytes: bytes) -> list[ET.Element]:
    root = ET.fromstring(
        package_entry(package_bytes, "ppt/notesSlides/notesSlide1.xml")
    )
    namespaces = {"a": DRAWING_NS, "p": PRESENTATION_NS}
    body_shape = next(
        shape
        for shape in root.findall("./p:cSld/p:spTree/p:sp", namespaces)
        if (
            placeholder := shape.find("./p:nvSpPr/p:nvPr/p:ph", namespaces)
        )
        is not None
        and placeholder.get("type") == "body"
    )
    return body_shape.findall("./p:txBody/a:p", namespaces)


def notes_package_with_body_style() -> bytes:
    source_bytes = IMPORT_FIDELITY_NOTES_FIXTURE.read_bytes()
    notes_part = "ppt/notesSlides/notesSlide1.xml"
    root = ET.fromstring(package_entry(source_bytes, notes_part))
    namespaces = {"a": DRAWING_NS, "p": PRESENTATION_NS}
    body_shape = next(
        shape
        for shape in root.findall("./p:cSld/p:spTree/p:sp", namespaces)
        if (
            placeholder := shape.find("./p:nvSpPr/p:nvPr/p:ph", namespaces)
        )
        is not None
        and placeholder.get("type") == "body"
    )
    paragraph = body_shape.findall("./p:txBody/a:p", namespaces)[2]
    paragraph.insert(0, ET.Element(f"{{{DRAWING_NS}}}pPr", {"algn": "ctr"}))
    run = paragraph.find(f"./{{{DRAWING_NS}}}r")
    assert run is not None
    run.insert(
        0,
        ET.Element(f"{{{DRAWING_NS}}}rPr", {"b": "1", "sz": "1800"}),
    )
    updated_notes = ET.tostring(root, encoding="utf-8", xml_declaration=True)

    output = BytesIO()
    with zipfile.ZipFile(BytesIO(source_bytes), "r") as source:
        with zipfile.ZipFile(output, "w") as destination:
            for info in source.infolist():
                destination.writestr(
                    info,
                    updated_notes if info.filename == notes_part else source.read(info),
                )
    return output.getvalue()


def zip_entry_hashes(package_bytes: bytes) -> dict[str, str]:
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as package:
        return {
            info.filename: hashlib.sha256(package.read(info.filename)).hexdigest()
            for info in package.infolist()
        }
