import base64
from collections.abc import Callable
from io import BytesIO
from pathlib import Path
from xml.etree import ElementTree as ET
from zipfile import ZIP_DEFLATED, ZipFile

import pytest
from pptx import Presentation
from pptx.util import Inches

from app.ai.deck_pptx_export import DeckPptxExportRequest, export_deck_pptx
from app.ai.pptx_ooxml_vector_importer import (
    import_pptx_ooxml_visual_tree,
    table_cell_fingerprint,
    text_vertical_align,
)


DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def test_generic_table_round_trip_preserves_tracks_cell_styles_and_locators(
    tmp_path: Path,
) -> None:
    deck = table_deck(
        rows=[
            [
                table_cell(
                    "Alpha",
                    fill="#123456",
                    text_color="#F8FAFC",
                    font_family="Arial",
                    font_size=30,
                    font_weight="bold",
                    align="center",
                    vertical_align="bottom",
                    border_color="#DC2626",
                    border_width=3,
                ),
                table_cell("Beta", fill="#E2E8F0"),
            ],
            [
                table_cell("Gamma", fill="#DCFCE7"),
                table_cell("Delta", fill="#FEF3C7"),
            ],
        ],
        column_widths=[320, 580],
        row_heights=[90, 150],
    )

    response = export_deck_pptx(DeckPptxExportRequest(deck=deck))
    assert response.warnings == []
    pptx_path = tmp_path / "table-round-trip.pptx"
    pptx_path.write_bytes(base64.b64decode(response.content_base64))

    imported = import_pptx_ooxml_visual_tree(pptx_path, "file_table")
    table = next(
        element
        for element in imported.blueprint["slides"][0]["elements"]
        if element["type"] == "table"
    )
    props = table["props"]
    assert props["columnWidths"] == pytest.approx([320, 580], abs=1)
    assert props["rowHeights"] == pytest.approx([90, 150], abs=1)
    assert props["rows"][0][1]["verticalAlign"] == "middle"
    assert props["rows"][0][0] == {
        "text": "Alpha",
        "fill": "#123456",
        "textColor": "#F8FAFC",
        "fontFamily": "Arial",
        "fontSize": 30,
        "fontWeight": "bold",
        "align": "center",
        "verticalAlign": "bottom",
        "borderColor": "#DC2626",
        "borderWidth": pytest.approx(3, abs=0.1),
        "colSpan": 1,
        "rowSpan": 1,
    }

    source = next(
        item
        for item in imported.template_blueprint["slides"][0]["elementSources"]
        if item["elementId"] == table["elementId"]
    )
    assert source["tableCellLocators"] == [
        {
            "rowIndex": row_index,
            "columnIndex": column_index,
            "fingerprint": source["tableCellLocators"][row_index * 2 + column_index][
                "fingerprint"
            ],
        }
        for row_index in range(2)
        for column_index in range(2)
    ]
    assert all(
        len(locator["fingerprint"]) == 64
        and locator["fingerprint"] == locator["fingerprint"].lower()
        for locator in source["tableCellLocators"]
    )


def test_generic_table_export_preserves_merged_cell_geometry_and_anchor_text(
    tmp_path: Path,
) -> None:
    rows = [
        [table_cell("Anchor", col_span=2, row_span=2), table_cell("Hidden B")],
        [table_cell("Hidden C"), table_cell("Hidden D")],
    ]
    response = export_deck_pptx(
        DeckPptxExportRequest(
            deck=table_deck(
                rows=rows,
                column_widths=[450, 450],
                row_heights=[120, 120],
            )
        )
    )

    assert response.warnings == []
    pptx_path = tmp_path / "merged-table-export.pptx"
    pptx_path.write_bytes(base64.b64decode(response.content_base64))
    presentation = Presentation(pptx_path)
    table = next(
        shape.table for shape in presentation.slides[0].shapes if shape.has_table
    )
    anchor = table.cell(0, 0)
    assert anchor.is_merge_origin
    assert anchor.span_width == 2
    assert anchor.span_height == 2
    assert anchor.text == "Anchor"

    imported = import_pptx_ooxml_visual_tree(pptx_path, "file_merged_export")
    table_element = next(
        element
        for element in imported.blueprint["slides"][0]["elements"]
        if element["type"] == "table"
    )
    assert table_element["props"]["rows"][0][0]["colSpan"] == 2
    assert table_element["props"]["rows"][0][0]["rowSpan"] == 2


def test_table_cell_fingerprint_excludes_text(tmp_path: Path) -> None:
    fingerprints: list[list[str]] = []
    for file_index, text in enumerate(("Before", "After"), start=1):
        response = export_deck_pptx(
            DeckPptxExportRequest(
                deck=table_deck(
                    rows=[[table_cell(text), table_cell("Same")]],
                    column_widths=[400, 500],
                    row_heights=[120],
                )
            )
        )
        pptx_path = tmp_path / f"fingerprint-{file_index}.pptx"
        pptx_path.write_bytes(base64.b64decode(response.content_base64))
        imported = import_pptx_ooxml_visual_tree(pptx_path, f"file_{file_index}")
        table_source = next(
            source
            for source in imported.template_blueprint["slides"][0]["elementSources"]
            if source.get("sourceType") == "table"
        )
        fingerprints.append(
            [locator["fingerprint"] for locator in table_source["tableCellLocators"]]
        )

    assert fingerprints[0] == fingerprints[1]


def test_table_cell_fingerprint_ignores_drawingml_text_whitespace(
    tmp_path: Path,
) -> None:
    fingerprints: list[str] = []
    for file_index, text in enumerate(("Alpha", " Alpha "), start=1):
        response = export_deck_pptx(
            DeckPptxExportRequest(
                deck=table_deck(
                    rows=[[table_cell(text)]],
                    column_widths=[900],
                    row_heights=[120],
                )
            )
        )
        pptx_path = tmp_path / f"whitespace-{file_index}.pptx"
        pptx_path.write_bytes(base64.b64decode(response.content_base64))
        imported = import_pptx_ooxml_visual_tree(pptx_path, f"file_{file_index}")
        table_source = next(
            source
            for source in imported.template_blueprint["slides"][0]["elementSources"]
            if source.get("sourceType") == "table"
        )
        fingerprints.append(table_source["tableCellLocators"][0]["fingerprint"])

    assert fingerprints[0] == fingerprints[1]


def test_table_cell_fingerprint_preserves_non_drawingml_text() -> None:
    first = ET.fromstring(
        f'<a:tc xmlns:a="{DRAWING_NS}" xmlns:x="urn:example">'
        "<a:txBody><a:bodyPr/><a:p><a:r><a:t>Cell</a:t></a:r></a:p></a:txBody>"
        "<a:tcPr><x:t>first</x:t></a:tcPr>"
        "</a:tc>"
    )
    second = ET.fromstring(
        f'<a:tc xmlns:a="{DRAWING_NS}" xmlns:x="urn:example">'
        "<a:txBody><a:bodyPr/><a:p><a:r><a:t>Cell</a:t></a:r></a:p></a:txBody>"
        "<a:tcPr><x:t>second</x:t></a:tcPr>"
        "</a:tc>"
    )

    assert table_cell_fingerprint(first) != table_cell_fingerprint(second)


def test_table_cell_fingerprint_ignores_drawingml_xml_space() -> None:
    without_space = ET.fromstring(
        f'<a:tc xmlns:a="{DRAWING_NS}">'
        "<a:txBody><a:bodyPr/><a:p><a:r><a:t>Cell</a:t></a:r></a:p></a:txBody>"
        "<a:tcPr/>"
        "</a:tc>"
    )
    with_space = ET.fromstring(
        f'<a:tc xmlns:a="{DRAWING_NS}">'
        '<a:txBody><a:bodyPr/><a:p><a:r><a:t xml:space="preserve"> Cell </a:t>'
        "</a:r></a:p></a:txBody><a:tcPr/>"
        "</a:tc>"
    )

    assert table_cell_fingerprint(without_space) == table_cell_fingerprint(with_space)


def test_generic_table_normalizes_tracks_to_preserve_element_frame(
    tmp_path: Path,
) -> None:
    deck = table_deck(
        rows=[[table_cell("A"), table_cell("B")], [table_cell("C"), table_cell("D")]],
        column_widths=[100, 200],
        row_heights=[20, 40],
        frame_width=900,
        frame_height=240,
    )

    response = export_deck_pptx(DeckPptxExportRequest(deck=deck))
    assert response.warnings == []
    pptx_path = tmp_path / "normalized-tracks.pptx"
    pptx_path.write_bytes(base64.b64decode(response.content_base64))
    imported = import_pptx_ooxml_visual_tree(pptx_path, "file_normalized")
    table = next(
        element
        for element in imported.blueprint["slides"][0]["elements"]
        if element["type"] == "table"
    )

    assert table["width"] == pytest.approx(900, abs=1)
    assert table["height"] == pytest.approx(240, abs=1)
    assert table["props"]["columnWidths"] == pytest.approx([300, 600], abs=1)
    assert table["props"]["rowHeights"] == pytest.approx([80, 160], abs=1)


def test_generic_transparent_table_cell_round_trips_as_no_fill(
    tmp_path: Path,
) -> None:
    response = export_deck_pptx(
        DeckPptxExportRequest(
            deck=table_deck(
                rows=[[table_cell("Transparent", fill="transparent")]],
                column_widths=[900],
                row_heights=[120],
            )
        )
    )
    pptx_path = tmp_path / "transparent-cell.pptx"
    pptx_path.write_bytes(base64.b64decode(response.content_base64))

    imported = import_pptx_ooxml_visual_tree(pptx_path, "file_transparent")
    table = next(
        element
        for element in imported.blueprint["slides"][0]["elements"]
        if element["type"] == "table"
    )

    assert table["props"]["rows"][0][0]["fill"] == "transparent"


def test_generic_empty_table_cell_round_trips_run_style(tmp_path: Path) -> None:
    response = export_deck_pptx(
        DeckPptxExportRequest(
            deck=table_deck(
                rows=[[
                    table_cell(
                        "",
                        text_color="#7C3AED",
                        font_family="Arial",
                        font_size=30,
                        font_weight="bold",
                        align="right",
                        vertical_align="bottom",
                    )
                ]],
                column_widths=[900],
                row_heights=[120],
            )
        )
    )
    pptx_path = tmp_path / "empty-cell-style.pptx"
    pptx_path.write_bytes(base64.b64decode(response.content_base64))

    imported = import_pptx_ooxml_visual_tree(pptx_path, "file_empty_cell")
    table = next(
        element
        for element in imported.blueprint["slides"][0]["elements"]
        if element["type"] == "table"
    )
    cell = table["props"]["rows"][0][0]

    assert cell["text"] == ""
    assert cell["textColor"] == "#7C3AED"
    assert cell["fontFamily"] == "Arial"
    assert cell["fontSize"] == 30
    assert cell["fontWeight"] == "bold"
    assert cell["align"] == "right"
    assert cell["verticalAlign"] == "bottom"


def test_imported_empty_table_cell_reads_paragraph_default_run_style(
    tmp_path: Path,
) -> None:
    response = export_deck_pptx(
        DeckPptxExportRequest(
            deck=table_deck(
                rows=[[table_cell("")]],
                column_widths=[900],
                row_heights=[120],
            )
        )
    )
    source_path = tmp_path / "empty-cell-default-source.pptx"
    styled_path = tmp_path / "empty-cell-default-style.pptx"
    source_path.write_bytes(base64.b64decode(response.content_base64))
    rewrite_slide_xml(source_path, styled_path, set_empty_cell_default_run_style)

    imported = import_pptx_ooxml_visual_tree(styled_path, "file_empty_default")
    table = next(
        element
        for element in imported.blueprint["slides"][0]["elements"]
        if element["type"] == "table"
    )
    cell = table["props"]["rows"][0][0]

    assert cell["text"] == ""
    assert cell["textColor"] == "#7C3AED"
    assert cell["fontFamily"] == "Arial"
    assert cell["fontSize"] == 36
    assert cell["fontWeight"] == "bold"


def test_generic_table_extreme_tracks_remain_positive_and_locatable(
    tmp_path: Path,
) -> None:
    response = export_deck_pptx(
        DeckPptxExportRequest(
            deck=table_deck(
                rows=[[table_cell("A"), table_cell("B"), table_cell("C")]],
                column_widths=[1_000_000_000_000, 1, 1],
                row_heights=[120],
                frame_width=900,
            )
        )
    )
    assert response.warnings == []
    pptx_path = tmp_path / "extreme-tracks.pptx"
    package_bytes = base64.b64decode(response.content_base64)
    pptx_path.write_bytes(package_bytes)

    with ZipFile(BytesIO(package_bytes)) as archive:
        root = ET.fromstring(archive.read("ppt/slides/slide1.xml"))
    widths = [
        int(column.attrib["w"])
        for column in root.findall(
            f".//{{{DRAWING_NS}}}tblGrid/{{{DRAWING_NS}}}gridCol"
        )
    ]
    assert len(widths) == 3
    assert all(width > 0 for width in widths)
    assert sum(widths) == 5_715_000

    imported = import_pptx_ooxml_visual_tree(pptx_path, "file_extreme_tracks")
    source = next(
        item
        for item in imported.template_blueprint["slides"][0]["elementSources"]
        if item.get("sourceType") == "table"
    )
    assert len(source["tableCellLocators"]) == 3


def test_generic_table_cell_lines_precede_fill_in_tcpr() -> None:
    response = export_deck_pptx(
        DeckPptxExportRequest(
            deck=table_deck(
                rows=[[table_cell("Ordered", fill="#F8FAFC")]],
                column_widths=[900],
                row_heights=[120],
            )
        )
    )

    with ZipFile(BytesIO(base64.b64decode(response.content_base64))) as archive:
        root = ET.fromstring(archive.read("ppt/slides/slide1.xml"))
    tc_pr = root.find(f".//{{{DRAWING_NS}}}tcPr")
    assert tc_pr is not None
    child_names = [child.tag.rsplit("}", maxsplit=1)[-1] for child in tc_pr]

    assert child_names[:5] == ["lnL", "lnR", "lnT", "lnB", "solidFill"]


def test_imported_merged_table_reads_actual_tc_span_and_omits_locators(
    tmp_path: Path,
) -> None:
    pptx_path = tmp_path / "merged-table.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = slide.shapes.add_table(
        2, 2, Inches(1), Inches(1), Inches(4), Inches(2)
    ).table
    table.cell(0, 0).text = "Merged"
    table.cell(0, 0).merge(table.cell(0, 1))
    presentation.save(pptx_path)

    imported = import_pptx_ooxml_visual_tree(pptx_path, "file_merged")
    table_element = next(
        element
        for element in imported.blueprint["slides"][0]["elements"]
        if element["type"] == "table"
    )
    source = next(
        item
        for item in imported.template_blueprint["slides"][0]["elementSources"]
        if item["elementId"] == table_element["elementId"]
    )

    assert table_element["props"]["rows"][0][0]["colSpan"] == 2
    assert "tableCellLocators" not in source
    assert any(
        warning.startswith(
            "PPTX_TABLE_STRUCTURE_UNSUPPORTED: slide=1; shape="
        )
        and warning.endswith("reason=merged-cell")
        for warning in imported.warnings
    )


def test_imported_jagged_and_track_mismatched_tables_omit_locators(
    tmp_path: Path,
) -> None:
    for reason, mutate in (
        ("jagged-grid", remove_last_cell),
        ("column-track-mismatch", remove_last_grid_column),
    ):
        source_path = tmp_path / f"source-{reason}.pptx"
        malformed_path = tmp_path / f"malformed-{reason}.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_table(
            2, 2, Inches(1), Inches(1), Inches(4), Inches(2)
        )
        presentation.save(source_path)
        rewrite_slide_xml(source_path, malformed_path, mutate)

        imported = import_pptx_ooxml_visual_tree(
            malformed_path, f"file_{reason.replace('-', '_')}"
        )
        table_source = next(
            source
            for source in imported.template_blueprint["slides"][0]["elementSources"]
            if source.get("sourceType") == "table"
        )
        assert "tableCellLocators" not in table_source
        assert any(
            warning.startswith("PPTX_TABLE_") and f"reason={reason}" in warning
            for warning in imported.warnings
        )


def test_generic_export_rejects_jagged_invalid_span_and_track_mismatched_tables() -> None:
    decks = (
        (
            table_deck(
                rows=[],
                column_widths=[],
                row_heights=[],
                frame_width=900,
                frame_height=120,
            ),
            "PPTX_TABLE_STRUCTURE_UNSUPPORTED: element=el_table; reason=empty-grid",
        ),
        (
            table_deck(
                rows=[[table_cell("A"), table_cell("B")], [table_cell("C")]],
                column_widths=[450, 450],
                row_heights=[120, 120],
            ),
            "PPTX_TABLE_STRUCTURE_UNSUPPORTED: element=el_table; reason=jagged-grid",
        ),
        (
            table_deck(
                rows=[[table_cell("A", col_span=3), table_cell("B")]],
                column_widths=[450, 450],
                row_heights=[120],
            ),
            "PPTX_TABLE_STRUCTURE_UNSUPPORTED: element=el_table; reason=invalid-cell-span",
        ),
        (
            table_deck(
                rows=[[table_cell("A"), table_cell("B")]],
                column_widths=[900],
                row_heights=[120],
            ),
            "PPTX_TABLE_TRACK_MISMATCH: element=el_table; track=columnWidths; expected=2; actual=1",
        ),
    )

    for deck, expected_warning in decks:
        response = export_deck_pptx(DeckPptxExportRequest(deck=deck))
        assert response.warnings == [expected_warning]
        presentation = Presentation(BytesIO(base64.b64decode(response.content_base64)))
        assert not any(shape.has_table for shape in presentation.slides[0].shapes)


def test_generic_export_reports_unsupported_table_cell_style() -> None:
    deck = table_deck(
        rows=[[table_cell("Gradient")]],
        column_widths=[900],
        row_heights=[120],
    )
    cell = deck["slides"][0]["elements"][0]["props"]["rows"][0][0]  # type: ignore[index]
    cell["fill"] = {  # type: ignore[index]
        "type": "linear-gradient",
        "angle": 0,
        "stops": [
            {"offset": 0, "color": "#FFFFFF"},
            {"offset": 1, "color": "#000000"},
        ],
    }

    response = export_deck_pptx(DeckPptxExportRequest(deck=deck))

    assert response.warnings == [
        "PPTX_TABLE_STYLE_UNSUPPORTED: element=el_table; row=0; column=0; property=fill"
    ]
    presentation = Presentation(BytesIO(base64.b64decode(response.content_base64)))
    assert not any(shape.has_table for shape in presentation.slides[0].shapes)


def test_text_vertical_align_maps_drawingml_center_anchor() -> None:
    body = ET.fromstring(
        f'<a:txBody xmlns:a="{DRAWING_NS}"><a:bodyPr anchor="ctr"/></a:txBody>'
    )

    assert text_vertical_align(body) == "middle"


def table_deck(
    *,
    rows: list[list[dict[str, object]]],
    column_widths: list[int],
    row_heights: list[int],
    frame_width: int | None = None,
    frame_height: int | None = None,
) -> dict[str, object]:
    return {
        "canvas": {"width": 1920, "height": 1080},
        "theme": {
            "backgroundColor": "#FFFFFF",
            "textColor": "#111827",
            "fontFamily": "Aptos",
        },
        "slides": [
            {
                "order": 1,
                "style": {"backgroundColor": "#FFFFFF"},
                "elements": [
                    {
                        "elementId": "el_table",
                        "type": "table",
                        "x": 100,
                        "y": 100,
                        "width": frame_width or sum(column_widths),
                        "height": frame_height or sum(row_heights),
                        "zIndex": 1,
                        "visible": True,
                        "props": {
                            "rows": rows,
                            "columnWidths": column_widths,
                            "rowHeights": row_heights,
                            "borderColor": "#CBD5E1",
                            "borderWidth": 1,
                        },
                    }
                ],
            }
        ],
    }


def table_cell(
    text: str,
    *,
    fill: str = "#FFFFFF",
    text_color: str = "#111827",
    font_family: str = "Aptos",
    font_size: int = 24,
    font_weight: str = "normal",
    align: str = "left",
    vertical_align: str = "middle",
    border_color: str = "#CBD5E1",
    border_width: int = 1,
    col_span: int = 1,
    row_span: int = 1,
) -> dict[str, object]:
    return {
        "text": text,
        "fill": fill,
        "textColor": text_color,
        "fontFamily": font_family,
        "fontSize": font_size,
        "fontWeight": font_weight,
        "align": align,
        "verticalAlign": vertical_align,
        "borderColor": border_color,
        "borderWidth": border_width,
        "colSpan": col_span,
        "rowSpan": row_span,
    }


def rewrite_slide_xml(
    source_path: Path,
    target_path: Path,
    mutate: Callable[[ET.Element], None],
) -> None:
    with ZipFile(source_path, "r") as source, ZipFile(
        target_path, "w", ZIP_DEFLATED
    ) as target:
        for name in source.namelist():
            payload = source.read(name)
            if name == "ppt/slides/slide1.xml":
                root = ET.fromstring(payload)
                mutate(root)
                payload = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            target.writestr(name, payload)


def remove_last_cell(root: ET.Element) -> None:
    rows = root.findall(f".//{{{DRAWING_NS}}}tbl/{{{DRAWING_NS}}}tr")
    cells = rows[-1].findall(f"{{{DRAWING_NS}}}tc")
    rows[-1].remove(cells[-1])


def remove_last_grid_column(root: ET.Element) -> None:
    grid = root.find(f".//{{{DRAWING_NS}}}tbl/{{{DRAWING_NS}}}tblGrid")
    assert grid is not None
    columns = grid.findall(f"{{{DRAWING_NS}}}gridCol")
    grid.remove(columns[-1])


def set_empty_cell_default_run_style(root: ET.Element) -> None:
    paragraph = root.find(
        f".//{{{DRAWING_NS}}}tc/{{{DRAWING_NS}}}txBody/{{{DRAWING_NS}}}p"
    )
    assert paragraph is not None
    for child in list(paragraph):
        paragraph.remove(child)
    paragraph_properties = ET.SubElement(
        paragraph,
        f"{{{DRAWING_NS}}}pPr",
    )
    default_properties = ET.SubElement(
        paragraph_properties,
        f"{{{DRAWING_NS}}}defRPr",
        {"b": "1", "sz": "1800"},
    )
    solid_fill = ET.SubElement(
        default_properties,
        f"{{{DRAWING_NS}}}solidFill",
    )
    ET.SubElement(
        solid_fill,
        f"{{{DRAWING_NS}}}srgbClr",
        {"val": "7C3AED"},
    )
    ET.SubElement(
        default_properties,
        f"{{{DRAWING_NS}}}latin",
        {"typeface": "Arial"},
    )
