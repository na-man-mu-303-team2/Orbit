import base64
from io import BytesIO

from pptx import Presentation

from app.ai.deck_pptx_export import DeckPptxExportRequest, export_deck_pptx


def test_static_activity_projection_exports_visible_copy_without_private_runtime_data() -> (
    None
):
    deck = {
        "deckId": "deck_activity_1",
        "projectId": "project_1",
        "title": "Activity export",
        "canvas": {"width": 1920, "height": 1080},
        "theme": {"textColor": "#111827", "backgroundColor": "#FFFFFF"},
        "slides": [
            projected_slide(
                "activity",
                1,
                [
                    "발표 만족도",
                    "발표가 유익했나요?",
                    "실시간 참여는 발표 중 제공됩니다.",
                ],
            ),
            projected_slide(
                "activity-results",
                2,
                ["발표 만족도", "7명 응답", "4.5 / 5", "승인된 의견입니다."],
            ),
        ],
    }

    response = export_deck_pptx(DeckPptxExportRequest(deck=deck))
    presentation = Presentation(BytesIO(base64.b64decode(response.content_base64)))
    exported_text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )

    assert "실시간 참여는 발표 중 제공됩니다." in exported_text
    assert "승인된 의견입니다." in exported_text
    assert "QR_PRIVATE_SENTINEL" not in exported_text
    assert "RAW_RESPONSE_SENTINEL" not in exported_text
    assert all(
        slide.notes_slide.notes_text_frame.text == "" for slide in presentation.slides
    )


def projected_slide(kind: str, order: int, texts: list[str]) -> dict:
    return {
        "kind": kind,
        "slideId": f"slide_{order}",
        "order": order,
        "speakerNotes": "",
        "style": {"backgroundColor": "#FFFFFF"},
        "elements": [
            {
                "elementId": f"element_{order}_{index}",
                "type": "text",
                "x": 120,
                "y": 120 + index * 150,
                "width": 1600,
                "height": 120,
                "zIndex": index,
                "visible": True,
                "props": {
                    "text": text,
                    "fontSize": 36,
                    "fontWeight": "normal",
                    "color": "#111827",
                    "lineHeight": 1.2,
                },
            }
            for index, text in enumerate(texts)
        ],
    }
