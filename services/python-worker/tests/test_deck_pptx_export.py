import base64
import math
from io import BytesIO
from typing import Any
from xml.etree import ElementTree as ET
from zipfile import ZipFile

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.ai.deck_pptx_export import DeckPptxExportRequest, export_deck_pptx
from app.ai.pptx_motion import parse_main_sequence


DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
ONE_PIXEL_PNG_DATA_URL = (
    "data:image/png;base64,"
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwC"
    "AAAAC0lEQVR42mNk+A8AAQUBAScY42YAAAAASUVORK5CYII="
)


def test_export_deck_pptx_serializes_and_reopens_all_image_crop_edges() -> None:
    crop = {"left": 0.2, "top": 0.1, "right": 0.15, "bottom": 0.05}

    binary = export_binary(image_deck(crop=crop))

    with ZipFile(BytesIO(binary)) as archive:
        slide_root = ET.fromstring(archive.read("ppt/slides/slide1.xml"))
    src_rect = slide_root.find(f".//{{{DRAWING_NS}}}srcRect")
    assert src_rect is not None
    assert src_rect.attrib == {
        "l": "20000",
        "t": "10000",
        "r": "15000",
        "b": "5000",
    }

    reopened = Presentation(BytesIO(binary))
    picture = next(
        shape
        for shape in reopened.slides[0].shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    )
    assert picture.crop_left == pytest.approx(crop["left"])
    assert picture.crop_top == pytest.approx(crop["top"])
    assert picture.crop_right == pytest.approx(crop["right"])
    assert picture.crop_bottom == pytest.approx(crop["bottom"])


def test_export_deck_pptx_clamps_rounded_crop_pairs_below_full_image() -> None:
    crop = {
        "left": 0.5000049,
        "top": 0.5000049,
        "right": 0.499995,
        "bottom": 0.499995,
    }

    binary = export_binary(image_deck(crop=crop))

    with ZipFile(BytesIO(binary)) as archive:
        slide_root = ET.fromstring(archive.read("ppt/slides/slide1.xml"))
    src_rect = slide_root.find(f".//{{{DRAWING_NS}}}srcRect")
    assert src_rect is not None
    assert src_rect.attrib == {
        "l": "50000",
        "t": "50000",
        "r": "49999",
        "b": "49999",
    }


@pytest.mark.parametrize("include_null", [False, True], ids=["omitted", "null"])
def test_export_deck_pptx_keeps_no_crop_images_without_src_rect(
    include_null: bool,
) -> None:
    deck = image_deck(crop=None) if include_null else image_deck()

    binary = export_binary(deck)

    with ZipFile(BytesIO(binary)) as archive:
        slide_root = ET.fromstring(archive.read("ppt/slides/slide1.xml"))
    assert slide_root.find(f".//{{{DRAWING_NS}}}srcRect") is None


@pytest.mark.parametrize(
    "crop",
    [
        {"left": math.nan},
        {"top": math.inf},
        {"right": -0.01},
        {"bottom": 1.01},
        {"left": True},
        {"left": "0.1"},
        {"left": 0.5, "right": 0.5},
        {"top": 0.6, "bottom": 0.4},
        "invalid",
    ],
)
def test_export_deck_pptx_rejects_crop_outside_shared_contract(
    crop: Any,
) -> None:
    with pytest.raises(ValueError, match="image crop"):
        export_deck_pptx(DeckPptxExportRequest(deck=image_deck(crop=crop)))


def test_export_deck_pptx_clips_text_to_its_fixed_frame() -> None:
    deck = image_deck()
    deck["slides"][0]["elements"] = [
        {
            "elementId": "el_overflow_text",
            "type": "text",
            "x": 120,
            "y": 180,
            "width": 320,
            "height": 80,
            "zIndex": 1,
            "visible": True,
            "props": {
                "text": "한글 English Supercalifragilisticexpialidocious",
                "fontSize": 28,
                "fontWeight": "normal",
                "align": "left",
                "verticalAlign": "top",
                "lineHeight": 1.2,
            },
        }
    ]

    binary = export_binary(deck)

    with ZipFile(BytesIO(binary)) as archive:
        slide_root = ET.fromstring(archive.read("ppt/slides/slide1.xml"))
    body_pr = slide_root.find(f".//{{{DRAWING_NS}}}bodyPr")
    assert body_pr is not None
    assert body_pr.get("horzOverflow") == "clip"
    assert body_pr.get("vertOverflow") == "clip"
    assert body_pr.get("wrap") == "square"


def test_export_deck_pptx_preserves_flattened_group_motion_semantics() -> None:
    response = export_deck_pptx(
        DeckPptxExportRequest(deck=group_motion_deck())
    )
    binary = base64.b64decode(response.content_base64)
    assert response.motion_diagnostics == [
        {
            "code": "PPTX_MOTION_TARGET_FLATTENED",
            "slideIndex": 1,
            "elementId": "el_group",
            "count": 2,
        }
    ]

    reopened = Presentation(BytesIO(binary))
    shape_targets = {
        str(shape.shape_id): "el_group"
        for shape in reopened.slides[0].shapes
    }
    with ZipFile(BytesIO(binary)) as archive:
        slide_root = ET.fromstring(archive.read("ppt/slides/slide1.xml"))
    parsed, coverage, diagnostics = parse_main_sequence(
        slide_root,
        slide_index=1,
        shape_targets=shape_targets,
    )

    assert coverage == "complete"
    assert diagnostics == []
    assert [
        (
            animation["elementId"],
            animation["type"],
            animation["startMode"],
            animation["durationMs"],
            animation["delayMs"],
        )
        for animation in parsed
    ] == [
        ("el_group", "fade-in", "on-click", 300, 0),
        ("el_group", "fade-in", "with-previous", 300, 0),
    ]


def export_binary(deck: dict[str, Any]) -> bytes:
    response = export_deck_pptx(DeckPptxExportRequest(deck=deck))
    return base64.b64decode(response.content_base64)


def image_deck(*, crop: Any = ...) -> dict[str, Any]:
    props: dict[str, Any] = {
        "src": ONE_PIXEL_PNG_DATA_URL,
        "alt": "Crop export fixture",
        "fit": "contain",
        "focusX": 0.5,
        "focusY": 0.5,
    }
    if crop is not ...:
        props["crop"] = crop
    return {
        "canvas": {"width": 1920, "height": 1080},
        "theme": {"backgroundColor": "#FFFFFF"},
        "slides": [
            {
                "order": 1,
                "style": {"backgroundColor": "#FFFFFF"},
                "elements": [
                    {
                        "elementId": "el_crop_image",
                        "type": "image",
                        "x": 120,
                        "y": 180,
                        "width": 640,
                        "height": 360,
                        "zIndex": 1,
                        "visible": True,
                        "props": props,
                    }
                ],
            }
        ],
    }


def group_motion_deck() -> dict[str, Any]:
    return {
        "canvas": {"width": 1920, "height": 1080},
        "theme": {"backgroundColor": "#FFFFFF"},
        "slides": [
            {
                "slideId": "slide_group",
                "order": 1,
                "style": {"backgroundColor": "#FFFFFF"},
                "elements": [
                    {
                        "elementId": "el_group_background",
                        "type": "rect",
                        "role": "decoration",
                        "x": 240,
                        "y": 240,
                        "width": 720,
                        "height": 400,
                        "zIndex": 1,
                        "visible": True,
                        "props": {"fill": "#DCEBFF"},
                    },
                    {
                        "elementId": "el_group_body",
                        "type": "text",
                        "role": "body",
                        "x": 300,
                        "y": 320,
                        "width": 600,
                        "height": 220,
                        "zIndex": 2,
                        "visible": True,
                        "props": {
                            "text": "그룹 전체 본문",
                            "fontSize": 28,
                        },
                    },
                    {
                        "elementId": "el_group",
                        "type": "group",
                        "role": "body",
                        "x": 240,
                        "y": 240,
                        "width": 720,
                        "height": 400,
                        "zIndex": 3,
                        "visible": True,
                        "props": {
                            "childElementIds": [
                                "el_group_background",
                                "el_group_body",
                            ]
                        },
                    },
                ],
                "animations": [
                    {
                        "animationId": "anim_group_reveal",
                        "elementId": "el_group",
                        "type": "fade-in",
                        "order": 1,
                        "startMode": "on-click",
                        "durationMs": 300,
                        "delayMs": 0,
                        "easing": "ease-out",
                    }
                ],
            }
        ],
    }
