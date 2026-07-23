import base64
import json
from collections import Counter
from io import BytesIO
from pathlib import Path
from types import SimpleNamespace
from typing import Any

import pytest
from pptx import Presentation

import app.ai.deck_generation.design_planning as design_planning_module
from app.ai.composition_library import (
    COMPOSITION_SPECS,
    compile_composition,
    normalize_design_program,
)
from app.ai.design_program import (
    ArtDirectorContext,
    DeckDesignProgram,
    DesignProgramError,
    art_director_prompt,
    create_design_program,
    design_program_response_format,
)
from app.ai.deck_pptx_export import DeckPptxExportRequest, export_deck_pptx
from app.ai.deck_generation.design_planning import (
    apply_program_v2_design_tokens,
    eligible_cover_compositions,
    plan_design,
    program_v2_slide_summary,
)
from app.ai.deck_generation.layout_compiler import compile_layout
from app.ai.deck_generation.models import (
    GenerateDeckRequest,
    CoverContent,
    GeneratedContentItem,
    MediaIntent,
    SlidePlan,
    SourceRecord,
    VisualIntent,
)
from app.ai.deck_generation.pipeline import (
    DeckGenerationOrchestrator,
    analyze_input,
)
from app.ai.deck_generation.quality import (
    is_expected_media_placeholder,
    validate_presentation,
)
from app.ai.deck_generation.source_grounding import (
    design_pack_source_ledgers,
    initial_source_records,
)
from app.ai.deck_generation.visual_requirements import (
    apply_visual_requirements,
    plan_visual_requirements,
    program_v2_visual_plan,
)


def context() -> ArtDirectorContext:
    return ArtDirectorContext(
        topic="Splatoon Raiders",
        presentationProfile="product-launch",
        brief={"presentationType": "제품 공개", "audience": "게임 팬"},
        designDirection="강한 잉크 색상과 명확한 시각적 중심",
        palette={"background": "#FFFFFF", "primary": "#6D28D9"},
        typography={"headingFont": "Pretendard", "bodyFont": "Pretendard"},
        forbiddenStyles=["gradient", "pastel"],
        mediaPolicy="hybrid",
        mediaBudget=4,
    )


def slides() -> list[dict[str, Any]]:
    return [
        {
            "title": "새로운 모험",
            "message": "스플래툰 레이더스가 새로운 경험을 연다",
            "contentItems": [{"text": "공식 공개 정보"}],
            "slideType": "cover",
            "mediaIntent": {
                "kind": "generate",
                "prompt": "ink island adventure",
                "alt": "게임 세계",
                "required": True,
            },
            "speakerNotes": "프롬프트에 포함되면 안 되는 전체 발표 메모",
            "sourceRecords": ["프롬프트에 포함되면 안 되는 연구 원문"],
        },
        {
            "title": "지금 확인하세요",
            "message": "공식 채널에서 다음 소식을 확인한다",
            "contentItems": [{"text": "공식 사이트"}],
            "slideType": "summary",
            "mediaIntent": {"kind": "none", "required": False},
        },
    ]


def valid_program() -> dict[str, Any]:
    return {
        "version": "program-v2",
        "visualConcept": "Energetic ink expedition",
        "paletteRoles": {
            "dominant": "#FFFFFF",
            "surface": "#F3F4F6",
            "text": "#111827",
            "focal": "#6D28D9",
            "secondary": "#22D3EE",
        },
        "typography": {
            "headingFont": "Pretendard",
            "bodyFont": "Pretendard",
            "typeScale": {"cover": 64, "title": 40, "body": 22, "caption": 14},
        },
        "backgroundSequence": ["image", "dark"],
        "imageStyle": "Official key art with clean crops",
        "surfaceStyle": "Flat ink color fields",
        "slides": [
            {
                "order": 1,
                "compositionId": "hero-full-bleed",
                "variant": "image",
                "backgroundMode": "image",
                "focalType": "hero-image",
                "assetRole": "atmosphere",
                "requiredAsset": True,
            },
            {
                "order": 2,
                "compositionId": "cta-closing",
                "variant": "dark",
                "backgroundMode": "dark",
                "focalType": "cta",
                "assetRole": "none",
                "requiredAsset": False,
            },
        ],
    }


def test_program_v2_typography_keeps_presentation_scale_floors() -> None:
    design_program = DeckDesignProgram.model_validate(valid_program())
    themed = apply_program_v2_design_tokens(
        design_program,
        {
            "backgroundColor": "#FFFFFF",
            "textColor": "#111827",
            "accentColor": "#6D28D9",
            "fontFamily": "Gmarket Sans",
            "palette": {"surface": "#F3F4F6", "secondary": "#22D3EE"},
            "typography": {
                "headingFontFamily": "Gmarket Sans",
                "bodyFontFamily": "Gmarket Sans",
                "titleSize": 40,
                "headingSize": 32,
                "bodySize": 22,
                "captionSize": 14,
            },
        },
    )

    assert themed.typography.type_scale == {
        "cover": 72,
        "title": 56,
        "body": 32,
        "caption": 24,
    }


def test_program_v2_palette_keeps_focal_and_secondary_roles_distinct() -> None:
    design_program = DeckDesignProgram.model_validate(valid_program())
    themed = apply_program_v2_design_tokens(
        design_program,
        {
            "backgroundColor": "#FFFFFF",
            "textColor": "#111827",
            "accentColor": "#2563EB",
            "palette": {
                "surface": "#FFFFFF",
                "primary": "#2563EB",
                "secondary": "#2563EB",
            },
            "typography": {},
        },
    )

    assert themed.palette_roles.focal == "#2563EB"
    assert themed.palette_roles.secondary == "#22D3EE"


def test_program_v2_compiles_empty_cover_and_closing_content_items() -> None:
    slide_summaries = slides()
    for slide in slide_summaries:
        slide["contentItems"] = []
    program = normalize_design_program(
        DeckDesignProgram.model_validate(valid_program()),
        slide_summaries,
        media_policy="minimal",
    )

    compiled = [
        compile_composition(direction, slide, program)
        for direction, slide in zip(
            program.slides,
            slide_summaries,
            strict=True,
        )
    ]

    assert [slide.composition_id for slide in program.slides] == [
        "cover-classic-corporate",
        "closing-centered-minimal",
    ]
    assert all(item.elements for item in compiled)


class FakeResponses:
    def __init__(self, payloads: list[dict[str, Any] | str]) -> None:
        self.payloads = payloads
        self.requests: list[dict[str, Any]] = []

    def create(self, **kwargs: Any) -> SimpleNamespace:
        self.requests.append(kwargs)
        payload = self.payloads[min(len(self.requests) - 1, len(self.payloads) - 1)]
        output_text = (
            payload
            if isinstance(payload, str)
            else json.dumps(payload, ensure_ascii=False)
        )
        return SimpleNamespace(output_text=output_text)


class PipelineFakeResponses:
    def __init__(self, payloads: dict[str, dict[str, Any]]) -> None:
        self.payloads = payloads
        self.requests: list[dict[str, Any]] = []

    def create(self, **kwargs: Any) -> SimpleNamespace:
        self.requests.append(kwargs)
        response_format_name = kwargs["text"]["format"]["name"]
        return SimpleNamespace(
            output_text=json.dumps(
                self.payloads[response_format_name],
                ensure_ascii=False,
            )
        )


def test_art_director_prompt_uses_only_compact_slide_summaries() -> None:
    prompt = art_director_prompt(context(), slides())

    assert "프롬프트에 포함되면 안 되는 전체 발표 메모" not in prompt
    assert "프롬프트에 포함되면 안 되는 연구 원문" not in prompt
    assert "hero-split" in prompt
    assert "mediaBudget" in prompt
    assert "강한 잉크 색상과 명확한 시각적 중심" in prompt
    assert "visualIntent" in prompt


def test_response_format_requires_exact_slide_count() -> None:
    schema = design_program_response_format(10)["format"]["schema"]

    assert schema["properties"]["slides"]["minItems"] == 10
    assert schema["properties"]["slides"]["maxItems"] == 10
    assert schema["properties"]["backgroundSequence"]["minItems"] == 10


def test_cover_eligibility_respects_media_and_verified_profile_requirements() -> None:
    minimal_input = analyze_input(
        GenerateDeckRequest(
            projectId="project_program_v2",
            topic="제품 출시 이벤트",
            design={"mediaPolicy": "minimal"},
        )
    )
    cover = SlidePlan(
        order=1,
        slide_type="cover",
        title="제품 출시 이벤트",
        message="새로운 경험을 공개합니다",
        speaker_notes="",
        keywords=[],
        evidence=[],
    )

    assert eligible_cover_compositions(minimal_input, cover) == [
        "cover-classic-corporate",
        "cover-modern-high-tech",
    ]

    provided_input = analyze_input(
        GenerateDeckRequest(
            projectId="project_program_v2",
            topic="연구 결과",
            design={"mediaPolicy": "provided-only"},
            officialAssetFileIds=["asset-profile"],
        )
    )
    research_cover = cover.model_copy(
        update={
            "cover_content": CoverContent(
                title="연구 결과",
                presenterName="김민지",
                profileImageAssetId="asset-profile",
            )
        }
    )

    assert "cover-research-author" in eligible_cover_compositions(
        provided_input,
        research_cover,
    )
    assert "cover-research-author" not in eligible_cover_compositions(
        provided_input,
        cover,
    )

    hybrid_input = analyze_input(
        GenerateDeckRequest(
            projectId="project_program_v2",
            topic="연구 결과",
            design={"mediaPolicy": "hybrid"},
            officialAssetFileIds=["asset-profile"],
        )
    )
    summary = program_v2_slide_summary(research_cover, hybrid_input)
    program_payload = valid_program()
    program_payload["backgroundSequence"] = ["light"]
    program_payload["slides"] = [
        {
            "order": 1,
            "compositionId": "cover-research-author",
            "variant": "light",
            "backgroundMode": "light",
            "focalType": "cover-image",
            "assetRole": "none",
            "requiredAsset": True,
        }
    ]
    normalized = normalize_design_program(
        DeckDesignProgram.model_validate(program_payload),
        [summary],
        media_policy="hybrid",
        preserve_slide_types=True,
    )

    assert summary["officialSourceAvailable"] is True
    assert normalized.slides[0].composition_id == "cover-research-author"
    assert normalized.slides[0].asset_role == "evidence"


def test_hybrid_official_asset_placeholder_is_expected_before_resolution() -> None:
    slide = {
        "aiNotes": {
            "visualPlan": {
                "imageNeeded": True,
                "imageSourcePolicy": "official-assets",
            }
        }
    }

    assert is_expected_media_placeholder(slide) is True


def test_program_v2_visual_plan_replaces_generic_media_prompt_with_slide_subject() -> None:
    raw_input = analyze_input(
        GenerateDeckRequest(
            projectId="project_program_v2",
            topic="Splatoon Raiders",
            design={
                "mediaPolicy": "hybrid",
                "constraints": {
                    "forbiddenStyles": ["gradient", "pastel"],
                },
            },
        )
    )
    slide = SlidePlan(
        order=2,
        slide_type="data",
        title="Nintendo Switch 2 전용 경험",
        message="공식 발표 사실을 확인한다",
        speaker_notes="공식 발표 내용을 설명합니다.",
        keywords=[],
        evidence=[],
        content_items=[
            GeneratedContentItem(contentItemId="item-1", text="전용 플랫폼")
        ],
        media_intent=MediaIntent(prompt="none", alt="공식 제품 이미지"),
        visual_intent=VisualIntent(mediaStyle="clean"),
    )
    design_program = DeckDesignProgram.model_validate(valid_program())
    direction = design_program.slides[0].model_copy(
        update={"asset_role": "evidence", "required_asset": True}
    )

    plan = program_v2_visual_plan(
        raw_input,
        slide,
        design_program,
        direction,
    )

    assert plan["imageSourcePolicy"] == "official-assets"
    assert "Splatoon Raiders" in plan["imagePrompt"]
    assert "Nintendo Switch 2 전용 경험" in plan["imagePrompt"]
    assert "official product evidence" in plan["imagePrompt"]
    assert "none" not in plan["imagePrompt"]
    assert "avoid gradient and pastel" in plan["imagePrompt"]


def test_provided_only_visual_plan_uses_worker_supported_official_asset_policy() -> None:
    raw_input = analyze_input(
        GenerateDeckRequest(
            projectId="project_program_v2",
            topic="연구 발표",
            design={"mediaPolicy": "provided-only"},
            officialAssetFileIds=["asset-profile"],
        )
    )
    slide = SlidePlan(
        order=1,
        slide_type="cover",
        title="연구 발표",
        message="검증된 프로필과 함께 소개합니다",
        speaker_notes="",
        keywords=[],
        evidence=[],
    )
    design_program = DeckDesignProgram.model_validate(valid_program())
    direction = design_program.slides[0].model_copy(
        update={"asset_role": "evidence", "required_asset": True}
    )

    plan = program_v2_visual_plan(raw_input, slide, design_program, direction)

    assert plan["imageSourcePolicy"] == "official-assets"


def test_program_v2_visual_plan_omits_icon_style_from_large_media_prompt() -> None:
    raw_input = analyze_input(
        GenerateDeckRequest(
            projectId="project_program_v2",
            topic="Splatoon Raiders",
            design={"mediaPolicy": "hybrid"},
        )
    )
    slide = SlidePlan(
        order=2,
        slide_type="solution",
        title="차별화된 협동 모험",
        message="협동 탐험 경험을 강조한다",
        speaker_notes="협동 탐험 경험을 설명합니다.",
        keywords=[],
        evidence=[],
        content_items=[
            GeneratedContentItem(contentItemId="item-1", text="협동 탐험")
        ],
        media_intent=MediaIntent(alt="협동 모험 장면"),
        visual_intent=VisualIntent(mediaStyle="icon"),
    )
    design_program = DeckDesignProgram.model_validate(valid_program())
    direction = design_program.slides[0].model_copy(
        update={"asset_role": "atmosphere", "required_asset": False}
    )

    plan = program_v2_visual_plan(raw_input, slide, design_program, direction)

    assert "atmospheric key visual" in plan["imagePrompt"]
    assert "icon" not in plan["imagePrompt"]


def test_program_v2_slide_summary_reports_official_source_availability() -> None:
    raw_input = analyze_input(
        GenerateDeckRequest(
            projectId="project_program_v2",
            topic="Splatoon Raiders",
        )
    )
    raw_input.source_records = [
        SourceRecord(
            sourceType="web",
            sourceId="web:official",
            url="https://example.com/official",
            title="Official announcement",
            content="Official product facts",
            authority="official",
        )
    ]
    slide = SlidePlan(
        order=1,
        slide_type="cover",
        title="Official reveal",
        message="The official reveal is available.",
        speaker_notes="Introduce the official reveal.",
        keywords=[],
        evidence=[],
        content_items=[
            GeneratedContentItem(contentItemId="item-1", text="Official reveal")
        ],
        source_refs=["web:official"],
    )

    summary = program_v2_slide_summary(slide, raw_input)

    assert summary["officialSourceAvailable"] is True


def test_program_v2_hybrid_cover_reserves_deck_official_source() -> None:
    raw_input = analyze_input(
        GenerateDeckRequest(
            projectId="project_program_v2",
            topic="Splatoon Raiders",
            design={"mediaPolicy": "hybrid"},
        )
    )
    raw_input.source_records = [
        SourceRecord(
            sourceType="web",
            sourceId="web:official",
            url="https://example.com/official",
            title="Official announcement",
            content="Official product facts",
            authority="official",
        ),
        SourceRecord(
            sourceType="web",
            sourceId="web:independent",
            url="https://example.com/news",
            title="Independent coverage",
            content="Independent product coverage",
            authority="independent",
        ),
    ]
    cover = SlidePlan(
        order=1,
        slide_type="cover",
        title="Official reveal",
        message="The official reveal is available.",
        speaker_notes="Introduce the official reveal.",
        keywords=[],
        evidence=[],
        source_refs=["web:independent"],
    )
    body = cover.model_copy(
        update={"order": 2, "slide_type": "context", "title": "Coverage"}
    )

    assert program_v2_slide_summary(cover, raw_input)[
        "officialSourceAvailable"
    ] is True
    assert program_v2_slide_summary(body, raw_input)[
        "officialSourceAvailable"
    ] is False


def test_program_v2_evidence_ledger_includes_deck_official_source() -> None:
    raw_input = analyze_input(
        GenerateDeckRequest(
            projectId="project_program_v2",
            topic="Splatoon Raiders",
            brief={"referencePolicy": "research-first"},
            design={"mediaPolicy": "hybrid"},
        )
    )
    raw_input.source_records = [
        SourceRecord(
            sourceType="web",
            sourceId="web:official",
            url="https://example.com/official",
            title="Official announcement",
            content="Official product facts",
            authority="official",
        ),
        SourceRecord(
            sourceType="web",
            sourceId="web:independent",
            url="https://example.com/news",
            title="Independent coverage",
            content="Independent product coverage",
            authority="independent",
        ),
    ]
    slide = SlidePlan(
        order=1,
        slide_type="cover",
        title="Official reveal",
        message="The official reveal is available.",
        speaker_notes="Introduce the official reveal.",
        keywords=[],
        evidence=[],
        content_items=[
            GeneratedContentItem(contentItemId="item-1", text="Official reveal")
        ],
        source_refs=["web:independent"],
    )

    ledgers = design_pack_source_ledgers(
        raw_input,
        slide,
        include_official_web=True,
    )

    assert ledgers[0]["sourceId"] == "web:official"
    assert ledgers[0]["authority"] == "official"
    assert {ledger["sourceId"] for ledger in ledgers} == {
        "web:official",
        "web:independent",
    }


def test_normalize_replaces_adjacent_comparison_silhouettes() -> None:
    plans = golden_slide_plans()
    plans[6].slide_type = "comparison"
    payload = golden_design_program()
    payload["slides"][6]["compositionId"] = "feature-comparison"
    program = DeckDesignProgram.model_validate(payload)

    normalized = normalize_design_program(
        program,
        [program_v2_slide_summary(plan) for plan in plans],
        media_policy="hybrid",
    )

    comparison_ids = {
        normalized.slides[5].composition_id,
        normalized.slides[6].composition_id,
    }
    silhouettes = [
        COMPOSITION_SPECS[direction.composition_id].silhouette
        for direction in normalized.slides
    ]

    assert comparison_ids == {"feature-comparison", "editorial-split"}
    assert all(left != right for left, right in zip(silhouettes, silhouettes[1:]))


def test_create_design_program_normalizes_background_sequence_without_retry() -> None:
    mismatched = {**valid_program(), "backgroundSequence": ["light", "light"]}
    responses = FakeResponses([mismatched, {**valid_program(), "slides": []}])
    client = SimpleNamespace(responses=responses)

    program = create_design_program(context(), slides(), client=client)

    assert program.background_sequence == ["image", "dark"]
    assert program.visual_concept == "Energetic ink expedition"
    assert "강한 잉크 색상과 명확한 시각적 중심" in program.image_style
    assert len(responses.requests) == 1


def test_create_design_program_retries_one_wrong_slide_count_response() -> None:
    invalid = valid_program()
    invalid["slides"] = invalid["slides"][:1]
    responses = FakeResponses([invalid, valid_program()])

    program = create_design_program(
        context(),
        slides(),
        client=SimpleNamespace(responses=responses),
    )

    assert program.background_sequence == ["image", "dark"]
    assert len(responses.requests) == 2


def test_create_design_program_fails_after_one_retry() -> None:
    invalid = {**valid_program(), "slides": []}
    responses = FakeResponses([invalid, invalid])

    with pytest.raises(DesignProgramError) as captured:
        create_design_program(
            context(),
            slides(),
            client=SimpleNamespace(responses=responses),
        )

    assert str(captured.value) == (
        "Art Director could not create a valid design plan. "
        "Please retry deck generation."
    )
    assert "validation error" not in str(captured.value)
    assert "input_value" not in str(captured.value)
    assert len(responses.requests) == 2


@pytest.mark.parametrize(
    "invalid",
    [
        "{not-json",
        {
            **valid_program(),
            "slides": [
                {**valid_program()["slides"][0], "backgroundMode": "transparent"},
                valid_program()["slides"][1],
            ],
        },
        {
            **valid_program(),
            "slides": [
                valid_program()["slides"][0],
                {**valid_program()["slides"][1], "order": 3},
            ],
        },
    ],
    ids=["json", "background-mode", "order"],
)
def test_create_design_program_keeps_irrecoverable_validation_errors(
    invalid: dict[str, Any] | str,
) -> None:
    responses = FakeResponses([invalid, invalid])

    with pytest.raises(DesignProgramError):
        create_design_program(
            context(),
            slides(),
            client=SimpleNamespace(responses=responses),
        )

    assert len(responses.requests) == 2


def test_program_v2_design_and_layout_stages_compile_canonical_backgrounds() -> None:
    mismatched = {**valid_program(), "backgroundSequence": ["light", "light"]}
    responses = FakeResponses([mismatched])
    orchestrator = DeckGenerationOrchestrator(
        GenerateDeckRequest(
            projectId="project_program_v2",
            topic="Splatoon Raiders",
            targetDurationMinutes=2,
            slideCountRange={"min": 2, "max": 2},
            design={"mediaPolicy": "hybrid"},
            visualPlanPolicy={"mediaPolicy": "hybrid"},
        ),
        client=SimpleNamespace(responses=responses),
    )
    raw_input = orchestrator.run_brief_agent()
    raw_input.source_records = initial_source_records(raw_input)
    slide_plans = [
        SlidePlan(
            order=1,
            slide_type="cover",
            title="새로운 모험",
            message="스플래툰 레이더스가 새로운 경험을 연다",
            speaker_notes="공식 공개 내용을 소개합니다.",
            keywords=["Splatoon Raiders"],
            evidence=[],
            content_items=[
                GeneratedContentItem(
                    contentItemId="item-1",
                    text="스플래툰 레이더스가 새로운 경험을 연다",
                )
            ],
            source_refs=["topic:brief"],
            media_intent=MediaIntent(
                kind="generate",
                prompt="ink island adventure",
                alt="게임 세계",
                required=True,
            ),
        ),
        SlidePlan(
            order=2,
            slide_type="summary",
            title="지금 확인하세요",
            message="공식 채널에서 다음 소식을 확인한다",
            speaker_notes="공식 채널 확인을 요청합니다.",
            keywords=["공식 채널"],
            evidence=[],
            content_items=[
                GeneratedContentItem(contentItemId="item-2", text="공식 사이트")
            ],
            source_refs=["topic:brief"],
        ),
    ]

    design_plan = plan_design(
        raw_input,
        slide_plans,
        client=orchestrator.client,
    )
    layout_result = compile_layout(
        raw_input,
        design_plan,
    )
    assert all(
        "visualPlan" not in slide["aiNotes"] for slide in layout_result.slides
    )
    expected_visual_plans = [
        program_v2_visual_plan(
            raw_input,
            slide_plan,
            design_plan.design_program,
            design_plan.design_program.slides[slide_plan.order - 1],
        )
        for slide_plan in design_plan.slide_plans
    ]
    requirements = plan_visual_requirements(
        raw_input,
        design_plan,
        layout_result,
    )
    assert [item.visual_plan for item in requirements.items] == expected_visual_plans
    visualized_slides = apply_visual_requirements(layout_result, requirements)
    assert all(
        "visualPlan" not in slide["aiNotes"] for slide in layout_result.slides
    )
    assert [
        slide["aiNotes"]["visualPlan"] for slide in visualized_slides
    ] == expected_visual_plans
    assert all(
        slide["aiNotes"]["visualPlan"] is not requirement.visual_plan
        for slide, requirement in zip(
            visualized_slides,
            requirements.items,
            strict=True,
        )
    )
    assert list(visualized_slides[0]["aiNotes"]) == [
        "emphasisPoints",
        "sourceEvidence",
        "visualPlan",
        "sourceLedger",
        "timingPlan",
        "compositionPlan",
    ]
    deck = orchestrator.build_deck(
        raw_input,
        type("Outline", (), {"title": "Splatoon Raiders"})(),
        design_plan,
        visualized_slides,
    )

    assert deck["metadata"]["designProgramSnapshot"]["version"] == "program-v2"
    assert design_plan.design_program.background_sequence == [
        direction.background_mode
        for direction in design_plan.design_program.slides
    ]
    assert deck["metadata"]["designProgramSnapshot"]["backgroundSequence"] == [
        slide["aiNotes"]["compositionPlan"]["backgroundMode"]
        for slide in deck["slides"]
    ]
    assert len(responses.requests) == 1
    assert deck["slides"][0]["aiNotes"]["compositionPlan"]["compositionId"] == (
        "cover-classic-corporate"
    )
    assert deck["slides"][0]["aiNotes"]["visualPlan"]["imageSourcePolicy"] == (
        "minimal"
    )
    assert sum(
        element.get("props", {}).get("text")
        == "스플래툰 레이더스가 새로운 경험을 연다"
        for element in deck["slides"][0]["elements"]
    ) == 1
    assert all(
        "_design_pack_" not in element["elementId"]
        for slide in deck["slides"]
        for element in slide["elements"]
    )
    assert all(
        element.get("role") != "background"
        for slide in deck["slides"]
        for element in slide["elements"]
    )
    assert all(slide["animations"] == [] for slide in deck["slides"])


def test_program_v2_golden_pipeline_contract() -> None:
    fixture_path = (
        Path(__file__).parent
        / "fixtures"
        / "splatoon_product_launch_golden_request.json"
    )
    request = GenerateDeckRequest.model_validate_json(fixture_path.read_text("utf-8"))
    request.reference_policy = "user-input-only"
    request.brief.reference_policy = "user-input-only"
    request.design.reference_policy = "user-input-only"
    responses = PipelineFakeResponses(
        {
            "design_pack_content_plan": golden_content_plan(),
            "deck_design_program": golden_design_program(),
        }
    )
    orchestrator = DeckGenerationOrchestrator(
        request,
        client=SimpleNamespace(responses=responses),
    )

    response = orchestrator.run()

    assert list(orchestrator.agent_outputs) == [
        "BriefAgent",
        "SourceGroundingAgent",
        "NarrativeAgent",
        "DesignDirectorAgent",
        "LayoutAgent",
        "ChartDataAgent",
        "MediaAgent",
        "QualityReviewerAgent",
        "RefinerAgent",
    ]
    assert [
        request_payload["text"]["format"]["name"]
        for request_payload in responses.requests
    ] == [
        "design_pack_content_plan",
        "design_pack_content_plan",
        "deck_design_program",
    ]

    deck = response.deck
    assert [slide["order"] for slide in deck["slides"]] == list(range(1, 11))
    assert [slide["slideId"] for slide in deck["slides"]] == [
        f"slide_{order}" for order in range(1, 11)
    ]
    actual_element_ids = [
        [element["elementId"] for element in slide["elements"]]
        for slide in deck["slides"]
    ]
    assert all(actual_element_ids)
    assert all(
        all(element_id.startswith(f"el_{order}_program_v2_") for element_id in ids)
        for order, ids in enumerate(actual_element_ids, start=1)
    )
    assert actual_element_ids[1] == golden_element_ids()[1]
    assert actual_element_ids[-1] == golden_element_ids()[-1]
    assert [
        slide["aiNotes"]["compositionPlan"]["compositionId"]
        for slide in deck["slides"]
    ] == [
        "cover-visual-impact",
        "agenda-numbered-list",
        "bento-focus",
        "editorial-split",
        "diagram-orbit",
        "editorial-split",
        "process-vertical-rail",
        "editorial-media-band",
        "diagram-hub",
        "closing-centered-minimal",
    ]
    title_element = next(
        element
        for element in deck["slides"][0]["elements"]
        if element["elementId"] == "el_1_program_v2_title"
    )
    assert {
        key: title_element[key]
        for key in ("elementId", "type", "role", "x", "y", "width", "height")
    } == {
        "elementId": "el_1_program_v2_title",
        "type": "text",
        "role": "title",
        "x": 120,
        "y": 250,
        "width": 720,
        "height": 350,
    }
    assert title_element["props"]["text"] == "미지의 군도로"
    assert all(slide["animations"] == [] for slide in deck["slides"])
    pptx_response = export_deck_pptx(DeckPptxExportRequest(deck=deck))
    presentation = Presentation(
        BytesIO(base64.b64decode(pptx_response.content_base64))
    )
    assert len(presentation.slides) == len(deck["slides"])

    assert response.validation.model_dump(by_alias=True) == {
        "passed": True,
        "layoutIssues": [],
        "contentIssues": [],
        "designIssues": [],
        "presentationIssues": [],
    }
    assert response.warnings == [
        "참고자료 없이 topic-only generation으로 생성했습니다."
    ]
    assert response.diagnostics.model_dump(by_alias=True) == {
        "referencePolicy": "user-input-only",
        "uploadedSourceCount": 0,
        "webSourceCount": 0,
        "researchAttempts": 0,
        "relevantWebSourceCount": 0,
        "officialWebSourceCount": 0,
        "independentWebSourceCount": 0,
        "researchQuality": "not-run",
        "researchIssueCodes": [],
        "researchFactCoverageSatisfied": False,
        "repairAttempted": True,
        "repairReasons": ["SPEAKER_NOTES_SHORT", "CONTENT_DUPLICATED"],
        "uniqueCoreLayoutCount": 8,
        "validationIssueCount": 0,
        "visualQaStatus": "not-run",
        "visualReviewAttempts": 0,
        "visualRepairAttempts": 0,
        "visualIssueCodes": [],
        "visualIssueSlideOrders": [],
        "warningCodes": [],
    }


def test_splatoon_product_launch_golden_composition_contract() -> None:
    fixture_path = (
        Path(__file__).parent
        / "fixtures"
        / "splatoon_product_launch_golden_request.json"
    )
    request = GenerateDeckRequest.model_validate_json(fixture_path.read_text("utf-8"))
    slide_plans = golden_slide_plans()
    responses = FakeResponses([golden_design_program()])
    orchestrator = DeckGenerationOrchestrator(
        request,
        client=SimpleNamespace(responses=responses),
    )
    raw_input = orchestrator.run_brief_agent()
    raw_input.source_records = [
        *initial_source_records(raw_input),
        SourceRecord(
            sourceType="web",
            sourceId="web:official",
            url="https://example.com/splatoon-official",
            title="Official Splatoon Raiders announcement",
            content="Official product reveal images and facts",
            authority="official",
        ),
    ]
    for slide_order in (2, 8):
        slide_plans[slide_order - 1].source_refs = ["web:official"]

    design_plan = orchestrator.run_design_director_agent(
        raw_input,
        slide_plans,
    )
    layout_result = orchestrator.run_layout_agent(
        raw_input,
        design_plan,
    )
    deck = orchestrator.build_deck(
        raw_input,
        type("Outline", (), {"title": "스플래툰 레이더스"})(),
        design_plan,
        layout_result.slides,
    )

    composition_ids = [
        slide["aiNotes"]["compositionPlan"]["compositionId"]
        for slide in deck["slides"]
    ]
    silhouettes = [COMPOSITION_SPECS[value].silhouette for value in composition_ids]
    asset_roles = [
        slide["aiNotes"]["compositionPlan"]["assetRole"]
        for slide in deck["slides"]
    ]
    backgrounds = {
        slide["aiNotes"]["compositionPlan"]["backgroundMode"]
        for slide in deck["slides"]
    }

    assert len(deck["slides"]) == 10
    assert deck["metadata"]["presentationProfile"] == "product-launch"
    assert max(Counter(composition_ids).values()) <= 2
    assert all(left != right for left, right in zip(silhouettes, silhouettes[1:]))
    assert len(backgrounds) >= 2
    assert 3 <= sum(role != "none" for role in asset_roles) <= 5
    for slide in deck["slides"]:
        focal_id = slide["aiNotes"]["compositionPlan"]["primaryFocalElementId"]
        assert focal_id in {element["elementId"] for element in slide["elements"]}
        assert all(
            "_program_v2_" in element["elementId"]
            for element in slide["elements"]
        )
    guarded_issue_codes = {
        "CONTENT_DUPLICATED",
        "VISUAL_HIERARCHY_WEAK",
        "GRID_ALIGNMENT_INCONSISTENT",
        "LINE_HEIGHT_OUT_OF_RANGE",
    }
    quality_issues = [
        issue
        for issue in validate_presentation(deck)
        if issue.code in guarded_issue_codes
    ]
    assert quality_issues == [], [
        (issue.code, issue.path, issue.message) for issue in quality_issues
    ]


@pytest.mark.parametrize(
    ("topic", "prompt", "expected_style_pack_id"),
    [
        ("Team alignment", "Create a concise strategy narrative.", "modern-editorial"),
        ("Product launch", "Show the new product features.", "product-showcase"),
        ("Quarterly KPI", "Create an executive metrics report.", "data-report"),
        ("Cloud architecture", "Explain the API security system.", "technical-system"),
    ],
)
def test_effective_style_pack_is_inferred_from_existing_request_fields(
    topic: str,
    prompt: str,
    expected_style_pack_id: str,
) -> None:
    raw_input = analyze_input(
        GenerateDeckRequest(
            projectId="project_demo_1",
            topic=topic,
            prompt=prompt,
            design={"visualRhythm": "auto"},
        )
    )

    assert (
        design_planning_module.effective_style_pack_id(raw_input)
        == expected_style_pack_id
    )
    assert (
        design_planning_module.select_style_pack(raw_input, [])["id"]
        == expected_style_pack_id
    )
    assert design_planning_module.effective_style_pack_prompt(raw_input)


def test_explicit_style_pack_stays_ahead_of_automatic_inference() -> None:
    raw_input = analyze_input(
        GenerateDeckRequest(
            projectId="project_demo_1",
            topic="Cloud architecture",
            prompt="Explain the API security system.",
            design={
                "stylePackId": "brandlogy-modern",
                "visualRhythm": "auto",
            },
        )
    )

    assert (
        design_planning_module.effective_style_pack_id(raw_input)
        == "brandlogy-modern"
    )
    context = design_planning_module.art_director_context(
        raw_input,
        design_planning_module.direct_design(raw_input, []),
        style_pack_id="brandlogy-modern",
        style_prompt=design_planning_module.effective_style_pack_prompt(raw_input),
    )
    assert "Effective style pack: brandlogy-modern." in context.design_direction
    assert "Brandlogy Modern Design Pack" in context.design_direction


def test_effective_style_pack_uses_generated_slide_types() -> None:
    raw_input = analyze_input(
        GenerateDeckRequest(
            projectId="project_demo_1",
            topic="Business update",
            prompt="Summarize the latest operating results.",
            design={"visualRhythm": "auto"},
        )
    )
    slide_plans = [
        SlidePlan(
            order=index,
            slide_type="data",
            title=f"Operating signal {index}",
            message="Evidence and interpretation",
            speaker_notes="Explain the evidence.",
            keywords=[],
            evidence=[],
        )
        for index in (1, 2)
    ]

    assert (
        design_planning_module.effective_style_pack_id(raw_input, slide_plans)
        == "data-report"
    )


def test_automatic_style_pack_preserves_palette_and_font_overrides() -> None:
    raw_input = analyze_input(
        GenerateDeckRequest(
            projectId="project_demo_1",
            topic="Product launch",
            prompt="Show the new product features.",
            design={
                "visualRhythm": "auto",
                "paletteOverride": {
                    "primary": "#123456",
                    "background": "#F7F7F2",
                    "text": "#111111",
                },
                "fontOverride": {
                    "fontId": "custom-sans",
                    "name": "Custom Sans",
                    "headingFontFamily": "Custom Sans",
                    "bodyFontFamily": "Custom Sans",
                    "recommendedTitleSize": 44,
                    "recommendedBodySize": 20,
                    "lineHeight": 1.24,
                },
            },
        )
    )

    theme = design_planning_module.direct_design(raw_input, [])
    theme = design_planning_module.apply_font_override(
        theme,
        raw_input.design.font_override,
    )

    assert theme["name"] == "product-showcase"
    assert theme["backgroundColor"] == "#F7F7F2"
    assert theme["textColor"] == "#111111"
    assert theme["palette"]["primary"] == "#123456"
    assert theme["fontFamily"] == "Custom Sans"
    assert theme["typography"]["headingFontFamily"] == "Custom Sans"
    assert theme["typography"]["bodyFontFamily"] == "Custom Sans"


def golden_slide_definitions() -> list[tuple[str, str, str, list[str]]]:
    return [
        ("cover", "미지의 군도로", "레이더스는 탐험 중심의 새 경험을 연다", ["공식 공개", "새로운 무대"]),
        ("solution", "혼자 떠나는 탐사", "익숙한 잉크 액션이 단독 탐험으로 확장된다", ["단독 플레이", "탐사 루프", "잉크 이동"]),
        ("feature-grid", "발견이 플레이가 된다", "섬의 발견과 수집이 진행 동기를 만든다", ["탐색", "수집", "성장"]),
        ("data", "공식 정보 한눈에", "확인된 정보와 미정 정보를 분리해 전달한다", ["플랫폼", "공개 상태"]),
        ("architecture", "탐험 루프", "이동과 발견, 대응, 귀환이 하나의 루프를 이룬다", ["이동", "발견", "대응", "귀환"]),
        ("comparison", "시리즈의 새 축", "대전 중심 경험과 다른 탐험 가치를 제안한다", ["기존 대전", "레이더스 탐험", "공통 잉크 액션"]),
        ("process", "한 번의 원정", "준비부터 귀환까지 선택이 이어진다", ["준비", "진입", "탐사", "귀환"]),
        ("data", "공식 장면이 증거다", "공식 키 아트와 트레일러 장면으로 변화를 보여준다", ["공식 키 아트", "트레일러 장면"]),
        ("solution", "출시 정보를 확인하세요", "공식 사이트에서 다음 공개와 출시 정보를 확인하세요", ["공식 사이트", "공식 채널", "다음 공개"]),
        ("summary", "다음 공개를 확인하세요", "공식 채널에서 출시 정보를 이어서 확인한다", ["공식 사이트", "공식 채널"]),
    ]


def golden_slide_plans() -> list[SlidePlan]:
    plans: list[SlidePlan] = []
    for order, (slide_type, title, message, items) in enumerate(
        golden_slide_definitions(),
        start=1,
    ):
        media = None
        if order in {1, 2, 8, 10}:
            media = MediaIntent(
                kind="generate",
                prompt=f"Splatoon Raiders official visual for {title}",
                alt=title,
                required=order in {1, 8},
            )
        plans.append(
            SlidePlan(
                order=order,
                slide_type=slide_type,
                title=title,
                message=message,
                speaker_notes=f"{order}번 슬라이드에서 {message}는 점을 설명합니다.",
                keywords=[title],
                evidence=[],
                content_items=[
                    GeneratedContentItem(
                        contentItemId=f"item-{order}-{index}",
                        text=value,
                    )
                    for index, value in enumerate(items, start=1)
                ],
                source_refs=["topic:brief"],
                **({"media_intent": media} if media is not None else {}),
            )
        )
    return plans


def golden_content_plan() -> dict[str, Any]:
    speaker_note_lengths = [205, 315, 315, 363, 363, 363, 363, 362, 315, 236]
    visual_intent = {
        "emphasis": "",
        "mood": "",
        "structure": "",
        "paletteHint": "",
        "emphasisStyle": "",
        "composition": "",
        "decorationDensity": "",
        "mediaStyle": "",
        "metricCardCaption": "",
    }
    slides: list[dict[str, Any]] = []
    for order, (definition, target) in enumerate(
        zip(
            golden_slide_definitions(),
            speaker_note_lengths,
            strict=True,
        ),
        start=1,
    ):
        slide_type, title, message, items = definition
        note_seed = "".join([title, message, *items])
        speaker_notes = (note_seed * (target // len(note_seed) + 1))[:target]
        if order == 10:
            speaker_notes = (
                "경청해 주셔서 감사합니다. 오늘 공유한 핵심 내용을 차분히 되짚으며 발표를 마무리하겠습니다. "
                "표지에서 제시한 주제와 본문에서 살펴본 근거가 하나의 흐름으로 연결되었음을 기억해 주시기 바랍니다. "
                "각 장표의 판단 기준은 확인된 정보의 범위 안에서 정리했으며, 추가 논의가 필요한 사항은 후속 자리에서 이어가겠습니다. "
                "함께 살펴본 내용이 앞으로의 대화와 검토에 유용한 출발점이 되기를 바랍니다. "
                "끝까지 집중해 주신 모든 분께 다시 한번 감사드립니다. "
                "발표 자료의 본문은 핵심 메시지와 근거를 중심으로 구성했으며, 마무리 장표에는 감사의 뜻만 담았습니다. "
                "오늘 나눈 관점이 각자의 상황에서 의미 있는 질문과 이해로 이어지기를 바랍니다. "
                "소중한 시간을 내어 함께해 주셔서 고맙습니다."
            )
        media_intent = {
            "kind": "none",
            "prompt": "",
            "alt": "",
            "caption": "",
            "rationale": "",
            "required": False,
            "placement": "auto",
            "src": "",
        }
        if order in {1, 2, 8, 10}:
            media_intent.update(
                {
                    "kind": "generate",
                    "prompt": f"Splatoon Raiders official visual for {title}",
                    "alt": title,
                    "required": order in {1, 8},
                }
            )
        slides.append(
            {
                "title": title,
                "message": message,
                "speakerNotes": speaker_notes,
                "keywords": [title],
                "slideType": slide_type,
                "visualIntent": visual_intent.copy(),
                "mediaIntent": media_intent,
                "contentItems": [
                    {
                        "contentItemId": f"item-{order}-{index}",
                        "text": item,
                    }
                    for index, item in enumerate(items, start=1)
                ],
                "sourceRefs": ["topic:brief"],
            }
        )
    return {"title": "스플래툰 레이더스", "slides": slides}


def golden_element_ids() -> list[list[str]]:
    return [
        [
            "el_1_program_v2_media_placeholder",
            "el_1_program_v2_media_caption",
            "el_1_program_v2_media_edge",
            "el_1_program_v2_title",
            "el_1_program_v2_subtitle",
        ],
        [
            "el_2_program_v2_title",
            "el_2_program_v2_agenda_rule",
            "el_2_program_v2_agenda_index_1",
            "el_2_program_v2_agenda_item_1",
            "el_2_program_v2_agenda_divider_1",
            "el_2_program_v2_agenda_index_2",
            "el_2_program_v2_agenda_item_2",
            "el_2_program_v2_agenda_divider_2",
            "el_2_program_v2_agenda_index_3",
            "el_2_program_v2_agenda_item_3",
            "el_2_program_v2_agenda_divider_3",
            "el_2_program_v2_agenda_index_4",
            "el_2_program_v2_agenda_item_4",
            "el_2_program_v2_agenda_divider_4",
            "el_2_program_v2_agenda_index_5",
            "el_2_program_v2_agenda_item_5",
            "el_2_program_v2_agenda_divider_5",
            "el_2_program_v2_agenda_index_6",
            "el_2_program_v2_agenda_item_6",
        ],
        [
            "el_3_program_v2_title",
            "el_3_program_v2_bento_1_field",
            "el_3_program_v2_bento_1_index",
            "el_3_program_v2_bento_1",
            "el_3_program_v2_bento_2_field",
            "el_3_program_v2_bento_2_index",
            "el_3_program_v2_bento_2",
            "el_3_program_v2_bento_3_field",
            "el_3_program_v2_bento_3_index",
            "el_3_program_v2_bento_3",
        ],
        [
            "el_4_program_v2_title",
            "el_4_program_v2_item_1_field",
            "el_4_program_v2_item_1_index",
            "el_4_program_v2_item_1",
            "el_4_program_v2_item_2_field",
            "el_4_program_v2_item_2_index",
            "el_4_program_v2_item_2",
        ],
        [
            "el_5_program_v2_title",
            "el_5_program_v2_orbit_connector_1",
            "el_5_program_v2_orbit_connector_2",
            "el_5_program_v2_orbit_connector_3",
            "el_5_program_v2_orbit_connector_4",
            "el_5_program_v2_orbit_hub_field",
            "el_5_program_v2_orbit_hub",
            "el_5_program_v2_orbit_node_1_field",
            "el_5_program_v2_orbit_node_1",
            "el_5_program_v2_orbit_node_2_field",
            "el_5_program_v2_orbit_node_2",
            "el_5_program_v2_orbit_node_3_field",
            "el_5_program_v2_orbit_node_3",
            "el_5_program_v2_orbit_node_4_field",
            "el_5_program_v2_orbit_node_4",
        ],
        [
            "el_6_program_v2_title",
            "el_6_program_v2_comparison_1_field",
            "el_6_program_v2_comparison_1_index",
            "el_6_program_v2_comparison_1",
            "el_6_program_v2_comparison_2_field",
            "el_6_program_v2_comparison_2_index",
            "el_6_program_v2_comparison_2",
            "el_6_program_v2_comparison_3_field",
            "el_6_program_v2_comparison_3_index",
            "el_6_program_v2_comparison_3",
        ],
        [
            "el_7_program_v2_title",
            "el_7_program_v2_vertical_rail",
            "el_7_program_v2_rail_marker_1",
            "el_7_program_v2_rail_marker_label_1",
            "el_7_program_v2_rail_rule_1",
            "el_7_program_v2_rail_step_1",
            "el_7_program_v2_rail_marker_2",
            "el_7_program_v2_rail_marker_label_2",
            "el_7_program_v2_rail_rule_2",
            "el_7_program_v2_rail_step_2",
            "el_7_program_v2_rail_marker_3",
            "el_7_program_v2_rail_marker_label_3",
            "el_7_program_v2_rail_rule_3",
            "el_7_program_v2_rail_step_3",
            "el_7_program_v2_rail_marker_4",
            "el_7_program_v2_rail_marker_label_4",
            "el_7_program_v2_rail_rule_4",
            "el_7_program_v2_rail_step_4",
        ],
        [
            "el_8_program_v2_title",
            "el_8_program_v2_editorial_band_rule_1",
            "el_8_program_v2_editorial_band_item_1",
            "el_8_program_v2_editorial_band_rule_2",
            "el_8_program_v2_editorial_band_item_2",
        ],
        [
            "el_9_program_v2_title",
            "el_9_program_v2_hub_field",
            "el_9_program_v2_hub",
            "el_9_program_v2_connector_left",
            "el_9_program_v2_connector_right",
            "el_9_program_v2_connector_bottom",
            "el_9_program_v2_node_1_field",
            "el_9_program_v2_node_1_index",
            "el_9_program_v2_node_1",
            "el_9_program_v2_node_2_field",
            "el_9_program_v2_node_2_index",
            "el_9_program_v2_node_2",
            "el_9_program_v2_node_3_field",
            "el_9_program_v2_node_3_index",
            "el_9_program_v2_node_3",
        ],
        [
            "el_10_program_v2_closing_mark",
            "el_10_program_v2_title",
            "el_10_program_v2_subtitle",
        ],
    ]


def golden_design_program() -> dict[str, Any]:
    composition_ids = [
        "hero-full-bleed",
        "agenda-numbered-list",
        "feature-comparison",
        "metric-poster",
        "diagram-hub",
        "feature-comparison",
        "process-horizontal",
        "image-evidence",
        "kpi-strip-evidence",
        "closing-centered-minimal",
    ]
    background_modes = [
        "image",
        "light",
        "dark",
        "light",
        "dark",
        "light",
        "dark",
        "light",
        "dark",
        "light",
    ]
    asset_roles = [
        "atmosphere",
        "evidence",
        "none",
        "none",
        "none",
        "none",
        "none",
        "evidence",
        "none",
        "atmosphere",
    ]
    return {
        "version": "program-v2",
        "visualConcept": "Playful ink expedition with editorial evidence",
        "paletteRoles": {
            "dominant": "#FFFFFF",
            "surface": "#F3F4F6",
            "text": "#111827",
            "focal": "#6D28D9",
            "secondary": "#22D3EE",
        },
        "typography": {
            "headingFont": "Pretendard",
            "bodyFont": "Pretendard",
            "typeScale": {"cover": 64, "title": 40, "body": 22, "caption": 14},
        },
        "backgroundSequence": background_modes,
        "imageStyle": "Official game art with bold clean crops",
        "surfaceStyle": "Flat ink fields without gradients",
        "slides": [
            {
                "order": order,
                "compositionId": composition_id,
                "variant": background,
                "backgroundMode": background,
                "focalType": "primary-message",
                "assetRole": asset_role,
                "requiredAsset": order in {1, 8},
            }
            for order, (composition_id, background, asset_role) in enumerate(
                zip(composition_ids, background_modes, asset_roles, strict=True),
                start=1,
            )
        ],
    }
