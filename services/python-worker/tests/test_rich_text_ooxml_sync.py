import base64
import copy
from io import BytesIO
from pathlib import Path
from unittest.mock import patch
from xml.etree import ElementTree as ET
from zipfile import ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from app.ai.pptx_ooxml_generation import (
    canonical_text_paragraphs,
    generate_pptx_ooxml,
    sync_pptx_ooxml,
    text_equal_spans,
)
from app.ai.pptx_design_importer import ImportedDesignAsset


PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
DML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def template_slide_id(generated: object) -> str:
    return generated.template_blueprint["slides"][0]["slideId"]


def test_imported_rich_text_capability_distinguishes_simple_and_hyperlink(
    tmp_path: Path,
) -> None:
    pptx_path = rich_text_source_pptx(tmp_path, add_hyperlink=True)

    generated = generate_pptx_ooxml(pptx_path, "file_rich_text", render=False)

    simple_source = source_for_text(generated, "Simple source")
    linked_source = source_for_text(generated, "Linked source")
    assert simple_source["ooxmlEditCapabilities"]["richText"] == "full"
    assert linked_source["ooxmlEditCapabilities"]["richText"] == "style-only"


def test_equal_plain_text_projection_preserves_full_mixed_run_bytes(
    tmp_path: Path,
) -> None:
    pptx_path = rich_text_source_pptx(tmp_path)
    original_package = pptx_path.read_bytes()
    generated = generate_pptx_ooxml(pptx_path, "file_rich_text", render=False)
    text = text_element(generated, "Simple source")

    result = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        operations=[
            {
                "type": "update_element_props",
                "slideId": template_slide_id(generated),
                "elementId": text["elementId"],
                "props": {"text": "Simple source"},
            }
        ],
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
    )

    assert result.unsupported_operations == []
    assert current_package_bytes(result.assets) == original_package


def test_equal_text_with_style_preserves_mixed_run_structure(
    tmp_path: Path,
) -> None:
    pptx_path = rich_text_source_pptx(tmp_path)
    generated = generate_pptx_ooxml(pptx_path, "file_rich_text", render=False)
    text = text_element(generated, "Simple source")
    source = source_for_text(generated, "Simple source")

    styled = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        operations=[
            {
                "type": "update_element_props",
                "slideId": template_slide_id(generated),
                "elementId": text["elementId"],
                "props": {"text": "Simple source", "underline": True},
            }
        ],
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
    )

    assert styled.unsupported_operations == []
    package = current_package_bytes(styled.assets)
    shape = shape_element(package, source["shapeId"])
    run_properties = list(shape.iter(f"{{{DML_NS}}}rPr"))
    assert len(run_properties) == 2
    assert run_properties[0].attrib["kumimoji"] == "1"
    assert run_properties[0].find(f"{{{DML_NS}}}effectLst") is not None
    assert run_properties[0].attrib["u"] == "sng"
    assert run_properties[1].attrib["u"] == "sng"
    assert run_properties[1].attrib["i"] == "1"

    original_first_rpr = run_property_xml(package, source["shapeId"], 0)
    original_second_rpr = run_property_xml(package, source["shapeId"], 1)
    styled_path = tmp_path / "mixed-run-styled.pptx"
    styled_path.write_bytes(package)
    body_only = sync_pptx_ooxml(
        styled_path,
        template_blueprint=generated.template_blueprint,
        operations=[
            {
                "type": "update_element_props",
                "slideId": template_slide_id(generated),
                "elementId": text["elementId"],
                "props": {"text": "Simple source", "verticalAlign": "middle"},
            }
        ],
        deck_canvas=generated.canvas,
        synced_deck_version=3,
        render=False,
    )
    assert body_only.unsupported_operations == []
    body_package = current_package_bytes(body_only.assets)
    assert run_property_xml(body_package, source["shapeId"], 0) == original_first_rpr
    assert run_property_xml(body_package, source["shapeId"], 1) == original_second_rpr
    body_pr = shape_element(body_package, source["shapeId"]).find(
        f"{{{PML_NS}}}txBody/{{{DML_NS}}}bodyPr"
    )
    assert body_pr is not None
    assert body_pr.attrib["anchor"] == "ctr"


def test_imported_letter_spacing_and_autofit_round_trip_through_targeted_sync(
    tmp_path: Path,
) -> None:
    pptx_path = rich_text_source_pptx(tmp_path)
    generated = generate_pptx_ooxml(pptx_path, "file_rich_text", render=False)
    text = text_element(generated, "Simple source")
    source = source_for_text(generated, "Simple source")
    props = copy.deepcopy(text["props"])
    props.update(
        {
            "autoFit": "shrink-text",
            "fontScale": 0.8,
            "lineSpaceReduction": 0.1,
        }
    )
    props["paragraphs"][0]["runs"][0]["letterSpacing"] = 2.4

    result = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        operations=[
            {
                "type": "update_element_props",
                "slideId": template_slide_id(generated),
                "elementId": text["elementId"],
                "props": props,
            }
        ],
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
    )

    assert result.unsupported_operations == []
    package = current_package_bytes(result.assets)
    shape = shape_element(package, source["shapeId"])
    body_properties = shape.find(f"{{{PML_NS}}}txBody/{{{DML_NS}}}bodyPr")
    assert body_properties is not None
    normal_autofit = body_properties.find(f"{{{DML_NS}}}normAutofit")
    assert normal_autofit is not None
    assert normal_autofit.attrib == {
        "fontScale": "80000",
        "lnSpcReduction": "10000",
    }
    run_properties = next(shape.iter(f"{{{DML_NS}}}rPr"))
    assert run_properties.attrib["spc"] == "120"

    synced_path = tmp_path / "autofit-synced.pptx"
    synced_path.write_bytes(package)
    reimported = generate_pptx_ooxml(
        synced_path,
        "file_rich_text_reimport",
        render=False,
    )
    imported = text_element(reimported, "Simple source")["props"]
    assert imported["autoFit"] == "shrink-text"
    assert imported["fontScale"] == 0.8
    assert imported["lineSpaceReduction"] == 0.1
    assert imported["paragraphs"][0]["runs"][0]["letterSpacing"] == 2.4


def test_imported_field_text_is_fail_closed_without_package_mutation(
    tmp_path: Path,
) -> None:
    pptx_path = rich_text_source_pptx(tmp_path, add_field=True)
    original_package = pptx_path.read_bytes()
    generated = generate_pptx_ooxml(pptx_path, "file_rich_text", render=False)
    text = text_element(generated, "Simple source")
    source = source_for_text(generated, "Simple source")
    assert source["ooxmlEditCapabilities"]["richText"] == "none"

    props = copy.deepcopy(text["props"])
    props["paragraphs"][0]["runs"][0]["fontWeight"] = "normal"
    result = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        operations=[
            {
                "type": "update_element_props",
                "slideId": template_slide_id(generated),
                "elementId": text["elementId"],
                "props": props,
            }
        ],
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
    )

    assert current_package_bytes(result.assets) == original_package
    assert result.applied_operations == []
    assert [item.reason_code for item in result.unsupported_operations] == [
        "RICH_TEXT_CAPABILITY_UNSAFE"
    ]


def test_imported_rich_text_style_sync_preserves_unknown_rpr_and_unselected_run(
    tmp_path: Path,
) -> None:
    pptx_path = rich_text_source_pptx(tmp_path)
    generated = generate_pptx_ooxml(pptx_path, "file_rich_text", render=False)
    text = text_element(generated, "Simple source")
    source = source_for_text(generated, "Simple source")
    original_package = pptx_path.read_bytes()
    original_second_rpr = run_property_xml(original_package, source["shapeId"], 1)
    props = copy.deepcopy(text["props"])
    props["paragraphs"][0]["runs"][0]["underline"] = True
    props["paragraphs"][0]["runs"][0]["color"] = "#16A34A"

    result = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        operations=[
            {
                "type": "update_element_props",
                "slideId": template_slide_id(generated),
                "elementId": text["elementId"],
                "props": props,
            }
        ],
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
    )

    assert result.unsupported_operations == []
    assert [item.operation_type for item in result.applied_operations] == [
        "update_element_props"
    ]
    package = current_package_bytes(result.assets)
    first_rpr = run_property_element(package, source["shapeId"], 0)
    assert first_rpr.attrib["kumimoji"] == "1"
    assert first_rpr.attrib["u"] == "sng"
    assert first_rpr.find(f"{{{DML_NS}}}effectLst") is not None
    assert first_rpr.find(f"{{{DML_NS}}}solidFill/{{{DML_NS}}}srgbClr").attrib == {
        "val": "16A34A"
    }
    assert run_property_xml(package, source["shapeId"], 1) == original_second_rpr

    round_trip_path = tmp_path / "rich-text-style-round-trip.pptx"
    round_trip_path.write_bytes(package)
    round_trip = generate_pptx_ooxml(
        round_trip_path,
        "file_rich_text_round_trip",
        render=False,
    )
    round_trip_text = text_element(round_trip, "Simple source")
    assert round_trip_text["props"]["paragraphs"][0]["runs"][0][
        "underline"
    ] is True
    assert round_trip_text["props"]["paragraphs"][0]["runs"][0]["color"] == (
        "#16A34A"
    )
    assert round_trip_text["props"]["paragraphs"][0]["runs"][1]["italic"] is (
        True
    )


def test_partial_style_props_preserve_existing_text_and_hyperlink(
    tmp_path: Path,
) -> None:
    simple_path = rich_text_source_pptx(tmp_path)
    simple_generated = generate_pptx_ooxml(
        simple_path,
        "file_rich_text",
        render=False,
    )
    simple_text = text_element(simple_generated, "Simple source")
    simple_source = source_for_text(simple_generated, "Simple source")
    simple_result = sync_pptx_ooxml(
        simple_path,
        template_blueprint=simple_generated.template_blueprint,
        operations=[
            {
                "type": "update_element_props",
                "slideId": template_slide_id(simple_generated),
                "elementId": simple_text["elementId"],
                "props": {"underline": True},
            }
        ],
        deck_canvas=simple_generated.canvas,
        synced_deck_version=2,
        render=False,
    )
    assert simple_result.unsupported_operations == []
    simple_package = current_package_bytes(simple_result.assets)
    simple_shape = shape_element(simple_package, simple_source["shapeId"])
    assert "".join(
        node.text or "" for node in simple_shape.iter(f"{{{DML_NS}}}t")
    ) == "Simple source"
    assert all(
        r_pr.attrib["u"] == "sng"
        for r_pr in simple_shape.iter(f"{{{DML_NS}}}rPr")
    )

    linked_path = rich_text_source_pptx(tmp_path, add_hyperlink=True)
    linked_generated = generate_pptx_ooxml(
        linked_path,
        "file_linked_text",
        render=False,
    )
    linked_text = text_element(linked_generated, "Linked source")
    linked_source = source_for_text(linked_generated, "Linked source")
    linked_result = sync_pptx_ooxml(
        linked_path,
        template_blueprint=linked_generated.template_blueprint,
        operations=[
            {
                "type": "update_element_props",
                "slideId": template_slide_id(linked_generated),
                "elementId": linked_text["elementId"],
                "props": {"fontWeight": "bold"},
            }
        ],
        deck_canvas=linked_generated.canvas,
        synced_deck_version=2,
        render=False,
    )
    assert linked_result.unsupported_operations == []
    linked_rpr = run_property_element(
        current_package_bytes(linked_result.assets),
        linked_source["shapeId"],
        0,
    )
    assert linked_rpr.attrib["b"] == "1"
    assert linked_rpr.find(f"{{{DML_NS}}}hlinkClick") is not None


def test_full_content_prefix_insertion_keeps_unknown_properties_on_unchanged_runs(
    tmp_path: Path,
) -> None:
    pptx_path = rich_text_source_pptx(tmp_path)
    generated = generate_pptx_ooxml(pptx_path, "file_rich_text", render=False)
    text = text_element(generated, "Simple source")
    source = source_for_text(generated, "Simple source")
    original_second_rpr = run_property_xml(
        pptx_path.read_bytes(),
        source["shapeId"],
        1,
    )
    props = copy.deepcopy(text["props"])
    original_runs = props["paragraphs"][0]["runs"]
    props["paragraphs"][0]["runs"] = [
        {"text": "🚀 inserted ", "baseline": "normal"},
        *original_runs,
    ]
    props["paragraphs"][0]["text"] = "🚀 inserted Simple source"
    props["text"] = "🚀 inserted Simple source"

    result = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        operations=[
            {
                "type": "update_element_props",
                "slideId": template_slide_id(generated),
                "elementId": text["elementId"],
                "props": props,
            }
        ],
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
    )

    assert result.unsupported_operations == []
    package = current_package_bytes(result.assets)
    unchanged_first_rpr = run_property_element(package, source["shapeId"], 1)
    assert unchanged_first_rpr.attrib["kumimoji"] == "1"
    assert unchanged_first_rpr.find(f"{{{DML_NS}}}effectLst") is not None
    assert run_property_xml(package, source["shapeId"], 2) == original_second_rpr


def test_hyperlink_style_only_sync_preserves_link_and_rejects_content_change(
    tmp_path: Path,
) -> None:
    pptx_path = rich_text_source_pptx(tmp_path, add_hyperlink=True)
    original_package = pptx_path.read_bytes()
    generated = generate_pptx_ooxml(pptx_path, "file_rich_text", render=False)
    text = text_element(generated, "Linked source")
    source = source_for_text(generated, "Linked source")
    style_props = copy.deepcopy(text["props"])
    style_props["paragraphs"][0]["runs"][0]["fontWeight"] = "bold"

    styled = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        operations=[
            {
                "type": "update_element_props",
                "slideId": template_slide_id(generated),
                "elementId": text["elementId"],
                "props": style_props,
            }
        ],
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
    )

    assert styled.unsupported_operations == []
    styled_rpr = run_property_element(
        current_package_bytes(styled.assets),
        source["shapeId"],
        0,
    )
    hyperlink = styled_rpr.find(f"{{{DML_NS}}}hlinkClick")
    assert hyperlink is not None
    assert hyperlink.attrib[f"{{{REL_NS}}}id"] == "rId999"

    segmented_props = copy.deepcopy(text["props"])
    original_run = segmented_props["paragraphs"][0]["runs"][0]
    segmented_props["paragraphs"][0]["runs"] = [
        {**original_run, "text": "Linked"},
        {**original_run, "text": " source"},
    ]
    segmented = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        operations=[
            {
                "type": "update_element_props",
                "slideId": template_slide_id(generated),
                "elementId": text["elementId"],
                "props": segmented_props,
            }
        ],
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
    )
    assert segmented.unsupported_operations == []
    segmented_package = current_package_bytes(segmented.assets)
    for run_index in (0, 1):
        split_link = run_property_element(
            segmented_package,
            source["shapeId"],
            run_index,
        ).find(f"{{{DML_NS}}}hlinkClick")
        assert split_link is not None
        assert split_link.attrib[f"{{{REL_NS}}}id"] == "rId999"

    content_props = copy.deepcopy(text["props"])
    content_props["paragraphs"][0]["runs"][0]["text"] = "Changed link"
    content_props["paragraphs"][0]["text"] = "Changed link"
    content_props["text"] = "Changed link"
    rejected = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        operations=[
            {
                "type": "update_element_props",
                "slideId": template_slide_id(generated),
                "elementId": text["elementId"],
                "props": content_props,
            }
        ],
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
    )
    assert current_package_bytes(rejected.assets) == original_package
    assert rejected.applied_operations == []
    assert [item.reason_code for item in rejected.unsupported_operations] == [
        "RICH_TEXT_CAPABILITY_UNSAFE"
    ]


def test_style_only_merge_keeps_distinct_hyperlink_run_boundaries(
    tmp_path: Path,
) -> None:
    pptx_path = rich_text_source_pptx(tmp_path, add_hyperlink=True)
    split_linked_run_relationships(pptx_path)
    generated = generate_pptx_ooxml(pptx_path, "file_rich_text", render=False)
    text = text_element(generated, "Linked source")
    source = source_for_text(generated, "Linked source")
    assert source["ooxmlEditCapabilities"]["richText"] == "style-only"
    props = copy.deepcopy(text["props"])
    props["paragraphs"][0]["runs"] = [
        {
            "text": "Linked source",
            "fontWeight": "bold",
            "baseline": "normal",
        }
    ]

    result = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        operations=[
            {
                "type": "update_element_props",
                "slideId": template_slide_id(generated),
                "elementId": text["elementId"],
                "props": props,
            }
        ],
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
    )

    assert result.unsupported_operations == []
    package = current_package_bytes(result.assets)
    for run_index, relationship_id in enumerate(("rId999", "rId998")):
        r_pr = run_property_element(package, source["shapeId"], run_index)
        assert r_pr.attrib["b"] == "1"
        hyperlink = r_pr.find(f"{{{DML_NS}}}hlinkClick")
        assert hyperlink is not None
        assert hyperlink.attrib[f"{{{REL_NS}}}id"] == relationship_id


def test_authored_canonical_text_add_and_same_batch_edit_round_trip(
    tmp_path: Path,
) -> None:
    pptx_path = rich_text_source_pptx(tmp_path)
    generated = generate_pptx_ooxml(pptx_path, "file_rich_text", render=False)
    slide_id = template_slide_id(generated)
    element = authored_text_element("Authored first\nSecond")
    edited_props = copy.deepcopy(element["props"])
    edited_props["paragraphs"][0]["runs"][1]["text"] = " edited"
    edited_props["paragraphs"][0]["text"] = "Authored edited"
    edited_props["text"] = "Authored edited\nSecond"

    result = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        operations=[
            {
                "type": "add_element",
                "slideId": slide_id,
                "element": element,
            },
            {
                "type": "update_element_props",
                "slideId": slide_id,
                "elementId": element["elementId"],
                "props": edited_props,
            },
        ],
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
    )

    assert result.unsupported_operations == []
    assert [item.operation_type for item in result.applied_operations] == [
        "add_element",
        "update_element_props",
    ]
    source = next(
        item for item in result.element_sources if item["elementId"] == element["elementId"]
    )
    assert source["ooxmlOrigin"] == "authored"
    assert source["ooxmlEditCapabilities"]["richText"] == "full"

    round_trip_path = tmp_path / "authored-rich-text-round-trip.pptx"
    round_trip_path.write_bytes(current_package_bytes(result.assets))
    round_trip = generate_pptx_ooxml(
        round_trip_path,
        "file_authored_round_trip",
        render=False,
    )
    imported = text_element(round_trip, "Authored edited\nSecond")
    assert imported["props"]["text"] == "Authored edited\nSecond"
    assert len(imported["props"]["paragraphs"]) == 2
    assert imported["props"]["paragraphs"][0]["runs"][0]["fontWeight"] == (
        "bold"
    )
    assert imported["props"]["paragraphs"][0]["runs"][1]["underline"] is (
        True
    )
    assert imported["props"]["paragraphs"][1]["bullet"]["enabled"] is True


def test_authored_empty_runs_follow_text_fallback_contract(tmp_path: Path) -> None:
    assert canonical_text_paragraphs(
        {
            "text": "Paragraph fallback",
            "paragraphs": [{"text": "Paragraph fallback", "runs": []}],
        }
    ) == [
        {
            "text": "Paragraph fallback",
            "runs": [{"text": "Paragraph fallback"}],
        }
    ]
    assert canonical_text_paragraphs({"text": "Top fallback", "runs": []}) == [
        {"text": "Top fallback", "runs": [{"text": "Top fallback"}]}
    ]
    assert canonical_text_paragraphs(
        {"text": "", "paragraphs": [{"text": "", "runs": []}]}
    ) == [{"text": "", "runs": []}]

    pptx_path = rich_text_source_pptx(tmp_path)
    generated = generate_pptx_ooxml(pptx_path, "file_rich_text", render=False)
    slide_id = template_slide_id(generated)
    elements: list[dict[str, object]] = []
    for index, props in enumerate(
        (
            {
                "text": "Paragraph fallback",
                "runs": [],
                "paragraphs": [{"text": "Paragraph fallback", "runs": []}],
            },
            {"text": "Top fallback", "runs": []},
            {"text": "", "runs": [], "paragraphs": [{"text": "", "runs": []}]},
            {
                "text": "Explicit runs",
                "runs": [{"text": "Explicit runs", "baseline": "normal"}],
            },
        )
    ):
        element = authored_text_element("Authored first\nSecond")
        element["elementId"] = f"el_empty_runs_{index}"
        element["y"] = 100 + index * 120
        element["props"] = {
            **props,
            "fontFamily": "Aptos",
            "fontSize": 28,
            "fontWeight": "normal",
            "align": "left",
            "verticalAlign": "top",
            "lineHeight": 1.2,
        }
        elements.append(element)

    result = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        operations=[
            {"type": "add_element", "slideId": slide_id, "element": element}
            for element in elements
        ],
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
    )

    assert result.unsupported_operations == []
    assert len(result.applied_operations) == 4
    package = current_package_bytes(result.assets)
    sources = {
        str(source["elementId"]): source for source in result.element_sources
    }
    expected_text_by_element = {
        "el_empty_runs_0": "Paragraph fallback",
        "el_empty_runs_1": "Top fallback",
        "el_empty_runs_3": "Explicit runs",
    }
    for element_id, expected in expected_text_by_element.items():
        assert expected in "".join(
            node.text or ""
            for node in shape_element(
                package,
                sources[element_id]["shapeId"],
            ).iter(f"{{{DML_NS}}}t")
        )
    assert list(
        shape_element(package, sources["el_empty_runs_2"]["shapeId"]).iter(
            f"{{{DML_NS}}}t"
        )
    ) == []


def test_large_repetitive_text_diff_uses_bounded_fallback() -> None:
    source = "ab" * 3_000
    target = "ba" * 3_000

    with patch(
        "app.ai.pptx_ooxml_generation.difflib.SequenceMatcher",
        side_effect=AssertionError("unbounded matcher invoked"),
    ):
        assert text_equal_spans(source, target) == []


def rich_text_source_pptx(
    tmp_path: Path,
    *,
    add_hyperlink: bool = False,
    add_field: bool = False,
) -> Path:
    pptx_path = tmp_path / "rich-text-sync-source.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    simple = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    simple.text_frame.clear()
    simple_first = simple.text_frame.paragraphs[0].add_run()
    simple_first.text = "Simple "
    simple_first.font.name = "Aptos"
    simple_first.font.size = Pt(24)
    simple_first.font.bold = True
    simple_first.font.color.rgb = RGBColor(0xDC, 0x26, 0x26)
    simple_second = simple.text_frame.paragraphs[0].add_run()
    simple_second.text = "source"
    simple_second.font.name = "Arial"
    simple_second.font.size = Pt(20)
    simple_second.font.italic = True

    linked = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(5), Inches(1))
    linked.text_frame.clear()
    linked_run = linked.text_frame.paragraphs[0].add_run()
    linked_run.text = "Linked source"
    linked_run.font.size = Pt(22)
    presentation.save(pptx_path)

    add_run_property_extensions(
        pptx_path,
        simple_shape_id=str(simple.shape_id),
        linked_shape_id=str(linked.shape_id) if add_hyperlink else None,
        add_field=add_field,
    )
    return pptx_path


def add_run_property_extensions(
    pptx_path: Path,
    *,
    simple_shape_id: str,
    linked_shape_id: str | None,
    add_field: bool,
) -> None:
    output = BytesIO()
    with ZipFile(pptx_path, "r") as source, ZipFile(output, "w") as target:
        root = ET.fromstring(source.read("ppt/slides/slide1.xml"))
        simple_shape = shape_by_id(root, simple_shape_id)
        simple_rpr = next(simple_shape.iter(f"{{{DML_NS}}}rPr"))
        simple_rpr.set("kumimoji", "1")
        ET.SubElement(simple_rpr, f"{{{DML_NS}}}effectLst")
        if add_field:
            simple_run = next(simple_shape.iter(f"{{{DML_NS}}}r"))
            paragraph = next(
                candidate
                for candidate in simple_shape.iter(f"{{{DML_NS}}}p")
                if simple_run in list(candidate)
            )
            field = ET.SubElement(
                paragraph,
                f"{{{DML_NS}}}fld",
                {"id": "{00000000-0000-0000-0000-000000000001}", "type": "slidenum"},
            )
            field.append(copy.deepcopy(simple_rpr))
            ET.SubElement(field, f"{{{DML_NS}}}t").text = "7"

        if linked_shape_id is not None:
            linked_shape = shape_by_id(root, linked_shape_id)
            linked_rpr = next(linked_shape.iter(f"{{{DML_NS}}}rPr"))
            ET.SubElement(
                linked_rpr,
                f"{{{DML_NS}}}hlinkClick",
                {f"{{{REL_NS}}}id": "rId999"},
            )

        slide_xml = ET.tostring(root, encoding="utf-8", xml_declaration=True)
        for info in source.infolist():
            target.writestr(
                info,
                slide_xml
                if info.filename == "ppt/slides/slide1.xml"
                else source.read(info.filename),
            )
    pptx_path.write_bytes(output.getvalue())


def split_linked_run_relationships(pptx_path: Path) -> None:
    output = BytesIO()
    with ZipFile(pptx_path, "r") as source, ZipFile(output, "w") as target:
        root = ET.fromstring(source.read("ppt/slides/slide1.xml"))
        linked_shape = next(
            shape
            for shape in root.iter(f"{{{PML_NS}}}sp")
            if "".join(
                node.text or ""
                for node in shape.iter()
                if node.tag == f"{{{DML_NS}}}t"
            )
            == "Linked source"
        )
        paragraph = next(linked_shape.iter(f"{{{DML_NS}}}p"))
        first_run = next(linked_shape.iter(f"{{{DML_NS}}}r"))
        first_text = first_run.find(f"{{{DML_NS}}}t")
        assert first_text is not None
        first_text.text = "Linked"
        second_run = copy.deepcopy(first_run)
        second_text = second_run.find(f"{{{DML_NS}}}t")
        assert second_text is not None
        second_text.text = " source"
        second_link = second_run.find(
            f"{{{DML_NS}}}rPr/{{{DML_NS}}}hlinkClick"
        )
        assert second_link is not None
        second_link.set(f"{{{REL_NS}}}id", "rId998")
        paragraph.insert(list(paragraph).index(first_run) + 1, second_run)

        slide_xml = ET.tostring(root, encoding="utf-8", xml_declaration=True)
        for info in source.infolist():
            target.writestr(
                info,
                slide_xml
                if info.filename == "ppt/slides/slide1.xml"
                else source.read(info.filename),
            )
    pptx_path.write_bytes(output.getvalue())


def authored_text_element(text: str) -> dict[str, object]:
    return {
        "elementId": "el_authored_rich_text",
        "type": "text",
        "x": 820,
        "y": 160,
        "width": 720,
        "height": 300,
        "rotation": 0,
        "opacity": 1,
        "visible": True,
        "locked": False,
        "zIndex": 20,
        "props": {
            "text": text,
            "fontFamily": "Aptos",
            "fontSize": 28,
            "fontWeight": "normal",
            "align": "left",
            "verticalAlign": "top",
            "lineHeight": 1.2,
            "paragraphs": [
                {
                    "text": "Authored first",
                    "runs": [
                        {
                            "text": "Authored",
                            "fontWeight": "bold",
                            "color": "#2563EB",
                            "baseline": "normal",
                        },
                        {
                            "text": " first",
                            "underline": True,
                            "baseline": "normal",
                        },
                    ],
                    "align": "center",
                    "lineHeight": 1.3,
                    "spaceBefore": 4,
                    "spaceAfter": 6,
                    "indent": 12,
                },
                {
                    "text": "Second",
                    "runs": [
                        {
                            "text": "Second",
                            "italic": True,
                            "baseline": "normal",
                        }
                    ],
                    "align": "left",
                    "lineHeight": 1.2,
                    "spaceBefore": 0,
                    "spaceAfter": 0,
                    "indent": 18,
                    "bullet": {"enabled": True, "character": "•", "indent": 18},
                },
            ],
        },
    }


def text_element(result: object, text: str) -> dict:
    blueprint = getattr(result, "blueprint")
    return next(
        element
        for element in blueprint["slides"][0]["elements"]
        if element["type"] == "text" and element["props"]["text"] == text
    )


def source_for_text(result: object, text: str) -> dict:
    element = text_element(result, text)
    blueprint = getattr(result, "template_blueprint")
    return next(
        source
        for source in blueprint["slides"][0]["elementSources"]
        if source["elementId"] == element["elementId"]
    )


def current_package_bytes(assets: list[ImportedDesignAsset]) -> bytes:
    package = next(asset for asset in assets if asset.asset_id == "current_package")
    return base64.b64decode(package.content_base64)


def shape_element(package: bytes, shape_id: str) -> ET.Element:
    with ZipFile(BytesIO(package), "r") as archive:
        root = ET.fromstring(archive.read("ppt/slides/slide1.xml"))
    return shape_by_id(root, shape_id)


def run_property_element(package: bytes, shape_id: str, run_index: int) -> ET.Element:
    with ZipFile(BytesIO(package), "r") as archive:
        root = ET.fromstring(archive.read("ppt/slides/slide1.xml"))
    shape = shape_by_id(root, shape_id)
    return list(shape.iter(f"{{{DML_NS}}}rPr"))[run_index]


def run_property_xml(package: bytes, shape_id: str, run_index: int) -> bytes:
    return ET.tostring(run_property_element(package, shape_id, run_index))


def shape_by_id(root: ET.Element, shape_id: str) -> ET.Element:
    for shape in root.iter(f"{{{PML_NS}}}sp"):
        c_nv_pr = next(shape.iter(f"{{{PML_NS}}}cNvPr"), None)
        if c_nv_pr is not None and c_nv_pr.get("id") == shape_id:
            return shape
    raise AssertionError(f"shape {shape_id} not found")
