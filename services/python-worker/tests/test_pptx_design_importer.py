import zipfile
from io import BytesIO
from pathlib import Path
from xml.etree import ElementTree as ET

from PIL import Image
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt

from app.ai.pptx_design_importer import (
    assign_text_roles,
    apply_repeated_text_roles,
    blip_fill_asset,
    build_template_blueprint,
    build_quality_report,
    import_pptx_design,
)
from app.ai.pptx_ooxml_vector_importer import (
    VECTOR_IMPORT_FLAG,
    import_pptx_design_with_optional_ooxml_vector,
    import_pptx_ooxml_visual_tree,
)


def test_import_pptx_design_extracts_editable_elements(tmp_path: Path) -> None:
    pptx_path = tmp_path / "design.pptx"
    image_path = tmp_path / "sample.png"
    Image.new("RGB", (24, 24), "#ff0000").save(image_path)

    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(240, 248, 255)

    title = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(5), Inches(1))
    title.text_frame.text = "Original Title"
    title_run = title.text_frame.paragraphs[0].runs[0]
    title_run.font.size = Pt(40)
    title_run.font.bold = True
    title_run.font.name = "Aptos Display"
    title_run.font.color.rgb = RGBColor(10, 20, 30)

    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1),
        Inches(2.4),
        Inches(3),
        Inches(1.2),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(17, 34, 51)
    box.line.color.rgb = RGBColor(68, 85, 102)
    slide.shapes.add_picture(
        str(image_path), Inches(7), Inches(2), Inches(2), Inches(2)
    )
    presentation.save(pptx_path)

    result = import_pptx_design(pptx_path, "file_design")
    slide_blueprint = result.blueprint["slides"][0]
    elements = slide_blueprint["elements"]

    assert slide_blueprint["style"]["backgroundColor"] == "#F0F8FF"
    assert slide_blueprint["style"]["textColor"] == "#0A141E"
    assert slide_blueprint["style"]["accentColor"] == "#112233"
    assert slide_blueprint["style"]["fontFamily"] == "Aptos Display"
    assert result.blueprint["theme"]["fontFamily"] == "Aptos Display"
    assert result.blueprint["theme"]["textColor"] == "#0A141E"
    assert result.blueprint["theme"]["accentColor"] == "#112233"
    assert result.blueprint["theme"]["typography"]["bodyFontFamily"] == "Aptos Display"
    assert any(
        element["type"] == "text"
        and element["role"] == "title"
        and element["props"]["text"] == "Original Title"
        and element["props"]["fontFamily"] == "Aptos Display"
        and element["props"]["color"] == "#0A141E"
        for element in elements
    )
    assert any(
        element["type"] == "rect" and element["props"]["fill"] == "#112233"
        for element in elements
    )
    assert any(
        element["type"] == "image" and element["props"]["src"] == "asset:image_1"
        for element in elements
    )
    title_slot = next(
        slot
        for slot in result.template_blueprint["slides"][0]["slots"]
        if slot["elementId"].endswith("_text")
    )
    title_source = next(
        source
        for source in result.template_blueprint["slides"][0]["elementSources"]
        if source["elementId"] == title_slot["elementId"]
    )
    image_slot = next(
        slot
        for slot in result.template_blueprint["slides"][0]["slots"]
        if slot["elementId"].endswith("_image")
    )
    image_source = next(
        source
        for source in result.template_blueprint["slides"][0]["elementSources"]
        if source["elementId"] == image_slot["elementId"]
    )
    assert result.template_blueprint["sourcePackageFileId"] == "file_design"
    assert result.template_blueprint["currentPackageFileId"] == "file_design"
    assert result.template_blueprint["ooxmlSyncedDeckVersion"] == 1
    assert title_slot["usage"] == "fixed-text"
    assert title_slot["replaceMode"] == "preserve"
    assert title_slot["confidence"] < 0.5
    assert title_source["slidePart"] == "ppt/slides/slide1.xml"
    assert title_source["writable"] is True
    assert image_slot["usage"] == "media-slot"
    assert image_slot["replaceMode"] == "preserve"
    assert image_source["relationshipId"].startswith("rId")
    assert result.quality_report["metrics"]["pixelSimilarity"] is None
    assert result.quality_report["slideReports"][0]["status"] == "not_evaluated"
    assert result.quality_report["compositeScore"] <= 100
    assert result.assets[0].mime_type == "image/png"


def test_import_pptx_design_marks_placeholders_as_replaceable_slots(
    tmp_path: Path,
) -> None:
    pptx_path = tmp_path / "placeholder.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text_frame.text = "Placeholder Title"
    slide.placeholders[1].text_frame.text = "Placeholder Subtitle"
    presentation.save(pptx_path)

    result = import_pptx_design(pptx_path, "file_design")
    slots = result.template_blueprint["slides"][0]["slots"]
    replaceable = [
        slot
        for slot in slots
        if slot["usage"] == "content-slot" and slot["replaceMode"] == "replace"
    ]

    assert len(replaceable) >= 2
    assert {slot["source"]["type"] for slot in replaceable} == {"placeholder"}
    assert any(slot["slotRole"] == "title" for slot in replaceable)


def test_assign_text_roles_uses_semantic_summary_for_plain_shapes() -> None:
    elements = [
        text_element("title", "Quarterly growth", 96, 80, 44),
        text_element("body", "Revenue improved across all regions.", 120, 260, 24),
        text_element("metric", "42%", 1200, 300, 52),
        text_element("footer", "ORBIT confidential", 100, 940, 12),
    ]

    assign_text_roles(elements, {}, slide_index=2, slide_count=8)

    assert [element["role"] for element in elements] == [
        "title",
        "body",
        "highlight",
        "caption",
    ]
    assert elements[2]["templateSlotRole"] == "metric"


def test_assign_text_roles_does_not_force_numeric_prefix_to_page_number() -> None:
    elements = [
        text_element("section", "01", 96, 120, 32),
        text_element("title", "Market overview", 220, 120, 40),
    ]

    assign_text_roles(elements, {}, slide_index=2, slide_count=8)

    assert elements[0]["role"] == "caption"
    assert elements[0]["templateSlotRole"] == "label"


def test_repeated_slide_text_stays_preserved_after_role_pass() -> None:
    slides = [
        {
            "sourceSlideIndex": index,
            "style": {"backgroundColor": "#ffffff"},
            "elements": [
                text_element(f"header_{index}", "ORBIT", 100, 60, 16),
                text_element(f"title_{index}", f"Slide {index}", 100, 180, 40),
            ],
        }
        for index in (1, 2)
    ]
    slot_sources = [
        {
            f"header_{index}": {
                "type": "slide",
                "slidePart": f"ppt/slides/slide{index}.xml",
                "shapeId": "1",
                "writable": True,
            },
            f"title_{index}": {
                "type": "slide",
                "slidePart": f"ppt/slides/slide{index}.xml",
                "shapeId": "2",
                "writable": True,
            },
        }
        for index in (1, 2)
    ]

    apply_repeated_text_roles(slides, slot_sources)
    blueprint = build_template_blueprint("file_design", slides, slot_sources)
    header_slot = next(
        slot
        for slot in blueprint["slides"][0]["slots"]
        if slot["elementId"] == "header_1"
    )

    assert header_slot["slotRole"] == "caption"
    assert header_slot["usage"] == "fixed-text"
    assert header_slot["replaceMode"] == "preserve"


def test_build_template_blueprint_preserves_semantic_slot_roles() -> None:
    slide = {
        "sourceSlideIndex": 1,
        "style": {"layout": "title-content"},
        "elements": [
            text_element("title", "Title", 100, 80, 44, role="title"),
            text_element("body", "Body", 100, 240, 24, role="body"),
            text_element("caption", "Caption", 100, 900, 14, role="caption"),
            {
                **text_element("metric", "42%", 1200, 240, 48, role="highlight"),
                "templateSlotRole": "metric",
            },
            {
                "elementId": "el_cell_1",
                "type": "rect",
                "role": "unknown",
                "x": 100,
                "y": 500,
                "width": 200,
                "height": 60,
                "props": {},
            },
            {"elementId": "chart_1", "type": "chart", "role": "chart", "x": 0, "y": 0, "width": 1, "height": 1, "props": {}},
            {"elementId": "image_1", "type": "image", "role": "image", "x": 0, "y": 0, "width": 1, "height": 1, "props": {}},
        ],
    }
    sources = {
        "image_1": {"type": "placeholder", "slidePart": "ppt/slides/slide1.xml", "shapeId": "7"},
    }

    blueprint = build_template_blueprint("file_design", [slide], [sources])
    slot_roles = {slot["elementId"]: slot["slotRole"] for slot in blueprint["slides"][0]["slots"]}

    assert blueprint["slides"][0]["slideRole"] == "metric"
    assert blueprint["slides"][0]["contentCapacity"] == "medium"
    assert slot_roles["title"] == "title"
    assert slot_roles["body"] == "body"
    assert slot_roles["caption"] == "caption"
    assert slot_roles["metric"] == "metric"
    assert slot_roles["el_cell_1"] == "table"
    assert slot_roles["chart_1"] == "chart"
    assert slot_roles["image_1"] == "image_placeholder"


def text_element(
    element_id: str,
    text: str,
    x: int,
    y: int,
    font_size: int,
    *,
    role: str = "body",
) -> dict[str, object]:
    return {
        "elementId": element_id,
        "type": "text",
        "role": role,
        "x": x,
        "y": y,
        "width": 520,
        "height": 90,
        "rotation": 0,
        "opacity": 1,
        "zIndex": 1,
        "locked": False,
        "visible": True,
        "props": {
            "text": text,
            "fontSize": font_size,
            "fontFamily": "Inter",
            "fontWeight": "normal",
            "color": "#111827",
            "align": "left",
            "verticalAlign": "top",
            "lineHeight": 1.2,
        },
    }


def test_import_pptx_design_flattens_group_shapes(tmp_path: Path) -> None:
    pptx_path = tmp_path / "group.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1),
        Inches(1),
        Inches(2),
        Inches(1),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0, 128, 128)
    textbox = slide.shapes.add_textbox(
        Inches(1.2), Inches(1.2), Inches(1.4), Inches(0.5)
    )
    textbox.text_frame.text = "Grouped text"
    slide.shapes.add_group_shape([box, textbox])
    presentation.save(pptx_path)

    result = import_pptx_design(pptx_path, "file_design")
    elements = result.blueprint["slides"][0]["elements"]

    assert not any("GROUP" in warning for warning in result.warnings)
    assert any(
        element["type"] == "rect" and element["props"]["fill"] == "#008080"
        for element in elements
    )
    text = next(element for element in elements if element["type"] == "text")
    assert text["props"]["text"] == "Grouped text"
    assert 160 <= text["x"] <= 190


def test_import_pptx_design_converts_freeform_to_custom_shape(tmp_path: Path) -> None:
    pptx_path = tmp_path / "freeform.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    builder = slide.shapes.build_freeform(Inches(1), Inches(1))
    builder.move_to(0, 0)
    builder.add_line_segments([(Inches(1), 0), (Inches(1), Inches(1))])
    freeform = builder.convert_to_shape()
    freeform.fill.solid()
    freeform.fill.fore_color.rgb = RGBColor(255, 128, 0)
    freeform.line.color.rgb = RGBColor(20, 20, 20)
    presentation.save(pptx_path)

    result = import_pptx_design(pptx_path, "file_design")
    elements = result.blueprint["slides"][0]["elements"]
    custom_shape = next(
        element for element in elements if element["type"] == "customShape"
    )
    custom_slot = next(
        slot
        for slot in result.template_blueprint["slides"][0]["slots"]
        if slot["elementId"] == custom_shape["elementId"]
    )

    assert "L" in custom_shape["props"]["pathData"]
    assert custom_shape["props"]["fill"] == "#FF8000"
    assert custom_shape["props"]["viewBoxWidth"] > 0
    assert custom_slot["usage"] == "decoration"
    assert custom_slot["replaceMode"] == "ignore"


def test_ooxml_visual_tree_converts_freeform_to_custom_shape(tmp_path: Path) -> None:
    pptx_path = tmp_path / "ooxml-freeform.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    builder = slide.shapes.build_freeform(Inches(1), Inches(2))
    builder.move_to(0, 0)
    builder.add_line_segments(
        [(Inches(1.2), Inches(0.2)), (Inches(2.2), Inches(1.4)), (0, Inches(1.8))]
    )
    freeform = builder.convert_to_shape()
    freeform.fill.solid()
    freeform.fill.fore_color.rgb = RGBColor(124, 58, 237)
    freeform.line.color.rgb = RGBColor(17, 24, 39)
    presentation.save(pptx_path)

    result = import_pptx_ooxml_visual_tree(pptx_path, "file_design")
    elements = result.blueprint["slides"][0]["elements"]
    fallback_images = [
        element
        for element in elements
        if element["type"] == "image"
        and str(element["props"]["src"]).startswith("asset:shape_render_1_slide_")
    ]
    custom_shape = next(
        element for element in elements if element["type"] == "customShape"
    )

    assert len(fallback_images) == 0
    assert custom_shape["props"]["fill"] == "#7C3AED"
    assert custom_shape["props"]["stroke"] == "#111827"
    assert "L" in custom_shape["props"]["pathData"]
    assert custom_shape["props"]["viewBoxWidth"] > 0
    assert custom_shape["props"]["viewBoxHeight"] > 0
    assert len(custom_shape["props"]["nodes"]) >= 4
    assert not any("unsupported custom geometry" in warning for warning in result.warnings)


def test_import_pptx_design_preserves_common_preset_shape_types(
    tmp_path: Path,
) -> None:
    pptx_path = tmp_path / "preset-shapes.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    presets = [
        (MSO_SHAPE.OVAL, 0.5, 0.5),
        (MSO_SHAPE.ISOSCELES_TRIANGLE, 2.0, 0.5),
        (MSO_SHAPE.DIAMOND, 3.5, 0.5),
        (MSO_SHAPE.DONUT, 5.0, 0.5),
        (MSO_SHAPE.STAR_5_POINT, 6.5, 0.5),
        (MSO_SHAPE.RIGHT_ARROW, 8.0, 0.5),
    ]
    for preset, left, top in presets:
        shape = slide.shapes.add_shape(
            preset,
            Inches(left),
            Inches(top),
            Inches(1),
            Inches(1),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(17, 34, 51)
        shape.line.color.rgb = RGBColor(68, 85, 102)
    presentation.save(pptx_path)

    result = import_pptx_design(pptx_path, "file_design")
    element_types = {
        element["type"]
        for element in result.blueprint["slides"][0]["elements"]
        if element["role"] == "decoration"
    }
    custom_paths = [
        element["props"]["pathData"]
        for element in result.blueprint["slides"][0]["elements"]
        if element["type"] == "customShape"
    ]

    assert "ellipse" in element_types
    assert "ring" in element_types
    assert "star" in element_types
    assert "customShape" in element_types
    assert len(custom_paths) >= 3
    assert all(path.startswith("M ") for path in custom_paths)


def test_ooxml_visual_tree_importer_converts_camel_case_preset_arrow(
    tmp_path: Path,
) -> None:
    pptx_path = tmp_path / "ooxml-preset-arrow.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(1),
        Inches(1),
        Inches(2),
        Inches(1),
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(37, 99, 235)
    presentation.save(pptx_path)

    result = import_pptx_ooxml_visual_tree(pptx_path, "file_design")
    elements = result.blueprint["slides"][0]["elements"]

    assert any(element["type"] == "customShape" for element in elements)
    assert not any(
        element["type"] == "image"
        and str(element["props"]["src"]).startswith("asset:shape_render_")
        for element in elements
    )
    assert not any("unsupported preset rightArrow" in warning for warning in result.warnings)


def test_ooxml_visual_tree_keeps_complex_text_editable(
    tmp_path: Path,
) -> None:
    pptx_path = tmp_path / "complex-text.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1.5))
    textbox.text_frame.text = "First paragraph"
    second = textbox.text_frame.add_paragraph()
    second.text = "Second paragraph"
    presentation.save(pptx_path)

    result = import_pptx_ooxml_visual_tree(pptx_path, "file_design")
    elements = result.blueprint["slides"][0]["elements"]
    text = next(element for element in elements if element["type"] == "text")

    assert text["props"]["text"] == "First paragraph\nSecond paragraph"
    assert not any(
        element["type"] == "image"
        and str(element["props"]["src"]).startswith("asset:shape_render_1_slide_")
        for element in elements
    )
    assert not any("multi-paragraph text layout" in warning for warning in result.warnings)


def test_ooxml_visual_tree_keeps_pattern_shape_editable_with_separate_text(
    tmp_path: Path,
) -> None:
    pptx_path = tmp_path / "shape-with-text.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(1.5),
    )
    shape.text_frame.text = "Editable label"
    replace_shape_fill_with_pattern(shape)
    presentation.save(pptx_path)

    result = import_pptx_ooxml_visual_tree(pptx_path, "file_design")
    elements = result.blueprint["slides"][0]["elements"]

    fallback_images = [
        element
        for element in elements
        if element["type"] == "image"
        and str(element["props"]["src"]).startswith("asset:shape_render_1_slide_")
    ]
    pattern_shape = next(
        element
        for element in elements
        if element["type"] == "rect"
        and isinstance(element["props"].get("fill"), dict)
        and element["props"]["fill"].get("type") == "pattern"
    )
    text = next(
        element
        for element in elements
        if element["type"] == "text" and element["props"]["text"] == "Editable label"
    )

    assert len(fallback_images) == 0
    assert pattern_shape["props"]["fill"] == {
        "type": "pattern",
        "preset": "pct20",
        "foreground": "#111827",
        "background": "#F59E0B",
    }
    assert pattern_shape["props"]["stroke"] != "transparent"
    assert pattern_shape["props"]["strokeWidth"] > 0
    assert pattern_shape["props"]["shadow"]["blur"] > 0
    assert pattern_shape["zIndex"] < text["zIndex"]


def test_ooxml_visual_tree_keeps_picture_fill_as_editable_image(
    tmp_path: Path,
) -> None:
    pptx_path = tmp_path / "picture-fill.pptx"
    image_path = tmp_path / "fill.png"
    Image.new("RGB", (32, 32), "#47604D").save(image_path)
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    picture = slide.shapes.add_picture(
        str(image_path),
        Inches(7),
        Inches(1),
        Inches(1),
        Inches(1),
    )
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(1.5),
    )
    shape.text_frame.text = "Picture fill label"
    replace_shape_fill_with_blip(shape, embedded_blip_relationship_id(picture))
    presentation.save(pptx_path)

    result = import_pptx_ooxml_visual_tree(pptx_path, "file_design")
    elements = result.blueprint["slides"][0]["elements"]
    fallback_images = [
        element
        for element in elements
        if element["type"] == "image"
        and str(element["props"]["src"]).startswith("asset:shape_render_1_slide_")
    ]
    picture_fill = next(
        element
        for element in elements
        if str(element["elementId"]).endswith("_picture_fill")
    )
    text = next(
        element
        for element in elements
        if element["type"] == "text"
        and element["props"]["text"] == "Picture fill label"
    )

    assert len(fallback_images) == 0
    assert picture_fill["type"] == "image"
    assert str(picture_fill["props"]["src"]).startswith("asset:image_")
    assert picture_fill["zIndex"] < text["zIndex"]


def test_ooxml_visual_tree_keeps_group_visual_children_editable(
    tmp_path: Path,
) -> None:
    pptx_path = tmp_path / "group-visual.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    box_a = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1),
        Inches(1),
        Inches(1),
        Inches(1),
    )
    box_a.fill.solid()
    box_a.fill.fore_color.rgb = RGBColor(245, 158, 11)
    box_b = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(2.2),
        Inches(1),
        Inches(1),
        Inches(1),
    )
    box_b.fill.solid()
    box_b.fill.fore_color.rgb = RGBColor(37, 99, 235)
    label = slide.shapes.add_textbox(Inches(1.1), Inches(2.2), Inches(2), Inches(0.5))
    label.text_frame.text = "Group text"
    slide.shapes.add_group_shape([box_a, box_b, label])
    presentation.save(pptx_path)

    result = import_pptx_ooxml_visual_tree(pptx_path, "file_design")
    elements = result.blueprint["slides"][0]["elements"]
    fallback_images = [
        element
        for element in elements
        if element["type"] == "image"
        and str(element["props"]["src"]).startswith("asset:shape_render_1_slide_")
    ]

    assert len(fallback_images) == 0
    assert any(element["type"] == "rect" for element in elements)
    assert any(element["type"] == "ellipse" for element in elements)
    assert any(
        element["type"] == "text" and element["props"]["text"] == "Group text"
        for element in elements
    )


def test_ooxml_visual_tree_keeps_supported_icon_group_shapes_editable(
    tmp_path: Path,
) -> None:
    pptx_path = tmp_path / "supported-icon-group.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(1),
        Inches(2),
        Inches(1),
        Inches(1),
    )
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(1.9),
        Inches(2.1),
        Inches(2),
        Inches(0.8),
    )
    label = slide.shapes.add_textbox(Inches(1.2), Inches(3.2), Inches(3), Inches(0.6))
    label.text_frame.text = "Grouped label"
    slide.shapes.add_group_shape([circle, arrow, label])
    presentation.save(pptx_path)

    result = import_pptx_ooxml_visual_tree(pptx_path, "file_design")
    elements = result.blueprint["slides"][0]["elements"]
    fallback_images = [
        element
        for element in elements
        if element["type"] == "image"
        and str(element["props"]["src"]).startswith("asset:shape_render_1_slide_")
    ]
    text = next(
        element
        for element in elements
        if element["type"] == "text" and element["props"]["text"] == "Grouped label"
    )

    assert len(fallback_images) == 0
    assert any(element["type"] == "ellipse" for element in elements)
    assert any(element["type"] == "customShape" for element in elements)
    assert text["zIndex"] > 0


def test_blip_fill_asset_extracts_embedded_image_and_color() -> None:
    image_buffer = BytesIO()
    Image.new("RGB", (4, 4), "#47604D").save(image_buffer, format="PNG")
    image_blob = image_buffer.getvalue()

    class FakeImagePart:
        content_type = "image/png"

        def __init__(self, blob: bytes) -> None:
            self.blob = blob

    class FakePart:
        def related_part(self, relationship_id: str) -> FakeImagePart:
            assert relationship_id == "rId2"
            return FakeImagePart(image_blob)

    class FakeShape:
        _element = parse_xml(
            """
            <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
                  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
              <p:spPr>
                <a:blipFill>
                  <a:blip r:embed="rId2"/>
                </a:blipFill>
              </p:spPr>
            </p:sp>
            """
        )
        part = FakePart()

    extracted = blip_fill_asset(FakeShape(), "image_1")

    assert extracted is not None
    asset, color = extracted
    assert asset.asset_id == "image_1"
    assert asset.mime_type == "image/png"
    assert color == "#47604D"


def test_import_pptx_design_expands_table_cells(tmp_path: Path) -> None:
    pptx_path = tmp_path / "table.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = slide.shapes.add_table(
        2, 2, Inches(1), Inches(1), Inches(4), Inches(2)
    ).table
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    table.cell(1, 0).text = "C"
    table.cell(1, 1).text = "D"
    presentation.save(pptx_path)

    result = import_pptx_design(pptx_path, "file_design")
    elements = result.blueprint["slides"][0]["elements"]
    cell_rects = [
        element
        for element in elements
        if element["type"] == "rect" and "_cell_" in element["elementId"]
    ]
    cell_texts = [
        element
        for element in elements
        if element["type"] == "text" and "_cell_" in element["elementId"]
    ]

    assert len(cell_rects) == 4
    assert [element["props"]["text"] for element in cell_texts] == ["A", "B", "C", "D"]
    table_slots = [
        slot
        for slot in result.template_blueprint["slides"][0]["slots"]
        if slot["source"]["type"] == "table"
    ]
    assert table_slots
    assert all(slot["replaceMode"] in {"ignore", "preserve"} for slot in table_slots)


def test_import_pptx_design_imports_layout_decorations(tmp_path: Path) -> None:
    pptx_path = tmp_path / "layout.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    layout = presentation.slide_layouts[6]
    layout.shapes._spTree.add_autoshape(
        99,
        "Decoration 1",
        "rect",
        Inches(0.2),
        Inches(0.2),
        Inches(1),
        Inches(0.25),
    )
    decoration = layout.shapes[-1]
    decoration.fill.solid()
    decoration.fill.fore_color.rgb = RGBColor(34, 197, 94)
    slide = presentation.slides.add_slide(layout)
    slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(4), Inches(1)
    ).text_frame.text = "Slide"
    presentation.save(pptx_path)

    result = import_pptx_design(pptx_path, "file_design")
    elements = result.blueprint["slides"][0]["elements"]
    layout_decoration = next(
        element
        for element in elements
        if element["type"] == "rect" and element["props"]["fill"] == "#22C55E"
    )

    assert layout_decoration["zIndex"] < 100
    assert layout_decoration["locked"] is True
    layout_slot = next(
        slot
        for slot in result.template_blueprint["slides"][0]["slots"]
        if slot["elementId"] == layout_decoration["elementId"]
    )
    assert layout_slot["usage"] == "decoration"
    assert layout_slot["source"]["type"] == "layout"


def test_ooxml_visual_tree_importer_preserves_vector_props(
    tmp_path: Path,
) -> None:
    pptx_path = tmp_path / "ooxml-vector.pptx"
    image_path = tmp_path / "sample.png"
    Image.new("RGB", (64, 64), "#336699").save(image_path)

    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    layout = presentation.slide_layouts[6]
    layout.shapes._spTree.add_autoshape(
        91,
        "Layout Decoration",
        "rect",
        Inches(0.2),
        Inches(0.2),
        Inches(1),
        Inches(0.25),
    )
    layout_decoration = layout.shapes[-1]
    layout_decoration.fill.solid()
    layout_decoration.fill.fore_color.rgb = RGBColor(34, 197, 94)

    slide = presentation.slides.add_slide(layout)
    textbox = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(5), Inches(1))
    paragraph = textbox.text_frame.paragraphs[0]
    first_run = paragraph.add_run()
    first_run.text = "Hello "
    first_run.font.name = "Aptos"
    first_run.font.size = Pt(36)
    first_run.font.bold = True
    first_run.font.color.rgb = RGBColor(17, 24, 39)
    second_run = paragraph.add_run()
    second_run.text = "World"
    second_run.font.name = "Aptos"
    second_run.font.size = Pt(28)
    second_run.font.color.rgb = RGBColor(37, 99, 235)

    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1),
        Inches(2.4),
        Inches(3),
        Inches(1.2),
    )
    box.line.color.rgb = RGBColor(17, 24, 39)
    replace_shape_fill_with_gradient(box)

    picture = slide.shapes.add_picture(
        str(image_path),
        Inches(5),
        Inches(2.4),
        Inches(2),
        Inches(2),
    )
    picture._element.blipFill.insert(
        0,
        parse_xml(
            """
            <a:srcRect xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
                       l="10000" t="5000" r="20000" b="15000"/>
            """
        ),
    )

    group_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(8),
        Inches(3),
        Inches(1.5),
        Inches(0.8),
    )
    group_box.fill.solid()
    group_box.fill.fore_color.rgb = RGBColor(245, 158, 11)
    group_text = slide.shapes.add_textbox(
        Inches(8.1),
        Inches(3.1),
        Inches(1.2),
        Inches(0.4),
    )
    group_text.text_frame.text = "Grouped vector"
    slide.shapes.add_group_shape([group_box, group_text])

    presentation.save(pptx_path)

    result = import_pptx_ooxml_visual_tree(pptx_path, "file_design")
    elements = result.blueprint["slides"][0]["elements"]

    text = next(
        element
        for element in elements
        if element["type"] == "text" and element["props"]["text"] == "Hello World"
    )
    assert text["props"]["runs"][0]["fontWeight"] == "bold"
    assert text["props"]["runs"][1]["fontSize"] == 56
    assert text["props"]["runs"][1]["color"] == "#2563EB"
    assert text["props"]["paragraphs"][0]["text"] == "Hello World"
    assert text["props"]["paragraphs"][0]["runs"][0]["fontWeight"] == "bold"
    assert text["props"]["paragraphs"][0]["runs"][1]["color"] == "#2563EB"
    assert text["props"]["bodyInset"] == {
        "left": 14,
        "right": 14,
        "top": 7,
        "bottom": 7,
    }

    gradient_shape = next(
        element
        for element in elements
        if isinstance(element.get("props", {}).get("fill"), dict)
    )
    assert gradient_shape["props"]["fill"]["type"] == "linear-gradient"
    assert gradient_shape["props"]["fill"]["stops"][1]["color"] == "#7C3AED"

    image = next(element for element in elements if element["type"] == "image")
    assert image["props"]["crop"] == {
        "left": 0.1,
        "top": 0.05,
        "right": 0.2,
        "bottom": 0.15,
    }

    layout_source = next(
        source
        for source in result.template_blueprint["slides"][0]["elementSources"]
        if source["sourceType"] == "layout"
    )
    assert layout_source["writable"] is False
    assert any(
        element["type"] == "text"
        and element["props"]["text"] == "Grouped vector"
        and element["x"] > 1000
        for element in elements
    )


def test_ooxml_visual_tree_importer_resolves_theme_scheme_colors(
    tmp_path: Path,
) -> None:
    pptx_path = tmp_path / "theme-colors.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    accent_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1),
        Inches(1),
        Inches(2),
        Inches(1),
    )
    replace_shape_fill_with_scheme(accent_shape, "accent4")

    transformed_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4),
        Inches(1),
        Inches(2),
        Inches(1),
    )
    replace_shape_fill_with_scheme(
        transformed_shape,
        "accent5",
        '<a:lumMod val="50000"/>',
    )
    presentation.save(pptx_path)
    replace_theme_colors(
        pptx_path,
        {
            "accent4": "FFEDA9",
            "accent5": "808080",
        },
    )

    result = import_pptx_ooxml_visual_tree(pptx_path, "file_design")
    fills = {
        element["props"]["fill"]
        for element in result.blueprint["slides"][0]["elements"]
        if element["type"] == "rect"
    }

    assert "#FFEDA9" in fills
    assert "#404040" in fills
    assert "#10B981" not in fills


def test_ooxml_visual_tree_importer_honors_hidden_master_shapes(
    tmp_path: Path,
) -> None:
    pptx_path = tmp_path / "hidden-master.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    master = presentation.slide_master
    master.shapes._spTree.add_autoshape(
        100,
        "Master Decoration",
        "rect",
        Inches(0.25),
        Inches(0.25),
        Inches(1),
        Inches(0.5),
    )
    decoration = master.shapes[-1]
    decoration.fill.solid()
    decoration.fill.fore_color.rgb = RGBColor(34, 51, 68)

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(
        Inches(1),
        Inches(1),
        Inches(3),
        Inches(1),
    ).text_frame.text = "Slide only"
    presentation.save(pptx_path)
    set_slide_show_master_shapes(pptx_path, 1, False)

    result = import_pptx_ooxml_visual_tree(pptx_path, "file_design")
    fills = [
        element["props"]["fill"]
        for element in result.blueprint["slides"][0]["elements"]
        if element["type"] == "rect"
    ]

    assert "#223344" not in fills


def test_ooxml_visual_tree_importer_preserves_text_box_geometry(
    tmp_path: Path,
) -> None:
    pptx_path = tmp_path / "text-geometry.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    textbox = slide.shapes.add_textbox(
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(2),
    )
    textbox.text_frame.text = "Inset text"
    body_pr = textbox._element.txBody.bodyPr
    body_pr.set("lIns", str(int(Inches(0.5))))
    body_pr.set("rIns", str(int(Inches(0.25))))
    body_pr.set("tIns", str(int(Inches(0.25))))
    body_pr.set("bIns", str(int(Inches(0.25))))
    body_pr.set("anchor", "mid")
    rotated_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6),
        Inches(1),
        Inches(2),
        Inches(1),
    )
    rotated_box.rotation = 30
    rotated_box.fill.solid()
    rotated_box.fill.fore_color.rgb = RGBColor(238, 68, 68)
    presentation.save(pptx_path)
    set_first_text_paragraph_line_spacing(pptx_path, 1, 130000)

    result = import_pptx_ooxml_visual_tree(pptx_path, "file_design")
    elements = result.blueprint["slides"][0]["elements"]
    text = next(
        element
        for element in elements
        if element["type"] == "text" and element["props"]["text"] == "Inset text"
    )
    shape = next(
        element
        for element in elements
        if element["type"] == "rect" and element["props"]["fill"] == "#EE4444"
    )

    assert text["x"] == 216
    assert text["y"] == 180
    assert text["width"] == 468
    assert text["height"] == 216
    assert text["props"]["verticalAlign"] == "middle"
    assert text["props"]["lineHeight"] == 1.3
    assert text["props"]["paragraphs"][0]["lineHeight"] == 1.3
    assert text["props"]["bodyInset"] == {
        "left": 72,
        "right": 36,
        "top": 36,
        "bottom": 36,
    }
    assert shape["rotation"] == 30


def test_ooxml_visual_tree_importer_preserves_vertical_text_mode(
    tmp_path: Path,
) -> None:
    pptx_path = tmp_path / "vertical-text.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(
        Inches(1),
        Inches(1.7),
        Inches(1),
        Inches(4.8),
    )
    textbox.text_frame.text = "VERTICAL"
    textbox._element.txBody.bodyPr.set("vert", "vert270")
    presentation.save(pptx_path)

    result = import_pptx_ooxml_visual_tree(pptx_path, "file_design")
    text = next(
        element
        for element in result.blueprint["slides"][0]["elements"]
        if element["type"] == "text" and element["props"]["text"] == "VERTICAL"
    )

    assert text["props"]["writingMode"] == "vertical-270"


def test_ooxml_visual_tree_importer_falls_back_graphic_frames(
    tmp_path: Path,
) -> None:
    pptx_path = tmp_path / "graphic-frame.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table_shape = slide.shapes.add_table(
        2,
        2,
        Inches(1),
        Inches(1),
        Inches(3),
        Inches(1.2),
    )
    table_shape.table.cell(0, 0).text = "A"
    table_shape.table.cell(1, 1).text = "B"
    chart_data = ChartData()
    chart_data.categories = ["A", "B"]
    chart_data.add_series("Series 1", (1, 2))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(5),
        Inches(1),
        Inches(3),
        Inches(2),
        chart_data,
    )
    presentation.save(pptx_path)

    result = import_pptx_ooxml_visual_tree(pptx_path, "file_design")
    elements = result.blueprint["slides"][0]["elements"]
    fallback_images = [
        element
        for element in elements
        if element["type"] == "image"
        and str(element["props"]["src"]).startswith("asset:shape_render_")
    ]
    table = next(element for element in elements if element["type"] == "table")
    chart = next(element for element in elements if element["type"] == "chart")
    chart_slot = next(
        slot
        for slot in result.template_blueprint["slides"][0]["slots"]
        if slot["elementId"] == chart["elementId"]
    )

    assert len(fallback_images) == 0
    assert table["props"]["rows"][0][0]["text"] == "A"
    assert table["props"]["rows"][0][0]["fill"] == "#4F81BD"
    assert table["props"]["rows"][0][0]["textColor"] == "#FFFFFF"
    assert table["props"]["rows"][1][1]["text"] == "B"
    assert chart["props"]["type"] == "bar"
    assert chart["props"]["data"] == [
        {"label": "A", "value": 1.0},
        {"label": "B", "value": 2.0},
    ]
    assert chart["props"]["style"]["colors"][0] == "#4F81BD"
    assert chart_slot["slotRole"] == "chart"
    assert chart_slot["source"]["type"] == "unknown"
    assert not any(
        "OOXML graphicFrame rendered as image fallback on slide 1: table"
        in warning
        for warning in result.warnings
    )
    assert not any(
        "OOXML graphicFrame rendered as image fallback on slide 1: chart"
        in warning
        for warning in result.warnings
    )


def test_ooxml_visual_tree_importer_preserves_svg_as_editable_media(
    tmp_path: Path,
) -> None:
    pptx_path = tmp_path / "svg-media.pptx"
    image_path = tmp_path / "placeholder.png"
    Image.new("RGB", (32, 32), "#2563EB").save(image_path)
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        str(image_path),
        Inches(1),
        Inches(1),
        Inches(2),
        Inches(2),
    )
    presentation.save(pptx_path)
    replace_first_picture_with_svg(pptx_path)

    result = import_pptx_ooxml_visual_tree(pptx_path, "file_design")
    elements = result.blueprint["slides"][0]["elements"]
    svg = next(element for element in elements if element["type"] == "svg")
    svg_slot = next(
        slot
        for slot in result.template_blueprint["slides"][0]["slots"]
        if slot["elementId"] == svg["elementId"]
    )

    assert svg["role"] == "media"
    assert svg["props"]["src"] == "asset:image_1"
    assert result.assets[0].file_name == "image_1.svg"
    assert result.assets[0].mime_type == "image/svg+xml"
    assert svg_slot["slotRole"] == "image"
    assert svg_slot["source"]["relationshipId"].startswith("rId")


def test_quality_report_counts_supported_deck_elements_as_editable() -> None:
    report = build_quality_report(
        [
            {
                "elements": [
                    editable_element("el_rect", "rect"),
                    editable_element("el_table", "table"),
                    editable_element("el_svg", "svg"),
                ]
            }
        ],
        [],
    )

    assert report["editabilityCoverage"] == 1
    assert report["metrics"]["editability"] == 100


def test_ooxml_visual_tree_importer_is_default(
    tmp_path: Path,
    monkeypatch: object,
) -> None:
    pptx_path = tmp_path / "default-vector.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(
        Inches(1),
        Inches(1),
        Inches(3),
        Inches(1),
    ).text_frame.text = "Default vector"
    presentation.save(pptx_path)

    monkeypatch.delenv(VECTOR_IMPORT_FLAG, raising=False)
    default_result = import_pptx_design_with_optional_ooxml_vector(
        pptx_path,
        "file_design",
    )
    assert any(
        str(element["elementId"]).startswith("el_ooxml_")
        for element in default_result.blueprint["slides"][0]["elements"]
    )

    monkeypatch.setenv(VECTOR_IMPORT_FLAG, "false")
    fallback_result = import_pptx_design_with_optional_ooxml_vector(
        pptx_path,
        "file_design",
    )
    assert any(
        str(element["elementId"]).startswith("el_imported_")
        for element in fallback_result.blueprint["slides"][0]["elements"]
    )


def replace_shape_fill_with_scheme(
    shape: object,
    scheme: str,
    transforms: str = "",
) -> None:
    sp_pr = shape._element.spPr
    for child in list(sp_pr):
        if child.tag.rsplit("}", maxsplit=1)[-1] in {"solidFill", "gradFill", "noFill"}:
            sp_pr.remove(child)
    sp_pr.insert(
        0,
        parse_xml(
            f"""
            <a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <a:schemeClr val="{scheme}">
                {transforms}
              </a:schemeClr>
            </a:solidFill>
            """
        ),
    )


def replace_theme_colors(pptx_path: Path, colors: dict[str, str]) -> None:
    with zipfile.ZipFile(pptx_path, "r") as package:
        entries = {item.filename: package.read(item.filename) for item in package.infolist()}

    theme_path = next(
        name
        for name in entries
        if name.startswith("ppt/theme/theme") and name.endswith(".xml")
    )
    root = ET.fromstring(entries[theme_path])
    for item in root.iter():
        key = item.tag.rsplit("}", maxsplit=1)[-1]
        if key not in colors:
            continue
        item.clear()
        ET.SubElement(
            item,
            "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr",
            {"val": colors[key]},
        )
    entries[theme_path] = ET.tostring(root, encoding="utf-8", xml_declaration=True)

    with zipfile.ZipFile(pptx_path, "w") as package:
        for name, content in entries.items():
            package.writestr(name, content)


def set_slide_show_master_shapes(
    pptx_path: Path,
    slide_index: int,
    show: bool,
) -> None:
    with zipfile.ZipFile(pptx_path, "r") as package:
        entries = {item.filename: package.read(item.filename) for item in package.infolist()}

    slide_path = f"ppt/slides/slide{slide_index}.xml"
    root = ET.fromstring(entries[slide_path])
    common_slide_data = next(
        item
        for item in root.iter()
        if item.tag.rsplit("}", maxsplit=1)[-1] == "cSld"
    )
    common_slide_data.set("showMasterSp", "1" if show else "0")
    entries[slide_path] = ET.tostring(root, encoding="utf-8", xml_declaration=True)

    with zipfile.ZipFile(pptx_path, "w") as package:
        for name, content in entries.items():
            package.writestr(name, content)


def set_first_text_paragraph_line_spacing(
    pptx_path: Path,
    slide_index: int,
    spacing_pct: int,
) -> None:
    with zipfile.ZipFile(pptx_path, "r") as package:
        entries = {item.filename: package.read(item.filename) for item in package.infolist()}

    slide_path = f"ppt/slides/slide{slide_index}.xml"
    root = ET.fromstring(entries[slide_path])
    paragraph = next(
        item for item in root.iter() if item.tag.rsplit("}", maxsplit=1)[-1] == "p"
    )
    namespace = paragraph.tag.rsplit("}", maxsplit=1)[0].lstrip("{")
    paragraph_props = next(
        (
            item
            for item in list(paragraph)
            if item.tag.rsplit("}", maxsplit=1)[-1] == "pPr"
        ),
        None,
    )
    if paragraph_props is None:
        paragraph_props = ET.Element(f"{{{namespace}}}pPr")
        paragraph.insert(0, paragraph_props)
    for child in list(paragraph_props):
        if child.tag.rsplit("}", maxsplit=1)[-1] == "lnSpc":
            paragraph_props.remove(child)
    line_spacing = ET.SubElement(paragraph_props, f"{{{namespace}}}lnSpc")
    ET.SubElement(line_spacing, f"{{{namespace}}}spcPct", {"val": str(spacing_pct)})
    entries[slide_path] = ET.tostring(root, encoding="utf-8", xml_declaration=True)

    with zipfile.ZipFile(pptx_path, "w") as package:
        for name, content in entries.items():
            package.writestr(name, content)


def replace_shape_fill_with_gradient(shape: object) -> None:
    sp_pr = shape._element.spPr
    for child in list(sp_pr):
        if child.tag.rsplit("}", maxsplit=1)[-1] in {"solidFill", "gradFill", "noFill"}:
            sp_pr.remove(child)
    sp_pr.insert(
        0,
        parse_xml(
            """
            <a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <a:gsLst>
                <a:gs pos="0">
                  <a:srgbClr val="2563EB"/>
                </a:gs>
                <a:gs pos="100000">
                  <a:srgbClr val="7C3AED">
                    <a:alpha val="75000"/>
                  </a:srgbClr>
                </a:gs>
              </a:gsLst>
              <a:lin ang="5400000" scaled="1"/>
            </a:gradFill>
            """
        ),
    )


def replace_shape_fill_with_pattern(shape: object) -> None:
    sp_pr = shape._element.spPr
    for child in list(sp_pr):
        if child.tag.rsplit("}", maxsplit=1)[-1] in {
            "solidFill",
            "gradFill",
            "noFill",
            "pattFill",
        }:
            sp_pr.remove(child)
    sp_pr.insert(
        0,
        parse_xml(
            """
            <a:pattFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
                        prst="pct20">
              <a:fgClr><a:srgbClr val="111827"/></a:fgClr>
              <a:bgClr><a:srgbClr val="F59E0B"/></a:bgClr>
            </a:pattFill>
            """
        ),
    )


def replace_shape_fill_with_blip(shape: object, relationship_id: str) -> None:
    sp_pr = shape._element.spPr
    for child in list(sp_pr):
        if child.tag.rsplit("}", maxsplit=1)[-1] in {
            "solidFill",
            "gradFill",
            "noFill",
            "pattFill",
            "blipFill",
        }:
            sp_pr.remove(child)
    sp_pr.insert(
        0,
        parse_xml(
            f"""
            <a:blipFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
                        xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
              <a:blip r:embed="{relationship_id}"/>
            </a:blipFill>
            """
        ),
    )


def embedded_blip_relationship_id(shape: object) -> str:
    for node in shape._element.iter():
        if node.tag.rsplit("}", maxsplit=1)[-1] != "blip":
            continue
        for key, value in node.attrib.items():
            if key.rsplit("}", maxsplit=1)[-1] == "embed":
                return str(value)
    raise AssertionError("shape has no embedded blip relationship")


def replace_first_picture_with_svg(pptx_path: Path) -> None:
    with zipfile.ZipFile(pptx_path, "r") as package:
        entries = {
            info.filename: package.read(info.filename)
            for info in package.infolist()
        }

    rels_path = "ppt/slides/_rels/slide1.xml.rels"
    rels_root = ET.fromstring(entries[rels_path])
    for relationship in rels_root:
        if str(relationship.get("Type", "")).endswith("/image"):
            relationship.set("Target", "../media/image1.svg")
            break
    entries[rels_path] = ET.tostring(
        rels_root,
        encoding="utf-8",
        xml_declaration=True,
    )

    content_types_root = ET.fromstring(entries["[Content_Types].xml"])
    namespace = str(content_types_root.tag).split("}", maxsplit=1)[0].strip("{")
    has_svg_default = any(
        child.get("Extension") == "svg"
        for child in list(content_types_root)
        if child.tag.endswith("Default")
    )
    if not has_svg_default:
        ET.SubElement(
            content_types_root,
            f"{{{namespace}}}Default",
            Extension="svg",
            ContentType="image/svg+xml",
        )
    entries["[Content_Types].xml"] = ET.tostring(
        content_types_root,
        encoding="utf-8",
        xml_declaration=True,
    )
    entries["ppt/media/image1.svg"] = (
        b'<svg xmlns="http://www.w3.org/2000/svg" width="32" height="32">'
        b'<rect width="32" height="32" fill="#2563EB"/>'
        b"</svg>"
    )

    rewritten_path = pptx_path.with_suffix(".rewritten.pptx")
    with zipfile.ZipFile(rewritten_path, "w", zipfile.ZIP_DEFLATED) as package:
        for filename, content in entries.items():
            package.writestr(filename, content)
    rewritten_path.replace(pptx_path)


def editable_element(element_id: str, element_type: str) -> dict[str, object]:
    return {
        "elementId": element_id,
        "type": element_type,
        "role": "media",
        "x": 10,
        "y": 10,
        "width": 100,
        "height": 80,
        "locked": False,
        "visible": True,
    }
