import base64
import copy
from io import BytesIO
from pathlib import Path
from typing import Any
import zipfile
from xml.etree import ElementTree as ET

import pytest
from pptx import Presentation
from pptx.util import Inches

from app.ai.pptx_ooxml_generation import (
    PML_NS,
    REL_NS,
    generate_pptx_ooxml,
    resolve_relationship_part,
    sync_pptx_ooxml,
)


def test_reorders_slide_id_list_without_changing_relationship_identity(
    tmp_path: Path,
) -> None:
    source_path = three_slide_pptx(tmp_path)
    generated = generate_pptx_ooxml(source_path, "file_reorder", render=False)
    operation = reorder_operation(generated, [2, 0, 1])
    original_bytes = source_path.read_bytes()

    result = sync_pptx_ooxml(
        source_path,
        template_blueprint=generated.template_blueprint,
        operations=[operation],
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
    )
    reordered_bytes = current_package_bytes(result)

    assert [item.operation_type for item in result.applied_operations] == [
        "reorder_slides"
    ]
    assert result.unsupported_operations == []
    before = slide_id_signatures(original_bytes)
    after = slide_id_signatures(reordered_bytes)
    assert after == [before[2], before[0], before[1]]
    assert {signature[:2] for signature in after} == {
        signature[:2] for signature in before
    }
    assert_package_entries_unchanged_except_presentation(
        original_bytes,
        reordered_bytes,
    )
    assert presentation_slide_texts(reordered_bytes) == [
        "Slide 3",
        "Slide 1",
        "Slide 2",
    ]

    round_trip_path = tmp_path / "round-trip.pptx"
    round_trip_path.write_bytes(reordered_bytes)
    round_trip = generate_pptx_ooxml(
        round_trip_path,
        "file_round_trip",
        render=False,
    )
    assert imported_slide_texts(round_trip.blueprint) == [
        "Slide 3",
        "Slide 1",
        "Slide 2",
    ]


def test_reorders_opaque_slide_ids_from_explicit_source_mapping(
    tmp_path: Path,
) -> None:
    source_path = three_slide_pptx(tmp_path)
    generated = generate_pptx_ooxml(source_path, "file_reorder", render=False)
    blueprint = copy.deepcopy(generated.template_blueprint)
    slide_ids = ["slide_cover", "slide_metrics", "slide_close"]
    for slide, slide_id in zip(blueprint["slides"], slide_ids, strict=True):
        slide["slideId"] = slide_id

    result = sync_pptx_ooxml(
        source_path,
        template_blueprint=blueprint,
        operations=[reorder_operation_from_blueprint(blueprint, [2, 0, 1])],
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
    )

    assert [item.operation_type for item in result.applied_operations] == [
        "reorder_slides"
    ]
    assert result.unsupported_operations == []
    assert presentation_slide_texts(current_package_bytes(result)) == [
        "Slide 3",
        "Slide 1",
        "Slide 2",
    ]


@pytest.mark.parametrize(
    ("mutation", "expected_reason"),
    [
        ("duplicate-id", "SLIDE_REORDER_PERMUTATION_INVALID"),
        ("missing-slide", "SLIDE_REORDER_PERMUTATION_INVALID"),
        ("unknown-id", "SLIDE_REORDER_LOCATOR_UNSAFE"),
        ("duplicate-locator", "SLIDE_REORDER_LOCATOR_UNSAFE"),
    ],
)
def test_invalid_reorder_returns_original_package(
    tmp_path: Path,
    mutation: str,
    expected_reason: str,
) -> None:
    source_path = three_slide_pptx(tmp_path)
    generated = generate_pptx_ooxml(source_path, "file_reorder", render=False)
    operation = reorder_operation(generated, [2, 0, 1])
    slide_orders = operation["slideOrders"]
    assert isinstance(slide_orders, list)
    if mutation == "duplicate-id":
        slide_orders[1]["slideId"] = slide_orders[0]["slideId"]
    elif mutation == "missing-slide":
        slide_orders.pop()
    elif mutation == "unknown-id":
        slide_orders[0]["slideId"] = "slide_unknown_99"
    else:
        slide_orders[1]["sourceSlidePart"] = slide_orders[0]["sourceSlidePart"]

    result = sync_pptx_ooxml(
        source_path,
        template_blueprint=generated.template_blueprint,
        operations=[operation],
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
    )

    assert current_package_bytes(result) == source_path.read_bytes()
    assert result.applied_operations == []
    assert [item.reason_code for item in result.unsupported_operations] == [
        expected_reason
    ]


def test_incomplete_blueprint_locator_returns_original_package(
    tmp_path: Path,
) -> None:
    source_path = three_slide_pptx(tmp_path)
    generated = generate_pptx_ooxml(source_path, "file_reorder", render=False)
    blueprint = copy.deepcopy(generated.template_blueprint)
    del blueprint["slides"][1]["sourceSlidePart"]

    result = sync_pptx_ooxml(
        source_path,
        template_blueprint=blueprint,
        operations=[reorder_operation(generated, [2, 0, 1])],
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
    )

    assert current_package_bytes(result) == source_path.read_bytes()
    assert result.applied_operations == []
    assert [item.reason_code for item in result.unsupported_operations] == [
        "SLIDE_REORDER_LOCATOR_UNSAFE"
    ]


def test_broken_presentation_relationship_returns_original_package(
    tmp_path: Path,
) -> None:
    source_path = three_slide_pptx(tmp_path)
    generated = generate_pptx_ooxml(source_path, "file_reorder", render=False)
    broken_bytes = without_first_presentation_relationship(source_path.read_bytes())
    broken_path = tmp_path / "broken.pptx"
    broken_path.write_bytes(broken_bytes)

    result = sync_pptx_ooxml(
        broken_path,
        template_blueprint=generated.template_blueprint,
        operations=[reorder_operation(generated, [2, 0, 1])],
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
    )

    assert current_package_bytes(result) == broken_bytes
    assert result.applied_operations == []
    assert [item.reason_code for item in result.unsupported_operations] == [
        "SLIDE_REORDER_RELATIONSHIP_UNSAFE"
    ]


def three_slide_pptx(tmp_path: Path) -> Path:
    presentation = Presentation()
    for index in range(1, 4):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text = f"Slide {index}"
    path = tmp_path / "three-slides.pptx"
    presentation.save(path)
    return path


def reorder_operation(generated: object, indices: list[int]) -> dict[str, object]:
    template_blueprint = getattr(generated, "template_blueprint")
    return reorder_operation_from_blueprint(template_blueprint, indices)


def reorder_operation_from_blueprint(
    template_blueprint: dict[str, Any],
    indices: list[int],
) -> dict[str, object]:
    return {
        "type": "reorder_slides",
        "slideOrders": [
            {
                "slideId": template_blueprint["slides"][index]["slideId"],
                "order": order,
                "sourceSlidePart": template_blueprint["slides"][index][
                    "sourceSlidePart"
                ],
            }
            for order, index in enumerate(indices, start=1)
        ],
    }


def current_package_bytes(result: object) -> bytes:
    assets = getattr(result, "assets")
    asset = next(item for item in assets if item.asset_id == "current_package")
    return base64.b64decode(asset.content_base64)


def slide_id_signatures(package_bytes: bytes) -> list[tuple[str, str, str]]:
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as package:
        presentation = ET.fromstring(package.read("ppt/presentation.xml"))
        relationships = ET.fromstring(
            package.read("ppt/_rels/presentation.xml.rels")
        )
    relationship_targets = {
        str(item.get("Id", "")): resolve_relationship_part(
            "ppt/presentation.xml",
            str(item.get("Target", "")),
        )
        for item in relationships
    }
    slide_id_list = presentation.find(f"{{{PML_NS}}}sldIdLst")
    assert slide_id_list is not None
    return [
        (
            str(item.get("id", "")),
            str(item.get(f"{{{REL_NS}}}id", "")),
            relationship_targets[str(item.get(f"{{{REL_NS}}}id", ""))],
        )
        for item in slide_id_list
    ]


def assert_package_entries_unchanged_except_presentation(
    before: bytes,
    after: bytes,
) -> None:
    with (
        zipfile.ZipFile(BytesIO(before), "r") as before_zip,
        zipfile.ZipFile(BytesIO(after), "r") as after_zip,
    ):
        assert before_zip.namelist() == after_zip.namelist()
        for name in before_zip.namelist():
            if name == "ppt/presentation.xml":
                assert before_zip.read(name) != after_zip.read(name)
            else:
                assert before_zip.read(name) == after_zip.read(name), name


def presentation_slide_texts(package_bytes: bytes) -> list[str]:
    presentation = Presentation(BytesIO(package_bytes))
    return [
        next(shape.text for shape in slide.shapes if hasattr(shape, "text_frame"))
        for slide in presentation.slides
    ]


def imported_slide_texts(blueprint: dict[str, object]) -> list[str]:
    slides = blueprint["slides"]
    assert isinstance(slides, list)
    values: list[str] = []
    for slide in slides:
        assert isinstance(slide, dict)
        elements = slide["elements"]
        assert isinstance(elements, list)
        text_values = [
            str(element["props"]["text"])
            for element in elements
            if isinstance(element, dict)
            and element.get("type") == "text"
            and isinstance(element.get("props"), dict)
            and str(element["props"].get("text", "")).startswith("Slide ")
        ]
        values.append(text_values[0])
    return values


def without_first_presentation_relationship(package_bytes: bytes) -> bytes:
    source = zipfile.ZipFile(BytesIO(package_bytes), "r")
    presentation = ET.fromstring(source.read("ppt/presentation.xml"))
    slide_id_list = presentation.find(f"{{{PML_NS}}}sldIdLst")
    assert slide_id_list is not None
    first_slide_relationship_id = str(
        list(slide_id_list)[0].get(f"{{{REL_NS}}}id", "")
    )
    rels_name = "ppt/_rels/presentation.xml.rels"
    rels = ET.fromstring(source.read(rels_name))
    relationship = next(
        item
        for item in rels
        if item.get("Id") == first_slide_relationship_id
    )
    rels.remove(relationship)
    output = BytesIO()
    with zipfile.ZipFile(output, "w") as target:
        for info in source.infolist():
            content = (
                ET.tostring(rels, encoding="utf-8", xml_declaration=True)
                if info.filename == rels_name
                else source.read(info.filename)
            )
            target.writestr(info, content)
    source.close()
    return output.getvalue()
