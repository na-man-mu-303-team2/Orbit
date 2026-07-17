from __future__ import annotations

import base64
from io import BytesIO
import zipfile

from pptx import Presentation
import pytest

from app.ai.pptx_design_importer import ImportedDesignAsset
from app.ai.pptx_png_zip_export import (
    PptxPngZipExportError,
    PptxPngZipExportRequest,
    export_pptx_png_zip,
)


def test_exports_every_rendered_slide_with_stable_names(monkeypatch: pytest.MonkeyPatch) -> None:
    png_one = b"\x89PNG\r\n\x1a\nfirst"
    png_two = b"\x89PNG\r\n\x1a\nsecond"

    def fake_render(_package: bytes, canvas: object) -> list[ImportedDesignAsset]:
        assert getattr(canvas, "width") == 1920
        return [
            _asset("slide_render_1", png_one),
            _asset("slide_render_2", png_two),
        ]

    monkeypatch.setattr(
        "app.ai.pptx_png_zip_export.render_pptx_to_png_assets",
        fake_render,
    )
    response = export_pptx_png_zip(
        PptxPngZipExportRequest(contentBase64=_pptx_base64())
    )

    with zipfile.ZipFile(BytesIO(base64.b64decode(response.content_base64))) as archive:
        assert archive.namelist() == ["slide-001.png", "slide-002.png"]
        assert archive.read("slide-001.png") == png_one
        assert archive.read("slide-002.png") == png_two


def test_rejects_invalid_pptx_content() -> None:
    with pytest.raises(PptxPngZipExportError, match="PPTX package is invalid"):
        export_pptx_png_zip(
            PptxPngZipExportRequest(
                contentBase64=base64.b64encode(b"not-a-pptx").decode("ascii")
            )
        )


def _pptx_base64() -> str:
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    output = BytesIO()
    presentation.save(output)
    return base64.b64encode(output.getvalue()).decode("ascii")


def _asset(asset_id: str, content: bytes) -> ImportedDesignAsset:
    return ImportedDesignAsset(
        assetId=asset_id,
        fileName=f"{asset_id}.png",
        mimeType="image/png",
        contentBase64=base64.b64encode(content).decode("ascii"),
    )
