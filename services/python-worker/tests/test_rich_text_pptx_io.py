import base64
from io import BytesIO
from pathlib import Path
from xml.etree import ElementTree as ET
from zipfile import ZipFile

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt

from app.ai.deck_pptx_export import DeckPptxExportRequest, export_deck_pptx
from app.ai.pptx_ooxml_vector_importer import import_pptx_ooxml_visual_tree


DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def test_vector_importer_records_rich_text_styles_and_stable_diagnostics(
    tmp_path: Path,
) -> None:
    pptx_path = tmp_path / "rich-text-import.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    paragraph = textbox.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.line_spacing = 1.3
    paragraph.space_before = Pt(2)
    paragraph.space_after = Pt(3)
    paragraph_properties = paragraph._p.get_or_add_pPr()
    paragraph_properties.set("marL", str(int(Inches(0.25))))
    paragraph_properties.append(
        parse_xml(
            f'<a:buChar xmlns:a="{DRAWING_NS}" char="▪"/>'
        )
    )

    run = paragraph.add_run()
    run.text = "Styled"
    run.font.name = "Aptos"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.italic = True
    run.font.underline = True
    run_properties = run._r.get_or_add_rPr()
    run_properties.set("spc", "120")
    run_properties.append(
        parse_xml(
            f'<a:hlinkClick xmlns:a="{DRAWING_NS}" xmlns:r="{REL_NS}" '
            'r:id="rId999"/>'
        )
    )
    presentation.save(pptx_path)

    result = import_pptx_ooxml_visual_tree(pptx_path, "file_design")

    text = next(
        element
        for element in result.blueprint["slides"][0]["elements"]
        if element["type"] == "text" and element["props"]["text"] == "Styled"
    )
    imported_paragraph = text["props"]["paragraphs"][0]
    imported_run = imported_paragraph["runs"][0]
    assert imported_run == {
        "text": "Styled",
        "fontFamily": "Aptos",
        "fontSize": 40,
        "fontWeight": "bold",
        "italic": True,
        "underline": True,
        "baseline": "normal",
    }
    assert imported_paragraph["align"] == "center"
    assert imported_paragraph["indent"] == 36
    assert imported_paragraph["spaceBefore"] == 4
    assert imported_paragraph["spaceAfter"] == 6
    assert imported_paragraph["lineHeight"] == 1.3
    assert imported_paragraph["bullet"] == {
        "enabled": True,
        "character": "▪",
        "indent": 36,
    }
    assert any(
        warning.startswith("PPTX_RICH_TEXT_UNSUPPORTED_HYPERLINK:")
        for warning in result.warnings
    )
    assert any(
        warning.startswith("PPTX_RICH_TEXT_UNSUPPORTED_LETTER_SPACING:")
        for warning in result.warnings
    )


def test_generic_export_round_trips_canonical_paragraphs_and_runs(
    tmp_path: Path,
) -> None:
    deck = rich_text_deck(
        {
            "text": "STALE TOP LEVEL TEXT",
            "runs": [{"text": "STALE TOP LEVEL RUN", "baseline": "normal"}],
            "fontFamily": "Fallback Sans",
            "fontSize": 24,
            "fontWeight": "normal",
            "align": "left",
            "lineHeight": 1.2,
            "verticalAlign": "top",
            "paragraphs": [
                {
                    "text": "Alpha Beta",
                    "runs": [
                        {
                            "text": "Alpha ",
                            "fontFamily": "Aptos",
                            "fontSize": 40,
                            "fontWeight": "bold",
                            "italic": True,
                            "color": "#DC2626",
                            "baseline": "normal",
                        },
                        {
                            "text": "Beta",
                            "fontFamily": "Arial",
                            "fontSize": 32,
                            "fontWeight": "normal",
                            "underline": True,
                            "color": "#2563EB",
                            "baseline": "normal",
                        },
                    ],
                    "align": "center",
                    "lineHeight": 1.35,
                    "spaceBefore": 8,
                    "spaceAfter": 10,
                    "indent": 24,
                    "bullet": {
                        "enabled": True,
                        "character": "▪",
                        "indent": 24,
                    },
                },
                {
                    "text": "Second",
                    "runs": [
                        {
                            "text": "Second",
                            "fontFamily": "Aptos",
                            "fontSize": 28,
                            "fontWeight": 600,
                            "italic": False,
                            "underline": False,
                            "color": "#111827",
                            "baseline": "normal",
                        }
                    ],
                    "align": "right",
                    "lineHeight": 1.1,
                    "spaceBefore": 0,
                    "spaceAfter": 4,
                    "indent": 0,
                },
                {
                    "text": "",
                    "runs": [],
                    "align": "left",
                    "lineHeight": 1.2,
                    "spaceBefore": 0,
                    "spaceAfter": 0,
                    "indent": 0,
                },
            ],
        }
    )

    response = export_deck_pptx(DeckPptxExportRequest(deck=deck))
    binary = base64.b64decode(response.content_base64)
    with ZipFile(BytesIO(binary)) as archive:
        slide_root = ET.fromstring(archive.read("ppt/slides/slide1.xml"))
    paragraphs = slide_root.findall(f".//{{{DRAWING_NS}}}p")
    assert len(paragraphs) == 3
    assert "".join(
        node.text or "" for node in paragraphs[0].iter(f"{{{DRAWING_NS}}}t")
    ) == "Alpha Beta"
    assert "▪" not in "".join(
        node.text or "" for node in paragraphs[0].iter(f"{{{DRAWING_NS}}}t")
    )
    first_properties = paragraphs[0].find(f"{{{DRAWING_NS}}}pPr")
    assert first_properties is not None
    assert first_properties.attrib["algn"] == "ctr"
    assert first_properties.attrib["marL"] == "152400"
    assert first_properties.find(f"{{{DRAWING_NS}}}buChar").attrib == {"char": "▪"}
    assert first_properties.find(
        f"{{{DRAWING_NS}}}lnSpc/{{{DRAWING_NS}}}spcPct"
    ).attrib == {"val": "135000"}
    assert first_properties.find(
        f"{{{DRAWING_NS}}}spcBef/{{{DRAWING_NS}}}spcPts"
    ).attrib == {"val": "400"}
    assert first_properties.find(
        f"{{{DRAWING_NS}}}spcAft/{{{DRAWING_NS}}}spcPts"
    ).attrib == {"val": "500"}

    first_run_properties = paragraphs[0].find(f"{{{DRAWING_NS}}}r/{{{DRAWING_NS}}}rPr")
    assert first_run_properties is not None
    assert first_run_properties.attrib["b"] == "1"
    assert first_run_properties.attrib["i"] == "1"
    assert first_run_properties.attrib.get("u") == "none"
    second_run_properties = paragraphs[0].findall(
        f"{{{DRAWING_NS}}}r/{{{DRAWING_NS}}}rPr"
    )[1]
    assert second_run_properties.attrib["u"] == "sng"

    exported_path = tmp_path / "rich-text-round-trip.pptx"
    exported_path.write_bytes(binary)
    imported = import_pptx_ooxml_visual_tree(exported_path, "file_round_trip")
    text = next(
        element
        for element in imported.blueprint["slides"][0]["elements"]
        if element["type"] == "text"
    )
    assert text["props"]["text"] == "Alpha Beta\nSecond\n"
    assert len(text["props"]["paragraphs"]) == 3
    first = text["props"]["paragraphs"][0]
    assert first["align"] == "center"
    assert first["lineHeight"] == 1.35
    assert first["spaceBefore"] == 8
    assert first["spaceAfter"] == 10
    assert first["indent"] == 24
    assert first["bullet"] == {
        "enabled": True,
        "character": "▪",
        "indent": 24,
    }
    assert first["runs"] == [
        {
            "text": "Alpha ",
            "fontFamily": "Aptos",
            "fontSize": 40,
            "fontWeight": "bold",
            "italic": True,
            "underline": False,
            "color": "#DC2626",
            "baseline": "normal",
        },
        {
            "text": "Beta",
            "fontFamily": "Arial",
            "fontSize": 32,
            "fontWeight": "normal",
            "italic": False,
            "underline": True,
            "color": "#2563EB",
            "baseline": "normal",
        },
    ]
    assert text["props"]["paragraphs"][1]["runs"][0]["fontWeight"] == "bold"
    assert text["props"]["paragraphs"][2]["text"] == ""


def test_generic_export_preserves_legacy_top_level_runs() -> None:
    deck = rich_text_deck(
        {
            "text": "stale",
            "runs": [
                {
                    "text": "Legacy ",
                    "fontFamily": "Aptos",
                    "fontSize": 30,
                    "italic": True,
                    "baseline": "normal",
                },
                {
                    "text": "runs",
                    "fontFamily": "Arial",
                    "fontSize": 30,
                    "underline": True,
                    "baseline": "normal",
                },
            ],
            "fontSize": 24,
            "fontWeight": "normal",
            "align": "left",
            "lineHeight": 1.2,
            "verticalAlign": "top",
        }
    )

    response = export_deck_pptx(DeckPptxExportRequest(deck=deck))
    binary = base64.b64decode(response.content_base64)
    with ZipFile(BytesIO(binary)) as archive:
        slide_root = ET.fromstring(archive.read("ppt/slides/slide1.xml"))
    runs = slide_root.findall(f".//{{{DRAWING_NS}}}p/{{{DRAWING_NS}}}r")

    assert [run.find(f"{{{DRAWING_NS}}}t").text for run in runs] == [
        "Legacy ",
        "runs",
    ]
    assert runs[0].find(f"{{{DRAWING_NS}}}rPr").attrib["i"] == "1"
    assert runs[1].find(f"{{{DRAWING_NS}}}rPr").attrib["u"] == "sng"


def test_generic_export_preserves_legacy_plain_text() -> None:
    deck = rich_text_deck(
        {
            "text": "Legacy plain text",
            "fontFamily": "Aptos",
            "fontSize": 30,
            "fontWeight": "bold",
            "italic": True,
            "underline": True,
            "color": "#2563EB",
            "align": "right",
            "lineHeight": 1.25,
            "verticalAlign": "middle",
        }
    )

    response = export_deck_pptx(DeckPptxExportRequest(deck=deck))
    presentation = Presentation(BytesIO(base64.b64decode(response.content_base64)))
    paragraph = presentation.slides[0].shapes[0].text_frame.paragraphs[0]

    assert paragraph.text == "Legacy plain text"
    assert paragraph.alignment == PP_ALIGN.RIGHT
    assert paragraph.runs[0].font.bold is True
    assert paragraph.runs[0].font.italic is True
    assert paragraph.runs[0].font.underline is True


def test_generic_export_reports_unsupported_run_properties_with_stable_codes() -> None:
    deck = rich_text_deck(
        {
            "text": "Linked tracking",
            "fontSize": 24,
            "fontWeight": "normal",
            "align": "left",
            "lineHeight": 1.2,
            "verticalAlign": "top",
            "paragraphs": [
                {
                    "text": "Linked tracking",
                    "runs": [
                        {
                            "text": "Linked tracking",
                            "hyperlink": "https://example.invalid",
                            "letterSpacing": 1.5,
                            "textShadow": "soft",
                            "baseline": "normal",
                        }
                    ],
                    "align": "left",
                    "lineHeight": 1.2,
                    "spaceBefore": 0,
                    "spaceAfter": 0,
                    "indent": 0,
                }
            ],
        }
    )

    response = export_deck_pptx(DeckPptxExportRequest(deck=deck))

    assert response.warnings == [
        "PPTX_RICH_TEXT_UNSUPPORTED_HYPERLINK: "
        "element=el_rich_text; paragraph=0; run=0",
        "PPTX_RICH_TEXT_UNSUPPORTED_LETTER_SPACING: "
        "element=el_rich_text; paragraph=0; run=0",
        "PPTX_RICH_TEXT_UNSUPPORTED_RUN_PROPERTY: "
        "property=textShadow; element=el_rich_text; paragraph=0; run=0",
    ]


def rich_text_deck(props: dict[str, object]) -> dict[str, object]:
    return {
        "canvas": {"width": 1920, "height": 1080},
        "theme": {
            "backgroundColor": "#FFFFFF",
            "textColor": "#111827",
            "fontFamily": "Pretendard",
        },
        "slides": [
            {
                "order": 1,
                "style": {"backgroundColor": "#FFFFFF"},
                "elements": [
                    {
                        "elementId": "el_rich_text",
                        "type": "text",
                        "x": 120,
                        "y": 120,
                        "width": 960,
                        "height": 400,
                        "zIndex": 1,
                        "visible": True,
                        "props": props,
                    }
                ],
            }
        ],
    }
