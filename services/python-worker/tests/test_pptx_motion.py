import base64
import copy
import re
import zipfile
from io import BytesIO
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from app.ai.deck_pptx_export import (
    DeckPptxExportRequest,
    bound_motion_diagnostics,
    export_deck_pptx,
)
from app.ai.pptx_motion import (
    DML_NS,
    MC_NS,
    P14_NS,
    PML_NS,
    main_sequence_node,
    parse_main_sequence,
    parse_slide_transition,
    replace_main_sequence,
    replace_slide_transition,
    serialize_slide_motion,
    supported_main_sequence_shape_ids,
)
from app.ai.pptx_ooxml_vector_importer import (
    animation_shape_targets,
    import_pptx_ooxml_visual_tree,
    motion_diagnostic_summary,
)
from app.ai.pptx_ooxml_generation import (
    excluded_timing_branch_bytes,
    generate_pptx_ooxml,
    preserve_excluded_timing_branch_bytes,
    preserve_xml_subtree_bytes,
    sync_pptx_ooxml,
)


def test_transition_import_selects_choice_once_and_prefers_p14_duration() -> None:
    slide = ET.fromstring(
        f'<p:sld xmlns:p="{PML_NS}" xmlns:p14="{P14_NS}" xmlns:mc="{MC_NS}">'
        '<p:cSld/><mc:AlternateContent><mc:Choice Requires="p14">'
        '<p:transition spd="slow" p14:dur="700"><p:fade/></p:transition>'
        "</mc:Choice><mc:Fallback>"
        '<p:transition spd="slow"><p:fade/></p:transition>'
        "</mc:Fallback></mc:AlternateContent></p:sld>"
    )

    assert parse_slide_transition(slide) == {"type": "fade", "durationMs": 700}


def test_transition_import_uses_speed_fallback_and_ignores_non_fade() -> None:
    medium = ET.fromstring(
        f'<p:sld xmlns:p="{PML_NS}"><p:cSld/>'
        '<p:transition spd="med"><p:fade/></p:transition></p:sld>'
    )
    unsupported = ET.fromstring(
        f'<p:sld xmlns:p="{PML_NS}"><p:cSld/>'
        '<p:transition spd="fast"><p:push/></p:transition></p:sld>'
    )

    assert parse_slide_transition(medium) == {"type": "fade", "durationMs": 500}
    assert parse_slide_transition(unsupported) is None


def test_transition_only_sync_preserves_self_closing_timing_bytes() -> None:
    original = b"<p:sld><p:cSld/><p:timing/><p:extLst/></p:sld>"
    rewritten = b"<p:sld><p:cSld /><p:timing /><p:extLst /></p:sld>"

    preserved = preserve_xml_subtree_bytes(original, rewritten, "timing")

    assert preserved is not None
    assert b"<p:timing/>" in preserved
    assert b"<p:timing />" not in preserved


def test_raw_timing_preservation_restores_root_scoped_namespace_prefixes() -> None:
    original = (
        f'<p:sld xmlns:p="{PML_NS}" xmlns:p15="urn:test:p15">'
        "<p:cSld/><p:timing><p15:future/></p:timing></p:sld>"
    ).encode()
    rewritten_root = ET.fromstring(original)
    rewritten_root.insert(1, ET.Element(f"{{{PML_NS}}}transition"))
    rewritten = ET.tostring(rewritten_root, encoding="utf-8")

    preserved = preserve_xml_subtree_bytes(original, rewritten, "timing")

    assert preserved is not None
    ET.fromstring(preserved)
    assert b'xmlns:p15="urn:test:p15"' in preserved
    assert b"<p15:future/>" in preserved


def test_excluded_timing_checksum_covers_interactive_and_media_wrappers() -> None:
    original = b"""<p:timing xmlns:p='urn:p'><p:tnLst>
<p:seq nextAc='seek' concurrent='1'><p:cTn id='90' nodeType='interactiveSeq'/></p:seq>
<p:par custom='keep'><p:cTn id='91' presetClass='mediacall'><p:childTnLst><p:audio/></p:childTnLst></p:cTn></p:par>
</p:tnLst></p:timing>"""
    rewritten = ET.tostring(ET.fromstring(original), encoding="utf-8")

    preserved = preserve_excluded_timing_branch_bytes(original, rewritten)

    assert preserved is not None
    assert excluded_timing_branch_bytes(preserved) == excluded_timing_branch_bytes(
        original
    )
    excluded = excluded_timing_branch_bytes(preserved)
    assert excluded[0].startswith(b"<p:seq nextAc='seek' concurrent='1'>")
    assert excluded[1].startswith(b"<p:par custom='keep'>")


def test_transition_import_ignores_unsupported_alternate_content_choice() -> None:
    slide = ET.fromstring(
        f'<p:sld xmlns:p="{PML_NS}" xmlns:mc="{MC_NS}"><p:cSld/>'
        '<mc:AlternateContent><mc:Choice Requires="p99">'
        '<p:transition spd="slow"><p:fade/></p:transition>'
        "</mc:Choice><mc:Fallback>"
        '<p:transition spd="fast"><p:fade/></p:transition>'
        "</mc:Fallback></mc:AlternateContent></p:sld>"
    )

    assert parse_slide_transition(slide) == {"type": "fade", "durationMs": 250}


def test_transition_serializer_writes_p14_choice_and_compatible_fallback() -> None:
    slide = ET.fromstring(f'<p:sld xmlns:p="{PML_NS}"><p:cSld/></p:sld>')

    assert replace_slide_transition(slide, {"type": "fade", "durationMs": 700})

    alternate = next(
        child for child in list(slide) if child.tag == f"{{{MC_NS}}}AlternateContent"
    )
    choice = next(
        child for child in list(alternate) if child.tag == f"{{{MC_NS}}}Choice"
    )
    fallback = next(
        child for child in list(alternate) if child.tag == f"{{{MC_NS}}}Fallback"
    )
    choice_transition = next(
        child for child in list(choice) if local_name(child) == "transition"
    )
    fallback_transition = next(
        child for child in list(fallback) if local_name(child) == "transition"
    )

    assert choice.get("Requires") == "p14"
    assert choice_transition.get(f"{{{P14_NS}}}dur") == "700"
    assert fallback_transition.get(f"{{{P14_NS}}}dur") is None
    assert first_local_child(choice_transition, "fade") is not None
    assert first_local_child(fallback_transition, "fade") is not None
    assert parse_slide_transition(slide) == {"type": "fade", "durationMs": 700}


def test_main_sequence_parser_dedupes_behaviors_and_uses_effect_duration() -> None:
    slide = synthetic_timing_slide(include_build=False)

    animations, coverage, diagnostics = parse_main_sequence(
        slide,
        slide_index=3,
        shape_targets={"7": "el_fade", "8": "el_appear", "9": "el_zoom"},
    )

    assert coverage == "complete"
    assert diagnostics == [
        {
            "code": "PPTX_MOTION_INTERACTIVE_EXCLUDED",
            "slideIndex": 3,
            "timingNodeId": "90",
        },
        {"code": "PPTX_MOTION_MEDIA_EXCLUDED", "slideIndex": 3},
    ]
    assert animations == [
        {
            "animationId": "anim_ooxml_3_3",
            "elementId": "el_fade",
            "type": "fade-in",
            "order": 1,
            "durationMs": 500,
            "delayMs": 25,
            "easing": "ease-out",
            "startMode": "on-click",
        },
        {
            "animationId": "anim_ooxml_3_6",
            "elementId": "el_appear",
            "type": "appear",
            "order": 2,
            "durationMs": 350,
            "delayMs": 0,
            "easing": "ease-out",
            "startMode": "with-previous",
        },
        {
            "animationId": "anim_ooxml_3_8",
            "elementId": "el_zoom",
            "type": "zoom-in",
            "order": 3,
            "durationMs": 650,
            "delayMs": 10,
            "easing": "ease-out",
            "startMode": "after-previous",
        },
    ]


def test_main_sequence_parser_marks_build_and_unresolved_targets_partial() -> None:
    slide = synthetic_timing_slide(include_build=True)

    animations, coverage, diagnostics = parse_main_sequence(
        slide,
        slide_index=2,
        shape_targets={"7": "el_group", "8": "el_appear"},
    )

    assert coverage == "partial"
    assert [animation["elementId"] for animation in animations] == [
        "el_group",
        "el_appear",
    ]
    assert {diagnostic["code"] for diagnostic in diagnostics} == {
        "PPTX_MOTION_INTERACTIVE_EXCLUDED",
        "PPTX_MOTION_MEDIA_EXCLUDED",
        "PPTX_MOTION_PARAGRAPH_BUILD_DOWNGRADED",
        "PPTX_MOTION_TARGET_UNRESOLVED",
    }


def test_main_sequence_parser_prunes_nested_interactive_entrance_effects() -> None:
    slide = ET.fromstring(
        f'<p:sld xmlns:p="{PML_NS}"><p:cSld/><p:timing><p:tnLst><p:par>'
        '<p:cTn id="1" nodeType="tmRoot"><p:childTnLst><p:seq>'
        '<p:cTn id="2" nodeType="mainSeq"><p:childTnLst><p:seq>'
        '<p:cTn id="90" nodeType="interactiveSeq"><p:childTnLst>'
        + effect_xml(
            outer_id=5,
            node_type="clickEffect",
            preset_id=10,
            target_id="7",
            duration=500,
            delay=0,
            behavior="fade",
        )
        + "<p:cmd/></p:childTnLst></p:cTn></p:seq></p:childTnLst></p:cTn>"
        "</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing></p:sld>"
    )

    animations, coverage, diagnostics = parse_main_sequence(
        slide,
        slide_index=1,
        shape_targets={"7": "el_interactive"},
    )

    assert animations == []
    assert coverage == "partial"
    assert {diagnostic["code"] for diagnostic in diagnostics} == {
        "PPTX_MOTION_INTERACTIVE_EXCLUDED"
    }

    elements = [
        {"elementId": "el_fill", "x": 0, "y": 0, "width": 100, "height": 100},
        {"elementId": "el_text", "x": 0, "y": 0, "width": 100, "height": 100},
    ]
    slot_sources = {
        "el_fill": {
            "slidePart": "ppt/slides/slide1.xml",
            "shapeId": "7",
            "writable": True,
        },
        "el_text": {
            "slidePart": "ppt/slides/slide1.xml",
            "shapeId": "7",
            "writable": True,
        },
    }

    targets = animation_shape_targets(
        slide,
        slide_index=1,
        slide_part="ppt/slides/slide1.xml",
        elements=elements,
        slot_sources=slot_sources,
    )

    assert targets == {}
    assert [element["elementId"] for element in elements] == [
        "el_fill",
        "el_text",
    ]


def test_media_exclusion_counts_logical_mediacall_branches_once() -> None:
    slide = ET.fromstring(
        f'<p:sld xmlns:p="{PML_NS}"><p:cSld/><p:timing><p:tnLst><p:par>'
        '<p:cTn id="1" nodeType="tmRoot"><p:childTnLst><p:seq>'
        '<p:cTn id="2" nodeType="mainSeq"><p:childTnLst>'
        '<p:par><p:cTn id="5" presetClass="mediacall"><p:childTnLst>'
        "<p:video/><p:cmd/></p:childTnLst></p:cTn></p:par>"
        '<p:par><p:cTn id="26" presetClass="mediacall"><p:childTnLst>'
        "<p:cmd/></p:childTnLst></p:cTn></p:par>"
        "</p:childTnLst></p:cTn></p:seq></p:childTnLst></p:cTn>"
        "</p:par></p:tnLst></p:timing></p:sld>"
    )

    animations, coverage, diagnostics = parse_main_sequence(
        slide,
        slide_index=8,
        shape_targets={},
    )

    assert animations == []
    assert coverage == "partial"
    assert diagnostics == [
        {
            "code": "PPTX_MOTION_MEDIA_EXCLUDED",
            "slideIndex": 8,
            "timingNodeId": "5",
        },
        {
            "code": "PPTX_MOTION_MEDIA_EXCLUDED",
            "slideIndex": 8,
            "timingNodeId": "26",
        },
    ]


def test_supported_shape_targets_exclude_multi_target_unresolved_effects() -> None:
    main_sequence = ET.fromstring(
        f'<p:cTn xmlns:p="{PML_NS}" id="2" nodeType="mainSeq">'
        '<p:childTnLst><p:par><p:cTn id="3" presetClass="entr" '
        'presetID="10" nodeType="clickEffect"><p:childTnLst>'
        '<p:animEffect filter="fade"><p:cBhvr><p:cTn id="4" dur="500"/>'
        '<p:tgtEl><p:spTgt spid="7"/></p:tgtEl>'
        '<p:tgtEl><p:spTgt spid="8"/></p:tgtEl>'
        "</p:cBhvr></p:animEffect></p:childTnLst></p:cTn></p:par>"
        "</p:childTnLst></p:cTn>"
    )

    assert supported_main_sequence_shape_ids(main_sequence) == set()


def test_main_sequence_parser_uses_nested_preset_nodes_as_logical_effects() -> None:
    slide = ET.fromstring(
        f'<p:sld xmlns:p="{PML_NS}"><p:cSld/><p:timing><p:tnLst><p:par>'
        '<p:cTn id="1" nodeType="tmRoot"><p:childTnLst><p:seq>'
        '<p:cTn id="2" nodeType="mainSeq"><p:childTnLst><p:par>'
        '<p:cTn id="3"><p:childTnLst>'
        + effect_xml(
            outer_id=5,
            node_type="clickEffect",
            preset_id=10,
            target_id="7",
            duration=500,
            delay=0,
            behavior="fade",
        )
        + effect_xml(
            outer_id=8,
            node_type="withEffect",
            preset_id=1,
            target_id="8",
            duration=350,
            delay=0,
            behavior="appear",
        )
        + "</p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn>"
        "</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing></p:sld>"
    )

    animations, coverage, diagnostics = parse_main_sequence(
        slide,
        slide_index=1,
        shape_targets={"7": "el_fade", "8": "el_appear"},
    )

    assert coverage == "complete"
    assert diagnostics == []
    assert [item["animationId"] for item in animations] == [
        "anim_ooxml_1_5",
        "anim_ooxml_1_8",
    ]
    assert [item["startMode"] for item in animations] == [
        "on-click",
        "with-previous",
    ]


def test_main_sequence_parser_rejects_unknown_declared_preset() -> None:
    slide = ET.fromstring(
        f'<p:sld xmlns:p="{PML_NS}"><p:cSld/><p:timing><p:tnLst><p:par>'
        '<p:cTn id="1" nodeType="tmRoot"><p:childTnLst><p:seq>'
        '<p:cTn id="2" nodeType="mainSeq"><p:childTnLst>'
        + effect_xml(
            outer_id=3,
            node_type="clickEffect",
            preset_id=22,
            target_id="7",
            duration=500,
            delay=0,
            behavior="appear",
        )
        + "</p:childTnLst></p:cTn></p:seq></p:childTnLst></p:cTn>"
        "</p:par></p:tnLst></p:timing></p:sld>"
    )

    animations, coverage, diagnostics = parse_main_sequence(
        slide,
        slide_index=1,
        shape_targets={"7": "el_target"},
    )

    assert animations == []
    assert coverage == "partial"
    assert diagnostics == [
        {
            "code": "PPTX_MOTION_PRESET_UNSUPPORTED",
            "slideIndex": 1,
            "timingNodeId": "3",
        }
    ]


def test_main_sequence_parser_uses_set_duration_for_appear_effect() -> None:
    slide = ET.fromstring(
        f'<p:sld xmlns:p="{PML_NS}"><p:cSld/><p:timing><p:tnLst><p:par>'
        '<p:cTn id="1" nodeType="tmRoot"><p:childTnLst><p:seq>'
        '<p:cTn id="2" nodeType="mainSeq"><p:childTnLst><p:par>'
        '<p:cTn id="3" presetClass="entr" presetID="1" '
        'nodeType="clickEffect"><p:childTnLst><p:set><p:cBhvr>'
        '<p:cTn id="4" dur="350"/><p:tgtEl><p:spTgt spid="7"/>'
        "</p:tgtEl></p:cBhvr></p:set></p:childTnLst></p:cTn></p:par>"
        "</p:childTnLst></p:cTn></p:seq></p:childTnLst></p:cTn>"
        "</p:par></p:tnLst></p:timing></p:sld>"
    )

    animations, coverage, diagnostics = parse_main_sequence(
        slide,
        slide_index=1,
        shape_targets={"7": "el_target"},
    )

    assert coverage == "complete"
    assert diagnostics == []
    assert animations[0]["type"] == "appear"
    assert animations[0]["durationMs"] == 350


def test_main_sequence_parser_fails_closed_for_nested_excluded_branch() -> None:
    slide = ET.fromstring(
        f'<p:sld xmlns:p="{PML_NS}"><p:cSld/><p:timing><p:tnLst><p:par>'
        '<p:cTn id="1" nodeType="tmRoot"><p:childTnLst><p:seq>'
        '<p:cTn id="2" nodeType="mainSeq"><p:childTnLst><p:par>'
        '<p:cTn id="3" dur="500" presetClass="entr" presetID="10" '
        'nodeType="clickEffect"><p:childTnLst>'
        '<p:animEffect filter="fade" transition="in"><p:cBhvr>'
        '<p:cTn id="4" dur="500"/><p:tgtEl><p:spTgt spid="7"/>'
        "</p:tgtEl></p:cBhvr></p:animEffect>"
        '<p:cmd type="call" cmd="fixture-command"/>'
        "</p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn>"
        "</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst>"
        "</p:timing></p:sld>"
    )
    before = canonical(slide)

    animations, coverage, diagnostics = parse_main_sequence(
        slide,
        slide_index=1,
        shape_targets={"7": "el_target"},
    )
    applied, replacement_diagnostics = replace_main_sequence(
        slide,
        animations,
        slide_index=1,
        element_targets={"el_target": ["7"]},
    )

    assert coverage == "partial"
    assert {item["code"] for item in diagnostics} == {"PPTX_MOTION_MEDIA_EXCLUDED"}
    assert applied is False
    assert {item["code"] for item in replacement_diagnostics} == {
        "PPTX_MOTION_MEDIA_EXCLUDED"
    }
    assert canonical(slide) == before


def test_motion_serializer_is_deterministic_and_flattens_group_targets() -> None:
    animations = [
        animation("anim_zoom", "el_group", "zoom-in", 2, "after-previous", 650, 10),
        animation("anim_fade", "el_title", "fade-in", 1, "on-click", 500, 25),
    ]

    first = serialize_slide_motion(
        animations,
        slide_index=1,
        element_targets={"el_title": ["7"], "el_group": ["8", "9"]},
    )
    second = serialize_slide_motion(
        copy.deepcopy(animations),
        slide_index=1,
        element_targets={"el_title": ["7"], "el_group": ["8", "9"]},
    )

    assert first.effect_count == 3
    assert first.diagnostics == [
        {
            "code": "PPTX_MOTION_TARGET_FLATTENED",
            "slideIndex": 1,
            "elementId": "el_group",
            "count": 2,
        }
    ]
    assert canonical(first.timing) == canonical(second.timing)
    assert first.timing is not None
    zoom_scale = next(
        node for node in first.timing.iter() if local_name(node) == "animScale"
    )
    zoom_names = [
        node.text for node in zoom_scale.iter() if local_name(node) == "attrName"
    ]
    zoom_from = first_local_child(zoom_scale, "from")
    zoom_to = first_local_child(zoom_scale, "to")
    assert zoom_names == ["ScaleX", "ScaleY"]
    assert zoom_from is not None and zoom_from.attrib == {"x": "0", "y": "0"}
    assert zoom_to is not None and zoom_to.attrib == {
        "x": "100000",
        "y": "100000",
    }
    assert first_local_child(zoom_scale, "by") is None
    zoom_parent = next(
        node
        for node in first.timing.iter()
        if local_name(node) == "cTn" and node.get("presetID") == "23"
    )
    assert any(local_name(node) == "set" for node in zoom_parent.iter())
    parsed, coverage, _diagnostics = parse_main_sequence(
        slide_with_timing(first.timing),
        slide_index=1,
        shape_targets={"7": "el_title", "8": "el_group_a", "9": "el_group_b"},
    )
    assert coverage == "complete"
    assert [item["startMode"] for item in parsed] == [
        "on-click",
        "after-previous",
        "with-previous",
    ]
    assert [item["durationMs"] for item in parsed] == [500, 650, 650]


def test_motion_serializer_orders_entry_roots_before_click_roots_natively() -> None:
    serialized = serialize_slide_motion(
        [
            animation("anim_click", "el_click", "fade-in", 1, "on-click", 400, 0),
            animation(
                "anim_entry",
                "el_entry",
                "fade-in",
                2,
                "on-slide-enter",
                300,
                0,
            ),
        ],
        slide_index=1,
        element_targets={"el_click": ["7"], "el_entry": ["8"]},
    )

    assert serialized.timing is not None
    main = main_sequence_node(serialized.timing)
    assert main is not None
    native_effects = [
        node
        for node in main.iter()
        if local_name(node) == "cTn" and node.get("presetClass") == "entr"
    ]
    assert [node.get("nodeType") for node in native_effects] == [
        "withEffect",
        "clickEffect",
    ]
    parsed, coverage, diagnostics = parse_main_sequence(
        slide_with_timing(serialized.timing),
        slide_index=1,
        shape_targets={"7": "el_click", "8": "el_entry"},
    )
    assert coverage == "complete"
    assert diagnostics == []
    assert [item["elementId"] for item in parsed] == ["el_click", "el_entry"]
    assert [item["startMode"] for item in parsed] == [
        "on-click",
        "on-slide-enter",
    ]


def test_transition_replacement_keeps_timing_subtree_byte_equivalent() -> None:
    slide = synthetic_timing_slide(include_build=False)
    timing = next(node for node in slide if local_name(node) == "timing")
    before = canonical(timing)

    assert replace_slide_transition(slide, {"type": "fade", "durationMs": 900})
    assert canonical(timing) == before
    assert parse_slide_transition(slide) == {"type": "fade", "durationMs": 900}
    assert replace_slide_transition(slide, None)
    assert canonical(timing) == before
    assert parse_slide_transition(slide) is None


def test_main_sequence_replacement_preserves_interactive_and_media_branches() -> None:
    slide = synthetic_timing_slide(include_build=False)
    timing = next(node for node in slide if local_name(node) == "timing")
    interactive = next(
        node
        for node in timing.iter()
        if local_name(node) == "cTn" and node.get("nodeType") == "interactiveSeq"
    )
    media = next(node for node in timing.iter() if local_name(node) == "audio")
    interactive_before = canonical(interactive)
    media_before = canonical(media)

    applied, diagnostics = replace_main_sequence(
        slide,
        [animation("anim_new", "el_fade", "appear", 1, "on-click", 400, 0)],
        slide_index=1,
        element_targets={"el_fade": ["7"]},
    )

    assert applied is True
    assert diagnostics == []
    assert canonical(interactive) == interactive_before
    assert canonical(media) == media_before
    main = main_sequence_node(timing)
    assert main is not None
    animations, coverage, _diagnostics = parse_main_sequence(
        slide,
        slide_index=1,
        shape_targets={"7": "el_fade"},
    )
    assert coverage == "complete"
    assert len(animations) == 1
    assert animations[0]["type"] == "appear"


def test_generic_export_serializes_transition_and_flattens_group_targets() -> None:
    deck = generic_motion_deck()

    first = export_deck_pptx(DeckPptxExportRequest(deck=deck))
    second = export_deck_pptx(DeckPptxExportRequest(deck=copy.deepcopy(deck)))
    first_package = base64.b64decode(first.content_base64)
    second_package = base64.b64decode(second.content_base64)
    first_slide = slide_root_from_package(first_package)
    second_slide = slide_root_from_package(second_package)

    assert first.motion_diagnostics == [
        {
            "code": "PPTX_MOTION_TARGET_FLATTENED",
            "slideIndex": 1,
            "elementId": "el_group",
            "count": 2,
        }
    ]
    assert parse_slide_transition(first_slide) == {
        "type": "fade",
        "durationMs": 700,
    }
    animations, coverage, diagnostics = parse_main_sequence(
        first_slide,
        slide_index=1,
        shape_targets={"2": "el_a", "3": "el_b"},
    )
    assert coverage == "complete"
    assert diagnostics == []
    assert [item["elementId"] for item in animations] == ["el_a", "el_b"]
    assert [item["startMode"] for item in animations] == [
        "on-click",
        "with-previous",
    ]
    assert canonical(first_local_child(first_slide, "timing")) == canonical(
        first_local_child(second_slide, "timing")
    )


def test_generic_export_bounds_large_motion_diagnostic_sets_deterministically() -> None:
    diagnostics = [
        {
            "code": "PPTX_MOTION_EFFECT_UNSUPPORTED",
            "slideIndex": (index % 3) + 1,
            "elementId": f"el_{index}",
        }
        for index in range(501)
    ]

    bounded = bound_motion_diagnostics(diagnostics)

    assert bounded == [
        {
            "code": "PPTX_MOTION_EFFECT_UNSUPPORTED",
            "slideIndex": 1,
            "count": 501,
        }
    ]


def test_vector_import_creates_deterministic_synthetic_animation_group(
    tmp_path: Path,
) -> None:
    pptx_path = tmp_path / "motion-group.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1),
        Inches(1),
        Inches(3),
        Inches(1),
    )
    shape.text = "Animation target"
    presentation.save(str(pptx_path))
    inject_fixture_motion(pptx_path, shape_id=str(shape.shape_id))

    first = import_pptx_ooxml_visual_tree(pptx_path, "file_motion")
    second = import_pptx_ooxml_visual_tree(pptx_path, "file_motion")
    first_slide = first.blueprint["slides"][0]
    second_slide = second.blueprint["slides"][0]
    animation = first_slide["animations"][0]
    group = next(
        element
        for element in first_slide["elements"]
        if element["elementId"] == animation["elementId"]
    )

    assert animation["type"] == "fade-in"
    assert group["type"] == "group"
    assert len(group["props"]["childElementIds"]) >= 2
    assert group == next(
        element
        for element in second_slide["elements"]
        if element["elementId"] == animation["elementId"]
    )
    assert first_slide["ooxmlMotionCapabilities"] == {
        "transitionWritable": True,
        "importedMainSequenceCoverage": "complete",
    }
    source = next(
        item
        for item in first.template_blueprint["slides"][0]["elementSources"]
        if item["elementId"] == group["elementId"]
    )
    assert source["shapeId"] == str(shape.shape_id)
    assert first.quality_report["motionDiagnostics"] == {
        "total": 0,
        "unsupported": 0,
        "downgraded": 0,
        "unresolved": 0,
        "excluded": 0,
        "details": [],
    }


def test_motion_diagnostic_details_fail_closed_at_shared_bound() -> None:
    diagnostics = [
        {
            "code": "PPTX_MOTION_PRESET_UNSUPPORTED",
            "slideIndex": slide_index,
        }
        for slide_index in range(1, 502)
    ]

    first = motion_diagnostic_summary(diagnostics)
    second = motion_diagnostic_summary(list(reversed(diagnostics)))

    assert first == second
    assert first["total"] == 501
    assert first["unsupported"] == 501
    assert first["details"] == []


def test_imported_transition_sync_preserves_timing_bytes(tmp_path: Path) -> None:
    pptx_path, _shape_id = synthetic_motion_pptx(tmp_path, include_excluded=True)
    generated = generate_pptx_ooxml(pptx_path, "file_motion", render=False)
    template_slide = generated.template_blueprint["slides"][0]
    original_package = pptx_path.read_bytes()
    original_timing = xml_subtree_bytes(
        slide_xml_from_package(original_package),
        "timing",
    )

    result = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        operations=[],
        slide_motion=[
            {
                "slideId": "slide_motion_1",
                "sourceSlidePart": "ppt/slides/slide1.xml",
                "transition": {"type": "fade", "durationMs": 900},
                "animations": generated.blueprint["slides"][0]["animations"],
                "capabilities": template_slide["ooxmlMotionCapabilities"],
                "touched": {"transition": True, "animations": False},
            }
        ],
        deck_canvas={"width": 1920, "height": 1080},
        synced_deck_version=2,
        render=False,
    )
    package = package_asset_bytes(result)

    assert [item.model_dump(by_alias=True) for item in result.applied_slide_motion] == [
        {
            "slideId": "slide_motion_1",
            "transition": True,
            "animations": False,
        }
    ]
    assert result.unsupported_slide_motion == []
    assert (
        xml_subtree_bytes(slide_xml_from_package(package), "timing") == original_timing
    )
    assert parse_slide_transition(slide_root_from_package(package)) == {
        "type": "fade",
        "durationMs": 900,
    }


def test_transition_sync_preserves_original_timing_with_element_edit(
    tmp_path: Path,
) -> None:
    pptx_path, _shape_id = synthetic_motion_pptx(tmp_path, include_excluded=True)
    compact_timing_empty_tags(pptx_path)
    generated = generate_pptx_ooxml(pptx_path, "file_motion", render=False)
    template_slide = generated.template_blueprint["slides"][0]
    text_source = next(
        source
        for source in template_slide["elementSources"]
        if source.get("writable") is True
        and source.get("elementType") == "text"
        and source.get("ooxmlEditCapabilities", {}).get("richText") == "full"
    )
    original_timing = xml_subtree_bytes(
        slide_xml_from_package(pptx_path.read_bytes()),
        "timing",
    )

    result = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        operations=[
            {
                "type": "update_element_props",
                "slideId": "slide_motion_1",
                "elementId": text_source["elementId"],
                "props": {"text": "Updated animation target"},
            }
        ],
        slide_motion=[
            {
                "slideId": "slide_motion_1",
                "sourceSlidePart": "ppt/slides/slide1.xml",
                "transition": {"type": "fade", "durationMs": 900},
                "animations": generated.blueprint["slides"][0]["animations"],
                "capabilities": template_slide["ooxmlMotionCapabilities"],
                "touched": {"transition": True, "animations": False},
            }
        ],
        deck_canvas={"width": 1920, "height": 1080},
        synced_deck_version=2,
        render=False,
    )
    package = package_asset_bytes(result)

    assert len(result.applied_operations) == 1
    assert len(result.applied_slide_motion) == 1
    assert xml_subtree_bytes(slide_xml_from_package(package), "timing") == (
        original_timing
    )


def test_imported_main_sequence_sync_preserves_excluded_branches(
    tmp_path: Path,
) -> None:
    pptx_path, shape_id = synthetic_motion_pptx(tmp_path, include_excluded=True)
    compact_timing_empty_tags(pptx_path)
    original_slide_xml = slide_xml_from_package(pptx_path.read_bytes())
    original_excluded_bytes = excluded_timing_branch_bytes(original_slide_xml)
    generated = generate_pptx_ooxml(pptx_path, "file_motion", render=False)
    template_slide = generated.template_blueprint["slides"][0]
    target_id = generated.blueprint["slides"][0]["animations"][0]["elementId"]
    before = slide_root_from_package(pptx_path.read_bytes())
    before_timing = first_local_child(before, "timing")
    assert before_timing is not None
    before_interactive = next(
        node
        for node in before_timing.iter()
        if local_name(node) == "cTn" and node.get("nodeType") == "interactiveSeq"
    )
    before_audio = next(
        node for node in before_timing.iter() if local_name(node) == "audio"
    )

    result = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        operations=[],
        slide_motion=[
            {
                "slideId": "slide_motion_1",
                "sourceSlidePart": "ppt/slides/slide1.xml",
                "transition": generated.blueprint["slides"][0].get("transition"),
                "animations": [
                    animation(
                        "anim_replaced",
                        target_id,
                        "appear",
                        1,
                        "after-previous",
                        400,
                        20,
                    )
                ],
                "capabilities": template_slide["ooxmlMotionCapabilities"],
                "touched": {"transition": False, "animations": True},
            }
        ],
        deck_canvas={"width": 1920, "height": 1080},
        synced_deck_version=2,
        render=False,
    )
    after_package = package_asset_bytes(result)
    after_slide_xml = slide_xml_from_package(after_package)
    after = slide_root_from_package(after_package)
    after_timing = first_local_child(after, "timing")
    assert after_timing is not None
    after_interactive = next(
        node
        for node in after_timing.iter()
        if local_name(node) == "cTn" and node.get("nodeType") == "interactiveSeq"
    )
    after_audio = next(
        node for node in after_timing.iter() if local_name(node) == "audio"
    )
    animations, coverage, diagnostics = parse_main_sequence(
        after,
        slide_index=1,
        shape_targets={shape_id: target_id},
    )

    assert canonical(after_interactive) == canonical(before_interactive)
    assert canonical(after_audio) == canonical(before_audio)
    assert excluded_timing_branch_bytes(after_slide_xml) == original_excluded_bytes
    after_root = slide_root_from_package(after_package)
    assert (
        "orbit"
        in after_root.get(
            "{http://schemas.openxmlformats.org/markup-compatibility/2006}Ignorable",
            "",
        ).split()
    )
    assert coverage == "complete"
    assert {item["code"] for item in diagnostics} == {
        "PPTX_MOTION_INTERACTIVE_EXCLUDED",
        "PPTX_MOTION_MEDIA_EXCLUDED",
    }
    assert [(item["type"], item["startMode"]) for item in animations] == [
        ("appear", "after-previous")
    ]
    assert result.applied_slide_motion[0].animations is True


def test_slide_motion_sync_failure_returns_original_package_atomically(
    tmp_path: Path,
) -> None:
    pptx_path, _shape_id = synthetic_motion_pptx(tmp_path, include_excluded=False)
    generated = generate_pptx_ooxml(pptx_path, "file_motion", render=False)
    original_package = pptx_path.read_bytes()
    stale_capabilities = {
        **generated.template_blueprint["slides"][0]["ooxmlMotionCapabilities"],
        "importedMainSequenceCoverage": "partial",
    }

    result = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        operations=[],
        slide_motion=[
            {
                "slideId": "slide_motion_1",
                "sourceSlidePart": "ppt/slides/slide1.xml",
                "transition": {"type": "fade", "durationMs": 900},
                "animations": generated.blueprint["slides"][0]["animations"],
                "capabilities": stale_capabilities,
                "touched": {"transition": True, "animations": True},
            }
        ],
        deck_canvas={"width": 1920, "height": 1080},
        synced_deck_version=2,
        render=False,
    )

    assert package_asset_bytes(result) == original_package
    assert result.applied_slide_motion == []
    assert [
        item.model_dump(by_alias=True) for item in result.unsupported_slide_motion
    ] == [
        {
            "slideId": "slide_motion_1",
            "scope": "animations",
            "reasonCode": "SLIDE_ANIMATION_CAPABILITY_UNSAFE",
        }
    ]


def test_sync_applies_delete_and_animation_replacement_atomically(
    tmp_path: Path,
) -> None:
    pptx_path, shape_id = synthetic_motion_pptx(
        tmp_path,
        include_excluded=False,
    )
    generated = generate_pptx_ooxml(pptx_path, "file_motion", render=False)
    template_slide = generated.template_blueprint["slides"][0]
    assert (
        template_slide["ooxmlMotionCapabilities"]["importedMainSequenceCoverage"]
        == "complete"
    )
    sources = [
        item for item in template_slide["elementSources"] if item["shapeId"] == shape_id
    ]
    assert sources

    result = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
        operations=[
            {
                "type": "delete_element",
                "slideId": "slide_motion_1",
                "sourceSlidePart": "ppt/slides/slide1.xml",
                "elementId": source["elementId"],
            }
            for source in sources
        ],
        slide_motion=[
            {
                "slideId": "slide_motion_1",
                "sourceSlidePart": "ppt/slides/slide1.xml",
                "transition": generated.blueprint["slides"][0].get("transition"),
                "animations": [],
                "capabilities": template_slide["ooxmlMotionCapabilities"],
                "touched": {"transition": False, "animations": True},
            }
        ],
    )
    package = package_asset_bytes(result)

    assert result.unsupported_operations == []
    assert result.unsupported_slide_motion == []
    assert len(result.applied_operations) == len(sources)
    assert result.applied_slide_motion[0].animations is True
    root = slide_root_from_package(package)
    assert not any(
        local_name(node) == "cNvPr" and node.get("id") == shape_id
        for node in root.iter()
    )
    animations, coverage, diagnostics = parse_main_sequence(
        root,
        slide_index=1,
        shape_targets={},
    )
    assert animations == []
    assert coverage == "complete"
    assert diagnostics == []


def test_sync_retargets_animation_after_same_batch_element_type_replacement(
    tmp_path: Path,
) -> None:
    pptx_path = tmp_path / "sync-motion-target-replacement.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1),
        Inches(1),
        Inches(3),
        Inches(1),
    )
    shape_id = str(shape.shape_id)
    slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(5),
        Inches(1),
        Inches(1),
        Inches(1),
    )
    presentation.save(str(pptx_path))
    inject_fixture_motion(pptx_path, shape_id=shape_id)
    generated = generate_pptx_ooxml(pptx_path, "file_motion", render=False)
    template_slide = generated.template_blueprint["slides"][0]
    source = next(
        item
        for item in template_slide["elementSources"]
        if item["shapeId"] == shape_id and item.get("writable") is True
    )
    element_id = source["elementId"]

    result = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=generated.template_blueprint,
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
        operations=[
            {
                "type": "delete_element",
                "slideId": "slide_motion_1",
                "sourceSlidePart": "ppt/slides/slide1.xml",
                "elementId": element_id,
            },
            {
                "type": "add_element",
                "slideId": "slide_motion_1",
                "sourceSlidePart": "ppt/slides/slide1.xml",
                "element": {
                    "elementId": element_id,
                    "type": "rect",
                    "x": 100,
                    "y": 100,
                    "width": 300,
                    "height": 100,
                    "rotation": 0,
                    "opacity": 1,
                    "locked": False,
                    "visible": True,
                    "props": {
                        "fill": "#2563eb",
                        "stroke": "transparent",
                        "strokeWidth": 0,
                        "borderRadius": 0,
                    },
                },
            },
        ],
        slide_motion=[
            {
                "slideId": "slide_motion_1",
                "sourceSlidePart": "ppt/slides/slide1.xml",
                "transition": generated.blueprint["slides"][0].get("transition"),
                "animations": [
                    animation(
                        "anim_retargeted",
                        element_id,
                        "fade-in",
                        1,
                        "on-click",
                        500,
                        0,
                    )
                ],
                "capabilities": template_slide["ooxmlMotionCapabilities"],
                "touched": {"transition": False, "animations": True},
            }
        ],
    )
    package = package_asset_bytes(result)
    assert result.unsupported_operations == []
    assert result.unsupported_slide_motion == []
    new_source = next(
        item for item in result.element_sources if item["elementId"] == element_id
    )
    new_shape_id = new_source["shapeId"]
    timing_targets = [
        node.get("spid")
        for node in slide_root_from_package(package).iter()
        if local_name(node) == "spTgt"
    ]

    assert new_shape_id != shape_id
    assert timing_targets
    assert set(timing_targets) == {new_shape_id}


def test_slide_motion_sync_rejects_ambiguous_source_part_atomically(
    tmp_path: Path,
) -> None:
    pptx_path, _shape_id = synthetic_motion_pptx(tmp_path, include_excluded=False)
    generated = generate_pptx_ooxml(pptx_path, "file_motion", render=False)
    original_package = pptx_path.read_bytes()
    template_blueprint = copy.deepcopy(generated.template_blueprint)
    duplicate_slide = copy.deepcopy(template_blueprint["slides"][0])
    duplicate_slide["slideIndex"] = 2
    duplicate_slide["sourceSlideIndex"] = 2
    template_blueprint["slides"].append(duplicate_slide)

    result = sync_pptx_ooxml(
        pptx_path,
        template_blueprint=template_blueprint,
        operations=[],
        slide_motion=[
            {
                "slideId": "slide_motion_1",
                "sourceSlidePart": "ppt/slides/slide1.xml",
                "transition": {"type": "fade", "durationMs": 900},
                "animations": generated.blueprint["slides"][0]["animations"],
                "capabilities": template_blueprint["slides"][0][
                    "ooxmlMotionCapabilities"
                ],
                "touched": {"transition": True, "animations": False},
            }
        ],
        deck_canvas={"width": 1920, "height": 1080},
        synced_deck_version=2,
        render=False,
    )

    assert package_asset_bytes(result) == original_package
    assert result.applied_slide_motion == []
    assert [
        item.model_dump(by_alias=True) for item in result.unsupported_slide_motion
    ] == [
        {
            "slideId": "slide_motion_1",
            "scope": "transition",
            "reasonCode": "SLIDE_MOTION_SOURCE_MISSING",
        }
    ]


def generic_motion_deck() -> dict[str, object]:
    return {
        "canvas": {"width": 1920, "height": 1080},
        "theme": {"backgroundColor": "#FFFFFF"},
        "slides": [
            {
                "order": 1,
                "style": {"backgroundColor": "#FFFFFF"},
                "transition": {"type": "fade", "durationMs": 700},
                "animations": [
                    animation(
                        "anim_group",
                        "el_group",
                        "fade-in",
                        1,
                        "on-click",
                        500,
                        0,
                    )
                ],
                "elements": [
                    motion_rect("el_a", 100, 100, 300, 200, 1),
                    motion_rect("el_b", 500, 100, 300, 200, 2),
                    {
                        "elementId": "el_group",
                        "type": "group",
                        "x": 100,
                        "y": 100,
                        "width": 700,
                        "height": 200,
                        "zIndex": 3,
                        "visible": True,
                        "props": {"childElementIds": ["el_a", "el_b"]},
                    },
                ],
            }
        ],
    }


def motion_rect(
    element_id: str,
    x: int,
    y: int,
    width: int,
    height: int,
    z_index: int,
) -> dict[str, object]:
    return {
        "elementId": element_id,
        "type": "rect",
        "x": x,
        "y": y,
        "width": width,
        "height": height,
        "zIndex": z_index,
        "visible": True,
        "props": {"fill": "#336699"},
    }


def slide_root_from_package(package_bytes: bytes) -> ET.Element:
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as package:
        return ET.fromstring(package.read("ppt/slides/slide1.xml"))


def inject_fixture_motion(path: Path, *, shape_id: str) -> None:
    package_bytes = path.read_bytes()
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as source:
        slide_root = ET.fromstring(source.read("ppt/slides/slide1.xml"))
        replace_slide_transition(
            slide_root,
            {"type": "fade", "durationMs": 700},
        )
        applied, diagnostics = replace_main_sequence(
            slide_root,
            [animation("anim_fixture", "el_fixture", "fade-in", 1, "on-click", 500, 0)],
            slide_index=1,
            element_targets={"el_fixture": [shape_id]},
        )
        assert applied is True
        assert diagnostics == []
        buffer = BytesIO()
        with zipfile.ZipFile(buffer, "w") as target:
            for info in source.infolist():
                target.writestr(
                    info,
                    ET.tostring(slide_root, encoding="utf-8", xml_declaration=True)
                    if info.filename == "ppt/slides/slide1.xml"
                    else source.read(info.filename),
                )
    path.write_bytes(buffer.getvalue())


def synthetic_motion_pptx(
    tmp_path: Path,
    *,
    include_excluded: bool,
) -> tuple[Path, str]:
    pptx_path = tmp_path / "sync-motion.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1),
        Inches(1),
        Inches(3),
        Inches(1),
    )
    shape.text = "Animation target"
    shape_id = str(shape.shape_id)
    presentation.save(str(pptx_path))
    inject_fixture_motion(pptx_path, shape_id=shape_id)
    if include_excluded:
        inject_excluded_timing_branches(pptx_path)
    return pptx_path, shape_id


def inject_excluded_timing_branches(path: Path) -> None:
    package_bytes = path.read_bytes()
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as source:
        slide_root = ET.fromstring(source.read("ppt/slides/slide1.xml"))
        timing = first_local_child(slide_root, "timing")
        assert timing is not None
        root_time = next(
            node
            for node in timing.iter()
            if local_name(node) == "cTn" and node.get("nodeType") == "tmRoot"
        )
        root_children = first_local_child(root_time, "childTnLst")
        assert root_children is not None
        excluded = ET.fromstring(
            f'<p:root xmlns:p="{PML_NS}">'
            '<p:seq><p:cTn id="90" nodeType="interactiveSeq"/></p:seq>'
            '<p:audio><p:cMediaNode><p:cTn id="91"/></p:cMediaNode></p:audio>'
            "</p:root>"
        )
        root_children.extend(list(excluded))
        buffer = BytesIO()
        with zipfile.ZipFile(buffer, "w") as target:
            for info in source.infolist():
                target.writestr(
                    info,
                    ET.tostring(slide_root, encoding="utf-8", xml_declaration=True)
                    if info.filename == "ppt/slides/slide1.xml"
                    else source.read(info.filename),
                )
    path.write_bytes(buffer.getvalue())


def compact_timing_empty_tags(path: Path) -> None:
    package_bytes = path.read_bytes()
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as source:
        slide_xml = source.read("ppt/slides/slide1.xml")
        timing = xml_subtree_bytes(slide_xml, "timing")
        compact_timing = timing.replace(b" />", b"/>")
        assert compact_timing != timing
        compact_slide_xml = slide_xml.replace(timing, compact_timing, 1)
        buffer = BytesIO()
        with zipfile.ZipFile(buffer, "w") as target:
            for info in source.infolist():
                target.writestr(
                    info,
                    compact_slide_xml
                    if info.filename == "ppt/slides/slide1.xml"
                    else source.read(info.filename),
                )
    path.write_bytes(buffer.getvalue())


def slide_xml_from_package(package_bytes: bytes) -> bytes:
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as package:
        return package.read("ppt/slides/slide1.xml")


def package_asset_bytes(result: object) -> bytes:
    assets = getattr(result, "assets")
    package = next(asset for asset in assets if asset.asset_id == "current_package")
    return base64.b64decode(package.content_base64)


def xml_subtree_bytes(content: bytes, name: str) -> bytes:
    escaped = re.escape(name.encode("ascii"))
    match = re.search(
        rb"<(?P<prefix>[A-Za-z_][A-Za-z0-9_.-]*:)?"
        + escaped
        + rb"\b[^>]*>.*?</(?P=prefix)"
        + escaped
        + rb"\s*>",
        content,
        re.DOTALL,
    )
    assert match is not None
    return match.group(0)


def first_local_child(element: ET.Element, name: str) -> ET.Element | None:
    return next(
        (child for child in list(element) if local_name(child) == name),
        None,
    )


def synthetic_timing_slide(*, include_build: bool) -> ET.Element:
    build = '<p:bldLst><p:bldP spid="7" grpId="0"/></p:bldLst>' if include_build else ""
    return ET.fromstring(
        f'<p:sld xmlns:p="{PML_NS}" xmlns:a="{DML_NS}"><p:cSld/>'
        "<p:timing><p:tnLst><p:par>"
        '<p:cTn id="1" dur="indefinite" nodeType="tmRoot"><p:childTnLst>'
        '<p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" '
        'nodeType="mainSeq"><p:childTnLst>'
        + effect_xml(
            outer_id=3,
            node_type="clickEffect",
            preset_id=10,
            target_id="7",
            duration=500,
            delay=25,
            behavior="fade",
        )
        + effect_xml(
            outer_id=6,
            node_type="withEffect",
            preset_id=1,
            target_id="8",
            duration=350,
            delay=0,
            behavior="appear",
        )
        + effect_xml(
            outer_id=8,
            node_type="afterEffect",
            preset_id=23,
            target_id="9",
            duration=650,
            delay=10,
            behavior="zoom",
        )
        + "</p:childTnLst></p:cTn></p:seq>"
        '<p:seq><p:cTn id="90" nodeType="interactiveSeq"/></p:seq>'
        '<p:audio><p:cMediaNode><p:cTn id="91"/></p:cMediaNode></p:audio>'
        "</p:childTnLst></p:cTn></p:par></p:tnLst>"
        f"{build}</p:timing></p:sld>"
    )


def effect_xml(
    *,
    outer_id: int,
    node_type: str,
    preset_id: int,
    target_id: str,
    duration: int,
    delay: int,
    behavior: str,
) -> str:
    if behavior == "fade":
        children = (
            '<p:set><p:cBhvr><p:cTn id="4" dur="1"/>'
            f'<p:tgtEl><p:spTgt spid="{target_id}"/></p:tgtEl>'
            "</p:cBhvr></p:set>"
            '<p:animEffect filter="fade" transition="in"><p:cBhvr>'
            f'<p:cTn id="5" dur="{duration}"/><p:tgtEl>'
            f'<p:spTgt spid="{target_id}"/></p:tgtEl>'
            "</p:cBhvr></p:animEffect>"
        )
    elif behavior == "zoom":
        children = (
            "<p:animScale><p:cBhvr>"
            f'<p:cTn id="9" dur="{duration}"/><p:tgtEl>'
            f'<p:spTgt spid="{target_id}"/></p:tgtEl>'
            "</p:cBhvr></p:animScale>"
        )
    else:
        children = (
            "<p:set><p:cBhvr>"
            f'<p:cTn id="7" dur="{duration}"/><p:tgtEl>'
            f'<p:spTgt spid="{target_id}"/></p:tgtEl>'
            "</p:cBhvr></p:set>"
        )
    return (
        "<p:par>"
        f'<p:cTn id="{outer_id}" dur="{duration}" presetClass="entr" '
        f'presetID="{preset_id}" nodeType="{node_type}">'
        f'<p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>'
        f"<p:childTnLst>{children}</p:childTnLst></p:cTn></p:par>"
    )


def animation(
    animation_id: str,
    element_id: str,
    animation_type: str,
    order: int,
    start_mode: str,
    duration: int,
    delay: int,
) -> dict[str, object]:
    return {
        "animationId": animation_id,
        "elementId": element_id,
        "type": animation_type,
        "order": order,
        "durationMs": duration,
        "delayMs": delay,
        "easing": "ease-out",
        "startMode": start_mode,
    }


def slide_with_timing(timing: ET.Element) -> ET.Element:
    slide = ET.Element(f"{{{PML_NS}}}sld")
    ET.SubElement(slide, f"{{{PML_NS}}}cSld")
    slide.append(copy.deepcopy(timing))
    return slide


def canonical(element: ET.Element | None) -> str:
    assert element is not None
    return ET.canonicalize(
        ET.tostring(element, encoding="unicode"),
        with_comments=False,
        rewrite_prefixes=True,
    )


def local_name(element: ET.Element) -> str:
    return str(element.tag).rsplit("}", maxsplit=1)[-1]
