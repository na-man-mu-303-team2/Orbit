from __future__ import annotations

import base64
import binascii
import copy
import difflib
import importlib
import json
import math
import posixpath
import re
import shutil
import subprocess
import zipfile
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from tempfile import TemporaryDirectory
from typing import Any, Literal, cast
from xml.etree import ElementTree as ET

from PIL import Image
from pptx import Presentation
from pydantic import BaseModel, ConfigDict, Field

from app.ai.pptx_design_importer import (
    CANVAS_HEIGHT,
    CANVAS_WIDTH,
    ImportedDesignAsset,
    build_quality_report,
)
from app.ai.pptx_ooxml_vector_importer import (
    OoxmlScale,
    direct_graphic_frame_table,
    import_pptx_design_with_optional_ooxml_vector,
    table_cell_locators,
    table_column_widths,
    table_row_heights,
    table_rows,
    theme_color_map,
)
from app.ai.pptx_motion import (
    parse_slide_motion,
    replace_main_sequence,
    replace_slide_transition,
)

PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
DML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
PKG_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
CONTENT_TYPES_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
IMAGE_REL_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
)
SLIDE_REL_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
)
SLIDE_LAYOUT_REL_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
)
SLIDE_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"
)
P_SP = f"{{{PML_NS}}}sp"
P_PIC = f"{{{PML_NS}}}pic"
P_GRAPHIC_FRAME = f"{{{PML_NS}}}graphicFrame"
A_T = f"{{{DML_NS}}}t"
A_BLIP = f"{{{DML_NS}}}blip"
XML_SPACE = "{http://www.w3.org/XML/1998/namespace}space"
TABLE_GRAPHIC_DATA_URI = "http://schemas.openxmlformats.org/drawingml/2006/table"
SUPPORTED_TABLE_PROPS = {
    "rows",
    "columnWidths",
    "rowHeights",
    "borderColor",
    "borderWidth",
}
SUPPORTED_TABLE_CELL_PROPS = {
    "text",
    "fill",
    "textColor",
    "fontFamily",
    "fontSize",
    "fontWeight",
    "align",
    "verticalAlign",
    "borderColor",
    "borderWidth",
    "colSpan",
    "rowSpan",
}

SUPPORTED_TEXT_PROPS = {
    "text", "runs", "paragraphs", "bodyInset", "fontFamily", "fontSize",
    "fontWeight", "italic", "underline", "color", "align", "verticalAlign",
    "writingMode", "lineHeight", "bullet",
}
SUPPORTED_TEXT_PARAGRAPH_PROPS = {
    "text", "runs", "fontFamily", "fontSize", "fontWeight", "italic",
    "underline", "color", "align", "lineHeight", "spaceBefore", "spaceAfter",
    "indent", "bullet",
}
SUPPORTED_TEXT_RUN_PROPS = {
    "text", "fontFamily", "fontSize", "fontWeight", "italic", "underline",
    "color", "baseline",
}
SUPPORTED_TEXT_STYLE_PROPS = {
    "fontFamily", "fontSize", "fontWeight", "italic", "underline", "color",
    "baseline",
}
MAX_TEXT_DIFF_MATRIX_CELLS = 250_000

ET.register_namespace("p", PML_NS)
ET.register_namespace("a", DML_NS)
ET.register_namespace("r", REL_NS)


class PptxOoxmlGenerationError(RuntimeError):
    pass


class UnsupportedPptxAspectRatioError(PptxOoxmlGenerationError):
    pass


class PptxRenderUnavailableError(PptxOoxmlGenerationError):
    pass


class PptxOoxmlGenerationResult(BaseModel):
    model_config = ConfigDict(populate_by_name=True)

    canvas: dict[str, Any]
    blueprint: dict[str, Any]
    template_blueprint: dict[str, Any] = Field(alias="templateBlueprint")
    quality_report: dict[str, Any] = Field(alias="qualityReport")
    assets: list[ImportedDesignAsset] = Field(default_factory=list)
    warnings: list[str] = Field(default_factory=list)


PptxOoxmlSyncOperationType = Literal[
    "add_slide",
    "delete_slide",
    "add_element",
    "update_element_frame",
    "update_element_props",
    "delete_element",
    "reorder_slides",
]

PptxOoxmlUnsupportedReasonCode = Literal[
    "ADD_SLIDE_FAILED",
    "ADD_SLIDE_LAYOUT_UNSAFE",
    "ADD_ELEMENT_FAILED",
    "ADD_ELEMENT_TYPE_UNSUPPORTED",
    "CROP_CAPABILITY_UNSAFE",
    "DELETE_SLIDE_FAILED",
    "DELETE_SLIDE_LOCATOR_UNSAFE",
    "DELETE_SLIDE_RELATIONSHIP_UNSAFE",
    "RICH_TEXT_CAPABILITY_UNSAFE",
    "ELEMENT_TYPE_MISMATCH",
    "FRAME_FIELDS_UNSUPPORTED",
    "GROUPED_FRAME_UNSUPPORTED",
    "MOTION_REFERENCE_COVERAGE_UNSAFE",
    "OPERATION_TYPE_UNSUPPORTED",
    "PROPS_FIELDS_UNSUPPORTED",
    "PROPS_UPDATE_FAILED",
    "SHAPE_MISSING",
    "SHARED_SHAPE_COHORT_UNSAFE",
    "SLIDE_PART_MISSING",
    "SLIDE_REORDER_LOCATOR_UNSAFE",
    "SLIDE_REORDER_PERMUTATION_INVALID",
    "SLIDE_REORDER_RELATIONSHIP_UNSAFE",
    "LAST_SLIDE_DELETE_FORBIDDEN",
    "SOURCE_MISSING",
    "SOURCE_NOT_WRITABLE",
    "SOURCE_PROVENANCE_UNSAFE",
    "SYNC_RESPONSE_INCOMPLETE",
    "TABLE_CELL_CAPABILITY_UNSAFE",
    "TABLE_STRUCTURE_UNSUPPORTED",
]

PptxOoxmlMotionCoverage = Literal["unknown", "absent", "partial", "complete"]
PptxOoxmlMotionScope = Literal["transition", "animations"]
PptxOoxmlMotionReasonCode = Literal[
    "SLIDE_MOTION_SOURCE_MISSING",
    "SLIDE_MOTION_PAYLOAD_INVALID",
    "SLIDE_TRANSITION_CAPABILITY_UNSAFE",
    "SLIDE_TRANSITION_UNSUPPORTED",
    "SLIDE_ANIMATION_CAPABILITY_UNSAFE",
    "SLIDE_ANIMATION_UNSUPPORTED",
    "SLIDE_ANIMATION_TARGET_UNRESOLVED",
    "SLIDE_MOTION_STRUCTURE_UNSUPPORTED",
]


class PptxOoxmlAppliedOperation(BaseModel):
    model_config = ConfigDict(populate_by_name=True, extra="forbid")

    operation_type: PptxOoxmlSyncOperationType = Field(alias="operationType")
    slide_id: str | None = Field(default=None, alias="slideId")
    element_id: str | None = Field(default=None, alias="elementId")


class PptxOoxmlUnsupportedOperation(PptxOoxmlAppliedOperation):
    reason_code: PptxOoxmlUnsupportedReasonCode = Field(alias="reasonCode")


class PptxOoxmlAppliedSlideMotion(BaseModel):
    model_config = ConfigDict(populate_by_name=True, extra="forbid")

    slide_id: str = Field(alias="slideId")
    transition: bool = False
    animations: bool = False


class PptxOoxmlUnsupportedSlideMotion(BaseModel):
    model_config = ConfigDict(populate_by_name=True, extra="forbid")

    slide_id: str = Field(alias="slideId")
    scope: PptxOoxmlMotionScope
    reason_code: PptxOoxmlMotionReasonCode = Field(alias="reasonCode")


class PptxOoxmlSyncResult(BaseModel):
    model_config = ConfigDict(populate_by_name=True)

    assets: list[ImportedDesignAsset] = Field(default_factory=list)
    element_sources: list[dict[str, Any]] = Field(
        default_factory=list,
        alias="elementSources",
    )
    applied_operations: list[PptxOoxmlAppliedOperation] = Field(
        default_factory=list,
        max_length=500,
        alias="appliedOperations",
    )
    unsupported_operations: list[PptxOoxmlUnsupportedOperation] = Field(
        default_factory=list,
        max_length=500,
        alias="unsupportedOperations",
    )
    applied_slide_motion: list[PptxOoxmlAppliedSlideMotion] = Field(
        default_factory=list,
        alias="appliedSlideMotion",
    )
    unsupported_slide_motion: list[PptxOoxmlUnsupportedSlideMotion] = Field(
        default_factory=list,
        alias="unsupportedSlideMotion",
    )
    warnings: list[str] = Field(default_factory=list)


@dataclass(frozen=True)
class CanvasSpec:
    preset: str
    width: int
    height: int
    aspect_ratio: str

    def payload(self) -> dict[str, Any]:
        return {
            "preset": self.preset,
            "width": self.width,
            "height": self.height,
            "aspectRatio": self.aspect_ratio,
        }


@dataclass(frozen=True)
class PackageFrameScale:
    canvas_width: int
    canvas_height: int
    slide_width_emu: int
    slide_height_emu: int


@dataclass(frozen=True)
class TextRunTemplate:
    start: int
    end: int
    run_properties: ET.Element[Any] | None


@dataclass(frozen=True)
class TextEqualSpan:
    target_start: int
    target_end: int
    source_start: int
    source_end: int


@dataclass(frozen=True)
class TextParagraphTemplate:
    start: int
    end: int
    paragraph: ET.Element[Any]


def generate_pptx_ooxml(
    path: Path,
    file_id: str,
    *,
    render: bool = True,
) -> PptxOoxmlGenerationResult:
    canvas = detect_canvas(path)
    imported = import_pptx_design_with_optional_ooxml_vector(
        path,
        file_id,
        canvas_width=canvas.width,
        canvas_height=canvas.height,
    )
    template_blueprint = prepare_template_blueprint(
        imported.template_blueprint,
        canvas,
        source_file_id=file_id,
        source_canvas=imported.blueprint.get("canvas", {}),
    )
    for index, slide in enumerate(template_blueprint.get("slides", [])):
        imported_slides = imported.blueprint.get("slides", [])
        imported_slide = (
            imported_slides[index]
            if isinstance(imported_slides, list) and index < len(imported_slides)
            else None
        )
        if (
            isinstance(slide, dict)
            and isinstance(imported_slide, dict)
            and isinstance(imported_slide.get("slideId"), str)
        ):
            slide["slideId"] = imported_slide["slideId"]
    warnings = list(imported.warnings)
    package_bytes = path.read_bytes()
    add_imported_ooxml_capabilities(
        imported.blueprint,
        template_blueprint,
        package_bytes,
    )

    assets = [
        package_asset("current_package", package_bytes, f"{safe_file_stem(path)}.pptx")
    ]
    assets.extend(imported.assets)
    if render:
        slide_render_assets = render_pptx_to_png_assets(package_bytes, canvas)
        assets.extend(slide_render_assets)
        fallback_render_assets = slide_render_assets
        if blueprint_has_shape_fallbacks(imported.blueprint):
            fallback_render_assets = render_pptx_to_png_assets(
                strip_text_from_pptx_package(package_bytes),
                canvas,
            )
        assets.extend(
            shape_fallback_assets(imported.blueprint, fallback_render_assets, warnings)
        )

    quality_report = dict(imported.quality_report)
    quality_report["notes"] = [
        note
        for note in quality_report.get("notes", [])
        if note != "pixel renderer unavailable"
    ]
    if render:
        metrics = dict(quality_report.get("metrics", {}))
        metrics["pixelSimilarity"] = None
        quality_report["metrics"] = metrics
        quality_report["notes"].append("OOXML package rendered to slide PNG")
    else:
        quality_report = build_quality_report(
            [{"elements": []} for _slide in template_blueprint["slides"]],
            warnings,
        )

    return PptxOoxmlGenerationResult(
        canvas=canvas.payload(),
        blueprint=imported.blueprint,
        templateBlueprint=template_blueprint,
        qualityReport=quality_report,
        assets=assets,
        warnings=warnings,
    )


def sync_pptx_ooxml(
    path: Path,
    *,
    template_blueprint: dict[str, Any],
    operations: list[dict[str, Any]],
    deck_canvas: dict[str, Any],
    synced_deck_version: int,
    slide_motion: list[dict[str, Any]] | None = None,
    render: bool = True,
) -> PptxOoxmlSyncResult:
    del synced_deck_version

    presentation = Presentation(str(path))
    scale = PackageFrameScale(
        canvas_width=int_value(deck_canvas.get("width"), CANVAS_WIDTH),
        canvas_height=int_value(deck_canvas.get("height"), CANVAS_HEIGHT),
        slide_width_emu=max(1, int(presentation.slide_width or 1)),
        slide_height_emu=max(1, int(presentation.slide_height or 1)),
    )
    (
        package_bytes,
        patch_warnings,
        updated_element_sources,
        applied_operations,
        unsupported_operations,
        applied_slide_motion,
        unsupported_slide_motion,
    ) = apply_patch_operations_to_package(
        path.read_bytes(),
        template_blueprint,
        operations,
        scale,
        slide_motion=slide_motion or [],
    )
    assets = [
        package_asset("current_package", package_bytes, f"{safe_file_stem(path)}.pptx")
    ]
    warnings: list[str] = patch_warnings

    if render:
        try:
            assets.extend(render_pptx_to_png_assets(package_bytes, detect_canvas(path)))
        except PptxRenderUnavailableError as error:
            warnings.append(str(error))

    return PptxOoxmlSyncResult(
        assets=assets,
        elementSources=updated_element_sources,
        appliedOperations=applied_operations,
        unsupportedOperations=unsupported_operations,
        appliedSlideMotion=applied_slide_motion,
        unsupportedSlideMotion=unsupported_slide_motion,
        warnings=warnings,
    )


def detect_canvas(path: Path) -> CanvasSpec:
    presentation = Presentation(str(path))
    width = max(1, int(presentation.slide_width or 1))
    height = max(1, int(presentation.slide_height or 1))
    ratio = width / height
    if abs(ratio - (16 / 9)) <= 0.02:
        return CanvasSpec("wide-16-9", 1920, 1080, "16:9")
    if abs(ratio - (4 / 3)) <= 0.02:
        return CanvasSpec("standard-4-3", 1024, 768, "4:3")
    raise UnsupportedPptxAspectRatioError(
        f"Unsupported PPTX aspect ratio {width}:{height}. Only 16:9 and 4:3 are supported."
    )


def prepare_template_blueprint(
    template_blueprint: dict[str, Any],
    canvas: CanvasSpec,
    *,
    source_file_id: str,
    source_canvas: dict[str, Any] | None = None,
) -> dict[str, Any]:
    prepared = cast(dict[str, Any], json.loads(json.dumps(template_blueprint)))
    prepared["sourcePackageFileId"] = source_file_id
    prepared["currentPackageFileId"] = "asset:current_package"
    source_canvas = source_canvas or {}
    scale_x = canvas.width / int_value(source_canvas.get("width"), CANVAS_WIDTH)
    scale_y = canvas.height / int_value(source_canvas.get("height"), CANVAS_HEIGHT)

    for slide in prepared.get("slides", []):
        if not isinstance(slide, dict):
            continue
        slide_index = int_value(
            slide.get("sourceSlideIndex"),
            int_value(slide.get("slideIndex"), 1),
        )
        slide_part = str(slide.get("sourceSlidePart", ""))
        if not slide_part:
            continue
        slide.setdefault(
            "slideId",
            f"slide_ooxml_{safe_id_component(source_file_id)}_{slide_index}",
        )
        slide["renderAssetFileId"] = f"asset:slide_render_{slide_index}"
        for slot_index, slot in enumerate(slide.get("slots", []), start=1):
            if not isinstance(slot, dict):
                continue
            scale_slot_bounds(slot, scale_x, scale_y)
            if slot.get("usage") == "media-slot":
                slot["replaceMode"] = "replace"
                slot["confidence"] = max(0.65, float(slot.get("confidence", 0)))
            source = slot.setdefault("source", {})
            if isinstance(source, dict):
                source.setdefault("slidePart", slide_part)
                source.setdefault("shapeId", str(slot_index))
    return prepared


def add_imported_ooxml_capabilities(
    blueprint: dict[str, Any],
    template_blueprint: dict[str, Any],
    package_bytes: bytes,
) -> None:
    blueprint_slides = {
        int_value(slide.get("sourceSlideIndex"), index + 1): slide
        for index, slide in enumerate(blueprint.get("slides", []))
        if isinstance(slide, dict)
    }

    try:
        package = zipfile.ZipFile(BytesIO(package_bytes), "r")
    except (OSError, zipfile.BadZipFile):
        package = None

    try:
        for index, slide in enumerate(template_blueprint.get("slides", [])):
            if not isinstance(slide, dict):
                continue
            source_slide_index = int_value(
                slide.get("sourceSlideIndex"),
                int_value(slide.get("slideIndex"), index + 1),
            )
            slide_part = str(slide.get("sourceSlidePart", ""))
            if not slide_part:
                continue
            slide["ooxmlOrigin"] = "imported"
            existing_motion_capabilities = dict_value(
                slide,
                "ooxmlMotionCapabilities",
            )
            existing_coverage = str(
                existing_motion_capabilities.get(
                    "importedMainSequenceCoverage",
                    "",
                )
            )
            motion_coverage = (
                existing_coverage
                if existing_coverage in {"unknown", "absent", "partial", "complete"}
                else imported_main_sequence_coverage(package, slide_part)
            )
            slide["ooxmlMotionCapabilities"] = {
                "transitionWritable": imported_slide_root(package, slide_part)
                is not None,
                "importedMainSequenceCoverage": motion_coverage,
            }

            blueprint_slide = blueprint_slides.get(source_slide_index, {})
            if isinstance(blueprint_slide, dict):
                blueprint_slide["ooxmlOrigin"] = "imported"
                blueprint_slide["ooxmlSourceSlidePart"] = slide_part
                blueprint_slide["ooxmlMotionCapabilities"] = copy.deepcopy(
                    slide["ooxmlMotionCapabilities"]
                )

            element_types = {
                str(element.get("elementId", "")): str(element.get("type", ""))
                for element in blueprint_slide.get("elements", [])
                if isinstance(element, dict)
            }
            element_sources = [
                source
                for source in slide.get("elementSources", [])
                if isinstance(source, dict)
            ]
            shape_cohort_sizes: dict[tuple[str, str], int] = {}
            for source in element_sources:
                cohort_key = (
                    str(source.get("slidePart", "")),
                    str(source.get("shapeId", "")),
                )
                if all(cohort_key):
                    shape_cohort_sizes[cohort_key] = (
                        shape_cohort_sizes.get(cohort_key, 0) + 1
                    )

            slide_root = imported_slide_root(package, slide_part)
            for source in element_sources:
                element_type = element_types.get(str(source.get("elementId", "")), "")
                if element_type:
                    source["elementType"] = element_type
                source["ooxmlOrigin"] = "imported"
                source["ooxmlEditCapabilities"] = imported_element_capabilities(
                    element_type,
                    source,
                    slide_root,
                    shape_cohort_sizes.get(
                        (
                            str(source.get("slidePart", "")),
                            str(source.get("shapeId", "")),
                        ),
                        0,
                    ),
                )
    finally:
        if package is not None:
            package.close()


def imported_slide_root(
    package: zipfile.ZipFile | None,
    slide_part: str,
) -> ET.Element[Any] | None:
    if package is None or not slide_part:
        return None
    try:
        return ET.fromstring(package.read(slide_part))
    except (KeyError, ET.ParseError, OSError):
        return None


def imported_main_sequence_coverage(
    package: zipfile.ZipFile | None,
    slide_part: str,
) -> PptxOoxmlMotionCoverage:
    root = imported_slide_root(package, slide_part)
    if root is None:
        return "unknown"
    motion = parse_slide_motion(root, slide_index=1, shape_targets={})
    if motion.coverage == "absent":
        return "absent"
    return "partial" if motion.coverage == "complete" else motion.coverage


def imported_element_capabilities(
    element_type: str,
    source: dict[str, Any],
    slide_root: ET.Element[Any] | None,
    shape_cohort_size: int,
) -> dict[str, Any]:
    frame_writable = False
    image_source_writable = False
    delete_writable = False
    crop_capability = "none"
    rich_text_capability = "none"
    table_cell_text_writable = False
    if slide_root is not None and bool(source.get("writable", False)):
        shape, _parent = find_shape_by_id(
            slide_root,
            str(source.get("shapeId", "")),
        )
        frame_writable = (
            shape_cohort_size == 1
            and shape is not None
            and not has_group_shape_ancestor(slide_root, shape)
            and element_type != "table"
        )
        delete_writable = (
            shape_cohort_size == 1
            and shape is not None
            and not has_group_shape_ancestor(slide_root, shape)
            and not source.get("fallbackReason")
            and element_type != "table"
        )
        crop_capability = imported_image_crop_capability(
            element_type,
            source,
            shape,
        )
        image_source_writable = (
            crop_capability == "picture"
            and direct_image_blip(shape, source) is not None
        )
        if (
            element_type == "text"
            and shape is not None
            and not source.get("fallbackReason")
        ):
            rich_text_capability = rich_text_capability_for_shape(shape)
        table_cell_text_writable = (
            shape_cohort_size == 1
            and element_type == "table"
            and shape is not None
            and not source.get("fallbackReason")
            and table_cell_text_capability_for_shape(shape, source)
        )
    return {
        "richText": rich_text_capability,
        "crop": crop_capability,
        "tableCellText": table_cell_text_writable,
        "frame": frame_writable,
        "delete": delete_writable,
        "imageSource": image_source_writable,
    }


def table_cell_text_capability_for_shape(
    shape: ET.Element[Any],
    source: dict[str, Any],
) -> bool:
    if (
        shape.tag != P_GRAPHIC_FRAME
        or str(source.get("sourceType", "")) != "table"
        or not bool(source.get("writable", False))
        or source.get("fallbackReason")
    ):
        return False
    locators, diagnostics = table_cell_locators(
        shape,
        slide_index=0,
        shape_id=str(source.get("shapeId", "")),
    )
    declared_locators = source.get("tableCellLocators")
    if diagnostics or not locators or declared_locators != locators:
        return False
    table = direct_graphic_frame_table(shape)
    return table is not None and all(
        table_cell_text_body_is_safe(cell)
        for row in direct_local_children(table, "tr")
        for cell in direct_local_children(row, "tc")
    )


def table_cell_text_body_is_safe(cell: ET.Element[Any]) -> bool:
    body = first_local_child(cell, "txBody")
    if body is None:
        return False
    paragraphs = direct_local_children(body, "p")
    if not paragraphs:
        return False
    if any(
        local_name(child) not in {"bodyPr", "lstStyle", "p", "extLst"}
        for child in list(body)
    ):
        return False
    for paragraph in paragraphs:
        if any(
            local_name(child) not in {"pPr", "r", "endParaRPr"}
            for child in list(paragraph)
        ):
            return False
        runs = direct_local_children(paragraph, "r")
        if len(runs) > 1:
            return False
        if not runs:
            continue
        run = runs[0]
        if any(local_name(child) not in {"rPr", "t"} for child in list(run)):
            return False
        if len(direct_local_children(run, "t")) != 1:
            return False
        if (
            first_local_descendant(run, "hlinkClick") is not None
            or first_local_descendant(run, "hlinkMouseOver") is not None
        ):
            return False
    return True


def rich_text_capability_for_shape(shape: ET.Element[Any]) -> str:
    if shape.tag != P_SP:
        return "none"
    body = first_local_child(shape, "txBody")
    if body is None:
        return "none"

    capability = "full"
    for child in list(body):
        child_name = local_name(child)
        if child_name in {"bodyPr", "lstStyle", "extLst"}:
            continue
        if child_name != "p":
            return "none"
        for paragraph_child in list(child):
            paragraph_child_name = local_name(paragraph_child)
            if paragraph_child_name == "fld":
                return "none"
            if paragraph_child_name not in {"pPr", "r", "br", "endParaRPr"}:
                return "none"
            if paragraph_child_name in {"r", "br"}:
                allowed = {"rPr", "t"} if paragraph_child_name == "r" else {"rPr"}
                if any(local_name(item) not in allowed for item in paragraph_child):
                    return "none"
                if paragraph_child_name == "r" and sum(
                    local_name(item) == "t" for item in paragraph_child
                ) != 1:
                    return "none"
            if first_local_descendant(paragraph_child, "hlinkClick") is not None:
                capability = "style-only"
            if first_local_descendant(paragraph_child, "hlinkMouseOver") is not None:
                capability = "style-only"
    return capability


def imported_image_crop_capability(
    element_type: str,
    source: dict[str, Any],
    shape: ET.Element[Any] | None,
) -> str:
    if (
        element_type != "image"
        or shape is None
        or not bool(source.get("writable", False))
        or source.get("fallbackReason")
    ):
        return "none"
    return image_crop_capability_for_shape(shape, source)


def apply_patch_operations_to_package(
    package_bytes: bytes,
    template_blueprint: dict[str, Any],
    operations: list[dict[str, Any]],
    scale: PackageFrameScale,
    *,
    slide_motion: list[dict[str, Any]] | None = None,
) -> tuple[
    bytes,
    list[str],
    list[dict[str, Any]],
    list[PptxOoxmlAppliedOperation],
    list[PptxOoxmlUnsupportedOperation],
    list[PptxOoxmlAppliedSlideMotion],
    list[PptxOoxmlUnsupportedSlideMotion],
]:
    slide_motion = slide_motion or []
    sources = element_source_map(template_blueprint)
    operations = route_operations_to_source_parts(template_blueprint, operations)
    warnings: list[str] = []
    updated_sources: dict[tuple[str, str], dict[str, Any]] = {}
    applied_operations: list[PptxOoxmlAppliedOperation] = []
    unsupported_operations: list[PptxOoxmlUnsupportedOperation] = []
    applied_slide_motion: list[PptxOoxmlAppliedSlideMotion] = []
    unsupported_slide_motion: list[PptxOoxmlUnsupportedSlideMotion] = []

    motion_failure = motion_reference_failure(
        operations,
        template_blueprint,
        slide_motion,
    )
    if motion_failure is not None:
        return package_bytes, warnings, [], [], [motion_failure], [], []

    redundant_shape_operations, cohort_failure = shared_shape_operation_plan(
        operations,
        sources,
    )
    if cohort_failure is not None:
        return package_bytes, warnings, [], [], [cohort_failure], [], []

    with zipfile.ZipFile(BytesIO(package_bytes), "r") as source:
        source_names = set(source.namelist())
        slide_parts = {
            str(slide.get("sourceSlidePart", ""))
            for slide in template_blueprint.get("slides", [])
            if isinstance(slide, dict) and slide.get("sourceSlidePart")
        }
        slide_parts.update(
            str(item.get("slidePart", ""))
            for item in sources.values()
            if str(item.get("slidePart", ""))
        )
        slide_parts.update(
            str(item.get("sourceSlidePart", ""))
            for item in slide_motion
            if isinstance(item, dict)
            and is_safe_slide_part(str(item.get("sourceSlidePart", "")))
        )
        package_entries = {
            slide_part: source.read(slide_part)
            for slide_part in slide_parts
            if slide_part in source_names
        }
        for slide_part in slide_parts:
            rels_part = rels_part_for_slide_part(slide_part)
            if rels_part in source_names:
                package_entries[rels_part] = source.read(rels_part)
        if "[Content_Types].xml" in source_names:
            package_entries["[Content_Types].xml"] = source.read(
                "[Content_Types].xml"
            )
        for presentation_part in (
            "ppt/presentation.xml",
            "ppt/_rels/presentation.xml.rels",
        ):
            if presentation_part in source_names:
                package_entries[presentation_part] = source.read(presentation_part)
        added_entries: dict[str, bytes] = {}

        for operation_index, operation in enumerate(operations):
            if operation_index in redundant_shape_operations:
                applied_operations.append(applied_operation(operation))
                continue
            reason_code = apply_sync_operation(
                operation,
                sources,
                package_entries,
                added_entries,
                updated_sources,
                scale,
                warnings,
                source,
                template_blueprint,
            )
            if reason_code is None:
                applied_operations.append(applied_operation(operation))
            else:
                unsupported_operations.append(
                    unsupported_operation(operation, reason_code)
                )

        if unsupported_operations:
            return package_bytes, warnings, [], [], unsupported_operations, [], []

        for motion_item in slide_motion:
            motion_result = apply_slide_motion_item(
                motion_item,
                template_blueprint,
                package_entries,
                sources,
            )
            if isinstance(motion_result, PptxOoxmlUnsupportedSlideMotion):
                unsupported_slide_motion.append(motion_result)
            else:
                applied_slide_motion.append(motion_result)

        if unsupported_slide_motion:
            return (
                package_bytes,
                warnings,
                [],
                [],
                [],
                [],
                unsupported_slide_motion,
            )

        animation_touched_parts = {
            str(item.get("sourceSlidePart", ""))
            for item in slide_motion
            if isinstance(item, dict)
            and isinstance(item.get("touched"), dict)
            and item["touched"].get("animations") is True
        }
        for part, content in list(package_entries.items()):
            if (
                part not in source_names
                or not is_safe_slide_part(part)
                or part in animation_touched_parts
            ):
                continue
            preserved_timing = preserve_xml_subtree_bytes(
                source.read(part),
                content,
                "timing",
            )
            if preserved_timing is None:
                first_operation = operations[0] if operations else {}
                failure = unsupported_operation(
                    first_operation,
                    "MOTION_REFERENCE_COVERAGE_UNSAFE",
                )
                return package_bytes, warnings, [], [], [failure], [], []
            package_entries[part] = preserved_timing

        changed_entries = {
            part: content
            for part, content in package_entries.items()
            if part not in source_names or content != source.read(part)
        }
        if not changed_entries and not added_entries:
            return (
                package_bytes,
                warnings,
                list(updated_sources.values()),
                applied_operations,
                unsupported_operations,
                applied_slide_motion,
                unsupported_slide_motion,
            )
        return (
            rewrite_zip(source, changed_entries, added_entries),
            warnings,
            list(updated_sources.values()),
            applied_operations,
            unsupported_operations,
            applied_slide_motion,
            unsupported_slide_motion,
        )


def apply_slide_motion_item(
    item: dict[str, Any],
    template_blueprint: dict[str, Any],
    package_entries: dict[str, bytes],
    current_sources: dict[tuple[str, str], dict[str, Any]] | None = None,
) -> PptxOoxmlAppliedSlideMotion | PptxOoxmlUnsupportedSlideMotion:
    slide_id = str(item.get("slideId", ""))
    slide_part = str(item.get("sourceSlidePart", ""))
    touched = item.get("touched")
    transition_touched = isinstance(touched, dict) and touched.get("transition") is True
    animations_touched = isinstance(touched, dict) and touched.get("animations") is True
    scope: PptxOoxmlMotionScope = (
        "transition" if transition_touched or not animations_touched else "animations"
    )
    matching_template_slides = [
        slide
        for slide in template_blueprint.get("slides", [])
        if isinstance(slide, dict) and source_slide_part(slide) == slide_part
    ]
    template_slide = (
        matching_template_slides[0] if len(matching_template_slides) == 1 else None
    )
    if (
        not slide_id
        or not is_safe_slide_part(slide_part)
        or template_slide is None
        or slide_part not in package_entries
    ):
        return unsupported_slide_motion(
            slide_id,
            scope,
            "SLIDE_MOTION_SOURCE_MISSING",
        )
    if not isinstance(touched, dict) or not (transition_touched or animations_touched):
        return unsupported_slide_motion(
            slide_id,
            scope,
            "SLIDE_MOTION_PAYLOAD_INVALID",
        )

    authoritative_capabilities = dict_value(
        template_slide,
        "ooxmlMotionCapabilities",
    )
    supplied_capabilities = item.get("capabilities")
    if not isinstance(supplied_capabilities, dict):
        return unsupported_slide_motion(
            slide_id,
            scope,
            "SLIDE_MOTION_PAYLOAD_INVALID",
        )
    authoritative_coverage = str(
        authoritative_capabilities.get("importedMainSequenceCoverage", "unknown")
    )
    if transition_touched and supplied_capabilities.get(
        "transitionWritable"
    ) != authoritative_capabilities.get("transitionWritable"):
        return unsupported_slide_motion(
            slide_id,
            "transition",
            "SLIDE_TRANSITION_CAPABILITY_UNSAFE",
        )
    if (
        animations_touched
        and supplied_capabilities.get("importedMainSequenceCoverage")
        != authoritative_coverage
    ):
        return unsupported_slide_motion(
            slide_id,
            "animations",
            "SLIDE_ANIMATION_CAPABILITY_UNSAFE",
        )

    original_slide_xml = package_entries[slide_part]
    try:
        root = ET.fromstring(original_slide_xml)
    except ET.ParseError:
        return unsupported_slide_motion(
            slide_id,
            scope,
            "SLIDE_MOTION_STRUCTURE_UNSUPPORTED",
        )

    if transition_touched:
        if authoritative_capabilities.get("transitionWritable") is not True:
            return unsupported_slide_motion(
                slide_id,
                "transition",
                "SLIDE_TRANSITION_CAPABILITY_UNSAFE",
            )
        if "transition" not in item or (
            item.get("transition") is not None
            and not isinstance(item.get("transition"), dict)
        ):
            return unsupported_slide_motion(
                slide_id,
                "transition",
                "SLIDE_MOTION_PAYLOAD_INVALID",
            )
        try:
            replace_slide_transition(root, item.get("transition"))
        except (TypeError, ValueError):
            return unsupported_slide_motion(
                slide_id,
                "transition",
                "SLIDE_TRANSITION_UNSUPPORTED",
            )

    if animations_touched:
        if authoritative_coverage not in {"absent", "complete"}:
            return unsupported_slide_motion(
                slide_id,
                "animations",
                "SLIDE_ANIMATION_CAPABILITY_UNSAFE",
            )
        animations = item.get("animations")
        if not isinstance(animations, list) or not all(
            isinstance(animation, dict) for animation in animations
        ):
            return unsupported_slide_motion(
                slide_id,
                "animations",
                "SLIDE_MOTION_PAYLOAD_INVALID",
            )
        try:
            applied, diagnostics = replace_main_sequence(
                root,
                cast(list[dict[str, Any]], animations),
                slide_index=int_value(template_slide.get("slideIndex"), 1),
                element_targets=slide_motion_element_targets(
                    template_slide,
                    current_sources,
                ),
            )
        except (TypeError, ValueError):
            return unsupported_slide_motion(
                slide_id,
                "animations",
                "SLIDE_MOTION_PAYLOAD_INVALID",
            )
        if not applied or diagnostics:
            codes = {str(diagnostic.get("code", "")) for diagnostic in diagnostics}
            if "PPTX_MOTION_TARGET_UNRESOLVED" in codes:
                reason = "SLIDE_ANIMATION_TARGET_UNRESOLVED"
            elif "PPTX_MOTION_STRUCTURE_UNSUPPORTED" in codes:
                reason = "SLIDE_MOTION_STRUCTURE_UNSUPPORTED"
            else:
                reason = "SLIDE_ANIMATION_UNSUPPORTED"
            return unsupported_slide_motion(
                slide_id,
                "animations",
                cast(PptxOoxmlMotionReasonCode, reason),
            )

    rewritten = preserve_root_namespace_declarations(
        original_slide_xml,
        xml_bytes(root),
    )
    if rewritten is None:
        return unsupported_slide_motion(
            slide_id,
            scope,
            "SLIDE_MOTION_STRUCTURE_UNSUPPORTED",
        )
    if animations_touched:
        preserved = preserve_excluded_timing_branch_bytes(
            original_slide_xml,
            rewritten,
        )
        if preserved is None:
            return unsupported_slide_motion(
                slide_id,
                "animations",
                "SLIDE_MOTION_STRUCTURE_UNSUPPORTED",
            )
        rewritten = preserved
    if transition_touched and not animations_touched:
        preserved_timing = preserve_xml_subtree_bytes(
            original_slide_xml,
            rewritten,
            "timing",
        )
        if preserved_timing is None:
            return unsupported_slide_motion(
                slide_id,
                "transition",
                "SLIDE_MOTION_STRUCTURE_UNSUPPORTED",
            )
        rewritten = preserved_timing
    package_entries[slide_part] = rewritten
    return PptxOoxmlAppliedSlideMotion(
        slideId=slide_id,
        transition=transition_touched,
        animations=animations_touched,
    )


def slide_motion_element_targets(
    template_slide: dict[str, Any],
    current_sources: dict[tuple[str, str], dict[str, Any]] | None = None,
) -> dict[str, list[str]]:
    targets: dict[str, list[str]] = {}
    slide_part = source_slide_part(template_slide)
    candidates = (
        current_sources.values()
        if current_sources is not None
        else template_slide.get("elementSources", [])
    )
    for source in candidates:
        if not isinstance(source, dict) or not bool(source.get("writable", False)):
            continue
        if (
            current_sources is not None
            and str(source.get("slidePart", "")) != slide_part
        ):
            continue
        element_id = str(source.get("elementId", ""))
        shape_id = str(source.get("shapeId", ""))
        if not element_id or not shape_id:
            continue
        values = targets.setdefault(element_id, [])
        if shape_id not in values:
            values.append(shape_id)
    return targets


def unsupported_slide_motion(
    slide_id: str,
    scope: PptxOoxmlMotionScope,
    reason_code: PptxOoxmlMotionReasonCode,
) -> PptxOoxmlUnsupportedSlideMotion:
    return PptxOoxmlUnsupportedSlideMotion(
        slideId=slide_id or "unknown",
        scope=scope,
        reasonCode=reason_code,
    )


def preserve_xml_subtree_bytes(
    original: bytes,
    rewritten: bytes,
    local_name_value: str,
) -> bytes | None:
    namespaced_rewritten = preserve_root_namespace_declarations(
        original,
        rewritten,
    )
    if namespaced_rewritten is None:
        return None
    original_match = xml_subtree_match(original, local_name_value)
    rewritten_match = xml_subtree_match(namespaced_rewritten, local_name_value)
    if original_match is None:
        return namespaced_rewritten
    if rewritten_match is None:
        return None
    return (
        namespaced_rewritten[: rewritten_match.start()]
        + original[original_match.start() : original_match.end()]
        + namespaced_rewritten[rewritten_match.end() :]
    )


def preserve_root_namespace_declarations(
    original: bytes,
    rewritten: bytes,
) -> bytes | None:
    original_root = xml_root_opening_tag(original)
    rewritten_root = xml_root_opening_tag(rewritten)
    if original_root is None or rewritten_root is None:
        return None
    original_namespaces = xml_namespace_declarations(original_root.group(0))
    rewritten_namespaces = xml_namespace_declarations(rewritten_root.group(0))
    missing: list[bytes] = []
    for prefix, (uri, declaration) in original_namespaces.items():
        current = rewritten_namespaces.get(prefix)
        if current is not None:
            if current[0] != uri:
                return None
            continue
        missing.append(declaration)
    if not missing:
        return rewritten
    insertion = b"".join(b" " + declaration.strip() for declaration in missing)
    insert_at = rewritten_root.end() - 1
    return rewritten[:insert_at] + insertion + rewritten[insert_at:]


def xml_root_opening_tag(content: bytes) -> re.Match[bytes] | None:
    return re.search(
        rb"<[A-Za-z_][A-Za-z0-9_.:-]*\b[^>]*>",
        content,
    )


def xml_namespace_declarations(
    opening_tag: bytes,
) -> dict[bytes, tuple[bytes, bytes]]:
    pattern = re.compile(
        rb"\s+xmlns(?::(?P<prefix>[A-Za-z_][A-Za-z0-9_.-]*))?"
        rb"\s*=\s*(?P<quote>['\"])(?P<uri>.*?)(?P=quote)"
    )
    return {
        match.group("prefix") or b"": (match.group("uri"), match.group(0))
        for match in pattern.finditer(opening_tag)
    }


def preserve_excluded_timing_branch_bytes(
    original: bytes,
    rewritten: bytes,
) -> bytes | None:
    original_ranges = excluded_timing_branch_ranges(original)
    rewritten_ranges = excluded_timing_branch_ranges(rewritten)
    if [item[0] for item in original_ranges] != [item[0] for item in rewritten_ranges]:
        return None
    result = rewritten
    for original_item, rewritten_item in zip(
        reversed(original_ranges),
        reversed(rewritten_ranges),
        strict=True,
    ):
        _, original_start, original_end = original_item
        _, rewritten_start, rewritten_end = rewritten_item
        result = (
            result[:rewritten_start]
            + original[original_start:original_end]
            + result[rewritten_end:]
        )
    return result


def excluded_timing_branch_bytes(content: bytes) -> list[bytes]:
    return [
        content[start:end] for _, start, end in excluded_timing_branch_ranges(content)
    ]


def excluded_timing_branch_ranges(
    content: bytes,
) -> list[tuple[str, int, int]]:
    tag_pattern = re.compile(
        rb"<(?P<closing>/)?(?P<name>[A-Za-z_][A-Za-z0-9_.:-]*)"
        rb"\b(?P<body>[^<>]*?)(?P<self_closing>/)?>",
        re.DOTALL,
    )
    nodes: list[XmlByteNode] = []
    stack: list[int] = []

    for match in tag_pattern.finditer(content):
        name = match.group("name")
        local = name.rsplit(b":", 1)[-1].decode("ascii")
        if match.group("closing"):
            if not stack:
                continue
            node_index = stack.pop()
            node = nodes[node_index]
            if node.name != name:
                return []
            node.end = match.end()
            continue

        body = match.group("body") or b""
        nodes.append(
            XmlByteNode(
                name=name,
                local=local,
                body=body,
                start=match.start(),
                end=match.end() if match.group("self_closing") else -1,
                parent=stack[-1] if stack else None,
            )
        )
        if not match.group("self_closing"):
            stack.append(len(nodes) - 1)

    if stack or any(node.end < 0 for node in nodes):
        return []

    selected: dict[int, str] = {}
    for node_index, node in enumerate(nodes):
        if node.local == "cTn" and xml_attribute_equals(
            node.body,
            "nodeType",
            "interactiveSeq",
        ):
            branch_index = nearest_ancestor(nodes, node_index, "seq")
            selected[branch_index if branch_index is not None else node_index] = (
                "interactiveSeq"
            )
            continue
        is_media_timeline = node.local == "cTn" and xml_attribute_equals(
            node.body,
            "presetClass",
            "mediacall",
        )
        if is_media_timeline:
            branch_index = nearest_ancestor(nodes, node_index, "par")
            selected[branch_index if branch_index is not None else node_index] = "media"
            continue
        if node.local in {"audio", "video", "cmd"}:
            media_timeline = nearest_ancestor_with_attribute(
                nodes,
                node_index,
                local="cTn",
                attribute="presetClass",
                value="mediacall",
            )
            if media_timeline is not None:
                branch_index = nearest_ancestor(nodes, media_timeline, "par")
                selected[
                    branch_index if branch_index is not None else media_timeline
                ] = "media"
            else:
                selected[node_index] = "media"

    selected_indexes = set(selected)
    outermost_indexes = [
        node_index
        for node_index in selected_indexes
        if not any(
            ancestor in selected_indexes
            for ancestor in ancestor_indexes(nodes, node_index)
        )
    ]
    return [
        (selected[node_index], nodes[node_index].start, nodes[node_index].end)
        for node_index in sorted(
            outermost_indexes,
            key=lambda index: nodes[index].start,
        )
    ]


@dataclass
class XmlByteNode:
    name: bytes
    local: str
    body: bytes
    start: int
    end: int
    parent: int | None


def xml_attribute_equals(body: bytes, name: str, value: str) -> bool:
    return (
        re.search(
            rb"\b"
            + re.escape(name.encode("ascii"))
            + rb"\s*=\s*['\"]"
            + re.escape(value.encode("ascii"))
            + rb"['\"]",
            body,
        )
        is not None
    )


def ancestor_indexes(nodes: list[XmlByteNode], node_index: int) -> list[int]:
    indexes: list[int] = []
    parent = nodes[node_index].parent
    while parent is not None:
        indexes.append(parent)
        parent = nodes[parent].parent
    return indexes


def nearest_ancestor(
    nodes: list[XmlByteNode],
    node_index: int,
    local: str,
) -> int | None:
    return next(
        (
            ancestor
            for ancestor in ancestor_indexes(nodes, node_index)
            if nodes[ancestor].local == local
        ),
        None,
    )


def nearest_ancestor_with_attribute(
    nodes: list[XmlByteNode],
    node_index: int,
    *,
    local: str,
    attribute: str,
    value: str,
) -> int | None:
    return next(
        (
            ancestor
            for ancestor in ancestor_indexes(nodes, node_index)
            if nodes[ancestor].local == local
            and xml_attribute_equals(nodes[ancestor].body, attribute, value)
        ),
        None,
    )


def xml_subtree_match(content: bytes, local_name_value: str) -> re.Match[bytes] | None:
    escaped = re.escape(local_name_value.encode("ascii"))
    pattern = re.compile(
        rb"<(?P<prefix>[A-Za-z_][A-Za-z0-9_.-]*:)?"
        + escaped
        + rb"\b(?:[^>]*/\s*>|[^>]*>.*?</(?P=prefix)"
        + escaped
        + rb"\s*>)",
        re.DOTALL,
    )
    return pattern.search(content)




def apply_sync_operation(
    operation: dict[str, Any],
    sources: dict[tuple[str, str], dict[str, Any]],
    package_entries: dict[str, bytes],
    added_entries: dict[str, bytes],
    updated_sources: dict[tuple[str, str], dict[str, Any]],
    scale: PackageFrameScale,
    warnings: list[str],
    source_package: zipfile.ZipFile,
    template_blueprint: dict[str, Any],
) -> PptxOoxmlUnsupportedReasonCode | None:
    operation_type = str(operation.get("type", ""))
    if operation_type == "add_slide":
        added_sources, reason_code = add_authored_slide_to_package(
            operation,
            package_entries,
            added_entries,
            scale,
            warnings,
            source_package,
            template_blueprint,
        )
        if reason_code is not None:
            return reason_code
        for added_source in added_sources:
            added_key = (
                str(added_source["slidePart"]),
                str(added_source["elementId"]),
            )
            sources[added_key] = added_source
            updated_sources[added_key] = added_source
        return None
    if operation_type == "reorder_slides":
        return reorder_presentation_slides(
            operation,
            package_entries,
            source_package,
            template_blueprint,
        )
    if operation_type == "delete_slide":
        return delete_presentation_slide(
            operation,
            sources,
            updated_sources,
            package_entries,
            source_package,
            template_blueprint,
        )
    element_id = operation_element_id(operation)
    operation_slide_part = slide_part_for_operation(
        operation,
        template_blueprint,
    )

    if operation_type == "add_element":
        element = operation.get("element")
        if not isinstance(element, dict):
            return "ADD_ELEMENT_FAILED"
        if not operation_slide_part or operation_slide_part not in package_entries:
            return "SLIDE_PART_MISSING"
        if str(element.get("type", "")) not in {"text", "rect", "image", "table"}:
            return "ADD_ELEMENT_TYPE_UNSUPPORTED"
        if add_element_has_unsupported_props(element):
            return "PROPS_FIELDS_UNSUPPORTED"
        element_source = add_element_to_slide_xml(
            operation_slide_part,
            element,
            package_entries,
            added_entries,
            scale,
            warnings,
        )
        if element_source is None:
            return "ADD_ELEMENT_FAILED"
        added_key = (
            str(element_source["slidePart"]),
            str(element_source["elementId"]),
        )
        sources[added_key] = element_source
        updated_sources[added_key] = element_source
        return None

    source_key = (operation_slide_part, element_id)
    source = sources.get(source_key)
    if not source:
        warnings.append(f"OOXML source missing for {element_id}.")
        return "SOURCE_MISSING"
    if not bool(source.get("writable", False)):
        warnings.append(f"OOXML source is locked for {element_id}.")
        return "SOURCE_NOT_WRITABLE"
    if source.get("ooxmlOrigin") not in {"imported", "authored"}:
        warnings.append(f"OOXML source provenance is unsafe for {element_id}.")
        return "SOURCE_PROVENANCE_UNSAFE"

    slide_part = str(source.get("slidePart", ""))
    slide_xml = package_entries.get(slide_part)
    if slide_xml is None:
        warnings.append(f"OOXML slide part missing for {element_id}.")
        return "SLIDE_PART_MISSING"

    root = ET.fromstring(slide_xml)
    shape, parent = find_shape_by_id(root, str(source.get("shapeId", "")))
    if shape is None:
        warnings.append(f"OOXML shape missing for {element_id}.")
        return "SHAPE_MISSING"

    shape_changed = False
    if operation_type == "update_element_props":
        props = operation.get("props")
        if not isinstance(props, dict) or not props:
            return "PROPS_FIELDS_UNSUPPORTED"
        source_shape_cohort_size = sum(
            1
            for candidate in sources.values()
            if str(candidate.get("slidePart", "")) == slide_part
            and str(candidate.get("shapeId", "")) == str(source.get("shapeId", ""))
        )
        props_reason = validate_source_props_update(
            source,
            shape,
            props,
            scale,
            source_package,
            source_shape_cohort_size,
        )
        if props_reason is not None:
            return props_reason
        shape_changed = update_shape_props(
            shape,
            props,
            source,
            scale,
            slide_part,
            package_entries,
            added_entries,
            updated_sources,
            source_key,
            warnings,
            element_id,
        )
        if not shape_changed:
            return "PROPS_UPDATE_FAILED"
    elif operation_type == "update_element_frame":
        frame = operation.get("frame")
        if not isinstance(frame, dict) or not frame:
            return "FRAME_FIELDS_UNSUPPORTED"
        if set(frame) - {
            "role",
            "x",
            "y",
            "width",
            "height",
            "rotation",
            "opacity",
            "zIndex",
            "locked",
            "visible",
        }:
            return "FRAME_FIELDS_UNSUPPORTED"
        opacity = frame.get("opacity", 1)
        if (
            isinstance(opacity, bool)
            or not isinstance(opacity, (int, float))
            or not math.isfinite(float(opacity))
            or not 0 <= float(opacity) <= 1
        ):
            return "FRAME_FIELDS_UNSUPPORTED"
        visible = frame.get("visible", True)
        if not isinstance(visible, bool):
            return "FRAME_FIELDS_UNSUPPORTED"
        geometry_fields = set(frame) & {"x", "y", "width", "height", "rotation"}
        if geometry_fields and has_group_shape_ancestor(root, shape):
            warnings.append(f"OOXML grouped frame sync skipped for {element_id}.")
            return "GROUPED_FRAME_UNSUPPORTED"
        capabilities = dict_value(source, "ooxmlEditCapabilities")
        source_shape_cohort_size = sum(
            1
            for candidate in sources.values()
            if str(candidate.get("slidePart", "")) == slide_part
            and str(candidate.get("shapeId", "")) == str(source.get("shapeId", ""))
        )
        safe_legacy_imported_frame = (
            source.get("ooxmlOrigin") == "imported"
            and source_shape_cohort_size == 1
            and source.get("elementType") != "table"
            and not has_group_shape_ancestor(root, shape)
        )
        if (
            source.get("ooxmlOrigin") == "imported"
            and not capabilities.get("frame")
            and not safe_legacy_imported_frame
        ):
            return "FRAME_FIELDS_UNSUPPORTED"
        if geometry_fields:
            update_shape_frame(shape, frame, scale)
            if (
                source.get("elementType") == "table"
                and source.get("ooxmlOrigin") == "authored"
                and not resize_authored_table_tracks_to_frame(shape)
            ):
                return "TABLE_STRUCTURE_UNSUPPORTED"
            shape_changed = True
        if "opacity" in frame and float(opacity) != 1:
            if source.get("elementType") != "image" or not set_picture_opacity(
                shape,
                source,
                float(opacity),
            ):
                return "FRAME_FIELDS_UNSUPPORTED"
            shape_changed = True
        elif "opacity" in frame and source.get("elementType") == "image":
            if not set_picture_opacity(shape, source, 1):
                return "FRAME_FIELDS_UNSUPPORTED"
            shape_changed = True
        if "visible" in frame:
            if not set_shape_visibility(shape, visible):
                return "FRAME_FIELDS_UNSUPPORTED"
            shape_changed = True
        if "zIndex" in frame:
            if parent is None or not reorder_visual_shape(
                parent,
                shape,
                frame["zIndex"],
            ):
                return "FRAME_FIELDS_UNSUPPORTED"
            shape_changed = True
    elif operation_type == "delete_element":
        if parent is not None:
            parent.remove(shape)
            remove_shape_sources(
                sources,
                updated_sources,
                slide_part,
                str(source.get("shapeId", "")),
            )
            shape_changed = True
    else:
        return "OPERATION_TYPE_UNSUPPORTED"

    if shape_changed:
        package_entries[slide_part] = xml_bytes(root)
    return None


def remove_shape_sources(
    sources: dict[tuple[str, str], dict[str, Any]],
    updated_sources: dict[tuple[str, str], dict[str, Any]],
    slide_part: str,
    shape_id: str,
) -> None:
    removed_keys = [
        source_key
        for source_key, candidate in sources.items()
        if str(candidate.get("slidePart", "")) == slide_part
        and str(candidate.get("shapeId", "")) == shape_id
    ]
    for source_key in removed_keys:
        sources.pop(source_key, None)
        updated_sources.pop(source_key, None)


def reorder_presentation_slides(
    operation: dict[str, Any],
    package_entries: dict[str, bytes],
    source_package: zipfile.ZipFile,
    template_blueprint: dict[str, Any],
) -> PptxOoxmlUnsupportedReasonCode | None:
    raw_slide_orders = operation.get("slideOrders")
    if not isinstance(raw_slide_orders, list) or not raw_slide_orders:
        return "SLIDE_REORDER_PERMUTATION_INVALID"

    blueprint_locators: dict[str, str] = {}
    blueprint_parts: set[str] = set()
    for raw_blueprint_slide in template_blueprint.get("slides", []):
        if not isinstance(raw_blueprint_slide, dict):
            return "SLIDE_REORDER_LOCATOR_UNSAFE"
        slide_id = raw_blueprint_slide.get("slideId")
        source_slide_part = raw_blueprint_slide.get("sourceSlidePart")
        if (
            not isinstance(slide_id, str)
            or not slide_id
            or not isinstance(source_slide_part, str)
            or not source_slide_part
        ):
            return "SLIDE_REORDER_LOCATOR_UNSAFE"
        if slide_id in blueprint_locators or source_slide_part in blueprint_parts:
            return "SLIDE_REORDER_LOCATOR_UNSAFE"
        blueprint_locators[slide_id] = source_slide_part
        blueprint_parts.add(source_slide_part)

    requested: list[tuple[int, str, str]] = []
    for raw_slide_order in raw_slide_orders:
        if not isinstance(raw_slide_order, dict):
            return "SLIDE_REORDER_PERMUTATION_INVALID"
        slide_id = raw_slide_order.get("slideId")
        order = raw_slide_order.get("order")
        source_slide_part = raw_slide_order.get("sourceSlidePart")
        if (
            not isinstance(slide_id, str)
            or not slide_id
            or not isinstance(order, int)
            or isinstance(order, bool)
        ):
            return "SLIDE_REORDER_PERMUTATION_INVALID"
        if (
            not isinstance(source_slide_part, str)
            or not source_slide_part
        ):
            return "SLIDE_REORDER_LOCATOR_UNSAFE"
        requested.append((order, slide_id, source_slide_part))

    requested_orders = [item[0] for item in requested]
    requested_ids = [item[1] for item in requested]
    expected_orders = set(range(1, len(requested) + 1))
    if len(set(requested_ids)) != len(requested_ids) or set(
        requested_orders
    ) != expected_orders:
        return "SLIDE_REORDER_PERMUTATION_INVALID"
    if any(blueprint_locators.get(item[1]) != item[2] for item in requested):
        return "SLIDE_REORDER_LOCATOR_UNSAFE"

    presentation_xml = package_entries.get("ppt/presentation.xml")
    presentation_rels_xml = package_entries.get(
        "ppt/_rels/presentation.xml.rels"
    )
    if presentation_xml is None or presentation_rels_xml is None:
        return "SLIDE_REORDER_RELATIONSHIP_UNSAFE"

    try:
        presentation_root = ET.fromstring(presentation_xml)
        relationships_root = ET.fromstring(presentation_rels_xml)
    except ET.ParseError:
        return "SLIDE_REORDER_RELATIONSHIP_UNSAFE"

    slide_id_list = presentation_root.find(f"{{{PML_NS}}}sldIdLst")
    if slide_id_list is None:
        return "SLIDE_REORDER_RELATIONSHIP_UNSAFE"
    slide_id_nodes = list(slide_id_list)
    if not slide_id_nodes or any(
        node.tag != f"{{{PML_NS}}}sldId" for node in slide_id_nodes
    ):
        return "SLIDE_REORDER_RELATIONSHIP_UNSAFE"

    relationships_by_id: dict[str, ET.Element[Any]] = {}
    for relationship in relationships_root:
        relationship_id = str(relationship.get("Id", ""))
        if not relationship_id or relationship_id in relationships_by_id:
            return "SLIDE_REORDER_RELATIONSHIP_UNSAFE"
        relationships_by_id[relationship_id] = relationship

    slide_nodes_by_part: dict[str, ET.Element[Any]] = {}
    source_names = set(source_package.namelist())
    for slide_id_node in slide_id_nodes:
        relationship_id = str(slide_id_node.get(f"{{{REL_NS}}}id", ""))
        mapped_relationship = relationships_by_id.get(relationship_id)
        if mapped_relationship is None:
            return "SLIDE_REORDER_RELATIONSHIP_UNSAFE"
        target = str(mapped_relationship.get("Target", ""))
        slide_part = resolve_relationship_part("ppt/presentation.xml", target)
        if (
            not slide_part.startswith("ppt/slides/")
            or (
                slide_part not in source_names
                and slide_part not in package_entries
            )
            or slide_part in slide_nodes_by_part
        ):
            return "SLIDE_REORDER_RELATIONSHIP_UNSAFE"
        slide_nodes_by_part[slide_part] = slide_id_node

    requested_parts = [item[2] for item in sorted(requested)]
    if len(requested_parts) != len(slide_id_nodes):
        return "SLIDE_REORDER_PERMUTATION_INVALID"
    if len(set(requested_parts)) != len(requested_parts) or set(
        requested_parts
    ) != set(slide_nodes_by_part):
        return "SLIDE_REORDER_LOCATOR_UNSAFE"

    slide_id_list[:] = [slide_nodes_by_part[part] for part in requested_parts]
    package_entries["ppt/presentation.xml"] = xml_bytes(presentation_root)
    return None


def delete_presentation_slide(
    operation: dict[str, Any],
    sources: dict[tuple[str, str], dict[str, Any]],
    updated_sources: dict[tuple[str, str], dict[str, Any]],
    package_entries: dict[str, bytes],
    source_package: zipfile.ZipFile,
    template_blueprint: dict[str, Any],
) -> PptxOoxmlUnsupportedReasonCode | None:
    slide_id = operation_slide_id(operation)
    slide_part = slide_part_for_operation(operation, template_blueprint)
    if not slide_id or not is_safe_slide_part(slide_part):
        return "DELETE_SLIDE_LOCATOR_UNSAFE"
    matching_blueprint_parts = [
        str(slide.get("sourceSlidePart", ""))
        for slide in template_blueprint.get("slides", [])
        if isinstance(slide, dict) and slide.get("slideId") == slide_id
    ]
    if matching_blueprint_parts != [slide_part]:
        return "DELETE_SLIDE_LOCATOR_UNSAFE"

    presentation_xml = package_entries.get("ppt/presentation.xml")
    presentation_rels_xml = package_entries.get(
        "ppt/_rels/presentation.xml.rels"
    )
    content_types_xml = package_entries.get("[Content_Types].xml")
    if presentation_xml is None or presentation_rels_xml is None:
        return "DELETE_SLIDE_RELATIONSHIP_UNSAFE"
    try:
        presentation_root = ET.fromstring(presentation_xml)
        relationships_root = ET.fromstring(presentation_rels_xml)
        content_types_root = (
            ET.fromstring(content_types_xml) if content_types_xml is not None else None
        )
    except ET.ParseError:
        return "DELETE_SLIDE_RELATIONSHIP_UNSAFE"

    slide_id_list = presentation_root.find(f"{{{PML_NS}}}sldIdLst")
    if slide_id_list is None:
        return "DELETE_SLIDE_RELATIONSHIP_UNSAFE"
    slide_id_nodes = list(slide_id_list)
    if len(slide_id_nodes) <= 1:
        return "LAST_SLIDE_DELETE_FORBIDDEN"
    relationships_by_id = {
        str(relationship.get("Id", "")): relationship
        for relationship in relationships_root
        if relationship.get("Id")
    }
    matching_nodes: list[tuple[ET.Element[Any], ET.Element[Any]]] = []
    for slide_id_node in slide_id_nodes:
        relationship_id = str(slide_id_node.get(f"{{{REL_NS}}}id", ""))
        relationship = relationships_by_id.get(relationship_id)
        if relationship is None:
            return "DELETE_SLIDE_RELATIONSHIP_UNSAFE"
        target = str(relationship.get("Target", ""))
        if resolve_relationship_part("ppt/presentation.xml", target) == slide_part:
            matching_nodes.append((slide_id_node, relationship))
    if len(matching_nodes) != 1 or slide_part not in source_package.namelist():
        return "DELETE_SLIDE_RELATIONSHIP_UNSAFE"

    slide_id_node, relationship = matching_nodes[0]
    slide_id_list.remove(slide_id_node)
    relationships_root.remove(relationship)
    if content_types_root is not None:
        part_name = f"/{slide_part}"
        for child in list(content_types_root):
            if child.tag.endswith("Override") and child.get("PartName") == part_name:
                content_types_root.remove(child)

    removed_source_keys = [
        key
        for key, source in sources.items()
        if str(source.get("slidePart", "")) == slide_part
    ]
    for key in removed_source_keys:
        sources.pop(key, None)
        updated_sources.pop(key, None)
    package_entries["ppt/presentation.xml"] = xml_bytes(presentation_root)
    package_entries["ppt/_rels/presentation.xml.rels"] = xml_bytes(
        relationships_root
    )
    if content_types_root is not None:
        package_entries["[Content_Types].xml"] = xml_bytes(content_types_root)
    return None


def resolve_relationship_part(source_part: str, target: str) -> str:
    if target.startswith("/"):
        return posixpath.normpath(target).lstrip("/")
    return posixpath.normpath(posixpath.join(posixpath.dirname(source_part), target))


def add_element_has_unsupported_props(element: dict[str, Any]) -> bool:
    props = dict_value(element, "props")
    element_type = str(element.get("type", ""))
    if (
        float(element.get("opacity", 1)) != 1
        or bool(element.get("locked", False))
        or not bool(element.get("visible", True))
        or float(element.get("rotation", 0)) != 0
    ):
        return True
    if element_type == "text":
        return not valid_text_props(props)
    if element_type == "rect":
        return not valid_rect_props(props)
    if element_type == "image":
        crop = props.get("crop")
        return (
            set(props) - {"src", "alt", "fit", "focusX", "focusY", "crop"}
            != set()
            or ("crop" in props and normalized_image_crop(crop) is None)
            or props.get("fit", "contain") != "contain"
            or float(props.get("focusX", 0.5)) != 0.5
            or float(props.get("focusY", 0.5)) != 0.5
        )
    if element_type == "table":
        return not valid_table_props(props)
    return True


def validate_source_props_update(
    source: dict[str, Any],
    shape: ET.Element[Any],
    props: dict[str, Any],
    scale: PackageFrameScale,
    source_package: zipfile.ZipFile,
    source_shape_cohort_size: int,
) -> PptxOoxmlUnsupportedReasonCode | None:
    prop_names = set(props)
    element_type = str(source.get("elementType", ""))
    if prop_names and prop_names.issubset(
        {"fill", "stroke", "strokeWidth", "borderRadius"}
    ):
        if (
            element_type != "rect"
            or shape.tag != P_SP
            or source.get("ooxmlOrigin") != "authored"
        ):
            return "ELEMENT_TYPE_MISMATCH"
        return None if valid_rect_props(props) else "PROPS_FIELDS_UNSUPPORTED"
    if prop_names and prop_names.issubset(SUPPORTED_TABLE_PROPS):
        if element_type != "table" or shape.tag != P_GRAPHIC_FRAME:
            return "ELEMENT_TYPE_MISMATCH"
        declared_capability = dict_value(source, "ooxmlEditCapabilities").get(
            "tableCellText"
        )
        if (
            declared_capability is not True
            or source_shape_cohort_size != 1
            or not table_cell_text_capability_for_shape(shape, source)
        ):
            return "TABLE_CELL_CAPABILITY_UNSAFE"
        if not valid_table_props(props):
            return "TABLE_STRUCTURE_UNSUPPORTED"
        if source.get("ooxmlOrigin") == "authored":
            return None
        if source.get("ooxmlOrigin") != "imported":
            return "TABLE_CELL_CAPABILITY_UNSAFE"
        return validate_imported_table_props_update(
            shape,
            props,
            scale,
            source_package,
        )
    if prop_names and prop_names.issubset(SUPPORTED_TEXT_PROPS):
        if element_type != "text" or shape.tag != P_SP:
            return "ELEMENT_TYPE_MISMATCH"
        if not valid_text_props(props):
            return "PROPS_FIELDS_UNSUPPORTED"
        capability = dict_value(source, "ooxmlEditCapabilities").get("richText")
        actual_capability = rich_text_capability_for_shape(shape)
        if capability not in {"full", "style-only"} or capability != actual_capability:
            return "RICH_TEXT_CAPABILITY_UNSAFE"
        if capability == "style-only":
            if text_props_has_content_projection(props):
                target_text = canonical_text_value(props)
                if target_text is None or target_text != text_body_value(shape):
                    return "RICH_TEXT_CAPABILITY_UNSAFE"
            if set(props) != {"text"} and text_props_has_content_projection(props):
                target = canonical_text_paragraphs(props)
                body = first_local_child(shape, "txBody")
                if body is None or target is None or not style_only_paragraphs_match(body, target):
                    return "RICH_TEXT_CAPABILITY_UNSAFE"
        return None
    if prop_names and prop_names.issubset({"src", "alt", "crop"}):
        capabilities = dict_value(source, "ooxmlEditCapabilities")
        if element_type != "image":
            return "ELEMENT_TYPE_MISMATCH"
        if {"src", "alt"}.intersection(prop_names) and (
            shape.tag != P_PIC
            or source.get("ooxmlOrigin") == "imported"
            and not capabilities.get("imageSource")
        ):
            return "ELEMENT_TYPE_MISMATCH"
        if "crop" in props:
            crop = props.get("crop")
            if crop is not None and normalized_image_crop(crop) is None:
                return "PROPS_FIELDS_UNSUPPORTED"
            capability = capabilities.get("crop")
            expected_capability = image_crop_capability_for_shape(shape, source)
            if capability != expected_capability or capability == "none":
                return "CROP_CAPABILITY_UNSAFE"
        return None
    return "PROPS_FIELDS_UNSUPPORTED"


def validate_imported_table_props_update(
    shape: ET.Element[Any],
    props: dict[str, Any],
    scale: PackageFrameScale,
    source_package: zipfile.ZipFile,
) -> PptxOoxmlUnsupportedReasonCode | None:
    table = direct_graphic_frame_table(shape)
    if table is None:
        return "TABLE_CELL_CAPABILITY_UNSAFE"
    ooxml_scale = OoxmlScale(
        canvas_width=scale.canvas_width,
        canvas_height=scale.canvas_height,
        slide_width_emu=scale.slide_width_emu,
        slide_height_emu=scale.slide_height_emu,
    )
    actual_rows = table_rows(table, ooxml_scale, theme_color_map(source_package))
    target_rows = props.get("rows")
    if (
        not isinstance(target_rows, list)
        or len(target_rows) != len(actual_rows)
        or any(
            not isinstance(target_row, list) or len(target_row) != len(actual_row)
            for target_row, actual_row in zip(target_rows, actual_rows, strict=True)
        )
    ):
        return "TABLE_STRUCTURE_UNSUPPORTED"

    changed_text_count = 0
    xml_rows = direct_local_children(table, "tr")
    for row_index, (target_row, actual_row) in enumerate(
        zip(target_rows, actual_rows, strict=True)
    ):
        xml_cells = direct_local_children(xml_rows[row_index], "tc")
        for column_index, (target_cell, actual_cell) in enumerate(
            zip(target_row, actual_row, strict=True)
        ):
            if not isinstance(target_cell, dict):
                return "TABLE_STRUCTURE_UNSUPPORTED"
            if not table_cell_non_text_equal(target_cell, actual_cell):
                return "TABLE_STRUCTURE_UNSUPPORTED"
            if str(target_cell.get("text", "")) != str(actual_cell.get("text", "")):
                changed_text_count += 1
                if not table_cell_text_can_set(
                    xml_cells[column_index],
                    str(target_cell.get("text", "")),
                ):
                    return "TABLE_STRUCTURE_UNSUPPORTED"
    if changed_text_count != 1:
        return "TABLE_STRUCTURE_UNSUPPORTED"

    if "columnWidths" in props and not numeric_track_values_equal(
        props.get("columnWidths"), table_column_widths(table, ooxml_scale)
    ):
        return "TABLE_STRUCTURE_UNSUPPORTED"
    if "rowHeights" in props and not numeric_track_values_equal(
        props.get("rowHeights"), table_row_heights(table, ooxml_scale)
    ):
        return "TABLE_STRUCTURE_UNSUPPORTED"
    if "borderColor" in props and not table_value_equal(
        props.get("borderColor"), "#CBD5E1"
    ):
        return "TABLE_STRUCTURE_UNSUPPORTED"
    if "borderWidth" in props and not table_value_equal(props.get("borderWidth"), 1):
        return "TABLE_STRUCTURE_UNSUPPORTED"
    return None


def table_cell_non_text_equal(
    target: dict[str, Any],
    actual: dict[str, Any],
) -> bool:
    if set(target) - SUPPORTED_TABLE_CELL_PROPS:
        return False
    defaults: dict[str, Any] = {
        "fill": "transparent",
        "fontSize": 18,
        "fontWeight": "normal",
        "align": "left",
        "verticalAlign": "middle",
        "borderColor": "#CBD5E1",
        "borderWidth": 1,
        "colSpan": 1,
        "rowSpan": 1,
    }
    for key in SUPPORTED_TABLE_CELL_PROPS - {"text"}:
        target_value = target.get(key, defaults.get(key))
        actual_value = actual.get(key, defaults.get(key))
        if not table_value_equal(target_value, actual_value):
            return False
    return True


def numeric_track_values_equal(target: Any, actual: list[int]) -> bool:
    return (
        isinstance(target, list)
        and len(target) == len(actual)
        and all(
            table_value_equal(left, right)
            for left, right in zip(target, actual, strict=True)
        )
    )


def table_value_equal(left: Any, right: Any) -> bool:
    if isinstance(left, bool) or isinstance(right, bool):
        return left is right
    if isinstance(left, (int, float)) and isinstance(right, (int, float)):
        return math.isclose(float(left), float(right), rel_tol=1e-4, abs_tol=0.1)
    if isinstance(left, str) and isinstance(right, str):
        if valid_hex_color(left) and valid_hex_color(right):
            return left.upper() == right.upper()
    return bool(left == right)


def valid_table_props(props: dict[str, Any]) -> bool:
    if not props or set(props) - SUPPORTED_TABLE_PROPS:
        return False
    rows = props.get("rows")
    if not isinstance(rows, list) or not 1 <= len(rows) <= 1000:
        return False
    if not isinstance(rows[0], list) or not 1 <= len(rows[0]) <= 1000:
        return False
    column_count = len(rows[0])
    if len(rows) * column_count > 10_000:
        return False
    if any(not isinstance(row, list) or len(row) != column_count for row in rows):
        return False
    if any(
        not isinstance(cell, dict) or not valid_table_cell_props(cell)
        for row in rows
        for cell in row
    ):
        return False
    if not valid_table_tracks(props.get("columnWidths"), column_count):
        return False
    if not valid_table_tracks(props.get("rowHeights"), len(rows)):
        return False
    border_color = props.get("borderColor", "#CBD5E1")
    border_width = props.get("borderWidth", 1)
    return (
        isinstance(border_color, str)
        and valid_hex_color(border_color)
        and finite_table_number(border_width, minimum=0)
    )


def valid_table_cell_props(cell: dict[str, Any]) -> bool:
    if set(cell) - SUPPORTED_TABLE_CELL_PROPS:
        return False
    text = cell.get("text", "")
    fill = cell.get("fill", "transparent")
    text_color = cell.get("textColor")
    font_family = cell.get("fontFamily")
    font_weight = cell.get("fontWeight", "normal")
    return (
        isinstance(text, str)
        and isinstance(fill, str)
        and (fill == "transparent" or valid_hex_color(fill))
        and (
            text_color is None
            or isinstance(text_color, str)
            and valid_hex_color(text_color)
        )
        and (font_family is None or isinstance(font_family, str) and bool(font_family))
        and finite_table_number(cell.get("fontSize", 18), minimum=0, strict=True)
        and valid_table_font_weight(font_weight)
        and cell.get("align", "left") in {"left", "center", "right", "justify"}
        and cell.get("verticalAlign", "middle") in {"top", "middle", "bottom"}
        and isinstance(cell.get("borderColor", "#CBD5E1"), str)
        and valid_hex_color(str(cell.get("borderColor", "#CBD5E1")))
        and finite_table_number(cell.get("borderWidth", 1), minimum=0)
        and valid_positive_integer(cell.get("colSpan", 1))
        and valid_positive_integer(cell.get("rowSpan", 1))
        and int(cell.get("colSpan", 1)) == 1
        and int(cell.get("rowSpan", 1)) == 1
    )


def valid_table_tracks(value: Any, count: int) -> bool:
    return value is None or (
        isinstance(value, list)
        and len(value) == count
        and all(finite_table_number(item, minimum=0, strict=True) for item in value)
    )


def valid_table_font_weight(value: Any) -> bool:
    return isinstance(value, str) and value in {"normal", "bold"}


def valid_positive_integer(value: Any) -> bool:
    return isinstance(value, int) and not isinstance(value, bool) and value > 0


def finite_table_number(
    value: Any,
    *,
    minimum: float,
    strict: bool = False,
) -> bool:
    return (
        isinstance(value, (int, float))
        and not isinstance(value, bool)
        and math.isfinite(float(value))
        and (float(value) > minimum if strict else float(value) >= minimum)
    )


def valid_text_props(props: dict[str, Any]) -> bool:
    if set(props) - SUPPORTED_TEXT_PROPS:
        return False
    if "text" in props and not isinstance(props.get("text"), str):
        return False
    if not valid_text_style_values(props):
        return False
    if props.get("align", "left") not in {"left", "center", "right", "justify"}:
        return False
    if props.get("verticalAlign", "top") not in {"top", "middle", "bottom"}:
        return False
    if props.get("writingMode", "horizontal") not in {"horizontal", "vertical-270"}:
        return False
    if not valid_positive_number(props.get("lineHeight", 1.2)):
        return False
    if not valid_text_body_inset(props.get("bodyInset")):
        return False
    if "bullet" in props and not valid_text_bullet(props.get("bullet")):
        return False
    runs = props.get("runs")
    if runs is not None and (
        not isinstance(runs, list) or any(not valid_text_run(run) for run in runs)
    ):
        return False
    paragraphs = props.get("paragraphs")
    if paragraphs is not None and (
        not isinstance(paragraphs, list)
        or any(not valid_text_paragraph(item) for item in paragraphs)
    ):
        return False
    return not text_props_has_content_projection(props) or canonical_text_paragraphs(props) is not None


def valid_text_paragraph(value: Any) -> bool:
    if not isinstance(value, dict) or set(value) - SUPPORTED_TEXT_PARAGRAPH_PROPS:
        return False
    if "text" in value and not isinstance(value.get("text"), str):
        return False
    if not valid_text_style_values(value):
        return False
    if value.get("align", "left") not in {"left", "center", "right", "justify"}:
        return False
    if not valid_positive_number(value.get("lineHeight", 1.2)):
        return False
    if any(key in value and not valid_nonnegative_number(value.get(key)) for key in ("spaceBefore", "spaceAfter")):
        return False
    if "indent" in value and not valid_finite_number(value.get("indent")):
        return False
    if "bullet" in value and not valid_text_bullet(value.get("bullet")):
        return False
    runs = value.get("runs")
    return runs is None or isinstance(runs, list) and all(valid_text_run(run) for run in runs)


def valid_text_run(value: Any) -> bool:
    return (
        isinstance(value, dict)
        and not set(value) - SUPPORTED_TEXT_RUN_PROPS
        and isinstance(value.get("text", ""), str)
        and valid_text_style_values(value)
        and value.get("baseline", "normal") in {"normal", "superscript", "subscript"}
    )


def valid_text_style_values(value: dict[str, Any]) -> bool:
    font_family = value.get("fontFamily")
    if font_family is not None and (not isinstance(font_family, str) or not font_family):
        return False
    if "fontSize" in value and not valid_positive_number(value.get("fontSize")):
        return False
    weight = value.get("fontWeight")
    if weight is not None and weight not in {"normal", "bold"}:
        return False
    if any(key in value and not isinstance(value.get(key), bool) for key in ("italic", "underline")):
        return False
    color = value.get("color")
    return color is None or isinstance(color, str) and valid_hex_color(color)


def valid_text_bullet(value: Any) -> bool:
    return (
        isinstance(value, dict)
        and not set(value) - {"enabled", "character", "indent"}
        and isinstance(value.get("enabled", False), bool)
        and isinstance(value.get("character", "\u2022"), str)
        and bool(value.get("character", "\u2022"))
        and valid_nonnegative_number(value.get("indent", 0))
    )


def valid_text_body_inset(value: Any) -> bool:
    return value is None or (
        isinstance(value, dict)
        and not set(value) - {"left", "right", "top", "bottom"}
        and all(valid_nonnegative_number(item) for item in value.values())
    )


def valid_finite_number(value: Any) -> bool:
    return isinstance(value, (int, float)) and not isinstance(value, bool) and math.isfinite(float(value))


def valid_positive_number(value: Any) -> bool:
    return valid_finite_number(value) and float(value) > 0


def valid_nonnegative_number(value: Any) -> bool:
    return valid_finite_number(value) and float(value) >= 0


def valid_rect_props(props: dict[str, Any]) -> bool:
    if set(props) - {"fill", "stroke", "strokeWidth", "borderRadius"}:
        return False
    for color_name in ("fill", "stroke"):
        if color_name not in props:
            continue
        color = props[color_name]
        if color != "transparent" and not (
            isinstance(color, str) and valid_hex_color(color)
        ):
            return False
    return all(
        name not in props or valid_nonnegative_number(props[name])
        for name in ("strokeWidth", "borderRadius")
    )


def canonical_text_value(props: dict[str, Any]) -> str | None:
    paragraphs = canonical_text_paragraphs(props)
    if paragraphs is None:
        return None
    return "\n".join(str(paragraph["text"]) for paragraph in paragraphs)


def text_props_has_content_projection(props: dict[str, Any]) -> bool:
    return any(key in props for key in ("text", "runs", "paragraphs"))


def canonical_text_paragraphs(props: dict[str, Any]) -> list[dict[str, Any]] | None:
    raw_paragraphs = props.get("paragraphs")
    if isinstance(raw_paragraphs, list):
        paragraphs: list[dict[str, Any]] = []
        for raw in raw_paragraphs:
            if not isinstance(raw, dict):
                return None
            paragraph = copy.deepcopy(raw)
            raw_runs = paragraph.get("runs")
            if isinstance(raw_runs, list) and raw_runs:
                if any(not isinstance(run, dict) for run in raw_runs):
                    return None
                runs = [copy.deepcopy(run) for run in raw_runs]
                text = "".join(str(run.get("text", "")) for run in runs)
                if "text" in paragraph and paragraph.get("text") != text:
                    return None
            else:
                text = str(paragraph.get("text", ""))
                runs = [{"text": text}] if text else []
            paragraph.update({"text": text, "runs": runs})
            paragraphs.append(paragraph)
        paragraphs = paragraphs or [{"text": "", "runs": []}]
        value = "\n".join(str(item["text"]) for item in paragraphs)
        return None if "text" in props and props.get("text") != value else paragraphs
    raw_runs = props.get("runs")
    if isinstance(raw_runs, list) and raw_runs:
        paragraphs = [{"text": "", "runs": []}]
        for raw in raw_runs:
            if not isinstance(raw, dict):
                return None
            pieces = str(raw.get("text", "")).split("\n")
            for index, piece in enumerate(pieces):
                if piece or len(pieces) == 1:
                    run = copy.deepcopy(raw)
                    run["text"] = piece
                    paragraphs[-1]["runs"].append(run)
                    paragraphs[-1]["text"] += piece
                if index < len(pieces) - 1:
                    paragraphs.append({"text": "", "runs": []})
        value = "\n".join(str(item["text"]) for item in paragraphs)
        return None if "text" in props and props.get("text") != value else paragraphs
    text = str(props.get("text", ""))
    return [{"text": part, "runs": [{"text": part}] if part else []} for part in text.split("\n")]


def text_body_value(shape: ET.Element[Any]) -> str:
    body = first_local_child(shape, "txBody")
    if body is None:
        return ""
    paragraphs: list[str] = []
    for paragraph in direct_local_children(body, "p"):
        parts: list[str] = []
        for child in paragraph:
            name = local_name(child)
            if name in {"r", "fld"}:
                parts.append("".join(node.text or "" for node in child.iter() if local_name(node) == "t"))
            elif name == "br":
                parts.append("\n")
        paragraphs.append("".join(parts))
    return "\n".join(paragraphs)


def normalized_image_crop(value: Any) -> dict[str, float] | None:
    if not isinstance(value, dict) or set(value) - {
        "left",
        "top",
        "right",
        "bottom",
    }:
        return None
    crop: dict[str, float] = {}
    for edge in ("left", "top", "right", "bottom"):
        raw_value = value.get(edge, 0)
        if (
            not isinstance(raw_value, (int, float))
            or isinstance(raw_value, bool)
            or not math.isfinite(float(raw_value))
            or not 0 <= float(raw_value) <= 1
        ):
            return None
        crop[edge] = float(raw_value)
    if crop["left"] + crop["right"] >= 1:
        return None
    if crop["top"] + crop["bottom"] >= 1:
        return None
    return crop


def operation_element_id(operation: dict[str, Any]) -> str:
    element_id = operation.get("elementId")
    if isinstance(element_id, str):
        return element_id
    element = operation.get("element")
    if isinstance(element, dict) and isinstance(element.get("elementId"), str):
        return str(element["elementId"])
    return ""


def applied_operation(operation: dict[str, Any]) -> PptxOoxmlAppliedOperation:
    return PptxOoxmlAppliedOperation(
        operationType=cast(PptxOoxmlSyncOperationType, operation.get("type")),
        slideId=operation_slide_id(operation) or None,
        elementId=operation_element_id(operation) or None,
    )


def unsupported_operation(
    operation: dict[str, Any],
    reason_code: PptxOoxmlUnsupportedReasonCode,
) -> PptxOoxmlUnsupportedOperation:
    return PptxOoxmlUnsupportedOperation(
        operationType=cast(PptxOoxmlSyncOperationType, operation.get("type")),
        slideId=operation_slide_id(operation) or None,
        elementId=operation_element_id(operation) or None,
        reasonCode=reason_code,
    )


def operation_slide_id(operation: dict[str, Any]) -> str:
    slide_id = operation.get("slideId")
    if isinstance(slide_id, str):
        return slide_id
    slide = operation.get("slide")
    if isinstance(slide, dict) and isinstance(slide.get("slideId"), str):
        return str(slide["slideId"])
    return ""


def element_source_map(
    template_blueprint: dict[str, Any],
) -> dict[tuple[str, str], dict[str, Any]]:
    return {
        (
            str(source.get("slidePart", "")),
            str(source.get("elementId", "")),
        ): dict(source)
        for slide in template_blueprint.get("slides", [])
        if isinstance(slide, dict)
        for source in slide.get("elementSources", [])
        if isinstance(source, dict) and source.get("elementId")
    }


def shared_shape_operation_plan(
    operations: list[dict[str, Any]],
    sources: dict[tuple[str, str], dict[str, Any]],
) -> tuple[set[int], PptxOoxmlUnsupportedOperation | None]:
    cohorts: dict[tuple[str, str], dict[str, dict[str, Any]]] = {}
    for source in sources.values():
        slide_part = str(source.get("slidePart", ""))
        shape_id = str(source.get("shapeId", ""))
        element_id = str(source.get("elementId", ""))
        if not slide_part or not shape_id or not element_id:
            continue
        cohorts.setdefault((slide_part, shape_id), {})[element_id] = source

    shared_cohorts = {
        cohort_key: members
        for cohort_key, members in cohorts.items()
        if len(members) > 1
    }
    operations_by_cohort: dict[tuple[str, str], list[tuple[int, dict[str, Any]]]] = {}
    for operation_index, operation in enumerate(operations):
        if operation.get("type") not in {"delete_element", "update_element_frame"}:
            continue
        source_key = (
            str(operation.get("sourceSlidePart", "")),
            operation_element_id(operation),
        )
        operation_source = sources.get(source_key)
        if operation_source is None:
            continue
        cohort_key = (
            str(operation_source.get("slidePart", "")),
            str(operation_source.get("shapeId", "")),
        )
        if cohort_key in shared_cohorts:
            operations_by_cohort.setdefault(cohort_key, []).append(
                (operation_index, operation)
            )

    redundant_indexes: set[int] = set()
    for cohort_key, indexed_operations in operations_by_cohort.items():
        member_ids = set(shared_cohorts[cohort_key])
        representative = indexed_operations[0][1]
        member_count = len(member_ids)
        if len(indexed_operations) % member_count != 0:
            return (
                set(),
                unsupported_operation(
                    representative,
                    "SHARED_SHAPE_COHORT_UNSAFE",
                ),
            )
        for round_start in range(0, len(indexed_operations), member_count):
            operation_round = indexed_operations[
                round_start : round_start + member_count
            ]
            round_representative = operation_round[0][1]
            representative_type = round_representative.get("type")
            representative_frame = round_representative.get("frame")
            round_member_ids = {
                operation_element_id(operation)
                for _operation_index, operation in operation_round
            }
            unsafe = round_member_ids != member_ids or any(
                operation.get("type") != representative_type
                or (
                    representative_type == "update_element_frame"
                    and operation.get("frame") != representative_frame
                )
                for _operation_index, operation in operation_round[1:]
            )
            if unsafe:
                return (
                    set(),
                    unsupported_operation(
                        round_representative,
                        "SHARED_SHAPE_COHORT_UNSAFE",
                    ),
                )
            redundant_indexes.update(
                operation_index for operation_index, _operation in operation_round[1:]
            )

    return redundant_indexes, None


def motion_reference_failure(
    operations: list[dict[str, Any]],
    template_blueprint: dict[str, Any],
    slide_motion: list[dict[str, Any]],
) -> PptxOoxmlUnsupportedOperation | None:
    coverage_by_slide_part = {
        source_slide_part(slide): str(
            dict_value(slide, "ooxmlMotionCapabilities").get(
                "importedMainSequenceCoverage",
                "unknown",
            )
        )
        for slide in template_blueprint.get("slides", [])
        if isinstance(slide, dict) and source_slide_part(slide)
    }
    animation_replacement_parts = {
        str(item.get("sourceSlidePart", ""))
        for item in slide_motion
        if isinstance(item, dict)
        and isinstance(item.get("touched"), dict)
        and item["touched"].get("animations") is True
    }
    for operation in operations:
        if operation.get("type") != "delete_element":
            continue
        slide_part = slide_part_for_operation(operation, template_blueprint)
        if slide_part in animation_replacement_parts:
            continue
        if coverage_by_slide_part.get(slide_part, "unknown") != "absent":
            return unsupported_operation(
                operation,
                "MOTION_REFERENCE_COVERAGE_UNSAFE",
            )
    return None


def find_shape_by_id(
    root: ET.Element[Any], shape_id: str
) -> tuple[ET.Element[Any] | None, ET.Element[Any] | None]:
    for parent in root.iter():
        for child in list(parent):
            if child.tag not in {P_SP, P_PIC, P_GRAPHIC_FRAME}:
                continue
            non_visual_name = {
                P_SP: "nvSpPr",
                P_PIC: "nvPicPr",
                P_GRAPHIC_FRAME: "nvGraphicFramePr",
            }[child.tag]
            non_visual = first_local_child(child, non_visual_name)
            c_nv_pr = (
                first_local_child(non_visual, "cNvPr")
                if non_visual is not None
                else None
            )
            if c_nv_pr is not None and c_nv_pr.get("id") == shape_id:
                return child, parent
    return None, None


def direct_image_blip_fill(shape: ET.Element[Any] | None) -> ET.Element[Any] | None:
    if shape is None:
        return None
    if shape.tag == P_PIC:
        return first_local_child(shape, "blipFill")
    if shape.tag == P_SP:
        shape_properties = first_local_child(shape, "spPr")
        if shape_properties is not None:
            return first_local_child(shape_properties, "blipFill")
    return None


def direct_image_blip(
    shape: ET.Element[Any] | None,
    source: dict[str, Any],
) -> ET.Element[Any] | None:
    blip_fill = direct_image_blip_fill(shape)
    if blip_fill is None:
        return None
    blip = first_local_child(blip_fill, "blip")
    expected_relationship_id = str(source.get("relationshipId", ""))
    current_relationship_id = (
        str(blip.get(f"{{{REL_NS}}}embed", "")) if blip is not None else ""
    )
    if (
        not expected_relationship_id
        or current_relationship_id != expected_relationship_id
    ):
        return None
    return blip


def set_picture_opacity(
    shape: ET.Element[Any],
    source: dict[str, Any],
    opacity: float,
) -> bool:
    blip = direct_image_blip(shape, source)
    if blip is None:
        return False
    for child in list(blip):
        if local_name(child) == "alphaModFix":
            blip.remove(child)
    if opacity < 1:
        alpha = ET.Element(
            f"{{{DML_NS}}}alphaModFix",
            {"amt": str(round(opacity * 100000))},
        )
        extension_list = first_local_child(blip, "extLst")
        if extension_list is None:
            blip.append(alpha)
        else:
            blip.insert(list(blip).index(extension_list), alpha)
    return True


def set_shape_visibility(shape: ET.Element[Any], visible: bool) -> bool:
    non_visual_name = {
        P_SP: "nvSpPr",
        P_PIC: "nvPicPr",
        P_GRAPHIC_FRAME: "nvGraphicFramePr",
    }.get(shape.tag)
    if non_visual_name is None:
        return False
    non_visual = first_local_child(shape, non_visual_name)
    c_nv_pr = (
        first_local_child(non_visual, "cNvPr")
        if non_visual is not None
        else None
    )
    if c_nv_pr is None:
        return False
    if visible:
        c_nv_pr.attrib.pop("hidden", None)
    else:
        c_nv_pr.set("hidden", "1")
    return True


def image_crop_capability_for_shape(
    shape: ET.Element[Any],
    source: dict[str, Any],
) -> str:
    if direct_image_blip(shape, source) is None:
        return "none"
    if shape.tag == P_PIC:
        return "picture"
    if shape.tag == P_SP:
        return "picture-fill"
    return "none"


def has_group_shape_ancestor(
    root: ET.Element[Any], shape: ET.Element[Any]
) -> bool:
    parents = {child: parent for parent in root.iter() for child in list(parent)}
    parent = parents.get(shape)
    while parent is not None:
        if local_name(parent) == "grpSp":
            return True
        parent = parents.get(parent)
    return False


def update_shape_props(
    shape: ET.Element[Any],
    props: dict[str, Any],
    source: dict[str, Any],
    scale: PackageFrameScale,
    slide_part: str,
    package_entries: dict[str, bytes],
    added_entries: dict[str, bytes],
    updated_sources: dict[tuple[str, str], dict[str, Any]],
    source_key: tuple[str, str],
    warnings: list[str],
    element_id: str,
) -> bool:
    changed = False
    if source.get("elementType") == "rect":
        return update_authored_rect_props(shape, props, scale)
    if source.get("elementType") == "table":
        if source.get("ooxmlOrigin") == "imported":
            changed = sync_imported_table_cell_text(shape, props)
        else:
            changed = replace_authored_table_subtree(shape, props, scale)
        if not changed or not refresh_table_source_locators(shape, source):
            warnings.append(f"OOXML table sync skipped for {element_id}.")
            return False
        updated_sources[source_key] = dict(source)
        return True
    if source.get("elementType") == "text":
        return sync_text_shape(shape, props, source, scale)
    if "src" in props:
        if source.get("fallbackReason"):
            warnings.append(f"OOXML fallback source preserved for {element_id}.")
            return False
        replacement = decode_image_data_url(props.get("src"))
        if isinstance(replacement, str):
            warnings.append(f"OOXML image sync skipped for {element_id}: {replacement}.")
            return False
        mime_type, image_blob = replacement
        relationship_id = replace_picture_media_relationship(
            shape,
            source,
            slide_part,
            mime_type,
            image_blob,
            package_entries,
            added_entries,
            warnings,
            element_id,
        )
        if relationship_id is None:
            return False
        source["relationshipId"] = relationship_id
        updated_sources[source_key] = dict(source)
        changed = True
    if "crop" in props:
        crop_value = props.get("crop")
        crop = normalized_image_crop(crop_value) if crop_value is not None else None
        if not set_picture_crop_source_rect(shape, crop):
            warnings.append(f"OOXML image crop target missing for {element_id}.")
            return False
        changed = True
    if changed:
        return True
    warnings.append(f"OOXML prop sync skipped for {element_id}.")
    return False


def sync_imported_table_cell_text(
    shape: ET.Element[Any],
    props: dict[str, Any],
) -> bool:
    table = direct_graphic_frame_table(shape)
    target_rows = props.get("rows")
    if table is None or not isinstance(target_rows, list):
        return False
    changed_cell: tuple[ET.Element[Any], str] | None = None
    rows = direct_local_children(table, "tr")
    for row_index, row in enumerate(rows):
        cells = direct_local_children(row, "tc")
        if row_index >= len(target_rows) or not isinstance(
            target_rows[row_index], list
        ):
            return False
        for column_index, cell in enumerate(cells):
            target_cell = target_rows[row_index][column_index]
            if not isinstance(target_cell, dict):
                return False
            target_text = str(target_cell.get("text", ""))
            if target_text == table_cell_text_value(cell):
                continue
            if changed_cell is not None or not table_cell_text_can_set(
                cell, target_text
            ):
                return False
            changed_cell = (cell, target_text)
    if changed_cell is None:
        return False
    cell, target_text = changed_cell
    return set_table_cell_text_value(cell, target_text)


def table_cell_text_value(cell: ET.Element[Any]) -> str:
    body = first_local_child(cell, "txBody")
    if body is None:
        return ""
    return "\n".join(
        "".join(text_run_value(run) for run in direct_local_children(paragraph, "r"))
        for paragraph in direct_local_children(body, "p")
    )


def table_cell_text_can_set(cell: ET.Element[Any], value: str) -> bool:
    body = first_local_child(cell, "txBody")
    return (
        body is not None
        and table_cell_text_body_is_safe(cell)
        and len(value.split("\n")) == len(direct_local_children(body, "p"))
    )


def set_table_cell_text_value(cell: ET.Element[Any], value: str) -> bool:
    if not table_cell_text_can_set(cell, value):
        return False
    body = first_local_child(cell, "txBody")
    if body is None:
        return False
    for paragraph, paragraph_text in zip(
        direct_local_children(body, "p"),
        value.split("\n"),
        strict=True,
    ):
        runs = direct_local_children(paragraph, "r")
        if runs:
            text_node = first_local_child(runs[0], "t")
            if text_node is None:
                return False
            set_text_node_value(text_node, paragraph_text)
            continue
        if not paragraph_text:
            continue
        run = ET.Element(f"{{{DML_NS}}}r")
        end_properties = first_local_child(paragraph, "endParaRPr")
        if end_properties is not None:
            run_properties = copy.deepcopy(end_properties)
            run_properties.tag = f"{{{DML_NS}}}rPr"
            run.append(run_properties)
        ET.SubElement(run, A_T)
        text_node = first_local_child(run, "t")
        if text_node is None:
            return False
        set_text_node_value(text_node, paragraph_text)
        insert_at = (
            list(paragraph).index(end_properties)
            if end_properties is not None
            else len(paragraph)
        )
        paragraph.insert(insert_at, run)
    return True


def refresh_table_source_locators(
    shape: ET.Element[Any],
    source: dict[str, Any],
) -> bool:
    locators, diagnostics = table_cell_locators(
        shape,
        slide_index=0,
        shape_id=str(source.get("shapeId", "")),
    )
    if diagnostics or not locators:
        return False
    source["tableCellLocators"] = locators
    capabilities = dict_value(source, "ooxmlEditCapabilities")
    capabilities["tableCellText"] = table_cell_text_capability_for_shape(shape, source)
    source["ooxmlEditCapabilities"] = capabilities
    return capabilities["tableCellText"] is True


def decode_image_data_url(value: Any) -> tuple[str, bytes] | str:
    if not isinstance(value, str):
        return "invalid data URL"
    header, separator, payload = value.partition(",")
    if not separator or not header.lower().startswith("data:"):
        return "invalid data URL"
    if not header.lower().endswith(";base64"):
        return "invalid data URL"

    mime_type = header[5:-7].lower()
    expected_format = {
        "image/png": "PNG",
        "image/jpeg": "JPEG",
        "image/gif": "GIF",
        "image/webp": "WEBP",
    }.get(mime_type)
    if expected_format is None:
        return f"unsupported MIME type {mime_type or 'unknown'}"

    try:
        image_blob = base64.b64decode(payload, validate=True)
        with Image.open(BytesIO(image_blob)) as image:
            image.verify()
            actual_format = str(image.format or "").upper()
    except (binascii.Error, OSError, SyntaxError, ValueError):
        return "invalid image data"
    if actual_format != expected_format:
        return f"image data does not match {mime_type}"
    return mime_type, image_blob


def replace_picture_media_relationship(
    shape: ET.Element[Any],
    source: dict[str, Any],
    slide_part: str,
    mime_type: str,
    image_blob: bytes,
    package_entries: dict[str, bytes],
    added_entries: dict[str, bytes],
    warnings: list[str],
    element_id: str,
) -> str | None:
    if shape.tag != P_PIC:
        warnings.append(f"OOXML image source is not a picture for {element_id}.")
        return None

    blip = next(shape.iter(A_BLIP), None)
    expected_relationship_id = str(source.get("relationshipId", ""))
    current_relationship_id = (
        str(blip.get(f"{{{REL_NS}}}embed", "")) if blip is not None else ""
    )
    if (
        blip is None
        or not expected_relationship_id
        or current_relationship_id != expected_relationship_id
    ):
        warnings.append(f"OOXML image relationship mismatch for {element_id}.")
        return None

    rels_part = rels_part_for_slide_part(slide_part)
    rels_xml = package_entries.get(rels_part)
    content_types_xml = package_entries.get("[Content_Types].xml")
    if rels_xml is None or not is_image_relationship(
        rels_xml, expected_relationship_id
    ):
        warnings.append(f"OOXML image relationship missing for {element_id}.")
        return None
    if content_types_xml is None:
        warnings.append(f"OOXML content types missing for {element_id}.")
        return None

    extension = extension_for_mime_type(mime_type)
    media_token = safe_package_token(
        f"{Path(slide_part).stem}_{source.get('shapeId', 'image')}"
    )
    media_name = f"orbit_sync_{media_token}.{extension}"
    media_part = f"ppt/media/{media_name}"
    relationship_id, next_rels_xml = append_image_relationship(
        rels_xml,
        f"../media/{media_name}",
    )

    blip.set(f"{{{REL_NS}}}embed", relationship_id)
    package_entries[rels_part] = next_rels_xml
    package_entries["[Content_Types].xml"] = ensure_content_type_default(
        content_types_xml,
        extension,
        mime_type,
    )
    added_entries[media_part] = image_blob
    return relationship_id


def is_image_relationship(rels_xml: bytes, relationship_id: str) -> bool:
    root = ET.fromstring(rels_xml)
    return any(
        child.get("Id") == relationship_id
        and child.get("Type") == IMAGE_REL_TYPE
        for child in list(root)
    )


def safe_package_token(value: str) -> str:
    token = "".join(
        char if char.isascii() and (char.isalnum() or char in "_-") else "_"
        for char in value
    )
    return token or "image"


def sync_text_shape(
    shape: ET.Element[Any],
    props: dict[str, Any],
    source: dict[str, Any],
    scale: PackageFrameScale,
) -> bool:
    if set(props) == {"text"} and str(props.get("text", "")) == text_body_value(shape):
        return True
    paragraphs = text_sync_paragraphs(shape, props)
    if paragraphs is None:
        return False
    body = ensure_text_body(shape)
    if dict_value(source, "ooxmlEditCapabilities").get("richText") == "style-only":
        paragraphs = preserve_existing_run_boundaries(body, paragraphs)
    equal_spans = text_equal_spans(
        text_body_value(shape),
        "\n".join(str(paragraph.get("text", "")) for paragraph in paragraphs),
    )
    apply_text_body_properties(body, props, scale)
    authored = source.get("ooxmlOrigin") == "authored"
    if text_structure_matches(body, paragraphs):
        patch_matching_text_structure(body, paragraphs, props, scale, authored)
    else:
        rebuild_text_structure(
            body,
            paragraphs,
            props,
            scale,
            authored,
            equal_spans,
        )
    return True


def text_sync_paragraphs(
    shape: ET.Element[Any],
    props: dict[str, Any],
) -> list[dict[str, Any]] | None:
    redundant_text_projection = (
        "text" in props
        and "runs" not in props
        and "paragraphs" not in props
        and str(props.get("text", "")) == text_body_value(shape)
    )
    if text_props_has_content_projection(props) and not redundant_text_projection:
        return canonical_text_paragraphs(props)
    body = first_local_child(shape, "txBody")
    if body is None:
        return [{"text": "", "runs": []}]
    paragraph_style = {
        key: props[key] for key in ("align", "lineHeight", "bullet") if key in props
    }
    paragraphs: list[dict[str, Any]] = []
    for paragraph in direct_local_children(body, "p"):
        runs: list[dict[str, Any]] = []
        for child in list(paragraph):
            name = local_name(child)
            if name == "r":
                runs.append({"text": text_run_value(child)})
            elif name == "br":
                runs.append({"text": "\n"})
        paragraphs.append(
            {
                "text": "".join(str(run["text"]) for run in runs),
                "runs": runs,
                **paragraph_style,
            }
        )
    return paragraphs or [{"text": "", "runs": [], **paragraph_style}]


def style_only_paragraphs_match(
    body: ET.Element[Any],
    paragraphs: list[dict[str, Any]],
) -> bool:
    existing = [
        text_paragraph_value(paragraph)
        for paragraph in direct_local_children(body, "p")
    ]
    target = [str(paragraph.get("text", "")) for paragraph in paragraphs]
    return existing == target


def preserve_existing_run_boundaries(
    body: ET.Element[Any],
    paragraphs: list[dict[str, Any]],
) -> list[dict[str, Any]]:
    existing_paragraphs = direct_local_children(body, "p")
    if len(existing_paragraphs) != len(paragraphs):
        return paragraphs
    result: list[dict[str, Any]] = []
    for existing, target in zip(existing_paragraphs, paragraphs, strict=True):
        text = str(target.get("text", ""))
        target_runs = target.get("runs", [])
        if not isinstance(target_runs, list):
            result.append(target)
            continue
        boundaries = {0, utf16_length(text)}
        offset = 0
        for child in list(existing):
            name = local_name(child)
            if name == "r":
                offset += utf16_length(text_run_value(child))
                boundaries.add(offset)
            elif name == "br":
                offset += 1
                boundaries.add(offset)
        target_intervals: list[tuple[int, int, dict[str, Any]]] = []
        offset = 0
        for target_run in target_runs:
            if not isinstance(target_run, dict):
                continue
            end = offset + utf16_length(str(target_run.get("text", "")))
            target_intervals.append((offset, end, target_run))
            boundaries.update({offset, end})
            offset = end
        ordered = sorted(boundaries)
        rebuilt_runs: list[dict[str, Any]] = []
        for start, end in zip(ordered, ordered[1:]):
            if start == end:
                continue
            target_run = next(
                (
                    run
                    for run_start, run_end, run in target_intervals
                    if run_start <= start < run_end
                ),
                {},
            )
            rebuilt_runs.append(
                {
                    **copy.deepcopy(target_run),
                    "text": utf16_slice(text, start, end),
                }
            )
        result.append({**copy.deepcopy(target), "runs": rebuilt_runs})
    return result


def apply_text_body_properties(
    body: ET.Element[Any],
    props: dict[str, Any],
    scale: PackageFrameScale,
) -> None:
    body_pr = first_local_child(body, "bodyPr")
    if body_pr is None:
        body_pr = ET.Element(f"{{{DML_NS}}}bodyPr")
        body.insert(0, body_pr)
    body_pr.set("horzOverflow", "clip")
    body_pr.set("vertOverflow", "clip")
    body_pr.set("wrap", "square")
    if "verticalAlign" in props:
        body_pr.set(
            "anchor",
            {"top": "t", "middle": "ctr", "bottom": "b"}.get(
                str(props.get("verticalAlign", "top")),
                "t",
            ),
        )
    if "writingMode" in props:
        body_pr.set(
            "vert",
            "vert270" if props.get("writingMode") == "vertical-270" else "horz",
        )
    if "bodyInset" in props:
        inset = dict_value(props, "bodyInset")
        for key, attribute, converter in (
            ("left", "lIns", canvas_x_to_emu),
            ("right", "rIns", canvas_x_to_emu),
            ("top", "tIns", canvas_y_to_emu),
            ("bottom", "bIns", canvas_y_to_emu),
        ):
            if key in inset:
                body_pr.set(attribute, str(converter(inset[key], scale)))


def text_structure_matches(
    body: ET.Element[Any],
    paragraphs: list[dict[str, Any]],
) -> bool:
    existing_paragraphs = direct_local_children(body, "p")
    if len(existing_paragraphs) != len(paragraphs):
        return False
    for existing, target in zip(existing_paragraphs, paragraphs, strict=True):
        content = [
            child
            for child in list(existing)
            if local_name(child) not in {"pPr", "endParaRPr"}
        ]
        if any(local_name(child) != "r" for child in content):
            return False
        target_runs = target.get("runs", [])
        if not isinstance(target_runs, list) or len(content) != len(target_runs):
            return False
        if any(
            text_run_value(run) != str(target_run.get("text", ""))
            for run, target_run in zip(content, target_runs, strict=True)
        ):
            return False
    return True


def patch_matching_text_structure(
    body: ET.Element[Any],
    paragraphs: list[dict[str, Any]],
    props: dict[str, Any],
    scale: PackageFrameScale,
    authored: bool,
) -> None:
    raw_paragraphs = props.get("paragraphs")
    for paragraph_index, (paragraph, target) in enumerate(
        zip(direct_local_children(body, "p"), paragraphs, strict=True)
    ):
        desired_paragraph = desired_paragraph_style(props, target, authored)
        patch_paragraph_properties(paragraph, desired_paragraph, scale)
        raw_paragraph = (
            raw_paragraphs[paragraph_index]
            if isinstance(raw_paragraphs, list)
            and paragraph_index < len(raw_paragraphs)
            and isinstance(raw_paragraphs[paragraph_index], dict)
            else None
        )
        raw_runs = raw_paragraph.get("runs") if raw_paragraph is not None else None
        has_explicit_runs = isinstance(raw_runs, list) and bool(raw_runs)
        existing_runs = direct_local_children(paragraph, "r")
        target_runs = target.get("runs", [])
        for run, target_run in zip(existing_runs, target_runs, strict=True):
            desired_run = desired_run_style(
                props,
                target,
                target_run,
                authored=authored,
                has_explicit_runs=has_explicit_runs,
            )
            patch_run_properties(run, desired_run, scale)


def rebuild_text_structure(
    body: ET.Element[Any],
    paragraphs: list[dict[str, Any]],
    props: dict[str, Any],
    scale: PackageFrameScale,
    authored: bool,
    equal_spans: list[TextEqualSpan],
) -> None:
    existing_paragraphs = direct_local_children(body, "p")
    run_templates = existing_text_run_templates(existing_paragraphs)
    paragraph_templates = existing_text_paragraph_templates(existing_paragraphs)
    raw_paragraphs = props.get("paragraphs")
    rebuilt: list[ET.Element[Any]] = []
    logical_offset = 0
    for paragraph_index, target in enumerate(paragraphs):
        paragraph_text_length = utf16_length(str(target.get("text", "")))
        source_start, source_end = map_target_interval_to_source(
            logical_offset,
            logical_offset + paragraph_text_length,
            equal_spans,
        )
        paragraph_template_record = nearest_text_paragraph_template(
            paragraph_templates,
            source_start,
            source_end,
        )
        paragraph_template = (
            paragraph_template_record.paragraph
            if paragraph_template_record is not None
            else None
        )
        paragraph = ET.Element(f"{{{DML_NS}}}p")
        if paragraph_template is not None:
            p_pr = first_local_child(paragraph_template, "pPr")
            if p_pr is not None:
                paragraph.append(copy.deepcopy(p_pr))
        patch_paragraph_properties(
            paragraph,
            desired_paragraph_style(props, target, authored),
            scale,
        )
        raw_paragraph = (
            raw_paragraphs[paragraph_index]
            if isinstance(raw_paragraphs, list)
            and paragraph_index < len(raw_paragraphs)
            and isinstance(raw_paragraphs[paragraph_index], dict)
            else None
        )
        raw_runs = raw_paragraph.get("runs") if raw_paragraph is not None else None
        has_explicit_runs = isinstance(raw_runs, list) and bool(raw_runs)
        target_runs = target.get("runs", [])
        for target_run in target_runs if isinstance(target_runs, list) else []:
            run_text = str(target_run.get("text", ""))
            desired_run = desired_run_style(
                props,
                target,
                target_run,
                authored=authored,
                has_explicit_runs=has_explicit_runs,
            )
            append_rebuilt_run_content(
                paragraph,
                run_text,
                logical_offset,
                desired_run,
                run_templates,
                equal_spans,
                scale,
            )
            logical_offset += utf16_length(run_text)
        if paragraph_template is not None:
            end_properties = first_local_child(paragraph_template, "endParaRPr")
            if end_properties is not None:
                paragraph.append(copy.deepcopy(end_properties))
        rebuilt.append(paragraph)
        if paragraph_index < len(paragraphs) - 1:
            logical_offset += 1

    children = list(body)
    paragraph_indexes = [
        index for index, child in enumerate(children) if local_name(child) == "p"
    ]
    insertion_index = paragraph_indexes[0] if paragraph_indexes else len(children)
    for paragraph in existing_paragraphs:
        body.remove(paragraph)
    for offset, paragraph in enumerate(rebuilt):
        body.insert(insertion_index + offset, paragraph)


def existing_text_run_templates(
    paragraphs: list[ET.Element[Any]],
) -> list[TextRunTemplate]:
    templates: list[TextRunTemplate] = []
    offset = 0
    for paragraph_index, paragraph in enumerate(paragraphs):
        for child in list(paragraph):
            name = local_name(child)
            if name == "r":
                text = text_run_value(child)
                length = utf16_length(text)
                r_pr = first_local_child(child, "rPr")
                templates.append(
                    TextRunTemplate(
                        start=offset,
                        end=offset + length,
                        run_properties=copy.deepcopy(r_pr)
                        if r_pr is not None
                        else None,
                    )
                )
                offset += length
            elif name == "br":
                r_pr = first_local_child(child, "rPr")
                templates.append(
                    TextRunTemplate(
                        start=offset,
                        end=offset + 1,
                        run_properties=copy.deepcopy(r_pr)
                        if r_pr is not None
                        else None,
                    )
                )
                offset += 1
        if paragraph_index < len(paragraphs) - 1:
            offset += 1
    return templates


def existing_text_paragraph_templates(
    paragraphs: list[ET.Element[Any]],
) -> list[TextParagraphTemplate]:
    templates: list[TextParagraphTemplate] = []
    offset = 0
    for paragraph_index, paragraph in enumerate(paragraphs):
        length = utf16_length(text_paragraph_value(paragraph))
        templates.append(
            TextParagraphTemplate(
                start=offset,
                end=offset + length,
                paragraph=paragraph,
            )
        )
        offset += length
        if paragraph_index < len(paragraphs) - 1:
            offset += 1
    return templates


def append_rebuilt_run_content(
    paragraph: ET.Element[Any],
    text: str,
    start: int,
    desired_style: dict[str, Any],
    templates: list[TextRunTemplate],
    equal_spans: list[TextEqualSpan],
    scale: PackageFrameScale,
) -> None:
    offset = start
    pieces = text.split("\n")
    for piece_index, piece in enumerate(pieces):
        if piece or len(pieces) == 1:
            run = ET.SubElement(paragraph, f"{{{DML_NS}}}r")
            source_start, source_end = map_target_interval_to_source(
                offset,
                offset + utf16_length(piece),
                equal_spans,
            )
            template = nearest_text_run_template(
                templates,
                source_start,
                source_end,
            )
            if template is not None and template.run_properties is not None:
                run.append(copy.deepcopy(template.run_properties))
            patch_run_properties(run, desired_style, scale)
            text_node = ET.SubElement(run, A_T)
            set_text_node_value(text_node, piece)
            offset += utf16_length(piece)
        if piece_index < len(pieces) - 1:
            line_break = ET.SubElement(paragraph, f"{{{DML_NS}}}br")
            source_start, source_end = map_target_interval_to_source(
                offset,
                offset + 1,
                equal_spans,
            )
            template = nearest_text_run_template(
                templates,
                source_start,
                source_end,
            )
            if template is not None and template.run_properties is not None:
                line_break.append(copy.deepcopy(template.run_properties))
            patch_run_properties(line_break, desired_style, scale)
            offset += 1


def nearest_text_run_template(
    templates: list[TextRunTemplate],
    start: int,
    end: int,
) -> TextRunTemplate | None:
    overlapping = [
        template
        for template in templates
        if max(start, template.start) < min(end, template.end)
    ]
    if overlapping:
        return max(
            overlapping,
            key=lambda template: min(end, template.end) - max(start, template.start),
        )
    return min(
        templates,
        key=lambda template: min(
            abs(start - template.start),
            abs(start - template.end),
        ),
        default=None,
    )


def nearest_text_paragraph_template(
    templates: list[TextParagraphTemplate],
    start: int,
    end: int,
) -> TextParagraphTemplate | None:
    overlapping = [
        template
        for template in templates
        if max(start, template.start) < min(end, template.end)
    ]
    if overlapping:
        return max(
            overlapping,
            key=lambda template: min(end, template.end) - max(start, template.start),
        )
    return min(
        templates,
        key=lambda template: min(
            abs(start - template.start),
            abs(start - template.end),
        ),
        default=None,
    )


def text_equal_spans(source: str, target: str) -> list[TextEqualSpan]:
    source_offsets = utf16_prefix_offsets(source)
    target_offsets = utf16_prefix_offsets(target)
    prefix_length = 0
    shared_length = min(len(source), len(target))
    while (
        prefix_length < shared_length and source[prefix_length] == target[prefix_length]
    ):
        prefix_length += 1

    suffix_length = 0
    source_remaining = len(source) - prefix_length
    target_remaining = len(target) - prefix_length
    while (
        suffix_length < min(source_remaining, target_remaining)
        and source[len(source) - suffix_length - 1]
        == target[len(target) - suffix_length - 1]
    ):
        suffix_length += 1

    spans: list[TextEqualSpan] = []
    if prefix_length:
        spans.append(
            TextEqualSpan(
                target_start=0,
                target_end=target_offsets[prefix_length],
                source_start=0,
                source_end=source_offsets[prefix_length],
            )
        )

    source_middle_end = len(source) - suffix_length
    target_middle_end = len(target) - suffix_length
    source_middle = source[prefix_length:source_middle_end]
    target_middle = target[prefix_length:target_middle_end]
    if (
        source_middle
        and target_middle
        and len(source_middle) * len(target_middle) <= MAX_TEXT_DIFF_MATRIX_CELLS
    ):
        matcher = difflib.SequenceMatcher(
            a=source_middle,
            b=target_middle,
            autojunk=False,
        )
        spans.extend(
            TextEqualSpan(
                target_start=target_offsets[prefix_length + match.b],
                target_end=target_offsets[prefix_length + match.b + match.size],
                source_start=source_offsets[prefix_length + match.a],
                source_end=source_offsets[prefix_length + match.a + match.size],
            )
            for match in matcher.get_matching_blocks()
            if match.size > 0
        )

    if suffix_length:
        spans.append(
            TextEqualSpan(
                target_start=target_offsets[target_middle_end],
                target_end=target_offsets[len(target)],
                source_start=source_offsets[source_middle_end],
                source_end=source_offsets[len(source)],
            )
        )
    return spans


def map_target_interval_to_source(
    start: int,
    end: int,
    equal_spans: list[TextEqualSpan],
) -> tuple[int, int]:
    overlaps: list[tuple[int, TextEqualSpan]] = []
    for span in equal_spans:
        overlap_start = max(start, span.target_start)
        overlap_end = min(end, span.target_end)
        if overlap_start < overlap_end:
            overlaps.append((overlap_end - overlap_start, span))
    if overlaps:
        _length, span = max(overlaps, key=lambda item: item[0])
        overlap_start = max(start, span.target_start)
        overlap_end = min(end, span.target_end)
        return (
            span.source_start + overlap_start - span.target_start,
            span.source_start + overlap_end - span.target_start,
        )

    preceding = [span for span in equal_spans if span.target_end <= start]
    if preceding:
        position = max(preceding, key=lambda span: span.target_end).source_end
        return position, position
    following = [span for span in equal_spans if span.target_start >= end]
    if following:
        position = min(following, key=lambda span: span.target_start).source_start
        return position, position
    return 0, 0


def desired_paragraph_style(
    props: dict[str, Any],
    paragraph: dict[str, Any],
    authored: bool,
) -> dict[str, Any]:
    keys = {"align", "lineHeight", "spaceBefore", "spaceAfter", "indent", "bullet"}
    desired: dict[str, Any] = {}
    if authored:
        for key in ("align", "lineHeight", "bullet"):
            if key in props:
                desired[key] = props[key]
    for key in keys:
        if key in paragraph:
            desired[key] = paragraph[key]
    return desired


def desired_run_style(
    props: dict[str, Any],
    paragraph: dict[str, Any],
    run: dict[str, Any],
    *,
    authored: bool,
    has_explicit_runs: bool,
) -> dict[str, Any]:
    desired: dict[str, Any] = {}
    if authored or not has_explicit_runs:
        for source in (props, paragraph):
            for key in SUPPORTED_TEXT_STYLE_PROPS:
                if key in source:
                    desired[key] = source[key]
    for key in SUPPORTED_TEXT_STYLE_PROPS:
        if key in run:
            desired[key] = run[key]
    return desired


def patch_run_properties(
    run: ET.Element[Any],
    desired: dict[str, Any],
    scale: PackageFrameScale,
) -> None:
    current = current_run_style(first_local_child(run, "rPr"), scale)
    differences = {
        key: value
        for key, value in desired.items()
        if key in SUPPORTED_TEXT_STYLE_PROPS
        and not text_style_values_equal(current.get(key), value)
    }
    if not differences:
        return
    r_pr = ensure_run_properties(run)
    if "fontFamily" in differences:
        for name in ("latin", "ea"):
            font = first_local_child(r_pr, name)
            if font is None:
                font = ET.SubElement(r_pr, f"{{{DML_NS}}}{name}")
            font.set("typeface", str(differences["fontFamily"]))
    if "fontSize" in differences:
        r_pr.set("sz", str(font_size_to_ooxml(differences["fontSize"], scale)))
    if "fontWeight" in differences:
        r_pr.set("b", "1" if is_bold_text_weight(differences["fontWeight"]) else "0")
    if "italic" in differences:
        r_pr.set("i", "1" if differences["italic"] else "0")
    if "underline" in differences:
        r_pr.set("u", "sng" if differences["underline"] else "none")
    if "baseline" in differences:
        baseline = differences["baseline"]
        if baseline == "superscript":
            r_pr.set("baseline", "30000")
        elif baseline == "subscript":
            r_pr.set("baseline", "-25000")
        else:
            r_pr.attrib.pop("baseline", None)
    if "color" in differences:
        for child in list(r_pr):
            if local_name(child) in {
                "solidFill",
                "gradFill",
                "noFill",
                "pattFill",
                "blipFill",
            }:
                r_pr.remove(child)
        color_fill = ET.Element(f"{{{DML_NS}}}solidFill")
        ET.SubElement(
            color_fill,
            f"{{{DML_NS}}}srgbClr",
            {"val": str(differences["color"])[1:].upper()},
        )
        r_pr.insert(0, color_fill)


def current_run_style(
    r_pr: ET.Element[Any] | None,
    scale: PackageFrameScale,
) -> dict[str, Any]:
    if r_pr is None:
        return {"baseline": "normal"}
    current: dict[str, Any] = {"baseline": "normal"}
    for name in ("latin", "ea", "cs"):
        font = first_local_child(r_pr, name)
        if font is not None and font.get("typeface"):
            current["fontFamily"] = str(font.get("typeface"))
            break
    size = int_value(r_pr.get("sz"), 0)
    if size > 0:
        current["fontSize"] = font_size_from_ooxml(size, scale)
    if r_pr.get("b") is not None:
        current["fontWeight"] = "bold" if r_pr.get("b") in {"1", "true"} else "normal"
    if r_pr.get("i") is not None:
        current["italic"] = r_pr.get("i") in {"1", "true"}
    if r_pr.get("u") is not None:
        current["underline"] = r_pr.get("u") not in {"0", "false", "none"}
    solid_fill = first_local_child(r_pr, "solidFill")
    srgb = first_local_child(solid_fill, "srgbClr") if solid_fill is not None else None
    if srgb is not None and srgb.get("val"):
        current["color"] = f"#{str(srgb.get('val')).upper()}"
    baseline = int_value(r_pr.get("baseline"), 0)
    if baseline > 0:
        current["baseline"] = "superscript"
    elif baseline < 0:
        current["baseline"] = "subscript"
    return current


def ensure_run_properties(run: ET.Element[Any]) -> ET.Element[Any]:
    r_pr = first_local_child(run, "rPr")
    if r_pr is not None:
        return r_pr
    r_pr = ET.Element(f"{{{DML_NS}}}rPr", {"lang": "ko-KR"})
    run.insert(0, r_pr)
    return r_pr


def patch_paragraph_properties(
    paragraph: ET.Element[Any],
    desired: dict[str, Any],
    scale: PackageFrameScale,
) -> None:
    current = current_paragraph_style(first_local_child(paragraph, "pPr"), scale)
    differences = {
        key: value
        for key, value in desired.items()
        if not text_style_values_equal(current.get(key), value)
    }
    if not differences:
        return
    p_pr = first_local_child(paragraph, "pPr")
    if p_pr is None:
        p_pr = ET.Element(f"{{{DML_NS}}}pPr")
        paragraph.insert(0, p_pr)
    if "align" in differences:
        p_pr.set(
            "algn",
            {"left": "l", "center": "ctr", "right": "r", "justify": "just"}.get(
                str(differences["align"]),
                "l",
            ),
        )
    if "indent" in differences:
        p_pr.set("marL", str(canvas_x_to_signed_emu(differences["indent"], scale)))
    if "lineHeight" in differences:
        set_paragraph_spacing_percent(p_pr, "lnSpc", differences["lineHeight"])
    if "spaceBefore" in differences:
        set_paragraph_spacing_points(p_pr, "spcBef", differences["spaceBefore"], scale)
    if "spaceAfter" in differences:
        set_paragraph_spacing_points(p_pr, "spcAft", differences["spaceAfter"], scale)
    if "bullet" in differences:
        for child in list(p_pr):
            if local_name(child) in {"buNone", "buChar", "buAutoNum"}:
                p_pr.remove(child)
        bullet = differences["bullet"]
        if isinstance(bullet, dict) and bullet.get("enabled"):
            ET.SubElement(
                p_pr,
                f"{{{DML_NS}}}buChar",
                {"char": str(bullet.get("character", "\u2022"))},
            )
            p_pr.set(
                "marL",
                str(canvas_x_to_signed_emu(bullet.get("indent", 0), scale)),
            )
        else:
            ET.SubElement(p_pr, f"{{{DML_NS}}}buNone")


def current_paragraph_style(
    p_pr: ET.Element[Any] | None,
    scale: PackageFrameScale,
) -> dict[str, Any]:
    if p_pr is None:
        return {
            "align": "left",
            "lineHeight": 1.15,
            "spaceBefore": 0,
            "spaceAfter": 0,
            "indent": 0,
        }
    current: dict[str, Any] = {
        "align": {
            "ctr": "center",
            "r": "right",
            "just": "justify",
        }.get(str(p_pr.get("algn", "l")), "left"),
        "lineHeight": paragraph_spacing_percent(p_pr, "lnSpc", 1.15),
        "spaceBefore": paragraph_spacing_canvas(p_pr, "spcBef", scale),
        "spaceAfter": paragraph_spacing_canvas(p_pr, "spcAft", scale),
        "indent": round(int_value(p_pr.get("marL"), 0) * canvas_x_scale(scale), 3),
    }
    bullet = first_local_child(p_pr, "buChar")
    if bullet is not None:
        current["bullet"] = {
            "enabled": True,
            "character": str(bullet.get("char", "\u2022")),
            "indent": max(0, current["indent"]),
        }
    elif first_local_child(p_pr, "buNone") is not None:
        current["bullet"] = {"enabled": False, "character": "\u2022", "indent": 0}
    return current


def set_paragraph_spacing_percent(
    p_pr: ET.Element[Any],
    name: str,
    value: Any,
) -> None:
    spacing = first_local_child(p_pr, name)
    if spacing is None:
        spacing = ET.SubElement(p_pr, f"{{{DML_NS}}}{name}")
    for child in list(spacing):
        spacing.remove(child)
    ET.SubElement(
        spacing,
        f"{{{DML_NS}}}spcPct",
        {"val": str(round(float(value) * 100000))},
    )


def set_paragraph_spacing_points(
    p_pr: ET.Element[Any],
    name: str,
    value: Any,
    scale: PackageFrameScale,
) -> None:
    spacing = first_local_child(p_pr, name)
    if spacing is None:
        spacing = ET.SubElement(p_pr, f"{{{DML_NS}}}{name}")
    for child in list(spacing):
        spacing.remove(child)
    ET.SubElement(
        spacing,
        f"{{{DML_NS}}}spcPts",
        {"val": str(canvas_spacing_to_ooxml(value, scale))},
    )


def paragraph_spacing_percent(
    p_pr: ET.Element[Any],
    name: str,
    fallback: float,
) -> float:
    spacing = first_local_child(p_pr, name)
    percentage = first_local_child(spacing, "spcPct") if spacing is not None else None
    return (
        int_value(percentage.get("val"), round(fallback * 100000)) / 100000
        if percentage is not None
        else fallback
    )


def paragraph_spacing_canvas(
    p_pr: ET.Element[Any],
    name: str,
    scale: PackageFrameScale,
) -> float:
    spacing = first_local_child(p_pr, name)
    points = first_local_child(spacing, "spcPts") if spacing is not None else None
    if points is None:
        return 0
    return round(
        int_value(points.get("val"), 0) / 100 * 12700 * canvas_average_scale(scale),
        3,
    )


def font_size_from_ooxml(size: int, scale: PackageFrameScale) -> float:
    return round(size / 100 * 12700 * canvas_average_scale(scale), 3)


def font_size_to_ooxml(value: Any, scale: PackageFrameScale) -> int:
    return max(1, round(float(value) / (12700 * canvas_average_scale(scale)) * 100))


def canvas_spacing_to_ooxml(value: Any, scale: PackageFrameScale) -> int:
    return max(0, round(float(value) / (12700 * canvas_average_scale(scale)) * 100))


def canvas_x_scale(scale: PackageFrameScale) -> float:
    return scale.canvas_width / scale.slide_width_emu


def canvas_average_scale(scale: PackageFrameScale) -> float:
    return (
        scale.canvas_width / scale.slide_width_emu
        + scale.canvas_height / scale.slide_height_emu
    ) / 2


def canvas_x_to_signed_emu(value: Any, scale: PackageFrameScale) -> int:
    return round(float(value) * scale.slide_width_emu / scale.canvas_width)


def text_style_values_equal(current: Any, desired: Any) -> bool:
    if isinstance(current, (int, float)) and isinstance(desired, (int, float)):
        return math.isclose(float(current), float(desired), rel_tol=1e-4, abs_tol=0.01)
    if isinstance(current, str) and isinstance(desired, str):
        if valid_hex_color(current) and valid_hex_color(desired):
            return current.upper() == desired.upper()
    return bool(current == desired)


def is_bold_text_weight(value: Any) -> bool:
    return value in {"semibold", "bold"} or (
        isinstance(value, int) and not isinstance(value, bool) and value >= 600
    )


def text_run_value(run: ET.Element[Any]) -> str:
    return "".join(node.text or "" for node in run.iter() if local_name(node) == "t")


def text_paragraph_value(paragraph: ET.Element[Any]) -> str:
    parts: list[str] = []
    for child in list(paragraph):
        name = local_name(child)
        if name in {"r", "fld"}:
            parts.append(text_run_value(child))
        elif name == "br":
            parts.append("\n")
    return "".join(parts)


def set_text_node_value(node: ET.Element[Any], value: str) -> None:
    node.text = value
    if value != value.strip():
        node.set(XML_SPACE, "preserve")
    else:
        node.attrib.pop(XML_SPACE, None)


def utf16_length(value: str) -> int:
    return len(value.encode("utf-16-le")) // 2


def utf16_prefix_offsets(value: str) -> list[int]:
    offsets = [0]
    for character in value:
        offsets.append(offsets[-1] + utf16_length(character))
    return offsets


def utf16_slice(value: str, start: int, end: int) -> str:
    encoded = value.encode("utf-16-le")
    return encoded[start * 2 : end * 2].decode("utf-16-le")


def canvas_x_to_emu(value: Any, scale: PackageFrameScale) -> int:
    return round(float(value) * scale.slide_width_emu / scale.canvas_width)


def canvas_y_to_emu(value: Any, scale: PackageFrameScale) -> int:
    return round(float(value) * scale.slide_height_emu / scale.canvas_height)


def update_shape_frame(
    shape: ET.Element[Any],
    frame: dict[str, Any],
    scale: PackageFrameScale,
) -> None:
    xfrm = ensure_xfrm(shape)
    off = first_local_child(xfrm, "off")
    if off is None:
        off = ET.SubElement(xfrm, f"{{{DML_NS}}}off")
    ext = first_local_child(xfrm, "ext")
    if ext is None:
        ext = ET.SubElement(xfrm, f"{{{DML_NS}}}ext")
    if "x" in frame:
        off.set(
            "x",
            str(round(float(frame["x"]) * scale.slide_width_emu / scale.canvas_width)),
        )
    if "y" in frame:
        off.set(
            "y",
            str(round(float(frame["y"]) * scale.slide_height_emu / scale.canvas_height)),
        )
    if "width" in frame:
        ext.set(
            "cx",
            str(
                max(
                    1,
                    round(
                        float(frame["width"])
                        * scale.slide_width_emu
                        / scale.canvas_width
                    ),
                )
            ),
        )
    if "height" in frame:
        ext.set(
            "cy",
            str(
                max(
                    1,
                    round(
                        float(frame["height"])
                        * scale.slide_height_emu
                        / scale.canvas_height
                    ),
                )
            ),
        )
    if "rotation" in frame:
        xfrm.set("rot", str(round(float(frame["rotation"]) * 60000)))


def table_graphic_frame_element(
    shape_id: int,
    element: dict[str, Any],
    scale: PackageFrameScale,
) -> ET.Element[Any]:
    frame = ET.Element(P_GRAPHIC_FRAME)
    non_visual = ET.SubElement(frame, f"{{{PML_NS}}}nvGraphicFramePr")
    ET.SubElement(
        non_visual,
        f"{{{PML_NS}}}cNvPr",
        {"id": str(shape_id), "name": "Orbit table"},
    )
    ET.SubElement(non_visual, f"{{{PML_NS}}}cNvGraphicFramePr")
    ET.SubElement(non_visual, f"{{{PML_NS}}}nvPr")
    update_shape_frame(frame, element, scale)
    graphic = ET.SubElement(frame, f"{{{DML_NS}}}graphic")
    graphic_data = ET.SubElement(
        graphic,
        f"{{{DML_NS}}}graphicData",
        {"uri": TABLE_GRAPHIC_DATA_URI},
    )
    _x, _y, width, height = frame_to_emu(element, scale)
    graphic_data.append(
        table_subtree_element(dict_value(element, "props"), width, height, scale)
    )
    return frame


def replace_authored_table_subtree(
    shape: ET.Element[Any],
    props: dict[str, Any],
    scale: PackageFrameScale,
) -> bool:
    if shape.tag != P_GRAPHIC_FRAME:
        return False
    graphic = first_local_child(shape, "graphic")
    graphic_data = (
        first_local_child(graphic, "graphicData") if graphic is not None else None
    )
    table = direct_graphic_frame_table(shape)
    frame_size = graphic_frame_size_emu(shape)
    if graphic_data is None or table is None or frame_size is None:
        return False
    width, height = frame_size
    rows = props.get("rows")
    if not isinstance(rows, list) or not rows or not isinstance(rows[0], list):
        return False
    column_count = len(rows[0])
    row_count = len(rows)
    preserved_column_widths = (
        table_column_tracks_emu(table, column_count)
        if "columnWidths" not in props
        else None
    )
    preserved_row_heights = (
        table_row_tracks_emu(table, row_count) if "rowHeights" not in props else None
    )
    if ("columnWidths" not in props and preserved_column_widths is None) or (
        "rowHeights" not in props and preserved_row_heights is None
    ):
        return False
    replacement = table_subtree_element(
        props,
        width,
        height,
        scale,
        column_widths_emu=preserved_column_widths,
        row_heights_emu=preserved_row_heights,
    )
    table_index = list(graphic_data).index(table)
    graphic_data.remove(table)
    graphic_data.insert(table_index, replacement)
    return True


def graphic_frame_size_emu(shape: ET.Element[Any]) -> tuple[int, int] | None:
    xfrm = first_local_child(shape, "xfrm")
    ext = first_local_child(xfrm, "ext") if xfrm is not None else None
    if ext is None:
        return None
    width = int_value(ext.get("cx"), 0)
    height = int_value(ext.get("cy"), 0)
    return (width, height) if width > 0 and height > 0 else None


def resize_authored_table_tracks_to_frame(shape: ET.Element[Any]) -> bool:
    table = direct_graphic_frame_table(shape)
    frame_size = graphic_frame_size_emu(shape)
    if table is None or frame_size is None:
        return False
    grid = first_local_child(table, "tblGrid")
    columns = direct_local_children(grid, "gridCol") if grid is not None else []
    rows = direct_local_children(table, "tr")
    if not columns or not rows:
        return False
    column_weights = [int_value(column.get("w"), 0) for column in columns]
    row_weights = [int_value(row.get("h"), 0) for row in rows]
    if any(value <= 0 for value in column_weights + row_weights):
        return False
    widths = normalized_table_tracks_emu(
        column_weights,
        total=frame_size[0],
        count=len(columns),
    )
    heights = normalized_table_tracks_emu(
        row_weights,
        total=frame_size[1],
        count=len(rows),
    )
    for column, width in zip(columns, widths, strict=True):
        column.set("w", str(width))
    for row, height in zip(rows, heights, strict=True):
        row.set("h", str(height))
    return True


def table_subtree_element(
    props: dict[str, Any],
    frame_width_emu: int,
    frame_height_emu: int,
    scale: PackageFrameScale,
    *,
    column_widths_emu: list[int] | None = None,
    row_heights_emu: list[int] | None = None,
) -> ET.Element[Any]:
    rows = cast(list[list[dict[str, Any]]], props["rows"])
    row_count = len(rows)
    column_count = len(rows[0])
    table = ET.Element(f"{{{DML_NS}}}tbl")
    ET.SubElement(table, f"{{{DML_NS}}}tblPr")
    grid = ET.SubElement(table, f"{{{DML_NS}}}tblGrid")
    column_widths = column_widths_emu or normalized_table_tracks_emu(
        props.get("columnWidths"),
        total=max(column_count, frame_width_emu),
        count=column_count,
    )
    row_heights = row_heights_emu or normalized_table_tracks_emu(
        props.get("rowHeights"),
        total=max(row_count, frame_height_emu),
        count=row_count,
    )
    for width in column_widths:
        ET.SubElement(grid, f"{{{DML_NS}}}gridCol", {"w": str(width)})
    for row_index, (row_payload, height) in enumerate(
        zip(rows, row_heights, strict=True)
    ):
        row = ET.SubElement(table, f"{{{DML_NS}}}tr", {"h": str(height)})
        for cell_payload in row_payload:
            row.append(table_cell_element(cell_payload, props, row_index, scale))
    return table


def table_column_tracks_emu(
    table: ET.Element[Any],
    expected_count: int,
) -> list[int] | None:
    grid = first_local_child(table, "tblGrid")
    if grid is None:
        return None
    tracks = [
        int_value(column.get("w"), 0)
        for column in direct_local_children(grid, "gridCol")
    ]
    return tracks if len(tracks) == expected_count and all(tracks) else None


def table_row_tracks_emu(
    table: ET.Element[Any],
    expected_count: int,
) -> list[int] | None:
    tracks = [int_value(row.get("h"), 0) for row in direct_local_children(table, "tr")]
    return tracks if len(tracks) == expected_count and all(tracks) else None


def normalized_table_tracks_emu(
    tracks: Any,
    *,
    total: int,
    count: int,
) -> list[int]:
    weights = (
        [float(value) for value in tracks]
        if isinstance(tracks, list) and tracks
        else [1.0] * count
    )
    maximum = max(weights)
    scaled = [weight / maximum for weight in weights]
    distributable = max(0, total - count)
    exact_extras = [distributable * weight / sum(scaled) for weight in scaled]
    floor_extras = [math.floor(value) for value in exact_extras]
    normalized = [1 + value for value in floor_extras]
    remainder = distributable - sum(floor_extras)
    order = sorted(
        range(count),
        key=lambda index: (-(exact_extras[index] - floor_extras[index]), index),
    )
    for index in order[:remainder]:
        normalized[index] += 1
    return normalized


def table_cell_element(
    cell_payload: dict[str, Any],
    table_props: dict[str, Any],
    row_index: int,
    scale: PackageFrameScale,
) -> ET.Element[Any]:
    del row_index
    cell = ET.Element(f"{{{DML_NS}}}tc")
    body = ET.SubElement(cell, f"{{{DML_NS}}}txBody")
    ET.SubElement(body, f"{{{DML_NS}}}bodyPr")
    ET.SubElement(body, f"{{{DML_NS}}}lstStyle")
    text = str(cell_payload.get("text", ""))
    for paragraph_text in text.split("\n"):
        paragraph = ET.SubElement(body, f"{{{DML_NS}}}p")
        align = {
            "center": "ctr",
            "right": "r",
            "justify": "just",
        }.get(str(cell_payload.get("align", "left")), "l")
        ET.SubElement(paragraph, f"{{{DML_NS}}}pPr", {"algn": align})
        run = ET.SubElement(paragraph, f"{{{DML_NS}}}r")
        run_properties = ET.SubElement(
            run,
            f"{{{DML_NS}}}rPr",
            {
                "lang": "ko-KR",
                "sz": str(font_size_to_ooxml(cell_payload.get("fontSize", 18), scale)),
                "b": "1"
                if is_bold_text_weight(cell_payload.get("fontWeight", "normal"))
                else "0",
            },
        )
        text_color = str(cell_payload.get("textColor") or "#000000")
        text_fill = ET.SubElement(run_properties, f"{{{DML_NS}}}solidFill")
        ET.SubElement(
            text_fill,
            f"{{{DML_NS}}}srgbClr",
            {"val": text_color[1:].upper()},
        )
        font_family = cell_payload.get("fontFamily")
        if isinstance(font_family, str) and font_family:
            ET.SubElement(
                run_properties,
                f"{{{DML_NS}}}latin",
                {"typeface": font_family},
            )
            ET.SubElement(
                run_properties,
                f"{{{DML_NS}}}ea",
                {"typeface": font_family},
            )
        text_node = ET.SubElement(run, A_T)
        set_text_node_value(text_node, paragraph_text)

    anchor = {
        "top": "t",
        "bottom": "b",
    }.get(str(cell_payload.get("verticalAlign", "middle")), "ctr")
    cell_properties = ET.SubElement(
        cell,
        f"{{{DML_NS}}}tcPr",
        {"anchor": anchor},
    )
    border_color = str(
        cell_payload.get("borderColor") or table_props.get("borderColor") or "#CBD5E1"
    )
    border_width = cell_payload.get(
        "borderWidth",
        table_props.get("borderWidth", 1),
    )
    for border_name in ("lnL", "lnR", "lnT", "lnB"):
        line = ET.SubElement(
            cell_properties,
            f"{{{DML_NS}}}{border_name}",
            {"w": str(table_border_width_to_emu(border_width, scale))},
        )
        if float(border_width) <= 0:
            ET.SubElement(line, f"{{{DML_NS}}}noFill")
        else:
            line_fill = ET.SubElement(line, f"{{{DML_NS}}}solidFill")
            ET.SubElement(
                line_fill,
                f"{{{DML_NS}}}srgbClr",
                {"val": border_color[1:].upper()},
            )
    fill = str(cell_payload.get("fill", "transparent"))
    if fill == "transparent":
        ET.SubElement(cell_properties, f"{{{DML_NS}}}noFill")
    else:
        solid_fill = ET.SubElement(cell_properties, f"{{{DML_NS}}}solidFill")
        ET.SubElement(
            solid_fill,
            f"{{{DML_NS}}}srgbClr",
            {"val": fill[1:].upper()},
        )
    return cell


def table_border_width_to_emu(value: Any, scale: PackageFrameScale) -> int:
    return max(0, round(float(value) / canvas_average_scale(scale)))


VISUAL_SHAPE_NAMES = {"cxnSp", "graphicFrame", "grpSp", "pic", "sp"}


def reorder_visual_shape(
    parent: ET.Element[Any],
    shape: ET.Element[Any],
    z_index_value: Any,
) -> bool:
    visual_children = [
        child for child in list(parent) if local_name(child) in VISUAL_SHAPE_NAMES
    ]
    if shape not in visual_children:
        return False
    target_index = normalized_z_index(z_index_value, len(visual_children))
    if target_index is None:
        return False
    current_index = visual_children.index(shape)
    if current_index == target_index:
        return True

    parent.remove(shape)
    remaining = [child for child in visual_children if child is not shape]
    if target_index >= len(remaining):
        insert_at = visual_insert_end_index(parent)
    else:
        insert_at = list(parent).index(remaining[target_index])
    parent.insert(insert_at, shape)
    return True


def normalized_z_index(value: Any, item_count: int) -> int | None:
    if isinstance(value, bool):
        return None
    try:
        numeric = float(value)
    except (TypeError, ValueError):
        return None
    if not math.isfinite(numeric) or not numeric.is_integer():
        return None
    return max(0, min(int(numeric), max(0, item_count - 1)))


def visual_insert_end_index(parent: ET.Element[Any]) -> int:
    children = list(parent)
    visual_indexes = [
        index
        for index, child in enumerate(children)
        if local_name(child) in VISUAL_SHAPE_NAMES
    ]
    if visual_indexes:
        return visual_indexes[-1] + 1
    for index, child in enumerate(children):
        if local_name(child) == "extLst":
            return index
    return len(children)


def add_authored_slide_to_package(
    operation: dict[str, Any],
    package_entries: dict[str, bytes],
    added_entries: dict[str, bytes],
    scale: PackageFrameScale,
    warnings: list[str],
    source_package: zipfile.ZipFile,
    template_blueprint: dict[str, Any],
) -> tuple[list[dict[str, Any]], PptxOoxmlUnsupportedReasonCode | None]:
    slide = operation.get("slide")
    slide_part = slide_part_for_operation(operation, template_blueprint)
    if not isinstance(slide, dict) or not slide_part:
        return [], "ADD_SLIDE_FAILED"
    if (
        slide_part in source_package.namelist()
        or slide_part in package_entries
        or not slide_part.startswith("ppt/slides/")
        or not slide_part.endswith(".xml")
    ):
        return [], "ADD_SLIDE_FAILED"

    layout_target = authored_slide_layout_target(
        package_entries,
        template_blueprint,
    )
    if not layout_target:
        return [], "ADD_SLIDE_LAYOUT_UNSAFE"
    presentation_xml = package_entries.get("ppt/presentation.xml")
    presentation_rels_xml = package_entries.get(
        "ppt/_rels/presentation.xml.rels"
    )
    content_types_xml = package_entries.get("[Content_Types].xml")
    if (
        presentation_xml is None
        or presentation_rels_xml is None
        or content_types_xml is None
    ):
        return [], "ADD_SLIDE_FAILED"

    try:
        presentation_root = ET.fromstring(presentation_xml)
        presentation_rels_root = ET.fromstring(presentation_rels_xml)
        content_types_root = ET.fromstring(content_types_xml)
    except ET.ParseError:
        return [], "ADD_SLIDE_FAILED"
    slide_id_list = presentation_root.find(f"{{{PML_NS}}}sldIdLst")
    if slide_id_list is None:
        return [], "ADD_SLIDE_FAILED"

    relationship_id = next_relationship_id(presentation_rels_root)
    ET.SubElement(
        presentation_rels_root,
        f"{{{PKG_REL_NS}}}Relationship",
        {
            "Id": relationship_id,
            "Type": SLIDE_REL_TYPE,
            "Target": posixpath.relpath(slide_part, "ppt"),
        },
    )
    current_slide_ids = [
        int(str(node.get("id", "0")))
        for node in list(slide_id_list)
        if str(node.get("id", "")).isdigit()
    ]
    slide_id_node = ET.Element(
        f"{{{PML_NS}}}sldId",
        {
            "id": str(max(current_slide_ids, default=255) + 1),
            f"{{{REL_NS}}}id": relationship_id,
        },
    )
    requested_order = slide.get("order")
    if not isinstance(requested_order, int) or isinstance(requested_order, bool):
        return [], "ADD_SLIDE_FAILED"
    insert_index = max(0, min(requested_order - 1, len(slide_id_list)))
    slide_id_list.insert(insert_index, slide_id_node)

    part_name = f"/{slide_part}"
    if not any(
        child.tag.endswith("Override") and child.get("PartName") == part_name
        for child in list(content_types_root)
    ):
        ET.SubElement(
            content_types_root,
            f"{{{CONTENT_TYPES_NS}}}Override",
            {"PartName": part_name, "ContentType": SLIDE_CONTENT_TYPE},
        )

    rels_part = rels_part_for_slide_part(slide_part)
    package_entries[slide_part] = empty_slide_xml()
    package_entries[rels_part] = slide_layout_relationships_xml(layout_target)
    package_entries["ppt/presentation.xml"] = xml_bytes(presentation_root)
    package_entries["ppt/_rels/presentation.xml.rels"] = xml_bytes(
        presentation_rels_root
    )
    package_entries["[Content_Types].xml"] = xml_bytes(content_types_root)

    added_sources: list[dict[str, Any]] = []
    elements = slide.get("elements", [])
    if not isinstance(elements, list):
        return [], "ADD_SLIDE_FAILED"
    for element in elements:
        if not isinstance(element, dict) or str(element.get("type", "")) not in {
            "text",
            "rect",
            "image",
            "table",
        }:
            return [], "ADD_SLIDE_FAILED"
        added_source = add_element_to_slide_xml(
            slide_part,
            element,
            package_entries,
            added_entries,
            scale,
            warnings,
        )
        if added_source is None:
            return [], "ADD_SLIDE_FAILED"
        added_sources.append(added_source)
    return added_sources, None


def authored_slide_layout_target(
    package_entries: dict[str, bytes],
    template_blueprint: dict[str, Any],
) -> str:
    for raw_slide in template_blueprint.get("slides", []):
        if not isinstance(raw_slide, dict) or raw_slide.get("ooxmlOrigin") == "authored":
            continue
        slide_part = str(raw_slide.get("sourceSlidePart", ""))
        if not slide_part:
            continue
        rels_xml = package_entries.get(rels_part_for_slide_part(slide_part))
        if rels_xml is None:
            continue
        try:
            relationships_root = ET.fromstring(rels_xml)
        except ET.ParseError:
            continue
        for relationship in list(relationships_root):
            if relationship.get("Type") == SLIDE_LAYOUT_REL_TYPE:
                target = str(relationship.get("Target", ""))
                if target:
                    return target
    return ""


def next_relationship_id(root: ET.Element[Any]) -> str:
    ids = [
        int(str(child.get("Id", "")).removeprefix("rId"))
        for child in list(root)
        if str(child.get("Id", "")).startswith("rId")
        and str(child.get("Id", "")).removeprefix("rId").isdigit()
    ]
    return f"rId{max(ids, default=0) + 1}"


def empty_slide_xml() -> bytes:
    slide = ET.Element(f"{{{PML_NS}}}sld")
    common = ET.SubElement(slide, f"{{{PML_NS}}}cSld")
    shape_tree = ET.SubElement(common, f"{{{PML_NS}}}spTree")
    non_visual = ET.SubElement(shape_tree, f"{{{PML_NS}}}nvGrpSpPr")
    ET.SubElement(
        non_visual,
        f"{{{PML_NS}}}cNvPr",
        {"id": "1", "name": ""},
    )
    ET.SubElement(non_visual, f"{{{PML_NS}}}cNvGrpSpPr")
    ET.SubElement(non_visual, f"{{{PML_NS}}}nvPr")
    group_properties = ET.SubElement(shape_tree, f"{{{PML_NS}}}grpSpPr")
    transform = ET.SubElement(group_properties, f"{{{DML_NS}}}xfrm")
    for name in ("off", "chOff"):
        ET.SubElement(transform, f"{{{DML_NS}}}{name}", {"x": "0", "y": "0"})
    for name in ("ext", "chExt"):
        ET.SubElement(transform, f"{{{DML_NS}}}{name}", {"cx": "0", "cy": "0"})
    color_map = ET.SubElement(slide, f"{{{PML_NS}}}clrMapOvr")
    ET.SubElement(color_map, f"{{{DML_NS}}}masterClrMapping")
    return xml_bytes(slide)


def slide_layout_relationships_xml(target: str) -> bytes:
    root = ET.Element(f"{{{PKG_REL_NS}}}Relationships")
    ET.SubElement(
        root,
        f"{{{PKG_REL_NS}}}Relationship",
        {"Id": "rId1", "Type": SLIDE_LAYOUT_REL_TYPE, "Target": target},
    )
    return xml_bytes(root)


def add_element_to_slide_xml(
    slide_part: str,
    element: dict[str, Any],
    package_entries: dict[str, bytes],
    added_entries: dict[str, bytes],
    scale: PackageFrameScale,
    warnings: list[str],
) -> dict[str, Any] | None:
    slide_xml = package_entries.get(slide_part)
    if slide_xml is None:
        warnings.append(
            f"OOXML slide part missing for added element {element.get('elementId')}."
        )
        return None

    root = ET.fromstring(slide_xml)
    shape_tree = first_local_descendant(root, "spTree")
    if shape_tree is None:
        warnings.append(
            f"OOXML shape tree missing for added element {element.get('elementId')}."
        )
        return None

    next_shape_id = next_c_nv_pr_id(root)
    element_id = str(element.get("elementId", ""))
    element_type = str(element.get("type", ""))
    if element_type == "text":
        shape = text_shape_element(next_shape_id, element, scale)
        source_type = "slide"
        relationship_id = None
    elif element_type == "rect":
        shape = rect_shape_element(next_shape_id, element, scale)
        source_type = "slide"
        relationship_id = None
    elif element_type == "image":
        replacement = decode_image_data_url(dict_value(element, "props").get("src"))
        if isinstance(replacement, str):
            warnings.append(
                f"OOXML add_element image skipped for {element_id}: {replacement}."
            )
            return None
        content_types_xml = package_entries.get("[Content_Types].xml")
        if content_types_xml is None:
            warnings.append(f"OOXML content types missing for {element_id}.")
            return None
        mime_type, image_blob = replacement
        extension = extension_for_mime_type(mime_type)
        rels_part = rels_part_for_slide_part(slide_part)
        rels_xml = package_entries.get(rels_part, empty_relationships_xml())
        media_name = f"orbit_sync_{Path(slide_part).stem}_{next_shape_id}.{extension}"
        media_part = f"ppt/media/{media_name}"
        try:
            relationship_id, next_rels_xml = append_image_relationship(
                rels_xml,
                f"../media/{media_name}",
            )
            next_content_types_xml = ensure_content_type_default(
                content_types_xml,
                extension,
                mime_type,
            )
        except (ET.ParseError, ValueError):
            warnings.append(f"OOXML image relationship invalid for {element_id}.")
            return None
        shape = picture_shape_element(
            next_shape_id,
            element,
            relationship_id,
            scale,
            image_dimensions(image_blob),
        )
        source_type = "image"
    elif element_type == "table":
        shape = table_graphic_frame_element(next_shape_id, element, scale)
        source_type = "table"
        relationship_id = None
    else:
        warnings.append(f"OOXML add_element skipped for {element_type}.")
        return None

    shape_tree.append(shape)
    package_entries[slide_part] = xml_bytes(root)
    if element_type == "image":
        package_entries[rels_part] = next_rels_xml
        package_entries["[Content_Types].xml"] = next_content_types_xml
        added_entries[media_part] = image_blob

    source: dict[str, Any] = {
        "elementId": element_id,
        "elementType": element_type,
        "ooxmlOrigin": "authored",
        "ooxmlEditCapabilities": {
            "richText": "full" if element_type == "text" else "none",
            "crop": "picture" if element_type == "image" else "none",
            "tableCellText": element_type == "table",
            "frame": True,
            "delete": True,
            "imageSource": element_type == "image",
        },
        "slidePart": slide_part,
        "shapeId": str(next_shape_id),
        "sourceType": source_type,
        "writable": True,
    }
    if relationship_id is not None:
        source["relationshipId"] = relationship_id
    if element_type == "table" and not refresh_table_source_locators(shape, source):
        warnings.append(
            f"OOXML authored table locator creation failed for {element_id}."
        )
        return None
    return source


def slide_part_for_operation(
    operation: dict[str, Any],
    template_blueprint: dict[str, Any],
) -> str:
    explicit_slide_part = str(operation.get("sourceSlidePart", ""))
    if explicit_slide_part:
        return explicit_slide_part
    operation_slide_id = str(operation.get("slideId", ""))
    matches = [
        str(slide.get("sourceSlidePart", ""))
        for slide in template_blueprint.get("slides", [])
        if isinstance(slide, dict) and slide.get("slideId") == operation_slide_id
    ]
    return matches[0] if len(matches) == 1 else ""


def source_slide_part(slide: dict[str, Any]) -> str:
    explicit = str(slide.get("sourceSlidePart", ""))
    if is_safe_slide_part(explicit):
        return explicit
    candidates = {
        str(source.get("slidePart", ""))
        for source in slide.get("elementSources", [])
        if isinstance(source, dict)
        and bool(source.get("writable", False))
        and is_safe_slide_part(str(source.get("slidePart", "")))
    }
    return next(iter(candidates)) if len(candidates) == 1 else ""


def route_operations_to_source_parts(
    template_blueprint: dict[str, Any],
    operations: list[dict[str, Any]],
) -> list[dict[str, Any]]:
    slides = [
        slide
        for slide in template_blueprint.get("slides", [])
        if isinstance(slide, dict)
    ]
    routed: list[dict[str, Any]] = []
    for operation in operations:
        if is_safe_slide_part(str(operation.get("sourceSlidePart", ""))):
            routed.append(operation)
            continue
        operation_slide_index = slide_index_from_id(str(operation.get("slideId", "")))
        slide = next(
            (
                candidate
                for candidate in slides
                if int_value(candidate.get("slideIndex"), 0) == operation_slide_index
            ),
            None,
        )
        if slide is None:
            slide = next(
                (
                    candidate
                    for candidate in slides
                    if int_value(candidate.get("sourceSlideIndex"), 0)
                    == operation_slide_index
                ),
                None,
            )
        slide_part = source_slide_part(slide) if slide is not None else ""
        routed.append(
            {**operation, **({"sourceSlidePart": slide_part} if slide_part else {})}
        )
    return routed


def is_safe_slide_part(value: str) -> bool:
    return (
        value.startswith("ppt/slides/slide")
        and value.endswith(".xml")
        and "/" not in value.removeprefix("ppt/slides/")
        and ".." not in value
    )


def slide_index_from_id(slide_id: str) -> int:
    suffix = slide_id.rsplit("_", maxsplit=1)[-1]
    return max(1, int_value(suffix, 1))


def text_shape_element(
    shape_id: int,
    element: dict[str, Any],
    scale: PackageFrameScale,
) -> ET.Element[Any]:
    shape = base_shape_element(shape_id, "Orbit text", element, scale)
    sync_text_shape(
        shape,
        dict_value(element, "props"),
        {
            "elementType": "text",
            "ooxmlOrigin": "authored",
            "ooxmlEditCapabilities": {"richText": "full"},
        },
        scale,
    )
    return shape


def rect_shape_element(
    shape_id: int,
    element: dict[str, Any],
    scale: PackageFrameScale,
) -> ET.Element[Any]:
    shape = base_shape_element(shape_id, "Orbit rect", element, scale)
    sp_pr = ensure_shape_properties(shape)
    props = dict_value(element, "props")
    border_radius = float(props.get("borderRadius", 0))
    geometry = ET.SubElement(
        sp_pr,
        f"{{{DML_NS}}}prstGeom",
        {"prst": "roundRect" if border_radius > 0 else "rect"},
    )
    adjustments = ET.SubElement(geometry, f"{{{DML_NS}}}avLst")
    if border_radius > 0:
        shortest_side = min(float(element["width"]), float(element["height"]))
        adjustment = round(min(0.5, border_radius / shortest_side) * 100000)
        ET.SubElement(
            adjustments,
            f"{{{DML_NS}}}gd",
            {"name": "adj", "fmla": f"val {adjustment}"},
        )
    fill = props.get("fill")
    if fill == "transparent":
        ET.SubElement(sp_pr, f"{{{DML_NS}}}noFill")
    elif isinstance(fill, str) and valid_hex_color(fill):
        solid_fill = ET.SubElement(sp_pr, f"{{{DML_NS}}}solidFill")
        ET.SubElement(solid_fill, f"{{{DML_NS}}}srgbClr", {"val": fill[1:]})
    stroke = props.get("stroke", "transparent")
    stroke_width = float(props.get("strokeWidth", 0))
    line = ET.SubElement(
        sp_pr,
        f"{{{DML_NS}}}ln",
        {"w": str(table_border_width_to_emu(stroke_width, scale))},
    )
    if stroke == "transparent" or stroke_width == 0:
        ET.SubElement(line, f"{{{DML_NS}}}noFill")
    elif isinstance(stroke, str) and valid_hex_color(stroke):
        line_fill = ET.SubElement(line, f"{{{DML_NS}}}solidFill")
        ET.SubElement(line_fill, f"{{{DML_NS}}}srgbClr", {"val": stroke[1:]})
    return shape


def update_authored_rect_props(
    shape: ET.Element[Any],
    props: dict[str, Any],
    scale: PackageFrameScale,
) -> bool:
    sp_pr = first_local_child(shape, "spPr")
    if sp_pr is None:
        return False
    if "borderRadius" in props:
        geometry = first_local_child(sp_pr, "prstGeom")
        xfrm = first_local_child(sp_pr, "xfrm")
        ext = first_local_child(xfrm, "ext") if xfrm is not None else None
        if geometry is None or ext is None:
            return False
        border_radius = float(props["borderRadius"])
        geometry.set("prst", "roundRect" if border_radius > 0 else "rect")
        adjustments = first_local_child(geometry, "avLst")
        if adjustments is None:
            adjustments = ET.SubElement(geometry, f"{{{DML_NS}}}avLst")
        adjustments.clear()
        if border_radius > 0:
            width = int_value(ext.get("cx"), 0) * scale.canvas_width / scale.slide_width_emu
            height = int_value(ext.get("cy"), 0) * scale.canvas_height / scale.slide_height_emu
            if width <= 0 or height <= 0:
                return False
            adjustment = round(min(0.5, border_radius / min(width, height)) * 100000)
            ET.SubElement(
                adjustments,
                f"{{{DML_NS}}}gd",
                {"name": "adj", "fmla": f"val {adjustment}"},
            )
    if "fill" in props:
        for child in list(sp_pr):
            if local_name(child) in {
                "blipFill",
                "gradFill",
                "grpFill",
                "noFill",
                "pattFill",
                "solidFill",
            }:
                sp_pr.remove(child)
        fill = props["fill"]
        fill_node = ET.Element(
            f"{{{DML_NS}}}{'noFill' if fill == 'transparent' else 'solidFill'}"
        )
        if fill != "transparent":
            ET.SubElement(fill_node, f"{{{DML_NS}}}srgbClr", {"val": fill[1:]})
        line = first_local_child(sp_pr, "ln")
        sp_pr.insert(list(sp_pr).index(line) if line is not None else len(sp_pr), fill_node)
    if "stroke" in props or "strokeWidth" in props:
        line = first_local_child(sp_pr, "ln")
        if line is None:
            line = ET.SubElement(sp_pr, f"{{{DML_NS}}}ln")
        if "strokeWidth" in props:
            line.set(
                "w",
                str(table_border_width_to_emu(props["strokeWidth"], scale)),
            )
        if "stroke" in props:
            for child in list(line):
                if local_name(child) in {
                    "gradFill",
                    "noFill",
                    "pattFill",
                    "solidFill",
                }:
                    line.remove(child)
            stroke = props["stroke"]
            line_fill = ET.Element(
                f"{{{DML_NS}}}{'noFill' if stroke == 'transparent' else 'solidFill'}"
            )
            if stroke != "transparent":
                ET.SubElement(
                    line_fill,
                    f"{{{DML_NS}}}srgbClr",
                    {"val": stroke[1:]},
                )
            line.insert(0, line_fill)
    return True


def picture_shape_element(
    shape_id: int,
    element: dict[str, Any],
    relationship_id: str,
    scale: PackageFrameScale,
    image_size: tuple[int, int] | None,
) -> ET.Element[Any]:
    picture = ET.Element(P_PIC)
    nv_pic_pr = ET.SubElement(picture, f"{{{PML_NS}}}nvPicPr")
    ET.SubElement(
        nv_pic_pr,
        f"{{{PML_NS}}}cNvPr",
        {"id": str(shape_id), "name": "Orbit image"},
    )
    ET.SubElement(nv_pic_pr, f"{{{PML_NS}}}cNvPicPr")
    ET.SubElement(nv_pic_pr, f"{{{PML_NS}}}nvPr")
    blip_fill = ET.SubElement(picture, f"{{{PML_NS}}}blipFill")
    ET.SubElement(
        blip_fill,
        A_BLIP,
        {f"{{{REL_NS}}}embed": relationship_id},
    )
    stretch = ET.SubElement(blip_fill, f"{{{DML_NS}}}stretch")
    ET.SubElement(stretch, f"{{{DML_NS}}}fillRect")
    sp_pr = ET.SubElement(picture, f"{{{PML_NS}}}spPr")
    update_shape_frame(picture, element, scale)
    frame_size = picture_frame_size(picture)
    if image_size is not None and frame_size is not None:
        set_picture_contain_source_rect(picture, image_size, frame_size)
    crop = normalized_image_crop(dict_value(element, "props").get("crop"))
    if crop is not None:
        set_picture_crop_source_rect(picture, crop)
    ET.SubElement(
        ET.SubElement(sp_pr, f"{{{DML_NS}}}prstGeom", {"prst": "rect"}),
        f"{{{DML_NS}}}avLst",
    )
    return picture


def image_dimensions(image_blob: bytes) -> tuple[int, int] | None:
    try:
        with Image.open(BytesIO(image_blob)) as image:
            width, height = image.size
    except (OSError, SyntaxError, ValueError):
        return None
    if width <= 0 or height <= 0:
        return None
    return int(width), int(height)


def picture_frame_size(shape: ET.Element[Any]) -> tuple[int, int] | None:
    xfrm = first_local_descendant(shape, "xfrm")
    if xfrm is None:
        return None
    ext = first_local_child(xfrm, "ext")
    if ext is None:
        return None
    width = int_value(ext.get("cx"), 0)
    height = int_value(ext.get("cy"), 0)
    if width <= 0 or height <= 0:
        return None
    return width, height


def set_picture_crop_source_rect(
    shape: ET.Element[Any],
    crop: dict[str, float] | None,
) -> bool:
    blip_fill = direct_image_blip_fill(shape)
    if blip_fill is None:
        return False
    children = list(blip_fill)
    blip_index = next(
        (index for index, child in enumerate(children) if local_name(child) == "blip"),
        -1,
    )
    if blip_index < 0:
        return False
    for child in children:
        if local_name(child) == "srcRect":
            blip_fill.remove(child)
    if crop is None:
        return True

    blip_fill.insert(
        blip_index + 1,
        ET.Element(
            f"{{{DML_NS}}}srcRect",
            crop_source_rect_attributes(crop),
        ),
    )
    return True


def crop_source_rect_attributes(crop: dict[str, float]) -> dict[str, str]:
    values = {
        edge: max(0, min(99_999, round(crop[name] * 100_000)))
        for edge, name in (
            ("l", "left"),
            ("t", "top"),
            ("r", "right"),
            ("b", "bottom"),
        )
    }
    for first, second in (("l", "r"), ("t", "b")):
        overflow = values[first] + values[second] - 99_999
        if overflow > 0:
            reduction = min(values[second], overflow)
            values[second] -= reduction
            values[first] -= overflow - reduction
    return {edge: str(values[edge]) for edge in ("l", "t", "r", "b")}


def set_picture_contain_source_rect(
    picture: ET.Element[Any],
    image_size: tuple[int, int],
    frame_size: tuple[int, int],
) -> None:
    blip_fill = direct_image_blip_fill(picture)
    if blip_fill is None:
        return

    image_width, image_height = image_size
    frame_width, frame_height = frame_size
    image_ratio = image_width / image_height
    frame_ratio = frame_width / frame_height
    attributes: dict[str, str] = {}
    if not math.isclose(image_ratio, frame_ratio, rel_tol=1e-6, abs_tol=1e-9):
        if image_ratio > frame_ratio:
            edge = -round((image_ratio / frame_ratio - 1) * 50_000)
            attributes = {"t": str(edge), "b": str(edge)}
        else:
            edge = -round((frame_ratio / image_ratio - 1) * 50_000)
            attributes = {"l": str(edge), "r": str(edge)}
    if not attributes:
        return

    children = list(blip_fill)
    blip_index = next(
        (index for index, child in enumerate(children) if local_name(child) == "blip"),
        -1,
    )
    blip_fill.insert(
        blip_index + 1,
        ET.Element(f"{{{DML_NS}}}srcRect", attributes),
    )


def valid_hex_color(value: str) -> bool:
    if len(value) != 7 or not value.startswith("#"):
        return False
    try:
        int(value[1:], 16)
    except ValueError:
        return False
    return True


def base_shape_element(
    shape_id: int,
    name: str,
    element: dict[str, Any],
    scale: PackageFrameScale,
) -> ET.Element[Any]:
    shape = ET.Element(P_SP)
    nv_sp_pr = ET.SubElement(shape, f"{{{PML_NS}}}nvSpPr")
    ET.SubElement(
        nv_sp_pr,
        f"{{{PML_NS}}}cNvPr",
        {"id": str(shape_id), "name": name},
    )
    ET.SubElement(nv_sp_pr, f"{{{PML_NS}}}cNvSpPr")
    ET.SubElement(nv_sp_pr, f"{{{PML_NS}}}nvPr")
    update_shape_frame(shape, element, scale)
    return shape


def ensure_shape_properties(shape: ET.Element[Any]) -> ET.Element[Any]:
    sp_pr = first_local_child(shape, "spPr")
    if sp_pr is None:
        sp_pr = ET.SubElement(shape, f"{{{PML_NS}}}spPr")
    return sp_pr


def ensure_xfrm(shape: ET.Element[Any]) -> ET.Element[Any]:
    if shape.tag == P_GRAPHIC_FRAME:
        xfrm = first_local_child(shape, "xfrm")
        if xfrm is None:
            non_visual = first_local_child(shape, "nvGraphicFramePr")
            insert_at = (
                list(shape).index(non_visual) + 1 if non_visual is not None else 0
            )
            xfrm = ET.Element(f"{{{PML_NS}}}xfrm")
            shape.insert(insert_at, xfrm)
        return xfrm
    sp_pr = ensure_shape_properties(shape)
    xfrm = first_local_child(sp_pr, "xfrm")
    if xfrm is None:
        xfrm = ET.SubElement(sp_pr, f"{{{DML_NS}}}xfrm")
    return xfrm


def ensure_text_body(shape: ET.Element[Any]) -> ET.Element[Any]:
    tx_body = first_local_child(shape, "txBody")
    if tx_body is None:
        tx_body = ET.SubElement(shape, f"{{{PML_NS}}}txBody")
        ET.SubElement(tx_body, f"{{{DML_NS}}}bodyPr")
        ET.SubElement(tx_body, f"{{{DML_NS}}}lstStyle")
    return tx_body


def frame_to_emu(
    frame: dict[str, Any],
    scale: PackageFrameScale,
) -> tuple[int, int, int, int]:
    return (
        round(float(frame.get("x", 0)) * scale.slide_width_emu / scale.canvas_width),
        round(float(frame.get("y", 0)) * scale.slide_height_emu / scale.canvas_height),
        max(
            1,
            round(
                float(frame.get("width", 1))
                * scale.slide_width_emu
                / scale.canvas_width
            ),
        ),
        max(
            1,
            round(
                float(frame.get("height", 1))
                * scale.slide_height_emu
                / scale.canvas_height
            ),
        ),
    )


def next_c_nv_pr_id(root: ET.Element[Any]) -> int:
    ids = [
        int(node.get("id", "0"))
        for node in root.iter()
        if local_name(node) == "cNvPr" and str(node.get("id", "")).isdigit()
    ]
    return max(ids, default=0) + 1


def dict_value(value: dict[str, Any], key: str) -> dict[str, Any]:
    item = value.get(key)
    return item if isinstance(item, dict) else {}


def first_local_child(element: ET.Element[Any], name: str) -> ET.Element[Any] | None:
    for child in list(element):
        if local_name(child) == name:
            return child
    return None


def direct_local_children(element: ET.Element[Any], name: str) -> list[ET.Element[Any]]:
    return [child for child in list(element) if local_name(child) == name]


def first_local_descendant(
    element: ET.Element[Any], name: str
) -> ET.Element[Any] | None:
    for child in element.iter():
        if local_name(child) == name:
            return child
    return None


def local_name(element: Any) -> str:
    tag = getattr(element, "tag", element)
    return str(tag).rsplit("}", maxsplit=1)[-1]


def scale_slot_bounds(slot: dict[str, Any], scale_x: float, scale_y: float) -> None:
    bounds = slot.get("bounds")
    if not isinstance(bounds, dict):
        return
    bounds["x"] = round(float(bounds.get("x", 0)) * scale_x, 3)
    bounds["y"] = round(float(bounds.get("y", 0)) * scale_y, 3)
    bounds["width"] = max(1, round(float(bounds.get("width", 1)) * scale_x, 3))
    bounds["height"] = max(1, round(float(bounds.get("height", 1)) * scale_y, 3))


def rels_part_for_slide_part(slide_part: str) -> str:
    path, name = slide_part.rsplit("/", maxsplit=1)
    return f"{path}/_rels/{name}.rels"


def append_image_relationship(rels_xml: bytes, target: str) -> tuple[str, bytes]:
    root = ET.fromstring(rels_xml)
    ids = [
        int(str(child.get("Id", "")).removeprefix("rId"))
        for child in list(root)
        if str(child.get("Id", "")).startswith("rId")
        and str(child.get("Id", "")).removeprefix("rId").isdigit()
    ]
    relationship_id = f"rId{max(ids, default=0) + 1}"
    ET.SubElement(
        root,
        f"{{{PKG_REL_NS}}}Relationship",
        {"Id": relationship_id, "Type": IMAGE_REL_TYPE, "Target": target},
    )
    return relationship_id, xml_bytes(root)


def ensure_content_type_default(
    content_types_xml: bytes,
    extension: str,
    mime_type: str,
) -> bytes:
    root = ET.fromstring(content_types_xml)
    for child in list(root):
        if child.tag.endswith("Default") and child.get("Extension") == extension:
            return content_types_xml
    ET.SubElement(
        root,
        f"{{{CONTENT_TYPES_NS}}}Default",
        {"Extension": extension, "ContentType": mime_type},
    )
    return xml_bytes(root)


def rewrite_zip(
    source: zipfile.ZipFile,
    changed_entries: dict[str, bytes],
    added_entries: dict[str, bytes] | None = None,
) -> bytes:
    buffer = BytesIO()
    added = added_entries or {}
    with zipfile.ZipFile(buffer, "w") as target:
        for info in source.infolist():
            if info.filename in added:
                continue
            target.writestr(
                info,
                changed_entries.get(info.filename, source.read(info.filename)),
            )
        for name, content in added.items():
            target.writestr(name, content)
        for name, content in changed_entries.items():
            if name not in source.namelist():
                target.writestr(name, content)
    return buffer.getvalue()


def render_pptx_to_png_assets(
    package_bytes: bytes,
    canvas: CanvasSpec,
) -> list[ImportedDesignAsset]:
    executable = shutil.which("libreoffice") or shutil.which("soffice")
    if executable is None:
        raise PptxRenderUnavailableError(
            "LibreOffice is required to render PPTX slides."
        )

    with TemporaryDirectory(prefix="orbit-ooxml-render-") as temp_dir:
        temp_path = Path(temp_dir)
        pptx_path = temp_path / "source.pptx"
        out_dir = temp_path / "out"
        out_dir.mkdir()
        pptx_path.write_bytes(package_bytes)
        try:
            subprocess.run(
                [
                    executable,
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(out_dir),
                    str(pptx_path),
                ],
                check=True,
                capture_output=True,
                text=True,
                timeout=120,
            )
        except (subprocess.CalledProcessError, subprocess.TimeoutExpired) as error:
            raise PptxRenderUnavailableError(
                "LibreOffice failed to render PPTX slides."
            ) from error
        pdf_path = out_dir / "source.pdf"
        if not pdf_path.exists():
            raise PptxRenderUnavailableError("LibreOffice did not produce a PDF.")
        return render_pdf_to_png_assets(pdf_path, canvas)


def render_pdf_to_png_assets(
    pdf_path: Path,
    canvas: CanvasSpec,
) -> list[ImportedDesignAsset]:
    fitz: Any = importlib.import_module("fitz")
    assets: list[ImportedDesignAsset] = []
    document = fitz.open(str(pdf_path))
    try:
        for page_index in range(document.page_count):
            page = document.load_page(page_index)
            matrix = fitz.Matrix(
                canvas.width / float(page.rect.width),
                canvas.height / float(page.rect.height),
            )
            pixmap = page.get_pixmap(matrix=matrix, alpha=False)
            assets.append(
                ImportedDesignAsset(
                    assetId=f"slide_render_{page_index + 1}",
                    fileName=f"slide-{page_index + 1:02d}.png",
                    mimeType="image/png",
                    contentBase64=base64.b64encode(pixmap.tobytes("png")).decode(
                        "ascii"
                    ),
                )
            )
    finally:
        document.close()
    if not assets:
        raise PptxRenderUnavailableError("Rendered PDF has no pages.")
    return assets


def shape_fallback_assets(
    blueprint: dict[str, Any],
    slide_render_assets: list[ImportedDesignAsset],
    warnings: list[str],
) -> list[ImportedDesignAsset]:
    render_assets_by_slide = slide_render_assets_by_index(slide_render_assets)
    assets: list[ImportedDesignAsset] = []
    seen_asset_ids: set[str] = set()
    slides = blueprint.get("slides", [])
    if not isinstance(slides, list):
        return assets

    for index, slide in enumerate(slides):
        if not isinstance(slide, dict):
            continue
        slide_index = int_value(slide.get("sourceSlideIndex"), index + 1)
        fallback_elements = [
            element
            for element in slide.get("elements", [])
            if isinstance(element, dict)
            and shape_fallback_asset_id_from_element(element) is not None
        ]
        if not fallback_elements:
            continue

        render_asset = render_assets_by_slide.get(slide_index)
        if render_asset is None:
            warnings.append(
                f"Shape image fallback skipped; slide render missing: {slide_index}"
            )
            continue

        try:
            image_bytes = base64.b64decode(render_asset.content_base64)
            with Image.open(BytesIO(image_bytes)) as source_image:
                rendered = source_image.convert("RGBA")
        except Exception:
            warnings.append(
                f"Shape image fallback skipped; slide render unreadable: {slide_index}"
            )
            continue

        for element in fallback_elements:
            asset_id = shape_fallback_asset_id_from_element(element)
            if asset_id is None or asset_id in seen_asset_ids:
                continue
            crop_box = element_crop_box(element, rendered.size)
            if crop_box is None:
                warnings.append(
                    f"Shape image fallback skipped; invalid frame: {asset_id}"
                )
                continue
            crop = rendered.crop(crop_box)
            buffer = BytesIO()
            crop.save(buffer, format="PNG")
            assets.append(
                ImportedDesignAsset(
                    assetId=asset_id,
                    fileName=f"{asset_id}.png",
                    mimeType="image/png",
                    contentBase64=base64.b64encode(buffer.getvalue()).decode("ascii"),
                )
            )
            seen_asset_ids.add(asset_id)

    return assets


def blueprint_has_shape_fallbacks(blueprint: dict[str, Any]) -> bool:
    slides = blueprint.get("slides", [])
    if not isinstance(slides, list):
        return False
    return any(
        isinstance(element, dict)
        and shape_fallback_asset_id_from_element(element) is not None
        for slide in slides
        if isinstance(slide, dict)
        for element in slide.get("elements", [])
    )


def strip_text_from_pptx_package(package_bytes: bytes) -> bytes:
    changed_entries: dict[str, bytes] = {}
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as package:
        for name in package.namelist():
            if not is_presentation_visual_part(name):
                continue
            root = ET.fromstring(package.read(name))
            if remove_text_bodies(root):
                changed_entries[name] = xml_bytes(root)
        if not changed_entries:
            return package_bytes
        return rewrite_zip(package, changed_entries)


def is_presentation_visual_part(name: str) -> bool:
    return (
        name.startswith("ppt/slides/slide")
        or name.startswith("ppt/slideLayouts/slideLayout")
        or name.startswith("ppt/slideMasters/slideMaster")
    ) and name.endswith(".xml")


def remove_text_bodies(element: ET.Element[Any]) -> bool:
    changed = False
    for child in list(element):
        if local_name(child) == "txBody":
            element.remove(child)
            changed = True
        elif remove_text_bodies(child):
            changed = True
    return changed


def slide_render_assets_by_index(
    slide_render_assets: list[ImportedDesignAsset],
) -> dict[int, ImportedDesignAsset]:
    assets: dict[int, ImportedDesignAsset] = {}
    for asset in slide_render_assets:
        prefix = "slide_render_"
        if not asset.asset_id.startswith(prefix):
            continue
        try:
            slide_index = int(asset.asset_id.removeprefix(prefix))
        except ValueError:
            continue
        assets[slide_index] = asset
    return assets


def shape_fallback_asset_id_from_element(element: dict[str, Any]) -> str | None:
    props = element.get("props")
    if not isinstance(props, dict):
        return None
    src = props.get("src")
    if not isinstance(src, str):
        return None
    prefix = "asset:shape_render_"
    if not src.startswith(prefix):
        return None
    return src.removeprefix("asset:")


def element_crop_box(
    element: dict[str, Any],
    image_size: tuple[int, int],
) -> tuple[int, int, int, int] | None:
    image_width, image_height = image_size
    x = math_floor_float(element.get("x"))
    y = math_floor_float(element.get("y"))
    width = math_ceil_float(element.get("width"))
    height = math_ceil_float(element.get("height"))
    if width <= 0 or height <= 0:
        return None
    left = max(0, min(image_width, x))
    top = max(0, min(image_height, y))
    right = max(left, min(image_width, x + width))
    bottom = max(top, min(image_height, y + height))
    if right <= left or bottom <= top:
        return None
    return left, top, right, bottom


def math_floor_float(value: Any) -> int:
    try:
        return math.floor(float(value))
    except (TypeError, ValueError):
        return 0


def math_ceil_float(value: Any) -> int:
    try:
        return math.ceil(float(value))
    except (TypeError, ValueError):
        return 0


def package_asset(
    asset_id: str, package_bytes: bytes, file_name: str
) -> ImportedDesignAsset:
    return ImportedDesignAsset(
        assetId=asset_id,
        fileName=file_name,
        mimeType="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        contentBase64=base64.b64encode(package_bytes).decode("ascii"),
    )


def empty_relationships_xml() -> bytes:
    return (
        b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        b'<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"/>'
    )


def extension_for_mime_type(mime_type: str) -> str:
    subtype = mime_type.rsplit("/", maxsplit=1)[-1].lower()
    if subtype == "jpeg":
        return "jpg"
    if subtype in {"png", "jpg", "gif", "webp"}:
        return subtype
    return "png"


def int_value(value: Any, fallback: int) -> int:
    try:
        return int(value)
    except (TypeError, ValueError):
        return fallback


def xml_bytes(element: ET.Element[Any]) -> bytes:
    namespace = namespace_for_tag(element.tag)
    if namespace in {CONTENT_TYPES_NS, PKG_REL_NS}:
        ET.register_namespace("", namespace)
    return cast(bytes, ET.tostring(element, encoding="utf-8", xml_declaration=True))


def namespace_for_tag(tag: str) -> str | None:
    if not tag.startswith("{"):
        return None
    return tag[1:].partition("}")[0]


def safe_file_stem(path: Path) -> str:
    stem = path.stem.strip() or "presentation"
    return "".join(
        char if char.isascii() and (char.isalnum() or char in "_-") else "_"
        for char in stem
    )


def safe_id_component(value: str) -> str:
    normalized = "".join(
        char if char.isascii() and (char.isalnum() or char in "_-") else "_"
        for char in value
    )
    return normalized or "pptx"
