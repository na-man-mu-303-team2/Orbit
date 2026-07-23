from __future__ import annotations

import base64
import math
import zipfile
from io import BytesIO
from typing import Any, Literal
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt
from pydantic import BaseModel, Field

from app.ai.pptx_motion import apply_generic_slide_motion


DECK_UNITS_PER_INCH = 144
POINTS_PER_INCH = 72
RICH_TEXT_UNSUPPORTED_HYPERLINK = "PPTX_RICH_TEXT_UNSUPPORTED_HYPERLINK"
RICH_TEXT_UNSUPPORTED_RUN_PROPERTY = "PPTX_RICH_TEXT_UNSUPPORTED_RUN_PROPERTY"
TABLE_STRUCTURE_UNSUPPORTED = "PPTX_TABLE_STRUCTURE_UNSUPPORTED"
TABLE_STYLE_UNSUPPORTED = "PPTX_TABLE_STYLE_UNSUPPORTED"
TABLE_TRACK_MISMATCH = "PPTX_TABLE_TRACK_MISMATCH"
SUPPORTED_RUN_PROPERTIES = {
    "baseline",
    "color",
    "fontFamily",
    "fontSize",
    "fontWeight",
    "italic",
    "letterSpacing",
    "text",
    "underline",
}
HYPERLINK_RUN_PROPERTIES = {"href", "hyperlink", "link"}


class DeckPptxExportRequest(BaseModel):
    deck: dict[str, Any]
    format: Literal["pptx"] = "pptx"


class DeckPptxExportResponse(BaseModel):
    content_base64: str = Field(alias="contentBase64")
    warnings: list[str] = Field(default_factory=list)
    motion_diagnostics: list[dict[str, Any]] = Field(
        default_factory=list,
        alias="motionDiagnostics",
    )


def export_deck_pptx(request: DeckPptxExportRequest) -> DeckPptxExportResponse:
    deck = request.deck
    presentation = Presentation()
    presentation.slide_width = Inches(float(deck["canvas"]["width"]) / 144)
    presentation.slide_height = Inches(float(deck["canvas"]["height"]) / 144)
    blank_layout = presentation.slide_layouts[6]
    warnings: list[str] = []
    element_targets_by_slide: list[dict[str, list[str]]] = []

    for slide_data in deck.get("slides", []):
        slide = presentation.slides.add_slide(blank_layout)
        apply_slide_background(slide, slide_data, deck)
        elements = sorted(
            slide_data.get("elements", []),
            key=lambda element: int(element.get("zIndex", 0)),
        )
        element_targets: dict[str, list[str]] = {}
        for element in elements:
            if not element.get("visible", True):
                continue
            before_ids = {int(shape.shape_id) for shape in slide.shapes}
            add_element(slide, element, deck, warnings)
            element_id = str(element.get("elementId", ""))
            if element_id:
                element_targets[element_id] = [
                    str(shape.shape_id)
                    for shape in slide.shapes
                    if int(shape.shape_id) not in before_ids
                ]
        resolve_group_targets(elements, element_targets)
        element_targets_by_slide.append(element_targets)
        add_speaker_notes(slide, slide_data, warnings)

    output = BytesIO()
    presentation.save(output)
    package_bytes, motion_diagnostics = inject_generic_motion(
        output.getvalue(),
        deck,
        element_targets_by_slide,
    )
    return DeckPptxExportResponse(
        contentBase64=base64.b64encode(package_bytes).decode("ascii"),
        warnings=dedupe(warnings),
        motionDiagnostics=bound_motion_diagnostics(motion_diagnostics),
    )


def bound_motion_diagnostics(
    diagnostics: list[dict[str, Any]],
) -> list[dict[str, Any]]:
    if len(diagnostics) <= 500:
        return diagnostics
    aggregates: dict[str, dict[str, Any]] = {}
    for diagnostic in diagnostics:
        code = str(diagnostic.get("code", "PPTX_MOTION_PAYLOAD_INVALID"))
        slide_index = max(1, int(diagnostic.get("slideIndex", 1)))
        count = max(1, int(diagnostic.get("count", 1)))
        current = aggregates.get(code)
        if current is None:
            aggregates[code] = {
                "code": code,
                "slideIndex": slide_index,
                "count": count,
            }
            continue
        current["slideIndex"] = min(int(current["slideIndex"]), slide_index)
        current["count"] = int(current["count"]) + count
    return [aggregates[code] for code in sorted(aggregates)]


def resolve_group_targets(
    elements: list[dict[str, Any]],
    element_targets: dict[str, list[str]],
) -> None:
    groups = {
        str(element.get("elementId", "")): element
        for element in elements
        if element.get("type") == "group" and element.get("elementId")
    }

    def targets_for(element_id: str, visiting: set[str]) -> list[str]:
        direct = element_targets.get(element_id, [])
        if direct or element_id not in groups or element_id in visiting:
            return direct
        visiting.add(element_id)
        result: list[str] = []
        child_ids = groups[element_id].get("props", {}).get("childElementIds", [])
        if isinstance(child_ids, list):
            for child_id in child_ids:
                for shape_id in targets_for(str(child_id), visiting):
                    if shape_id not in result:
                        result.append(shape_id)
        visiting.remove(element_id)
        element_targets[element_id] = result
        return result

    for group_id in groups:
        targets_for(group_id, set())


def inject_generic_motion(
    package_bytes: bytes,
    deck: dict[str, Any],
    element_targets_by_slide: list[dict[str, list[str]]],
) -> tuple[bytes, list[dict[str, Any]]]:
    diagnostics: list[dict[str, Any]] = []
    changed_entries: dict[str, bytes] = {}
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as source:
        source_names = set(source.namelist())
        for slide_index, slide_data in enumerate(deck.get("slides", []), start=1):
            slide_part = f"ppt/slides/slide{slide_index}.xml"
            if slide_part not in source_names or not isinstance(slide_data, dict):
                continue
            try:
                slide_root = ET.fromstring(source.read(slide_part))
                slide_diagnostics = apply_generic_slide_motion(
                    slide_root,
                    slide_data,
                    slide_index=slide_index,
                    element_targets=(
                        element_targets_by_slide[slide_index - 1]
                        if slide_index <= len(element_targets_by_slide)
                        else {}
                    ),
                )
            except (ET.ParseError, TypeError, ValueError):
                diagnostics.append(
                    {
                        "code": "PPTX_MOTION_PAYLOAD_INVALID",
                        "slideIndex": slide_index,
                    }
                )
                continue
            diagnostics.extend(slide_diagnostics)
            changed_entries[slide_part] = ET.tostring(
                slide_root,
                encoding="utf-8",
                xml_declaration=True,
            )
        if not changed_entries:
            return package_bytes, diagnostics
        target_buffer = BytesIO()
        with zipfile.ZipFile(target_buffer, "w") as target:
            for info in source.infolist():
                target.writestr(
                    info,
                    changed_entries.get(info.filename, source.read(info.filename)),
                )
        return target_buffer.getvalue(), diagnostics


def add_speaker_notes(
    slide: Any,
    slide_data: dict[str, Any],
    warnings: list[str],
) -> None:
    speaker_notes = str(slide_data.get("speakerNotes", "")).strip()
    if not speaker_notes:
        return
    notes_text_frame = slide.notes_slide.notes_text_frame
    if notes_text_frame is None:
        warnings.append(
            f"Skipped speaker notes for slide {slide_data.get('order', '?')}: "
            "notes placeholder is unavailable."
        )
        return
    notes_text_frame.text = speaker_notes


def apply_slide_background(slide: Any, slide_data: dict[str, Any], deck: dict[str, Any]) -> None:
    color = (
        slide_data.get("style", {}).get("backgroundColor")
        or deck.get("theme", {}).get("backgroundColor")
        or "#FFFFFF"
    )
    if not is_hex_color(color):
        return
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_element(
    slide: Any,
    element: dict[str, Any],
    deck: dict[str, Any],
    warnings: list[str],
) -> None:
    element_type = element.get("type")
    if element_type == "text":
        add_text(slide, element, deck, warnings)
    elif element_type in {"rect", "ellipse"}:
        add_shape(slide, element, element_type)
    elif element_type in {"line", "arrow"}:
        add_line(slide, element)
        if element_type == "arrow":
            warnings.append("Arrowheads are exported as plain lines in phase 1.")
    elif element_type == "image":
        add_image(slide, element, warnings)
    elif element_type == "chart":
        add_chart(slide, element, warnings)
    elif element_type == "table":
        add_table(slide, element, deck, warnings)
    else:
        warnings.append(f"Skipped unsupported element type: {element_type}")


def add_text(
    slide: Any,
    element: dict[str, Any],
    deck: dict[str, Any],
    warnings: list[str],
) -> None:
    props = element.get("props", {})
    shape = slide.shapes.add_textbox(
        emu_x(element),
        emu_y(element),
        emu_width(element),
        emu_height(element),
    )
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    body_pr = text_frame._txBody.bodyPr
    body_pr.set("horzOverflow", "clip")
    body_pr.set("vertOverflow", "clip")
    body_pr.set("wrap", "square")
    inset = props.get("bodyInset") or {}
    text_frame.margin_left = Inches(float(inset.get("left", 0)) / 144)
    text_frame.margin_right = Inches(float(inset.get("right", 0)) / 144)
    text_frame.margin_top = Inches(float(inset.get("top", 0)) / 144)
    text_frame.margin_bottom = Inches(float(inset.get("bottom", 0)) / 144)
    text_frame.vertical_anchor = vertical_anchor(props.get("verticalAlign", "top"))
    if "autoFit" in props:
        for child in list(body_pr):
            if child.tag.rsplit("}", maxsplit=1)[-1] in {
                "noAutofit",
                "normAutofit",
                "spAutoFit",
            }:
                body_pr.remove(child)
        auto_fit = str(props.get("autoFit"))
        if auto_fit == "none":
            body_pr.append(OxmlElement("a:noAutofit"))
        elif auto_fit == "resize-shape":
            body_pr.append(OxmlElement("a:spAutoFit"))
        elif auto_fit == "shrink-text":
            normal_autofit = OxmlElement("a:normAutofit")
            normal_autofit.set(
                "fontScale",
                str(round(float(props.get("fontScale", 1)) * 100000)),
            )
            normal_autofit.set(
                "lnSpcReduction",
                str(round(float(props.get("lineSpaceReduction", 0)) * 100000)),
            )
            body_pr.append(normal_autofit)

    paragraphs = props.get("paragraphs")
    element_id = str(element.get("elementId", "unknown"))
    if isinstance(paragraphs, list):
        for index, paragraph_payload in enumerate(paragraphs):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            apply_paragraph(
                paragraph,
                paragraph_payload,
                props,
                deck,
                warnings,
                element_id=element_id,
                paragraph_index=index,
            )
        return

    paragraph = text_frame.paragraphs[0]
    runs = props.get("runs")
    if isinstance(runs, list) and runs:
        apply_paragraph(
            paragraph,
            {"text": props.get("text", ""), "runs": runs},
            props,
            deck,
            warnings,
            element_id=element_id,
            paragraph_index=0,
        )
    else:
        apply_paragraph(
            paragraph,
            {"text": props.get("text", "")},
            props,
            deck,
            warnings,
            element_id=element_id,
            paragraph_index=0,
        )


def apply_paragraph(
    paragraph: Any,
    paragraph_payload: dict[str, Any],
    fallback: dict[str, Any],
    deck: dict[str, Any],
    warnings: list[str],
    *,
    element_id: str,
    paragraph_index: int,
) -> None:
    paragraph.text = ""
    paragraph_props = {**fallback, **paragraph_payload}
    apply_paragraph_style(paragraph, paragraph_props)
    apply_font(paragraph.font, paragraph_props, deck)
    runs = paragraph_payload.get("runs")
    if isinstance(runs, list) and runs:
        for run_index, run_payload in enumerate(runs):
            append_run_diagnostics(
                warnings,
                run_payload,
                element_id=element_id,
                paragraph_index=paragraph_index,
                run_index=run_index,
            )
            run = paragraph.add_run()
            run.text = str(run_payload.get("text", ""))
            apply_font(run.font, {**fallback, **paragraph_payload, **run_payload}, deck)
        return
    run = paragraph.add_run()
    run.text = str(paragraph_payload.get("text", ""))
    apply_font(run.font, {**fallback, **paragraph_payload}, deck)


def apply_paragraph_style(paragraph: Any, props: dict[str, Any]) -> None:
    paragraph.alignment = paragraph_alignment(props.get("align", "left"))
    line_height = float(props.get("lineHeight", 1.2))
    paragraph.line_spacing = line_height
    paragraph.space_before = Pt(canvas_units_to_points(props.get("spaceBefore", 0)))
    paragraph.space_after = Pt(canvas_units_to_points(props.get("spaceAfter", 0)))

    bullet = props.get("bullet") or {}
    indent = float(props.get("indent", 0))
    if bullet.get("enabled"):
        indent = max(indent, float(bullet.get("indent", 0)))

    paragraph_properties = paragraph._p.get_or_add_pPr()
    paragraph_properties.set("marL", str(canvas_units_to_emu(indent)))
    for child in list(paragraph_properties):
        if child.tag.rsplit("}", maxsplit=1)[-1] in {"buAutoNum", "buChar", "buNone"}:
            paragraph_properties.remove(child)
    if bullet.get("enabled"):
        bullet_character = OxmlElement("a:buChar")
        bullet_character.set("char", str(bullet.get("character", "•")))
        paragraph_properties.append(bullet_character)


def apply_font(font: Any, props: dict[str, Any], deck: dict[str, Any]) -> None:
    theme = deck.get("theme", {})
    font.name = str(
        props.get("fontFamily")
        or theme.get("fontFamily")
        or theme.get("typography", {}).get("bodyFontFamily")
        or "Pretendard"
    )
    font.size = Pt(
        float(props.get("fontSize", 24))
        * POINTS_PER_INCH
        / DECK_UNITS_PER_INCH
    )
    font.bold = is_bold(props.get("fontWeight", "normal"))
    font.italic = bool(props.get("italic", False))
    font.underline = bool(props.get("underline", False))
    color = props.get("color") or theme.get("textColor") or "#111827"
    if is_hex_color(color):
        font.color.rgb = rgb(color)
    baseline = props.get("baseline", "normal")
    run_properties = font._element
    if baseline == "superscript":
        run_properties.set("baseline", "30000")
    elif baseline == "subscript":
        run_properties.set("baseline", "-25000")
    else:
        run_properties.attrib.pop("baseline", None)
    if "letterSpacing" in props:
        run_properties.set(
            "spc",
            str(round(canvas_units_to_points(props["letterSpacing"]) * 100)),
        )


def append_run_diagnostics(
    warnings: list[str],
    run_payload: dict[str, Any],
    *,
    element_id: str,
    paragraph_index: int,
    run_index: int,
) -> None:
    location = (
        f"element={element_id}; paragraph={paragraph_index}; run={run_index}"
    )
    if any(key in run_payload for key in HYPERLINK_RUN_PROPERTIES):
        warnings.append(f"{RICH_TEXT_UNSUPPORTED_HYPERLINK}: {location}")
    diagnostic_properties = (
        set(run_payload)
        - SUPPORTED_RUN_PROPERTIES
        - HYPERLINK_RUN_PROPERTIES
    )
    for property_name in sorted(diagnostic_properties):
        warnings.append(
            f"{RICH_TEXT_UNSUPPORTED_RUN_PROPERTY}: property={property_name}; "
            f"{location}"
        )


def canvas_units_to_points(value: Any) -> float:
    return float(value) * POINTS_PER_INCH / DECK_UNITS_PER_INCH


def canvas_units_to_emu(value: Any) -> int:
    return round(float(value) * int(Inches(1)) / DECK_UNITS_PER_INCH)


def add_shape(slide: Any, element: dict[str, Any], element_type: str) -> None:
    shape_type = MSO_SHAPE.OVAL if element_type == "ellipse" else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type,
        emu_x(element),
        emu_y(element),
        emu_width(element),
        emu_height(element),
    )
    apply_fill(shape, element.get("props", {}).get("fill", "transparent"))
    apply_line(shape, element.get("props", {}))


def add_line(slide: Any, element: dict[str, Any]) -> None:
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        emu_x(element),
        emu_y(element),
        emu_x(element) + emu_width(element),
        emu_y(element) + emu_height(element),
    )
    apply_line(shape, element.get("props", {}))


def add_image(slide: Any, element: dict[str, Any], warnings: list[str]) -> None:
    props = element.get("props", {})
    crop = validated_image_crop(props.get("crop"))
    src = str(props.get("src", ""))
    if not src.startswith("data:image/") or ";base64," not in src:
        warnings.append("Skipped image without embedded data URL.")
        return
    _, encoded = src.split(";base64,", 1)
    picture = slide.shapes.add_picture(
        BytesIO(base64.b64decode(encoded)),
        emu_x(element),
        emu_y(element),
        width=emu_width(element),
        height=emu_height(element),
    )
    if crop is not None:
        picture.crop_left = crop["left"]
        picture.crop_top = crop["top"]
        picture.crop_right = crop["right"]
        picture.crop_bottom = crop["bottom"]


def validated_image_crop(value: Any) -> dict[str, float] | None:
    if value is None:
        return None
    if not isinstance(value, dict):
        raise ValueError("Invalid image crop: expected an object.")

    crop: dict[str, float] = {}
    for edge in ("left", "top", "right", "bottom"):
        raw_value = value.get(edge, 0)
        if (
            isinstance(raw_value, bool)
            or not isinstance(raw_value, (int, float))
            or not math.isfinite(raw_value)
            or raw_value < 0
            or raw_value > 1
        ):
            raise ValueError(f"Invalid image crop {edge} fraction.")
        crop[edge] = float(raw_value)

    if crop["left"] + crop["right"] >= 1:
        raise ValueError("Invalid image crop: left and right hide the full image.")
    if crop["top"] + crop["bottom"] >= 1:
        raise ValueError("Invalid image crop: top and bottom hide the full image.")
    units = image_crop_units(crop)
    return {
        "left": units["left"] / 100_000,
        "top": units["top"] / 100_000,
        "right": units["right"] / 100_000,
        "bottom": units["bottom"] / 100_000,
    }


def image_crop_units(crop: dict[str, float]) -> dict[str, int]:
    units = {
        edge: max(0, min(99_999, round(value * 100_000)))
        for edge, value in crop.items()
    }
    for first, second in (("left", "right"), ("top", "bottom")):
        overflow = units[first] + units[second] - 99_999
        if overflow > 0:
            reduction = min(units[second], overflow)
            units[second] -= reduction
            units[first] -= overflow - reduction
    return units


def add_chart(slide: Any, element: dict[str, Any], warnings: list[str]) -> None:
    props = element.get("props", {})
    chart_type = props.get("type")
    if chart_type == "scatter":
        warnings.append("Skipped unsupported scatter chart.")
        return
    data = props.get("data", [])
    if not data:
        warnings.append("Skipped chart without data.")
        return

    chart_data = CategoryChartData()  # type: ignore[no-untyped-call]
    if chart_type == "line":
        categories = list(dict.fromkeys(str(item.get("label", "")) for item in data))
        series_names = list(
            dict.fromkeys(str(item.get("series") or "Series 1") for item in data)
        )
        chart_data.categories = categories
        for series_name in series_names:
            values_by_category = {
                str(item.get("label", "")): float(item.get("value", 0))
                for item in data
                if str(item.get("series") or "Series 1") == series_name
            }
            chart_data.add_series(  # type: ignore[no-untyped-call]
                series_name,
                [values_by_category.get(category, 0.0) for category in categories],
            )
    else:
        chart_data.categories = [str(item.get("label", "")) for item in data]
        chart_data.add_series(  # type: ignore[no-untyped-call]
            str(props.get("title") or "Series"),
            [float(item.get("value", 0)) for item in data],
        )
    chart_shape_type = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
    }.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
    slide.shapes.add_chart(
        chart_shape_type,
        emu_x(element),
        emu_y(element),
        emu_width(element),
        emu_height(element),
        chart_data,
    )


def add_table(
    slide: Any,
    element: dict[str, Any],
    deck: dict[str, Any],
    warnings: list[str],
) -> None:
    props = element.get("props", {})
    rows = props.get("rows") or []
    element_id = str(element.get("elementId", "unknown"))
    issue = table_export_issue(props, element_id)
    if issue is not None:
        warnings.append(issue)
        return
    if not rows:
        return
    row_count = len(rows)
    col_count = len(rows[0])
    merge_anchors, covered_cells, merge_issue = table_merge_layout(
        rows, element_id
    )
    if merge_issue is not None:
        warnings.append(merge_issue)
        return
    table_shape = slide.shapes.add_table(
        row_count,
        col_count,
        emu_x(element),
        emu_y(element),
        emu_width(element),
        emu_height(element),
    )
    table = table_shape.table
    column_widths = normalized_table_track_emu(
        props.get("columnWidths"),
        total=float(element["width"]),
        count=col_count,
    )
    row_heights = normalized_table_track_emu(
        props.get("rowHeights"),
        total=float(element["height"]),
        count=row_count,
    )
    for column_index, width in enumerate(column_widths):
        table.columns[column_index].width = width
    for row_index, height in enumerate(row_heights):
        table.rows[row_index].height = height

    for row_index, row in enumerate(rows):
        for col_index, cell_payload in enumerate(row):
            cell = table.cell(row_index, col_index)
            text_frame = cell.text_frame
            text_frame.clear()
            cell.vertical_anchor = vertical_anchor(
                str(cell_payload.get("verticalAlign", "middle"))
            )
            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = paragraph_alignment(
                str(cell_payload.get("align", "left"))
            )
            run = paragraph.add_run()
            run.text = (
                ""
                if (row_index, col_index) in covered_cells
                else str(cell_payload.get("text", ""))
            )
            apply_font(
                run.font,
                {
                    **cell_payload,
                    "color": cell_payload.get("textColor"),
                },
                deck,
            )
            fill = cell_payload.get("fill")
            if is_hex_color(fill):
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(fill)
            elif fill == "transparent":
                cell.fill.background()
            apply_table_cell_border(
                cell,
                str(
                    cell_payload.get("borderColor")
                    or props.get("borderColor")
                    or "#CBD5E1"
                ),
                cell_payload.get("borderWidth", props.get("borderWidth", 1)),
            )

    for row_index, col_index, row_span, col_span in merge_anchors:
        if row_span == 1 and col_span == 1:
            continue
        table.cell(row_index, col_index).merge(
            table.cell(row_index + row_span - 1, col_index + col_span - 1)
        )


def table_merge_layout(
    rows: list[list[dict[str, Any]]],
    element_id: str,
) -> tuple[
    list[tuple[int, int, int, int]],
    set[tuple[int, int]],
    str | None,
]:
    row_count = len(rows)
    column_count = len(rows[0])
    occupied: set[tuple[int, int]] = set()
    covered: set[tuple[int, int]] = set()
    anchors: list[tuple[int, int, int, int]] = []

    for row_index, row in enumerate(rows):
        for column_index, cell in enumerate(row):
            row_span = cell.get("rowSpan", 1)
            col_span = cell.get("colSpan", 1)
            if (
                isinstance(row_span, bool)
                or not isinstance(row_span, int)
                or isinstance(col_span, bool)
                or not isinstance(col_span, int)
                or row_span < 1
                or col_span < 1
                or row_index + row_span > row_count
                or column_index + col_span > column_count
            ):
                return (
                    [],
                    set(),
                    f"{TABLE_STRUCTURE_UNSUPPORTED}: element={element_id}; "
                    "reason=invalid-cell-span",
                )

            origin = (row_index, column_index)
            if origin in occupied:
                if row_span != 1 or col_span != 1:
                    return (
                        [],
                        set(),
                        f"{TABLE_STRUCTURE_UNSUPPORTED}: element={element_id}; "
                        "reason=overlapping-cell-span",
                    )
                covered.add(origin)
                continue

            cells = {
                (candidate_row, candidate_column)
                for candidate_row in range(row_index, row_index + row_span)
                for candidate_column in range(
                    column_index, column_index + col_span
                )
            }
            if occupied.intersection(cells):
                return (
                    [],
                    set(),
                    f"{TABLE_STRUCTURE_UNSUPPORTED}: element={element_id}; "
                    "reason=overlapping-cell-span",
                )
            occupied.update(cells)
            covered.update(cells - {origin})
            anchors.append((row_index, column_index, row_span, col_span))

    return anchors, covered, None


def table_export_issue(props: Any, element_id: str) -> str | None:
    rows = props.get("rows") if isinstance(props, dict) else None
    if not isinstance(rows, list) or not rows:
        return (
            f"{TABLE_STRUCTURE_UNSUPPORTED}: element={element_id}; "
            "reason=empty-grid"
        )
    if not isinstance(rows[0], list) or not rows[0]:
        return f"{TABLE_STRUCTURE_UNSUPPORTED}: element={element_id}; reason=jagged-grid"
    column_count = len(rows[0])
    if any(not isinstance(row, list) or len(row) != column_count for row in rows):
        return f"{TABLE_STRUCTURE_UNSUPPORTED}: element={element_id}; reason=jagged-grid"
    if any(not isinstance(cell, dict) for row in rows for cell in row):
        return (
            f"{TABLE_STRUCTURE_UNSUPPORTED}: element={element_id}; "
            "reason=invalid-cell"
        )
    _, _, merge_issue = table_merge_layout(rows, element_id)
    if merge_issue is not None:
        return merge_issue

    for row_index, row in enumerate(rows):
        for column_index, cell in enumerate(row):
            fill = cell.get("fill", "transparent")
            if fill != "transparent" and not is_hex_color(fill):
                return (
                    f"{TABLE_STYLE_UNSUPPORTED}: element={element_id}; "
                    f"row={row_index}; column={column_index}; property=fill"
                )

    column_widths = props.get("columnWidths")
    if isinstance(column_widths, list) and len(column_widths) != column_count:
        return (
            f"{TABLE_TRACK_MISMATCH}: element={element_id}; track=columnWidths; "
            f"expected={column_count}; actual={len(column_widths)}"
        )
    row_heights = props.get("rowHeights")
    if isinstance(row_heights, list) and len(row_heights) != len(rows):
        return (
            f"{TABLE_TRACK_MISMATCH}: element={element_id}; track=rowHeights; "
            f"expected={len(rows)}; actual={len(row_heights)}"
        )
    return None


def normalized_table_track_emu(
    tracks: Any,
    *,
    total: float,
    count: int,
) -> list[int]:
    weights = (
        [float(value) for value in tracks]
        if isinstance(tracks, list) and tracks
        else [1.0] * count
    )
    total_emu = max(count, canvas_units_to_emu(total))
    max_weight = max(weights)
    scaled_weights = [weight / max_weight for weight in weights]
    weight_total = sum(scaled_weights)
    distributable = total_emu - count
    exact_extras = [
        distributable * weight / weight_total for weight in scaled_weights
    ]
    floor_extras = [math.floor(value) for value in exact_extras]
    normalized = [1 + value for value in floor_extras]
    remainder = distributable - sum(floor_extras)
    remainder_order = sorted(
        range(count),
        key=lambda index: (-(exact_extras[index] - floor_extras[index]), index),
    )
    for index in remainder_order[:remainder]:
        normalized[index] += 1
    return normalized


def apply_table_cell_border(cell: Any, color: str, width: Any) -> None:
    cell_properties = cell._tc.get_or_add_tcPr()
    border_width = max(0.0, float(width))
    borders: list[Any] = []
    for border_name in ("lnL", "lnR", "lnT", "lnB"):
        for existing in list(cell_properties):
            if existing.tag.rsplit("}", maxsplit=1)[-1] == border_name:
                cell_properties.remove(existing)
        border = OxmlElement(f"a:{border_name}")
        border.set("w", str(canvas_units_to_emu(border_width)))
        if not is_hex_color(color) or border_width <= 0:
            border.append(OxmlElement("a:noFill"))
            borders.append(border)
            continue
        solid_fill = OxmlElement("a:solidFill")
        color_value = OxmlElement("a:srgbClr")
        color_value.set("val", color[1:])
        solid_fill.append(color_value)
        border.append(solid_fill)
        dash = OxmlElement("a:prstDash")
        dash.set("val", "solid")
        border.append(dash)
        borders.append(border)

    fill_names = {"noFill", "solidFill", "gradFill", "blipFill", "pattFill", "grpFill"}
    insertion_index = next(
        (
            index
            for index, child in enumerate(cell_properties)
            if child.tag.rsplit("}", maxsplit=1)[-1] in fill_names
        ),
        0,
    )
    for offset, border in enumerate(borders):
        cell_properties.insert(insertion_index + offset, border)


def apply_fill(shape: Any, fill: Any) -> None:
    if fill == "transparent":
        shape.fill.background()
        return
    if is_hex_color(fill):
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)


def apply_line(shape: Any, props: dict[str, Any]) -> None:
    stroke = props.get("stroke", "transparent")
    width = float(props.get("strokeWidth", 0))
    if width <= 0 or stroke == "transparent":
        shape.line.fill.background()
        return
    if is_hex_color(stroke):
        shape.line.color.rgb = rgb(stroke)
        shape.line.width = Pt(width)


def emu_x(element: dict[str, Any]) -> Any:
    return Inches(float(element.get("x", 0)) / 144)


def emu_y(element: dict[str, Any]) -> Any:
    return Inches(float(element.get("y", 0)) / 144)


def emu_width(element: dict[str, Any]) -> Any:
    return Inches(float(element.get("width", 1)) / 144)


def emu_height(element: dict[str, Any]) -> Any:
    return Inches(float(element.get("height", 1)) / 144)


def paragraph_alignment(value: str) -> Any:
    return {
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }.get(value, PP_ALIGN.LEFT)


def vertical_anchor(value: str) -> Any:
    return {
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(value, MSO_ANCHOR.TOP)


def is_bold(value: Any) -> bool:
    return value in {"bold", "semibold"} or (isinstance(value, int) and value >= 600)


def is_hex_color(value: Any) -> bool:
    return isinstance(value, str) and len(value) == 7 and value.startswith("#")


def rgb(value: str) -> RGBColor:
    return RGBColor(  # type: ignore[no-untyped-call]
        int(value[1:3], 16),
        int(value[3:5], 16),
        int(value[5:7], 16),
    )


def dedupe(values: list[str]) -> list[str]:
    result: list[str] = []
    seen: set[str] = set()
    for value in values:
        if value in seen:
            continue
        seen.add(value)
        result.append(value)
    return result
