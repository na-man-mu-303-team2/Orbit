from __future__ import annotations

import base64
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pydantic import BaseModel, ConfigDict, Field

from app.ai.pptx_quality import not_evaluated_slide_reports


CANVAS_WIDTH = 1920
CANVAS_HEIGHT = 1080


class ImportedDesignAsset(BaseModel):
    model_config = ConfigDict(populate_by_name=True)

    asset_id: str = Field(alias="assetId")
    file_name: str = Field(alias="fileName")
    mime_type: str = Field(alias="mimeType")
    content_base64: str = Field(alias="contentBase64")


class PptxDesignImportResult(BaseModel):
    model_config = ConfigDict(populate_by_name=True)

    blueprint: dict[str, Any]
    template_blueprint: dict[str, Any] = Field(alias="templateBlueprint")
    quality_report: dict[str, Any] = Field(alias="qualityReport")
    assets: list[ImportedDesignAsset] = Field(default_factory=list)
    warnings: list[str] = Field(default_factory=list)


class ImportedElementBlueprint(BaseModel):
    model_config = ConfigDict(populate_by_name=True)

    element_id: str = Field(alias="elementId")
    type: str
    role: str
    x: int
    y: int
    width: int
    height: int
    rotation: float = 0
    opacity: float = 1
    z_index: int = Field(alias="zIndex")
    locked: bool = False
    visible: bool = True
    props: dict[str, Any] = Field(default_factory=dict)


class ImportedSlideBlueprint(BaseModel):
    model_config = ConfigDict(populate_by_name=True)

    source_file_id: str = Field(default="", alias="sourceFileId")
    source_slide_index: int = Field(default=1, alias="sourceSlideIndex")
    style: dict[str, Any]
    elements: list[ImportedElementBlueprint]


class ImportedDesignBlueprint(BaseModel):
    model_config = ConfigDict(populate_by_name=True)

    source_file_id: str = Field(default="", alias="sourceFileId")
    canvas: dict[str, int] = Field(
        default_factory=lambda: {"width": CANVAS_WIDTH, "height": CANVAS_HEIGHT}
    )
    theme: dict[str, Any]
    slides: list[ImportedSlideBlueprint]
    warnings: list[str] = Field(default_factory=list)


@dataclass(frozen=True)
class ShapeTransform:
    scale_x: float = 1
    scale_y: float = 1
    translate_x: float = 0
    translate_y: float = 0

    def rect(
        self, x: int, y: int, width: int, height: int
    ) -> tuple[float, float, float, float]:
        return (
            self.scale_x * x + self.translate_x,
            self.scale_y * y + self.translate_y,
            self.scale_x * width,
            self.scale_y * height,
        )

    def for_group(self, group_shape: Any) -> ShapeTransform:
        off_x, off_y, ext_x, ext_y, child_x, child_y, child_width, child_height = (
            group_transform_values(group_shape)
        )
        ratio_x = ext_x / max(1, child_width)
        ratio_y = ext_y / max(1, child_height)
        local = ShapeTransform(
            scale_x=ratio_x,
            scale_y=ratio_y,
            translate_x=off_x - child_x * ratio_x,
            translate_y=off_y - child_y * ratio_y,
        )
        return ShapeTransform(
            scale_x=self.scale_x * local.scale_x,
            scale_y=self.scale_y * local.scale_y,
            translate_x=self.scale_x * local.translate_x + self.translate_x,
            translate_y=self.scale_y * local.translate_y + self.translate_y,
        )


def import_pptx_design(
    path: Path,
    file_id: str,
    *,
    canvas_width: int = CANVAS_WIDTH,
    canvas_height: int = CANVAS_HEIGHT,
) -> PptxDesignImportResult:
    presentation = Presentation(str(path))
    source_width_value = presentation.slide_width
    source_height_value = presentation.slide_height
    canvas_width = max(1, int(canvas_width))
    canvas_height = max(1, int(canvas_height))
    source_width = max(
        1,
        int(source_width_value) if source_width_value is not None else canvas_width,
    )
    source_height = max(
        1,
        int(source_height_value) if source_height_value is not None else canvas_height,
    )
    scale_x = canvas_width / source_width
    scale_y = canvas_height / source_height
    assets: list[ImportedDesignAsset] = []
    asset_colors: dict[str, str] = {}
    warnings: list[str] = []
    slides: list[dict[str, Any]] = []
    slot_sources_by_slide: list[dict[str, dict[str, Any]]] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        elements: list[dict[str, Any]] = []
        slot_sources: dict[str, dict[str, Any]] = {}
        z_cursor = [1]
        background_color = slide_background_color(slide) or "#ffffff"
        elements.append(
            shape_element(
                element_id=f"el_imported_{slide_index}_background",
                role="background",
                x=0,
                y=0,
                width=canvas_width,
                height=canvas_height,
                z_index=0,
                fill=background_color,
                stroke="transparent",
                locked=True,
            )
        )

        for source_name, source_shapes in inherited_decoration_sources(slide):
            append_shape_collection(
                source_shapes,
                slide_index=slide_index,
                path_prefix=source_name,
                elements=elements,
                assets=assets,
                asset_colors=asset_colors,
                warnings=warnings,
                slot_sources=slot_sources,
                scale_x=scale_x,
                scale_y=scale_y,
                canvas_width=canvas_width,
                canvas_height=canvas_height,
                z_cursor=z_cursor,
                transform=ShapeTransform(),
                decoration_only=True,
            )

        z_cursor[0] = max(z_cursor[0], 100)
        append_shape_collection(
            slide.shapes,
            slide_index=slide_index,
            path_prefix="slide",
            elements=elements,
            assets=assets,
            asset_colors=asset_colors,
            warnings=warnings,
            slot_sources=slot_sources,
            scale_x=scale_x,
            scale_y=scale_y,
            canvas_width=canvas_width,
            canvas_height=canvas_height,
            z_cursor=z_cursor,
            transform=ShapeTransform(),
            decoration_only=False,
        )

        assign_text_roles(elements)
        background_color = background_color_from_elements(
            elements,
            asset_colors,
            background_color,
        )
        apply_imported_background_color(elements, background_color)
        slides.append(
            {
                "sourceFileId": file_id,
                "sourceSlideIndex": slide_index,
                "style": imported_slide_style(elements, background_color),
                "elements": elements,
            }
        )
        slot_sources_by_slide.append(slot_sources)

    blueprint = ImportedDesignBlueprint.model_validate(
        {
            "sourceFileId": file_id,
            "canvas": {
                "width": canvas_width,
                "height": canvas_height,
            },
            "theme": imported_theme(slides),
            "slides": slides,
            "warnings": warnings,
        }
    )
    return PptxDesignImportResult(
        blueprint=blueprint.model_dump(by_alias=True),
        templateBlueprint=build_template_blueprint(
            file_id,
            slides,
            slot_sources_by_slide,
        ),
        qualityReport=build_quality_report(slides, warnings),
        assets=assets,
        warnings=warnings,
    )


def append_shape_collection(
    shapes: Any,
    *,
    slide_index: int,
    path_prefix: str,
    elements: list[dict[str, Any]],
    assets: list[ImportedDesignAsset],
    asset_colors: dict[str, str],
    warnings: list[str],
    slot_sources: dict[str, dict[str, Any]],
    scale_x: float,
    scale_y: float,
    canvas_width: int,
    canvas_height: int,
    z_cursor: list[int],
    transform: ShapeTransform,
    decoration_only: bool,
) -> None:
    for shape_index, shape in enumerate(shapes, start=1):
        append_shape_elements(
            shape,
            slide_index=slide_index,
            element_path=f"{path_prefix}_{shape_index}",
            elements=elements,
            assets=assets,
            asset_colors=asset_colors,
            warnings=warnings,
            slot_sources=slot_sources,
            scale_x=scale_x,
            scale_y=scale_y,
            canvas_width=canvas_width,
            canvas_height=canvas_height,
            z_cursor=z_cursor,
            transform=transform,
            decoration_only=decoration_only,
        )


def append_shape_elements(
    shape: Any,
    *,
    slide_index: int,
    element_path: str,
    elements: list[dict[str, Any]],
    assets: list[ImportedDesignAsset],
    asset_colors: dict[str, str],
    warnings: list[str],
    slot_sources: dict[str, dict[str, Any]],
    scale_x: float,
    scale_y: float,
    canvas_width: int,
    canvas_height: int,
    z_cursor: list[int],
    transform: ShapeTransform,
    decoration_only: bool,
) -> None:
    if decoration_only and bool(getattr(shape, "is_placeholder", False)):
        return

    shape_type = getattr(shape, "shape_type", None)
    if shape_type == MSO_SHAPE_TYPE.GROUP:
        append_shape_collection(
            getattr(shape, "shapes", []),
            slide_index=slide_index,
            path_prefix=element_path,
            elements=elements,
            assets=assets,
            asset_colors=asset_colors,
            warnings=warnings,
            slot_sources=slot_sources,
            scale_x=scale_x,
            scale_y=scale_y,
            canvas_width=canvas_width,
            canvas_height=canvas_height,
            z_cursor=z_cursor,
            transform=transform.for_group(shape),
            decoration_only=decoration_only,
        )
        return

    frame = normalized_frame(
        shape,
        scale_x,
        scale_y,
        canvas_width,
        canvas_height,
        transform,
    )
    element_id = f"el_imported_{slide_index}_{element_path}"
    locked = decoration_only
    role = "decoration" if decoration_only else "media"

    if shape_type == MSO_SHAPE_TYPE.PICTURE:
        asset_id = f"image_{len(assets) + 1}"
        asset = image_asset(shape, asset_id)
        assets.append(asset)
        color = average_image_color(shape.image.blob)
        if color:
            asset_colors[asset_id] = color
        element = {
            **element_base(
                element_id=f"{element_id}_image",
                role=role,
                frame=frame,
                z_index=next_z(z_cursor),
                locked=locked,
            ),
            "type": "image",
            "props": {
                "src": f"asset:{asset_id}",
                "alt": str(getattr(shape, "name", "Imported image")),
                "fit": "contain",
                "focusX": 0.5,
                "focusY": 0.5,
            },
        }
        elements.append(element)
        slot_sources[element["elementId"]] = shape_slot_source(
            shape,
            slide_index,
            element_path,
            decoration_only,
            fallback_type="image",
        )
        return

    blip_asset = blip_fill_asset(shape, f"image_{len(assets) + 1}")
    if blip_asset is not None:
        asset, color = blip_asset
        assets.append(asset)
        if color:
            asset_colors[asset.asset_id] = color
        image_role = (
            "background"
            if is_full_canvas_frame(frame, canvas_width, canvas_height)
            else role
        )
        element = {
            **element_base(
                element_id=f"{element_id}_image_fill",
                role=image_role,
                frame=frame,
                z_index=next_z(z_cursor),
                locked=locked or image_role == "background",
            ),
            "type": "image",
            "props": {
                "src": f"asset:{asset.asset_id}",
                "alt": str(getattr(shape, "name", "Imported image fill")),
                "fit": "stretch",
                "focusX": 0.5,
                "focusY": 0.5,
            },
        }
        elements.append(element)
        slot_sources[element["elementId"]] = shape_slot_source(
            shape,
            slide_index,
            element_path,
            decoration_only,
            fallback_type="image",
        )
        return

    if shape_type == MSO_SHAPE_TYPE.TABLE and not decoration_only:
        table_items = table_elements(
            shape,
            element_id,
            frame,
            z_cursor,
            canvas_width,
            canvas_height,
        )
        elements.extend(table_items)
        for element in table_items:
            slot_sources[element["elementId"]] = shape_slot_source(
                shape,
                slide_index,
                element_path,
                decoration_only,
                fallback_type="table",
            )
        return

    fill = shape_fill_color(shape)
    stroke = shape_line_color(shape)
    if shape_type == MSO_SHAPE_TYPE.FREEFORM:
        if not fill and not stroke:
            return
        custom_shape = freeform_element(
            shape,
            element_id=f"{element_id}_custom",
            frame=frame,
            z_index=next_z(z_cursor),
            fill=fill or "transparent",
            stroke=stroke or "transparent",
            locked=locked,
        )
        if custom_shape is not None:
            elements.append(custom_shape)
            slot_sources[custom_shape["elementId"]] = shape_slot_source(
                shape,
                slide_index,
                element_path,
                decoration_only,
                fallback_type="shape",
            )
        else:
            warnings.append(
                f"Unsupported PPTX freeform path on slide {slide_index}: {getattr(shape, 'name', 'freeform')}"
            )
            fallback = append_fallback_shape(
                elements, element_id, frame, z_cursor, fill, stroke, locked
            )
            slot_sources[fallback["elementId"]] = shape_slot_source(
                shape,
                slide_index,
                element_path,
                decoration_only,
                fallback_type="shape",
            )
    elif fill or stroke:
        element = append_pptx_shape(
            elements, shape, element_id, frame, z_cursor, fill, stroke, locked
        )
        slot_sources[element["elementId"]] = shape_slot_source(
            shape,
            slide_index,
            element_path,
            decoration_only,
            fallback_type="shape",
        )

    text = "" if decoration_only else shape_text(shape)
    if text:
        element = {
            **element_base(
                element_id=f"{element_id}_text",
                role="body",
                frame=frame,
                z_index=next_z(z_cursor),
                locked=False,
            ),
            "type": "text",
            "props": {
                "text": text,
                **shape_text_props(shape),
            },
        }
        elements.append(element)
        slot_sources[element["elementId"]] = shape_slot_source(
            shape,
            slide_index,
            element_path,
            decoration_only,
            fallback_type="shape",
        )
    elif is_unsupported_complex_shape(shape):
        warnings.append(f"Unsupported PPTX shape on slide {slide_index}: {shape_type}")


def append_fallback_shape(
    elements: list[dict[str, Any]],
    element_id: str,
    frame: dict[str, int],
    z_cursor: list[int],
    fill: str | None,
    stroke: str | None,
    locked: bool,
) -> dict[str, Any]:
    element = shape_element(
        element_id=f"{element_id}_shape",
        role="decoration",
        x=frame["x"],
        y=frame["y"],
        width=frame["width"],
        height=frame["height"],
        z_index=next_z(z_cursor),
        fill=fill or "transparent",
        stroke=stroke or "transparent",
        locked=locked,
    )
    elements.append(element)
    return element


def append_pptx_shape(
    elements: list[dict[str, Any]],
    shape: Any,
    element_id: str,
    frame: dict[str, int],
    z_cursor: list[int],
    fill: str | None,
    stroke: str | None,
    locked: bool,
) -> dict[str, Any]:
    element = pptx_shape_element(
        shape,
        element_id=f"{element_id}_shape",
        frame=frame,
        z_index=next_z(z_cursor),
        fill=fill or "transparent",
        stroke=stroke or "transparent",
        locked=locked,
    )
    elements.append(element)
    return element


def pptx_shape_element(
    shape: Any,
    *,
    element_id: str,
    frame: dict[str, int],
    z_index: int,
    fill: str,
    stroke: str,
    locked: bool,
) -> dict[str, Any]:
    token = pptx_shape_token(shape)
    if is_line_preset(shape, token):
        element_type = "arrow" if has_arrow_head(shape) else "line"
        return shape_element_of_type(
            element_type,
            element_id=element_id,
            role="decoration",
            frame=frame,
            z_index=z_index,
            fill=fill,
            stroke=stroke,
            locked=locked,
        )
    if "oval" in token:
        return shape_element_of_type(
            "ellipse",
            element_id=element_id,
            role="decoration",
            frame=frame,
            z_index=z_index,
            fill=fill,
            stroke=stroke,
            locked=locked,
        )
    if "donut" in token:
        return shape_element_of_type(
            "ring",
            element_id=element_id,
            role="decoration",
            frame=frame,
            z_index=z_index,
            fill=fill,
            stroke=stroke,
            locked=locked,
        )
    if "star" in token:
        return shape_element_of_type(
            "star",
            element_id=element_id,
            role="decoration",
            frame=frame,
            z_index=z_index,
            fill=fill,
            stroke=stroke,
            locked=locked,
        )

    custom_path = preset_custom_shape_path(token)
    if custom_path:
        path_data, closed = custom_path
        return {
            **element_base(
                element_id=element_id,
                role="decoration",
                frame=frame,
                z_index=z_index,
                locked=locked,
            ),
            "type": "customShape",
            "props": {
                "pathData": path_data,
                "viewBoxWidth": 100,
                "viewBoxHeight": 100,
                "fill": fill,
                "stroke": stroke,
                "strokeWidth": 1 if stroke != "transparent" else 0,
                "closed": closed,
                "nodes": [],
            },
        }

    element = shape_element(
        element_id=element_id,
        role="decoration",
        x=frame["x"],
        y=frame["y"],
        width=frame["width"],
        height=frame["height"],
        z_index=z_index,
        fill=fill,
        stroke=stroke,
        locked=locked,
    )
    if "round" in token:
        element["props"]["borderRadius"] = round(
            min(frame["width"], frame["height"]) * 0.16
        )
    return element


def shape_element_of_type(
    element_type: str,
    *,
    element_id: str,
    role: str,
    frame: dict[str, int],
    z_index: int,
    fill: str,
    stroke: str,
    locked: bool,
    sides: int | None = None,
) -> dict[str, Any]:
    props: dict[str, Any] = {
        "fill": fill,
        "stroke": stroke,
        "strokeWidth": 1 if stroke != "transparent" else 0,
        "borderRadius": 0,
    }
    if sides is not None:
        props["sides"] = sides
    return {
        **element_base(
            element_id=element_id,
            role=role,
            frame=frame,
            z_index=z_index,
            locked=locked,
        ),
        "type": element_type,
        "props": props,
    }


def pptx_shape_token(shape: Any) -> str:
    raw = getattr(shape, "auto_shape_type", None)
    return enum_token(raw)


def enum_token(value: Any) -> str:
    token = str(value or "").split("(", maxsplit=1)[0].strip().lower()
    return token.replace(" ", "_").replace("-", "_")


def is_line_preset(shape: Any, token: str) -> bool:
    return (
        getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.LINE
        or "connector" in token
    )


def has_arrow_head(shape: Any) -> bool:
    line = first_descendant(shape._element, "ln")
    return (
        first_child(line, "headEnd") is not None
        or first_child(line, "tailEnd") is not None
    )


def preset_custom_shape_path(token: str) -> tuple[str, bool] | None:
    token = normalized_preset_token(token)
    if "triangle" in token:
        if "right" in token:
            return "M 0 0 L 100 100 L 0 100 Z", True
        return "M 50 0 L 100 100 L 0 100 Z", True
    if "diamond" in token:
        return "M 50 0 L 100 50 L 50 100 L 0 50 Z", True
    if "parallelogram" in token:
        return "M 25 0 L 100 0 L 75 100 L 0 100 Z", True
    if "trapezoid" in token:
        return "M 25 0 L 75 0 L 100 100 L 0 100 Z", True
    if "hexagon" in token:
        return "M 25 0 L 75 0 L 100 50 L 75 100 L 25 100 L 0 50 Z", True
    if "pentagon" in token:
        return "M 50 0 L 100 38 L 81 100 L 19 100 L 0 38 Z", True
    if "chevron" in token:
        return "M 0 0 L 65 0 L 100 50 L 65 100 L 0 100 L 35 50 Z", True
    if "left_right_arrow" in token:
        return (
            "M 0 50 L 25 15 L 25 35 L 75 35 L 75 15 L 100 50 L 75 85 L 75 65 L 25 65 L 25 85 Z",
            True,
        )
    if "up_down_arrow" in token:
        return (
            "M 50 0 L 85 25 L 65 25 L 65 75 L 85 75 L 50 100 L 15 75 L 35 75 L 35 25 L 15 25 Z",
            True,
        )
    if "right_arrow" in token:
        return "M 0 25 L 65 25 L 65 0 L 100 50 L 65 100 L 65 75 L 0 75 Z", True
    if "left_arrow" in token:
        return "M 100 25 L 35 25 L 35 0 L 0 50 L 35 100 L 35 75 L 100 75 Z", True
    if "up_arrow" in token:
        return "M 25 100 L 25 35 L 0 35 L 50 0 L 100 35 L 75 35 L 75 100 Z", True
    if "down_arrow" in token:
        return "M 25 0 L 25 65 L 0 65 L 50 100 L 100 65 L 75 65 L 75 0 Z", True
    return None


def normalized_preset_token(token: str) -> str:
    chars: list[str] = []
    for index, char in enumerate(token):
        if char.isupper() and index > 0:
            chars.append("_")
        chars.append(char.lower() if char != "-" else "_")
    return "".join(chars)


def shape_slot_source(
    shape: Any,
    slide_index: int,
    element_path: str,
    decoration_only: bool,
    *,
    fallback_type: str,
) -> dict[str, Any]:
    if bool(getattr(shape, "is_placeholder", False)):
        source_type = "placeholder"
    elif element_path.startswith("master_"):
        source_type = "master"
    elif element_path.startswith("layout_"):
        source_type = "layout"
    elif fallback_type in {"table", "image"}:
        source_type = fallback_type
    elif decoration_only:
        source_type = "layout"
    else:
        source_type = "slide"

    source = {
        "type": source_type,
        "name": str(getattr(shape, "name", "") or "").strip() or "shape",
        "slidePart": f"ppt/slides/slide{slide_index}.xml",
        "shapeId": str(getattr(shape, "shape_id", "") or "0"),
        "writable": source_type not in {"master", "layout"} and not decoration_only,
    }
    placeholder = placeholder_type(shape)
    if placeholder:
        source["placeholderType"] = placeholder
    relationship_id = blip_relationship_id(shape)
    if relationship_id:
        source["relationshipId"] = relationship_id
    return source


def blip_relationship_id(shape: Any) -> str:
    element = getattr(shape, "_element", None)
    if element is None:
        return ""
    blip = first_descendant(element, "blip")
    return attr_by_local_name(blip, "embed") or ""


def placeholder_type(shape: Any) -> str:
    try:
        raw = str(shape.placeholder_format.type)
    except Exception:
        return ""
    normalized = (
        raw.lower().replace("(", " ").replace(")", " ").replace("_", "-").split()
    )
    return normalized[0] if normalized else ""


def build_template_blueprint(
    file_id: str,
    slides: list[dict[str, Any]],
    slot_sources_by_slide: list[dict[str, dict[str, Any]]],
) -> dict[str, Any]:
    repeated_texts = repeated_text_values(slides)
    return {
        "templateId": f"template_{safe_id(file_id)}",
        "sourceFileId": file_id,
        "sourcePackageFileId": file_id,
        "currentPackageFileId": file_id,
        "ooxmlSyncedDeckVersion": 1,
        "slides": [
            {
                "slideIndex": index + 1,
                "sourceSlideIndex": int(slide.get("sourceSlideIndex", index + 1)),
                "elementSources": [
                    source
                    for source in (
                        template_element_source_for_element(
                            element,
                            slot_sources_by_slide[index].get(
                                str(element.get("elementId", "")),
                                infer_element_source(element),
                            ),
                        )
                        for element in slide.get("elements", [])
                        if isinstance(element, dict)
                    )
                    if source is not None
                ],
                "slots": [
                    slot
                    for slot in (
                        template_slot_for_element(
                            element,
                            slot_sources_by_slide[index].get(
                                str(element.get("elementId", "")),
                                infer_element_source(element),
                            ),
                            repeated_texts,
                        )
                        for element in slide.get("elements", [])
                        if isinstance(element, dict)
                    )
                    if slot is not None
                ],
            }
            for index, slide in enumerate(slides)
        ],
    }


def template_element_source_for_element(
    element: dict[str, Any],
    source: dict[str, Any],
) -> dict[str, Any] | None:
    element_id = str(element.get("elementId", ""))
    slide_part = str(source.get("slidePart", ""))
    shape_id = str(source.get("shapeId", ""))
    if not element_id or not slide_part or not shape_id:
        return None

    source_type = str(source.get("type", "unknown"))
    if source_type not in {
        "placeholder",
        "slide",
        "layout",
        "master",
        "table",
        "image",
        "shape",
        "unknown",
    }:
        source_type = "unknown"

    element_source = {
        "elementId": element_id,
        "slidePart": slide_part,
        "shapeId": shape_id,
        "sourceType": source_type,
        "writable": bool(source.get("writable", False)),
    }
    relationship_id = str(source.get("relationshipId", ""))
    if relationship_id:
        element_source["relationshipId"] = relationship_id
    fallback_reason = str(source.get("fallbackReason", ""))
    if fallback_reason:
        element_source["fallbackReason"] = fallback_reason
    return element_source


def template_slot_for_element(
    element: dict[str, Any],
    source: dict[str, Any],
    repeated_texts: set[str],
) -> dict[str, Any] | None:
    element_id = str(element.get("elementId", ""))
    if not element_id:
        return None

    element_type = str(element.get("type", ""))
    role = slot_role_for_element(element)
    source_type = str(source.get("type", "unknown"))
    locked = bool(element.get("locked", False))

    if element_type == "text":
        text_key = normalized_text_key(element_text(element))
        if source_type == "placeholder":
            usage, replace_mode, confidence = "content-slot", "replace", 0.95
        elif source_type in {"master", "layout"} or text_key in repeated_texts:
            usage, replace_mode, confidence = "fixed-text", "preserve", 0.8
        else:
            usage, replace_mode, confidence = "fixed-text", "preserve", 0.45
    elif element_type == "image":
        if source_type == "placeholder":
            usage, replace_mode, confidence = "media-slot", "replace", 0.95
        elif str(element.get("role", "")) == "background" or locked:
            usage, replace_mode, confidence = "decoration", "ignore", 0.9
        else:
            usage, replace_mode, confidence = "media-slot", "preserve", 0.55
    else:
        usage, replace_mode, confidence = (
            "decoration",
            "ignore",
            0.85 if locked else 0.6,
        )

    return {
        "elementId": element_id,
        "usage": usage,
        "slotRole": role,
        "replaceMode": replace_mode,
        "confidence": confidence,
        "bounds": {
            "x": number_or_zero(element.get("x")),
            "y": number_or_zero(element.get("y")),
            "width": max(1, number_or_zero(element.get("width"))),
            "height": max(1, number_or_zero(element.get("height"))),
        },
        "source": source,
    }


def slot_role_for_element(element: dict[str, Any]) -> str:
    role = str(element.get("role", "unknown"))
    if role in {"title", "subtitle", "body", "caption", "background"}:
        return role
    if element.get("type") in {"image", "svg"}:
        return "image"
    if element.get("type") == "chart":
        return "chart"
    if "_cell_" in str(element.get("elementId", "")):
        return "table"
    return "unknown"


def infer_element_source(element: dict[str, Any]) -> dict[str, Any]:
    element_id = str(element.get("elementId", ""))
    if "_master_" in element_id:
        return {"type": "master"}
    if "_layout_" in element_id:
        return {"type": "layout"}
    if element.get("type") in {"image", "svg"}:
        return {"type": "image"}
    if "_cell_" in element_id:
        return {"type": "table"}
    return {"type": "slide"}


QUALITY_WEIGHTS = {
    "geometry": 25,
    "text": 15,
    "color": 10,
    "layer": 10,
    "editability": 10,
    "pixelSimilarity": 30,
}


def build_quality_report(
    slides: list[dict[str, Any]],
    warnings: list[str],
) -> dict[str, Any]:
    editability_coverage = editable_element_coverage(slides)
    metrics: dict[str, float | None] = {
        "geometry": geometry_score(slides),
        "text": text_score(slides, warnings),
        "color": color_score(slides),
        "layer": layer_score(slides),
        "editability": round(editability_coverage * 100),
        "pixelSimilarity": None,
    }
    composite = weighted_quality_score(metrics)
    applied_cap: int | None = None
    notes = ["pixel renderer unavailable"]

    if editability_coverage < 0.2:
        applied_cap = 50
    elif editability_coverage < 0.5:
        applied_cap = 70

    if applied_cap is not None:
        composite = min(composite, applied_cap)
        notes.append(f"editability coverage cap {applied_cap}")

    return {
        "compositeScore": composite,
        "metrics": metrics,
        "weights": QUALITY_WEIGHTS,
        "editabilityCoverage": editability_coverage,
        "appliedCap": applied_cap,
        "slideReports": not_evaluated_slide_reports(
            len(slides),
            "candidate renderer unavailable",
        ),
        "notes": notes,
    }


def weighted_quality_score(metrics: dict[str, float | None]) -> int:
    weighted_total = 0.0
    weight_total = 0
    for key, weight in QUALITY_WEIGHTS.items():
        value = metrics.get(key)
        if value is None:
            continue
        weighted_total += value * weight
        weight_total += weight
    return round(weighted_total / max(1, weight_total))


def editable_element_coverage(slides: list[dict[str, Any]]) -> float:
    elements = [
        element
        for slide in slides
        for element in slide.get("elements", [])
        if isinstance(element, dict)
        and element.get("visible", True)
        and element.get("role") != "background"
    ]
    if not elements:
        return 0
    editable = [
        element
        for element in elements
        if not bool(element.get("locked", False))
        and element.get("type")
        in {
            "arrow",
            "chart",
            "customShape",
            "ellipse",
            "group",
            "image",
            "line",
            "polygon",
            "rect",
            "ring",
            "star",
            "svg",
            "table",
            "text",
        }
    ]
    return round(len(editable) / len(elements), 3)


def geometry_score(slides: list[dict[str, Any]]) -> int:
    elements = [
        element
        for slide in slides
        for element in slide.get("elements", [])
        if isinstance(element, dict)
    ]
    if not elements:
        return 50
    invalid = [
        element
        for element in elements
        if number_or_zero(element.get("width")) <= 0
        or number_or_zero(element.get("height")) <= 0
    ]
    return max(0, 100 - len(invalid) * 20)


def text_score(slides: list[dict[str, Any]], warnings: list[str]) -> int:
    text_count = sum(
        1
        for slide in slides
        for element in slide.get("elements", [])
        if isinstance(element, dict) and element.get("type") == "text"
    )
    warning_penalty = min(40, len(warnings) * 5)
    return max(0, (95 if text_count else 75) - warning_penalty)


def color_score(slides: list[dict[str, Any]]) -> int:
    colors = imported_visible_colors(slides)
    return 90 if colors else 70


def layer_score(slides: list[dict[str, Any]]) -> int:
    score = 100
    for slide in slides:
        z_indexes = [
            element.get("zIndex")
            for element in slide.get("elements", [])
            if isinstance(element, dict)
        ]
        if len(z_indexes) != len(set(z_indexes)):
            score -= 15
    return max(0, score)


def repeated_text_values(slides: list[dict[str, Any]]) -> set[str]:
    counts: dict[str, int] = {}
    for slide in slides:
        for element in slide.get("elements", []):
            if not isinstance(element, dict) or element.get("type") != "text":
                continue
            key = normalized_text_key(element_text(element))
            if key:
                counts[key] = counts.get(key, 0) + 1
    return {key for key, count in counts.items() if count > 1}


def normalized_text_key(value: Any) -> str:
    return " ".join(str(value or "").lower().split())


def element_text(element: dict[str, Any]) -> str:
    props = element.get("props", {})
    if not isinstance(props, dict):
        return ""
    return str(props.get("text", ""))


def number_or_zero(value: Any) -> float:
    try:
        return float(value)
    except Exception:
        return 0


def safe_id(value: str) -> str:
    return (
        "".join(
            char if char.isascii() and (char.isalnum() or char in "_-") else "_"
            for char in value
        )
        or "pptx"
    )


def next_z(z_cursor: list[int]) -> int:
    z_index = z_cursor[0]
    z_cursor[0] += 1
    return z_index


def normalized_frame(
    shape: Any,
    scale_x: float,
    scale_y: float,
    canvas_width: int,
    canvas_height: int,
    transform: ShapeTransform | None = None,
) -> dict[str, int]:
    transformer = transform or ShapeTransform()
    raw_x, raw_y, raw_width, raw_height = transformer.rect(
        int(getattr(shape, "left", 0)),
        int(getattr(shape, "top", 0)),
        int(getattr(shape, "width", 1)),
        int(getattr(shape, "height", 1)),
    )
    x = max(0, round(raw_x * scale_x))
    y = max(0, round(raw_y * scale_y))
    width = max(1, round(raw_width * scale_x))
    height = max(1, round(raw_height * scale_y))
    return {
        "x": min(x, canvas_width - 1),
        "y": min(y, canvas_height - 1),
        "width": min(width, canvas_width - min(x, canvas_width - 1)),
        "height": min(height, canvas_height - min(y, canvas_height - 1)),
    }


def group_transform_values(
    group_shape: Any,
) -> tuple[int, int, int, int, int, int, int, int]:
    fallback_x = int(getattr(group_shape, "left", 0))
    fallback_y = int(getattr(group_shape, "top", 0))
    fallback_width = max(1, int(getattr(group_shape, "width", 1)))
    fallback_height = max(1, int(getattr(group_shape, "height", 1)))
    xfrm = first_descendant(group_shape._element, "xfrm")
    if xfrm is None:
        return (
            fallback_x,
            fallback_y,
            fallback_width,
            fallback_height,
            fallback_x,
            fallback_y,
            fallback_width,
            fallback_height,
        )

    off = first_child(xfrm, "off")
    ext = first_child(xfrm, "ext")
    child_off = first_child(xfrm, "chOff")
    child_ext = first_child(xfrm, "chExt")
    return (
        int_attr(off, "x", fallback_x),
        int_attr(off, "y", fallback_y),
        int_attr(ext, "cx", fallback_width),
        int_attr(ext, "cy", fallback_height),
        int_attr(child_off, "x", fallback_x),
        int_attr(child_off, "y", fallback_y),
        int_attr(child_ext, "cx", fallback_width),
        int_attr(child_ext, "cy", fallback_height),
    )


def element_base(
    *,
    element_id: str,
    role: str,
    frame: dict[str, int],
    z_index: int,
    locked: bool,
) -> dict[str, Any]:
    return {
        "elementId": element_id,
        "role": role,
        "x": frame["x"],
        "y": frame["y"],
        "width": frame["width"],
        "height": frame["height"],
        "rotation": 0,
        "opacity": 1,
        "zIndex": z_index,
        "locked": locked,
        "visible": True,
    }


def shape_element(
    *,
    element_id: str,
    role: str,
    x: int,
    y: int,
    width: int,
    height: int,
    z_index: int,
    fill: str,
    stroke: str,
    locked: bool,
) -> dict[str, Any]:
    return {
        **element_base(
            element_id=element_id,
            role=role,
            frame={"x": x, "y": y, "width": width, "height": height},
            z_index=z_index,
            locked=locked,
        ),
        "type": "rect",
        "props": {
            "fill": fill,
            "stroke": stroke,
            "strokeWidth": 1 if stroke != "transparent" else 0,
            "borderRadius": 0,
        },
    }


def image_asset(shape: Any, asset_id: str) -> ImportedDesignAsset:
    image = shape.image
    extension = str(getattr(image, "ext", "png") or "png")
    mime_type = str(getattr(image, "content_type", f"image/{extension}"))
    return image_asset_from_blob(asset_id, image.blob, mime_type)


def image_asset_from_blob(
    asset_id: str,
    blob: bytes,
    mime_type: str,
) -> ImportedDesignAsset:
    extension = extension_for_mime_type(mime_type)
    return ImportedDesignAsset(
        assetId=asset_id,
        fileName=f"{asset_id}.{extension}",
        mimeType=mime_type,
        contentBase64=base64.b64encode(blob).decode("ascii"),
    )


def blip_fill_asset(
    shape: Any,
    asset_id: str,
) -> tuple[ImportedDesignAsset, str | None] | None:
    blip = first_descendant(shape._element, "blip")
    relationship_id = attr_by_local_name(blip, "embed") if blip is not None else None
    if not relationship_id:
        return None
    try:
        image_part = shape.part.related_part(relationship_id)
        blob = bytes(image_part.blob)
        mime_type = str(getattr(image_part, "content_type", "image/png"))
    except Exception:
        return None
    return image_asset_from_blob(asset_id, blob, mime_type), average_image_color(blob)


def extension_for_mime_type(mime_type: str) -> str:
    subtype = mime_type.rsplit("/", maxsplit=1)[-1].lower()
    if subtype == "jpeg":
        return "jpg"
    if subtype in {"svg", "svg+xml"}:
        return "svg"
    if subtype in {"png", "jpg", "gif", "webp"}:
        return subtype
    return "png"


def average_image_color(blob: bytes) -> str | None:
    try:
        from PIL import Image

        with Image.open(BytesIO(blob)) as image:
            image.thumbnail((24, 24))
            pixels = list(image.convert("RGB").getdata())
    except Exception:
        return None
    if not pixels:
        return None
    red = round(sum(pixel[0] for pixel in pixels) / len(pixels))
    green = round(sum(pixel[1] for pixel in pixels) / len(pixels))
    blue = round(sum(pixel[2] for pixel in pixels) / len(pixels))
    return f"#{red:02X}{green:02X}{blue:02X}"


def is_full_canvas_frame(
    frame: dict[str, int],
    canvas_width: int,
    canvas_height: int,
) -> bool:
    return (
        frame["x"] <= 4
        and frame["y"] <= 4
        and frame["width"] >= canvas_width - 8
        and frame["height"] >= canvas_height - 8
    )


def background_color_from_elements(
    elements: list[dict[str, Any]],
    asset_colors: dict[str, str],
    fallback: str,
) -> str:
    for element in elements:
        if element.get("type") != "image" or element.get("role") != "background":
            continue
        asset_id = asset_id_from_src(str(element.get("props", {}).get("src", "")))
        color = asset_colors.get(asset_id)
        if color:
            return color
    return fallback


def apply_imported_background_color(
    elements: list[dict[str, Any]],
    background_color: str,
) -> None:
    for element in elements:
        if element.get("role") == "background" and element.get("type") == "rect":
            props = element.get("props", {})
            if isinstance(props, dict):
                props["fill"] = background_color
            return


def asset_id_from_src(src: str) -> str:
    return src.removeprefix("asset:")


def table_elements(
    shape: Any,
    element_id: str,
    frame: dict[str, int],
    z_cursor: list[int],
    canvas_width: int,
    canvas_height: int,
) -> list[dict[str, Any]]:
    table = shape.table
    column_widths = [max(1, int(column.width)) for column in table.columns]
    row_heights = [max(1, int(row.height)) for row in table.rows]
    total_width = max(1, sum(column_widths))
    total_height = max(1, sum(row_heights))
    elements: list[dict[str, Any]] = []

    y = frame["y"]
    for row_index, row in enumerate(table.rows):
        row_height = max(
            1, round(frame["height"] * row_heights[row_index] / total_height)
        )
        x = frame["x"]
        for column_index, cell in enumerate(row.cells):
            column_width = max(
                1,
                round(frame["width"] * column_widths[column_index] / total_width),
            )
            if not bool(getattr(cell, "is_spanned", False)):
                cell_frame = {
                    "x": min(x, canvas_width - 1),
                    "y": min(y, canvas_height - 1),
                    "width": min(
                        column_width,
                        canvas_width - min(x, canvas_width - 1),
                    ),
                    "height": min(
                        row_height,
                        canvas_height - min(y, canvas_height - 1),
                    ),
                }
                fill = cell_fill_color(cell) or "#ffffff"
                elements.append(
                    shape_element(
                        element_id=f"{element_id}_cell_{row_index + 1}_{column_index + 1}",
                        role="decoration",
                        x=cell_frame["x"],
                        y=cell_frame["y"],
                        width=cell_frame["width"],
                        height=cell_frame["height"],
                        z_index=next_z(z_cursor),
                        fill=fill,
                        stroke="#d1d5db",
                        locked=True,
                    )
                )
                text = text_frame_text(cell.text_frame)
                if text:
                    elements.append(
                        {
                            **element_base(
                                element_id=f"{element_id}_cell_{row_index + 1}_{column_index + 1}_text",
                                role="body",
                                frame={
                                    "x": cell_frame["x"] + 8,
                                    "y": cell_frame["y"] + 6,
                                    "width": max(1, cell_frame["width"] - 16),
                                    "height": max(1, cell_frame["height"] - 12),
                                },
                                z_index=next_z(z_cursor),
                                locked=False,
                            ),
                            "type": "text",
                            "props": {
                                "text": text,
                                **text_frame_text_props(cell.text_frame),
                            },
                        }
                    )
            x += column_width
        y += row_height
    return elements


def freeform_element(
    shape: Any,
    *,
    element_id: str,
    frame: dict[str, int],
    z_index: int,
    fill: str,
    stroke: str,
    locked: bool,
) -> dict[str, Any] | None:
    path = custom_geometry_path(shape)
    if path is None:
        return None
    path_data, view_box_width, view_box_height, nodes = path
    return {
        **element_base(
            element_id=element_id,
            role="decoration",
            frame=frame,
            z_index=z_index,
            locked=locked,
        ),
        "type": "customShape",
        "props": {
            "pathData": path_data,
            "viewBoxWidth": view_box_width,
            "viewBoxHeight": view_box_height,
            "fill": fill,
            "stroke": stroke,
            "strokeWidth": 1 if stroke != "transparent" else 0,
            "closed": path_data.rstrip().endswith("Z"),
            "nodes": nodes,
        },
    }


def custom_geometry_path(
    shape: Any,
) -> tuple[str, int, int, list[dict[str, Any]]] | None:
    custom_geometry = first_descendant(shape._element, "custGeom")
    if custom_geometry is None:
        return None
    path_list = first_descendant(custom_geometry, "pathLst")
    if path_list is None:
        return None

    segments: list[str] = []
    nodes: list[dict[str, Any]] = []
    view_box_width = 0
    view_box_height = 0
    uses_curves = False
    for path in direct_children(path_list, "path"):
        view_box_width = max(view_box_width, int_attr(path, "w", 0))
        view_box_height = max(view_box_height, int_attr(path, "h", 0))
        for command in list(path):
            name = local_name(command)
            points = [point_xy(point) for point in direct_children(command, "pt")]
            if name == "moveTo" and len(points) == 1:
                x, y = points[0]
                segments.append(f"M {x} {y}")
                if not uses_curves:
                    nodes.append({"x": x, "y": y, "mode": "corner"})
            elif name == "lnTo" and len(points) == 1:
                x, y = points[0]
                segments.append(f"L {x} {y}")
                if not uses_curves:
                    nodes.append({"x": x, "y": y, "mode": "corner"})
            elif name == "quadBezTo" and len(points) == 2:
                uses_curves = True
                (x1, y1), (x, y) = points
                segments.append(f"Q {x1} {y1} {x} {y}")
            elif name == "cubicBezTo" and len(points) == 3:
                uses_curves = True
                (x1, y1), (x2, y2), (x, y) = points
                segments.append(f"C {x1} {y1} {x2} {y2} {x} {y}")
            elif name == "close":
                segments.append("Z")
            else:
                return None

    path_data = " ".join(segments).strip()
    if not path_data:
        return None
    if uses_curves:
        nodes = []
    return (
        path_data,
        max(1, view_box_width or int(getattr(shape, "width", 1))),
        max(1, view_box_height or int(getattr(shape, "height", 1))),
        nodes,
    )


def inherited_decoration_sources(slide: Any) -> list[tuple[str, Any]]:
    sources: list[tuple[str, Any]] = []
    layout = getattr(slide, "slide_layout", None)
    master = getattr(layout, "slide_master", None) if layout is not None else None
    if master is not None:
        sources.append(("master", master.shapes))
    if layout is not None:
        sources.append(("layout", layout.shapes))
    return sources


def slide_background_color(slide: Any) -> str | None:
    try:
        return color_to_hex(slide.background.fill.fore_color)
    except Exception:
        return None


def shape_fill_color(shape: Any) -> str | None:
    try:
        return color_to_hex(shape.fill.fore_color)
    except Exception:
        return None


def shape_line_color(shape: Any) -> str | None:
    try:
        return color_to_hex(shape.line.color)
    except Exception:
        return None


def cell_fill_color(cell: Any) -> str | None:
    try:
        return color_to_hex(cell.fill.fore_color)
    except Exception:
        return None


def color_to_hex(color: Any) -> str | None:
    try:
        rgb = color.rgb
    except Exception:
        return None
    if rgb is None:
        return None
    return f"#{rgb}"


def shape_text(shape: Any) -> str:
    if not bool(getattr(shape, "has_text_frame", False)):
        return ""
    return text_frame_text(shape.text_frame)


def text_frame_text(text_frame: Any) -> str:
    return "\n".join(
        paragraph.text.strip()
        for paragraph in text_frame.paragraphs
        if paragraph.text.strip()
    )


def shape_text_props(shape: Any) -> dict[str, Any]:
    font = first_font(
        shape.text_frame if bool(getattr(shape, "has_text_frame", False)) else None
    )
    return text_props_for_font(font)


def text_frame_text_props(text_frame: Any) -> dict[str, Any]:
    return text_props_for_font(first_font(text_frame))


def text_props_for_font(font: Any) -> dict[str, Any]:
    return {
        "fontFamily": str(getattr(font, "name", None) or "Inter"),
        "fontSize": font_size(font),
        "fontWeight": "bold" if bool(getattr(font, "bold", False)) else "normal",
        "color": font_color(font) or "#111827",
        "align": "left",
        "verticalAlign": "top",
        "lineHeight": 1.15,
    }


def first_font(text_frame: Any) -> Any:
    if text_frame is None:
        return None
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            return run.font
        return paragraph.font
    return None


def font_size(font: Any) -> int:
    size = getattr(font, "size", None)
    if size is None:
        return 24
    try:
        return max(8, round(float(size.pt)))
    except Exception:
        return 24


def font_color(font: Any) -> str | None:
    try:
        return color_to_hex(font.color)
    except Exception:
        return None


def assign_text_roles(elements: list[dict[str, Any]]) -> None:
    text_elements = [element for element in elements if element.get("type") == "text"]
    if not text_elements:
        return

    title = max(
        text_elements,
        key=lambda element: (
            int(element.get("props", {}).get("fontSize", 0)),
            -int(element.get("y", 0)),
        ),
    )
    title["role"] = "title"
    for element in text_elements:
        if element is not title:
            element["role"] = "body"


def imported_theme(slides: list[dict[str, Any]]) -> dict[str, Any]:
    first_slide = slides[0] if slides else {}
    background = str(
        first_slide.get("style", {}).get("backgroundColor", "#ffffff")
        if isinstance(first_slide.get("style"), dict)
        else "#ffffff"
    )
    colors = imported_visible_colors(slides)
    text_color = imported_text_color(slides) or text_color_for_background(background)
    accent = (
        first_non_background_color(
            colors,
            background,
            skip={text_color},
        )
        or text_color
    )
    font_family = imported_font_family(slides)
    muted = "#f3f4f6" if is_light_color(background) else "#1f3a2e"
    border = "#d1d5db" if is_light_color(background) else "#7fa38c"
    return {
        "name": "Imported PPTX",
        "fontFamily": font_family,
        "backgroundColor": background,
        "textColor": text_color,
        "accentColor": accent,
        "palette": {
            "primary": accent,
            "secondary": first_non_background_color(colors, background, skip={accent})
            or accent,
            "surface": background,
            "muted": muted,
            "border": border,
        },
        "typography": {
            "headingFontFamily": font_family,
            "bodyFontFamily": font_family,
            "titleSize": 56,
            "headingSize": 40,
            "bodySize": 24,
            "captionSize": 16,
        },
        "effects": {"borderRadius": 8},
    }


def imported_slide_style(
    elements: list[dict[str, Any]],
    background_color: str,
) -> dict[str, Any]:
    slide = {
        "style": {"backgroundColor": background_color},
        "elements": elements,
    }
    text_color = imported_text_color([slide]) or text_color_for_background(
        background_color
    )
    colors = imported_visible_colors([slide])
    accent = (
        first_non_background_color(
            colors,
            background_color,
            skip={text_color},
        )
        or text_color
    )
    return {
        "layout": "title-content",
        "backgroundColor": background_color,
        "textColor": text_color,
        "accentColor": accent,
        "fontFamily": imported_font_family([slide]),
    }


def imported_font_family(slides: list[dict[str, Any]]) -> str:
    counts: dict[str, int] = {}
    first_seen: dict[str, int] = {}
    for slide in slides:
        for element in slide.get("elements", []):
            if element.get("type") != "text":
                continue
            family = str(element.get("props", {}).get("fontFamily", "")).strip()
            if not family:
                continue
            first_seen.setdefault(family, len(first_seen))
            counts[family] = counts.get(family, 0) + 1
    if not counts:
        return "Inter"
    return min(counts, key=lambda family: (-counts[family], first_seen[family]))


def imported_visible_colors(slides: list[dict[str, Any]]) -> list[str]:
    colors: list[str] = []
    for slide in slides:
        style = slide.get("style", {})
        if isinstance(style, dict) and is_hex_color(style.get("backgroundColor")):
            colors.append(str(style["backgroundColor"]))
        for element in slide.get("elements", []):
            props = element.get("props", {})
            if not isinstance(props, dict):
                continue
            for key in ("fill", "stroke", "color"):
                value = props.get(key)
                if is_hex_color(value) and value != "transparent":
                    colors.append(str(value))
    return colors


def imported_text_color(slides: list[dict[str, Any]]) -> str | None:
    text_colors: list[str] = []
    for slide in slides:
        for element in slide.get("elements", []):
            if element.get("type") != "text":
                continue
            color = element.get("props", {}).get("color")
            if is_hex_color(color):
                text_colors.append(str(color))
    return text_colors[0] if text_colors else None


def first_non_background_color(
    colors: list[str],
    background: str,
    *,
    skip: set[str] | None = None,
) -> str | None:
    blocked = {background, "transparent", *(skip or set())}
    for color in colors:
        if color not in blocked:
            return color
    return None


def text_color_for_background(color: str) -> str:
    if not is_hex_color(color):
        return "#111827"
    return "#111827" if is_light_color(color) else "#f8fafc"


def is_light_color(color: str) -> bool:
    if not is_hex_color(color):
        return True
    red = int(color[1:3], 16)
    green = int(color[3:5], 16)
    blue = int(color[5:7], 16)
    return (0.2126 * red + 0.7152 * green + 0.0722 * blue) > 150


def is_hex_color(value: Any) -> bool:
    if not isinstance(value, str) or len(value) != 7 or not value.startswith("#"):
        return False
    try:
        int(value[1:], 16)
    except ValueError:
        return False
    return True


def is_unsupported_complex_shape(shape: Any) -> bool:
    shape_type = getattr(shape, "shape_type", None)
    unsupported = {MSO_SHAPE_TYPE.CHART}
    diagram = getattr(MSO_SHAPE_TYPE, "DIAGRAM", None)
    if diagram is not None:
        unsupported.add(diagram)
    return shape_type in unsupported


def first_descendant(element: Any, name: str) -> Any | None:
    for candidate in element.iter():
        if local_name(candidate) == name:
            return candidate
    return None


def first_child(element: Any, name: str) -> Any | None:
    for candidate in list(element):
        if local_name(candidate) == name:
            return candidate
    return None


def direct_children(element: Any, name: str) -> list[Any]:
    return [candidate for candidate in list(element) if local_name(candidate) == name]


def local_name(element: Any) -> str:
    tag = getattr(element, "tag", element)
    return str(tag).rsplit("}", maxsplit=1)[-1]


def int_attr(element: Any | None, name: str, fallback: int) -> int:
    if element is None:
        return fallback
    try:
        return int(element.get(name))
    except Exception:
        return fallback


def attr_by_local_name(element: Any | None, name: str) -> str | None:
    if element is None:
        return None
    for key, value in element.attrib.items():
        if local_name(key) == name:
            return str(value)
    return None


def point_xy(point: Any) -> tuple[int, int]:
    return int_attr(point, "x", 0), int_attr(point, "y", 0)
