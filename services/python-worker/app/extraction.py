from __future__ import annotations

import importlib.util
import json
import shutil
import subprocess
import tempfile
from collections.abc import Iterator
from contextlib import contextmanager
from dataclasses import dataclass, field
from enum import Enum
from io import BytesIO
from pathlib import Path
from typing import Any


class FileKind(str, Enum):
    IMAGE = "image"
    PDF = "pdf"
    DOCX = "docx"
    PPTX = "pptx"
    UNSUPPORTED = "unsupported"


class ResultStatus(str, Enum):
    SUCCEEDED = "succeeded"
    SKIPPED = "skipped"
    FAILED = "failed"


class PageStatus(str, Enum):
    TEXT_PAGE = "text_page"
    OCR_NEEDED_PAGE = "ocr_needed_page"
    BLANK_PAGE = "blank_page"
    SUSPICIOUS_PAGE = "suspicious_page"


class SlideStatus(str, Enum):
    TEXT_SLIDE = "text_slide"
    MIXED_SLIDE = "mixed_slide"
    OCR_NEEDED_SLIDE = "ocr_needed_slide"
    BLANK_SLIDE = "blank_slide"


@dataclass(frozen=True)
class ExtractedSection:
    title: str
    text: str = ""
    status: str = "text"
    index: int | None = None
    notes: list[str] = field(default_factory=list)
    metadata: dict[str, str | int | float] = field(default_factory=dict)


@dataclass(frozen=True)
class ExtractionResult:
    source_path: Path
    kind: FileKind
    status: ResultStatus
    sections: list[ExtractedSection] = field(default_factory=list)
    message: str = ""


@dataclass(frozen=True)
class ExtractConfig:
    ocr_lang: str = "kor+eng"
    pdf_render_scale: float = 2.0
    pdf_text_page_min_chars: int = 40
    pdf_ocr_max_chars: int = 20
    pdf_image_ratio_threshold: float = 0.35
    blank_image_ratio_threshold: float = 0.01
    docx_min_text_chars: int = 80
    pptx_text_slide_min_chars: int = 30
    pptx_image_ratio_threshold: float = 0.35
    tesseract_cmd: str | None = None
    soffice_path: Path | None = None


@dataclass(frozen=True)
class AICleanupResult:
    text: str
    status: str
    message: str = ""


@dataclass(frozen=True)
class PresentationKeyword:
    keyword: str
    reason: str
    priority: str


@dataclass(frozen=True)
class KeywordExtractionResult:
    keywords: list[PresentationKeyword]
    status: str
    message: str = ""


@dataclass(frozen=True)
class SlideText:
    index: int
    text: str
    image_area_ratio: float
    status: SlideStatus


class DependencyError(RuntimeError):
    pass


IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff", ".gif"}
DEFAULT_OPENAI_MODEL = "gpt-4o-mini"

CLEANUP_INSTRUCTIONS = """
You clean OCR text for a Korean presentation-generation assistant.
Return only the cleaned reference text, not commentary.

Rules:
- Preserve meaningful Korean and English content.
- Do not summarize, shorten, or collapse the material into only high-level points.
- Keep all recoverable facts, headings, lists, and slide-level topics in source order.
- Remove OCR noise, repeated UI artifacts, stray icons, and duplicated fragments.
- Keep useful document titles, names, table/list relationships, and section order.
- Correct obvious OCR mistakes when context is clear, such as AL -> AI.
- Do not invent facts that are not supported by the OCR text.
""".strip()

KEYWORD_INSTRUCTIONS = """
You extract presentation keywords from cleaned Korean reference material.
Return only a JSON object with a "keywords" array. Each item must have:
- keyword: short Korean or English phrase
- reason: why this matters for the presentation
- priority: one of "high", "medium", "low"

Choose keywords that help structure PPT slides and speaker notes.
Do not include duplicate or vague keywords.
Do not invent facts that are not supported by the cleaned text.
""".strip()

KEYWORD_RESPONSE_FORMAT: dict[str, Any] = {
    "format": {
        "type": "json_schema",
        "name": "presentation_keywords",
        "strict": True,
        "schema": {
            "type": "object",
            "additionalProperties": False,
            "properties": {
                "keywords": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "additionalProperties": False,
                        "properties": {
                            "keyword": {"type": "string"},
                            "reason": {"type": "string"},
                            "priority": {
                                "type": "string",
                                "enum": ["high", "medium", "low"],
                            },
                        },
                        "required": ["keyword", "reason", "priority"],
                    },
                },
            },
            "required": ["keywords"],
        },
    }
}


def extract_file(source_path: Path, config: ExtractConfig | None = None) -> ExtractionResult:
    selected_config = config or ExtractConfig()
    source_path = Path(source_path)
    kind = detect_file_kind(source_path)

    if kind == FileKind.UNSUPPORTED:
        return ExtractionResult(
            source_path=source_path,
            kind=kind,
            status=ResultStatus.SKIPPED,
            sections=[
                ExtractedSection(
                    title="Unsupported file",
                    status="skipped",
                    notes=["Supported files are images, PDF, DOCX, and PPTX."],
                )
            ],
            message=f"Unsupported extension: {source_path.suffix.lower()}",
        )

    try:
        if kind == FileKind.IMAGE:
            return extract_image(source_path, selected_config)
        if kind == FileKind.PDF:
            return extract_pdf(source_path, selected_config)
        if kind == FileKind.DOCX:
            return extract_docx(source_path, selected_config)
        if kind == FileKind.PPTX:
            return extract_pptx(source_path, selected_config)
    except DependencyError as error:
        return failed_extraction(source_path, kind, "Dependency failure", str(error))
    except Exception as error:
        return failed_extraction(source_path, kind, "Extraction failure", str(error))

    raise ValueError(f"Unhandled file kind: {kind}")


def detect_file_kind(path: Path) -> FileKind:
    suffix = path.suffix.lower()
    if suffix in IMAGE_EXTENSIONS:
        return FileKind.IMAGE
    if suffix == ".pdf":
        return FileKind.PDF
    if suffix == ".docx":
        return FileKind.DOCX
    if suffix == ".pptx":
        return FileKind.PPTX
    return FileKind.UNSUPPORTED


def failed_extraction(
    source_path: Path,
    kind: FileKind,
    title: str,
    message: str,
) -> ExtractionResult:
    return ExtractionResult(
        source_path=source_path,
        kind=kind,
        status=ResultStatus.FAILED,
        sections=[ExtractedSection(title=title, status="failed", notes=[message])],
        message=message,
    )


def require_python_package(import_name: str, package_name: str | None = None) -> None:
    if importlib.util.find_spec(import_name) is None:
        display_name = package_name or import_name
        raise DependencyError(f"Missing Python package '{display_name}'.")


def find_executable(candidates: list[str], configured_path: Path | str | None = None) -> str:
    if configured_path:
        path = Path(configured_path)
        if path.exists():
            return str(path)
        raise DependencyError(f"Configured executable was not found: {path}.")

    for candidate in candidates:
        resolved = shutil.which(candidate)
        if resolved:
            return resolved

    raise DependencyError(f"Missing external executable: {', '.join(candidates)}.")


def extract_text_from_image(image: Any, config: ExtractConfig) -> str:
    require_python_package("pytesseract")
    tesseract_path = find_executable(["tesseract"], config.tesseract_cmd)

    import pytesseract  # type: ignore[import-untyped]

    pytesseract.pytesseract.tesseract_cmd = tesseract_path
    text = pytesseract.image_to_string(image, lang=config.ocr_lang)
    return str(text).strip()


def extract_image(source_path: Path, config: ExtractConfig) -> ExtractionResult:
    require_python_package("PIL", "Pillow")

    from PIL import Image

    with Image.open(source_path) as image:
        text = extract_text_from_image(image, config)

    return ExtractionResult(
        source_path=source_path,
        kind=FileKind.IMAGE,
        status=ResultStatus.SUCCEEDED,
        sections=[ExtractedSection(title="Image OCR", text=text, status="ocr")],
        message="Image OCR completed.",
    )


def meaningful_text_chars(text: str) -> int:
    return sum(1 for char in text if not char.isspace())


def classify_pdf_page(
    text: str,
    image_area_ratio: float,
    config: ExtractConfig,
) -> PageStatus:
    char_count = meaningful_text_chars(text)

    if char_count >= config.pdf_text_page_min_chars:
        return PageStatus.TEXT_PAGE
    if (
        char_count <= config.pdf_ocr_max_chars
        and image_area_ratio >= config.pdf_image_ratio_threshold
    ):
        return PageStatus.OCR_NEEDED_PAGE
    if char_count == 0 and image_area_ratio <= config.blank_image_ratio_threshold:
        return PageStatus.BLANK_PAGE
    return PageStatus.SUSPICIOUS_PAGE


def page_image_area_ratio(page: Any) -> float:
    page_area = float(page.rect.width * page.rect.height)
    if page_area <= 0:
        return 0.0

    image_area = 0.0
    for image in page.get_images(full=True):
        xref = image[0]
        for rect in page.get_image_rects(xref):
            image_area += max(0.0, float(rect.width * rect.height))

    return min(image_area / page_area, 1.0)


def render_page_image(page: Any, config: ExtractConfig) -> Any:
    require_python_package("PIL", "Pillow")

    import fitz  # type: ignore[import-untyped]
    from PIL import Image

    matrix = fitz.Matrix(config.pdf_render_scale, config.pdf_render_scale)
    pixmap = page.get_pixmap(matrix=matrix, alpha=False)
    image = Image.open(BytesIO(pixmap.tobytes("png")))
    image.load()
    return image


def ocr_pdf_pages(
    pdf_path: Path,
    page_indexes: set[int],
    config: ExtractConfig,
) -> dict[int, str]:
    require_python_package("fitz", "PyMuPDF")

    import fitz

    texts: dict[int, str] = {}
    with fitz.open(pdf_path) as document:
        for page_index in sorted(page_indexes):
            if page_index < 0 or page_index >= document.page_count:
                continue
            page = document.load_page(page_index)
            with render_page_image(page, config) as image:
                texts[page_index] = extract_text_from_image(image, config)
    return texts


def extract_pdf(source_path: Path, config: ExtractConfig) -> ExtractionResult:
    require_python_package("fitz", "PyMuPDF")

    import fitz

    sections: list[ExtractedSection] = []
    ocr_count = 0

    with fitz.open(source_path) as document:
        for page_index in range(document.page_count):
            page = document.load_page(page_index)
            text = str(page.get_text()).strip()
            image_ratio = page_image_area_ratio(page)
            page_status = classify_pdf_page(text, image_ratio, config)
            notes: list[str] = []
            output_text = text

            if page_status == PageStatus.OCR_NEEDED_PAGE:
                with render_page_image(page, config) as image:
                    output_text = extract_text_from_image(image, config)
                ocr_count += 1
                notes.append("Rendered page image and applied OCR.")
            elif page_status == PageStatus.BLANK_PAGE:
                notes.append("No meaningful text or image content detected.")
            elif page_status == PageStatus.SUSPICIOUS_PAGE:
                notes.append("Page has limited text and did not meet OCR image threshold.")

            sections.append(
                ExtractedSection(
                    title=f"Page {page_index + 1}",
                    text=output_text,
                    status=page_status.value,
                    index=page_index + 1,
                    notes=notes,
                    metadata={
                        "meaningful_text_chars": meaningful_text_chars(text),
                        "image_area_ratio": round(image_ratio, 4),
                    },
                )
            )

    return ExtractionResult(
        source_path=source_path,
        kind=FileKind.PDF,
        status=ResultStatus.SUCCEEDED,
        sections=sections,
        message=f"PDF extraction completed. OCR pages: {ocr_count}.",
    )


def extract_docx_text(source_path: Path) -> str:
    require_python_package("docx", "python-docx")

    from docx import Document

    document = Document(str(source_path))
    chunks: list[str] = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            chunks.append(text)

    for table_index, table in enumerate(document.tables, start=1):
        rows: list[str] = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            chunks.append(f"Table {table_index}\n" + "\n".join(rows))

    return "\n\n".join(chunks).strip()


def extract_docx(source_path: Path, config: ExtractConfig) -> ExtractionResult:
    text = extract_docx_text(source_path)
    char_count = meaningful_text_chars(text)

    if char_count >= config.docx_min_text_chars:
        return ExtractionResult(
            source_path=source_path,
            kind=FileKind.DOCX,
            status=ResultStatus.SUCCEEDED,
            sections=[
                ExtractedSection(
                    title="DOCX text",
                    text=text,
                    status="text",
                    metadata={"meaningful_text_chars": char_count},
                )
            ],
            message="DOCX structural text extraction completed.",
        )

    with converted_to_pdf(source_path, config) as pdf_path:
        pdf_result = extract_pdf(pdf_path, config)

    return ExtractionResult(
        source_path=source_path,
        kind=FileKind.DOCX,
        status=ResultStatus.SUCCEEDED,
        sections=[
            ExtractedSection(
                title="DOCX structural text",
                text=text,
                status="sparse_text",
                notes=["Text was below threshold; PDF OCR fallback was used."],
                metadata={"meaningful_text_chars": char_count},
            ),
            *pdf_result.sections,
        ],
        message="DOCX text was sparse; extracted via PDF fallback.",
    )


@contextmanager
def converted_to_pdf(source_path: Path, config: ExtractConfig) -> Iterator[Path]:
    soffice = find_executable(["soffice", "libreoffice"], config.soffice_path)

    with tempfile.TemporaryDirectory(prefix="orbit-convert-") as temp_dir:
        output_dir = Path(temp_dir)
        profile_dir = output_dir / "lo-profile"
        profile_dir.mkdir()
        command = [
            soffice,
            f"-env:UserInstallation={profile_dir.as_uri()}",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir),
            str(source_path),
        ]
        completed = subprocess.run(
            command,
            check=False,
            capture_output=True,
            text=True,
        )

        if completed.returncode != 0:
            details = (completed.stderr or completed.stdout).strip()
            raise DependencyError(
                f"LibreOffice failed to convert '{source_path.name}' to PDF. {details}"
            )

        pdf_path = output_dir / f"{source_path.stem}.pdf"
        if not pdf_path.exists():
            matches = list(output_dir.glob("*.pdf"))
            if not matches:
                raise DependencyError(
                    f"LibreOffice did not produce a PDF for '{source_path.name}'."
                )
            pdf_path = matches[0]

        yield pdf_path


def classify_pptx_slide(
    text: str,
    image_area_ratio: float,
    config: ExtractConfig,
) -> SlideStatus:
    char_count = meaningful_text_chars(text)

    if char_count == 0 and image_area_ratio <= config.blank_image_ratio_threshold:
        return SlideStatus.BLANK_SLIDE
    if (
        char_count < config.pptx_text_slide_min_chars
        and image_area_ratio >= config.pptx_image_ratio_threshold
    ):
        return SlideStatus.OCR_NEEDED_SLIDE
    if (
        char_count >= config.pptx_text_slide_min_chars
        and image_area_ratio >= config.pptx_image_ratio_threshold
    ):
        return SlideStatus.MIXED_SLIDE
    return SlideStatus.TEXT_SLIDE


def extract_pptx_slides(source_path: Path, config: ExtractConfig) -> list[SlideText]:
    require_python_package("pptx", "python-pptx")

    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    presentation = Presentation(str(source_path))
    slide_area = float(presentation.slide_width * presentation.slide_height)
    slides: list[SlideText] = []

    for slide_index, slide in enumerate(presentation.slides):
        chunks: list[str] = []
        image_area = 0.0

        for shape in slide.shapes:
            chunks.extend(extract_shape_text(shape))
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_area += float(shape.width * shape.height)

        text = "\n".join(chunk for chunk in chunks if chunk.strip()).strip()
        image_ratio = min(image_area / slide_area, 1.0) if slide_area > 0 else 0.0
        slides.append(
            SlideText(
                index=slide_index + 1,
                text=text,
                image_area_ratio=image_ratio,
                status=classify_pptx_slide(text, image_ratio, config),
            )
        )

    return slides


def extract_shape_text(shape: Any) -> list[str]:
    chunks: list[str] = []

    if getattr(shape, "has_text_frame", False):
        text = str(shape.text).strip()
        if text:
            chunks.append(text)

    if getattr(shape, "has_table", False):
        table_rows: list[str] = []
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                table_rows.append(" | ".join(cells))
        if table_rows:
            chunks.append("\n".join(table_rows))

    if getattr(shape, "shape_type", None) == 6 and hasattr(shape, "shapes"):
        for child in shape.shapes:
            chunks.extend(extract_shape_text(child))

    return chunks


def extract_pptx(source_path: Path, config: ExtractConfig) -> ExtractionResult:
    slides = extract_pptx_slides(source_path, config)
    ocr_targets = {
        slide.index - 1
        for slide in slides
        if slide.status == SlideStatus.OCR_NEEDED_SLIDE
    }
    ocr_text_by_index: dict[int, str] = {}

    if ocr_targets:
        with converted_to_pdf(source_path, config) as pdf_path:
            ocr_text_by_index = ocr_pdf_pages(pdf_path, ocr_targets, config)

    sections: list[ExtractedSection] = []
    for slide in slides:
        page_index = slide.index - 1
        text_parts = [slide.text.strip()]
        notes: list[str] = []

        if page_index in ocr_text_by_index:
            text_parts.append(ocr_text_by_index[page_index].strip())
            notes.append("Rendered converted PDF page and applied OCR.")
        elif slide.status == SlideStatus.OCR_NEEDED_SLIDE:
            notes.append("Slide was marked for OCR, but no OCR text was produced.")

        sections.append(
            ExtractedSection(
                title=f"Slide {slide.index}",
                text="\n\n".join(part for part in text_parts if part),
                status=slide.status.value,
                index=slide.index,
                notes=notes,
                metadata={
                    "meaningful_text_chars": meaningful_text_chars(slide.text),
                    "image_area_ratio": round(slide.image_area_ratio, 4),
                },
            )
        )

    return ExtractionResult(
        source_path=source_path,
        kind=FileKind.PPTX,
        status=ResultStatus.SUCCEEDED,
        sections=sections,
        message=f"PPTX extraction completed. OCR slides: {len(ocr_targets)}.",
    )


def clean_reference_text(
    raw_text: str,
    *,
    client: Any | None = None,
    model: str | None = None,
    api_key: str | None = None,
) -> AICleanupResult:
    text = raw_text.strip()
    if not text:
        return AICleanupResult(text="", status="skipped", message="No raw text to clean.")

    api_client: Any = client
    if api_client is None:
        if not api_key:
            return AICleanupResult(
                text="",
                status="unavailable",
                message="OPENAI_API_KEY is not configured.",
            )

        from openai import OpenAI

        api_client = OpenAI(api_key=api_key)

    try:
        response = api_client.responses.create(
            model=model or DEFAULT_OPENAI_MODEL,
            instructions=CLEANUP_INSTRUCTIONS,
            input=f"Clean this OCR/extracted reference text:\n\n{text}",
        )
    except Exception as error:
        return AICleanupResult(text="", status="failed", message=str(error))

    cleaned_text = str(getattr(response, "output_text", "")).strip()
    if not cleaned_text:
        return AICleanupResult(text="", status="failed", message="OpenAI returned empty text.")

    return AICleanupResult(text=cleaned_text, status="succeeded")


def extract_presentation_keywords(
    cleaned_text: str,
    *,
    client: Any | None = None,
    model: str | None = None,
    api_key: str | None = None,
) -> KeywordExtractionResult:
    text = cleaned_text.strip()
    if not text:
        return KeywordExtractionResult(
            keywords=[],
            status="skipped",
            message="No cleaned text to analyze.",
        )

    api_client: Any = client
    if api_client is None:
        if not api_key:
            return KeywordExtractionResult(
                keywords=[],
                status="unavailable",
                message="OPENAI_API_KEY is not configured.",
            )

        from openai import OpenAI

        api_client = OpenAI(api_key=api_key)

    try:
        response = api_client.responses.create(
            model=model or DEFAULT_OPENAI_MODEL,
            instructions=KEYWORD_INSTRUCTIONS,
            input=f"Extract presentation keywords from this cleaned text:\n\n{text}",
            text=KEYWORD_RESPONSE_FORMAT,
        )
    except Exception as error:
        return KeywordExtractionResult(keywords=[], status="failed", message=str(error))

    output_text = str(getattr(response, "output_text", "")).strip()
    if not output_text:
        return KeywordExtractionResult(
            keywords=[],
            status="failed",
            message="OpenAI returned empty keywords.",
        )

    try:
        payload = json.loads(output_text)
    except json.JSONDecodeError as error:
        return KeywordExtractionResult(
            keywords=[],
            status="failed",
            message=f"OpenAI returned invalid keyword JSON: {error}",
        )

    if isinstance(payload, dict):
        payload = payload.get("keywords")

    if not isinstance(payload, list):
        return KeywordExtractionResult(
            keywords=[],
            status="failed",
            message="OpenAI keyword response did not include a keywords array.",
        )

    keywords: list[PresentationKeyword] = []
    for item in payload:
        if not isinstance(item, dict):
            continue

        keyword = str(item.get("keyword", "")).strip()
        reason = str(item.get("reason", "")).strip()
        priority = str(item.get("priority", "medium")).strip().lower()

        if not keyword or not reason:
            continue
        if priority not in {"high", "medium", "low"}:
            priority = "medium"

        keywords.append(PresentationKeyword(keyword, reason, priority))

    if not keywords:
        return KeywordExtractionResult(
            keywords=[],
            status="failed",
            message="OpenAI returned no usable keywords.",
        )

    return KeywordExtractionResult(keywords=keywords, status="succeeded")
