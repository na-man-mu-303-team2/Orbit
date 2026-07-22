from __future__ import annotations

import io
import zipfile
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[2]
FIXTURE = (
    ROOT
    / "services"
    / "python-worker"
    / "tests"
    / "fixtures"
    / "pptx"
    / "import-fidelity-notes.pptx"
)
FIXED_ZIP_TIMESTAMP = (2026, 1, 1, 0, 0, 0)
NOTES_MASTER_MARKER = "NOTES_MASTER_DECORATION_DO_NOT_IMPORT"
NOTES_NON_BODY_MARKER = "NOTES_NON_BODY_DO_NOT_IMPORT"


def main() -> None:
    FIXTURE.parent.mkdir(parents=True, exist_ok=True)
    image_bytes = build_fixture_image()
    with io.BytesIO() as image_stream:
        image_stream.write(image_bytes)
        image_stream.seek(0)
        presentation = build_presentation(image_stream)
        with io.BytesIO() as output:
            presentation.save(output)
            package = inject_notes_master_decoration(output.getvalue())
    FIXTURE.write_bytes(normalize_zip(package))
    print(FIXTURE.relative_to(ROOT).as_posix())


def build_fixture_image() -> bytes:
    image = Image.new("RGB", (240, 160), "#0F172A")
    draw = ImageDraw.Draw(image)
    draw.rectangle((16, 16, 224, 144), fill="#2563EB")
    draw.ellipse((72, 32, 168, 128), fill="#F59E0B")
    output = io.BytesIO()
    image.save(output, format="PNG", optimize=False)
    return output.getvalue()


def build_presentation(image_stream: io.BytesIO) -> Presentation:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])

    title = slide.shapes.title
    title.text_frame.clear()
    title_paragraph = title.text_frame.paragraphs[0]
    title_run = title_paragraph.add_run()
    title_run.text = "Inherited title style"
    title_run.font.name = "Pretendard SemiBold"
    title_run._r.get_or_add_rPr().set("spc", "120")

    subtitle = slide.placeholders[1]
    subtitle.text_frame.text = "Letter spacing and inherited sizing baseline"

    callout = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
        Inches(8.6),
        Inches(1.7),
        Inches(3.6),
        Inches(1.6),
    )
    callout.text_frame.text = "Unsupported callout"

    picture = slide.shapes.add_picture(
        image_stream,
        Inches(1.0),
        Inches(3.2),
        Inches(3.0),
        Inches(2.0),
    )
    label = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.5),
        Inches(3.55),
        Inches(2.6),
        Inches(1.1),
    )
    label.text_frame.text = "Grouped image"
    slide.shapes.add_group_shape([picture, label])

    notes = slide.notes_slide
    notes_text_frame = notes.notes_text_frame
    notes_text_frame.clear()
    notes_text_frame.paragraphs[0].text = "첫 번째 문단"
    notes_text_frame.add_paragraph().text = ""
    manual_break_paragraph = notes_text_frame.add_paragraph()
    first_run = manual_break_paragraph.add_run()
    first_run.text = "수동"
    second_run = manual_break_paragraph.add_run()
    second_run.text = "줄바꿈"
    first_run._r.addnext(OxmlElement("a:br"))
    return presentation


def inject_notes_master_decoration(package_bytes: bytes) -> bytes:
    source = io.BytesIO(package_bytes)
    output = io.BytesIO()
    with zipfile.ZipFile(source, "r") as package, zipfile.ZipFile(
        output, "w", compression=zipfile.ZIP_DEFLATED
    ) as updated:
        for item in package.infolist():
            content = package.read(item.filename)
            if item.filename == "ppt/notesMasters/notesMaster1.xml":
                closing = b"</p:spTree>"
                if closing not in content:
                    raise RuntimeError("notes master shape tree is missing")
                content = content.replace(
                    closing,
                    notes_master_shape_xml() + closing,
                    1,
                )
            if item.filename == "ppt/notesSlides/notesSlide1.xml":
                closing = b"</p:spTree>"
                if closing not in content:
                    raise RuntimeError("notes slide shape tree is missing")
                content = content.replace(
                    closing,
                    notes_non_body_shape_xml() + closing,
                    1,
                )
            updated.writestr(item, content)
    return output.getvalue()


def notes_master_shape_xml() -> bytes:
    return f"""
<p:sp>
  <p:nvSpPr>
    <p:cNvPr id="99" name="Notes master fixture marker"/>
    <p:cNvSpPr txBox="1"/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="457200" y="457200"/><a:ext cx="3657600" cy="365760"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/>
  </p:spPr>
  <p:txBody>
    <a:bodyPr/>
    <a:lstStyle/>
    <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>{NOTES_MASTER_MARKER}</a:t></a:r><a:endParaRPr lang="en-US" sz="1200"/></a:p>
  </p:txBody>
</p:sp>
""".strip().encode("utf-8")


def notes_non_body_shape_xml() -> bytes:
    return f"""
<p:sp>
  <p:nvSpPr>
    <p:cNvPr id="98" name="Notes non-body fixture marker"/>
    <p:cNvSpPr txBox="1"/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="457200" y="5943600"/><a:ext cx="3657600" cy="365760"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/>
  </p:spPr>
  <p:txBody>
    <a:bodyPr/>
    <a:lstStyle/>
    <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>{NOTES_NON_BODY_MARKER}</a:t></a:r><a:endParaRPr lang="en-US" sz="1200"/></a:p>
  </p:txBody>
</p:sp>
""".strip().encode("utf-8")


def normalize_zip(package_bytes: bytes) -> bytes:
    source = io.BytesIO(package_bytes)
    output = io.BytesIO()
    with zipfile.ZipFile(source, "r") as package, zipfile.ZipFile(
        output, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=9
    ) as normalized:
        for item in sorted(package.infolist(), key=lambda entry: entry.filename):
            info = zipfile.ZipInfo(item.filename, FIXED_ZIP_TIMESTAMP)
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = item.external_attr
            info.create_system = item.create_system
            normalized.writestr(info, package.read(item.filename))
    return output.getvalue()


if __name__ == "__main__":
    main()
