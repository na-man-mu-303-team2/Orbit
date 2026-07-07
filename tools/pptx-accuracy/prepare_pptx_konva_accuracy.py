from __future__ import annotations

# ruff: noqa: E402

import base64
import json
import os
import sys
import tempfile
from pathlib import Path
from typing import Any

ROOT = Path(__file__).resolve().parents[2]
WORKER = ROOT / "services" / "python-worker"
sys.path.insert(0, str(WORKER))

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt

from app.ai.pptx_ooxml_generation import generate_pptx_ooxml

OUT = ROOT / "tmp" / "pptx-konva-accuracy" / "run"
SAMPLES = OUT / "samples"
GOLDEN = OUT / "golden"
PAYLOADS = OUT / "payloads"
CANDIDATE = OUT / "candidate"
ASSETS = OUT / "assets"
TEMP = OUT / "tmp"

SLIDE_WIDE = 13.333333
SLIDE_HIGH = 7.5


def main() -> None:
    for directory in (SAMPLES, GOLDEN, PAYLOADS, CANDIDATE, ASSETS, TEMP):
        directory.mkdir(parents=True, exist_ok=True)
    os.environ["TMP"] = str(TEMP)
    os.environ["TEMP"] = str(TEMP)
    tempfile.tempdir = str(TEMP)

    image_assets = create_image_assets()
    builders = [
        ("01_text_basic", sample_text_basic),
        ("02_rich_text_cjk", sample_rich_text_cjk),
        ("03_vertical_multiline_text", sample_vertical_multiline_text),
        ("04_basic_shapes", sample_basic_shapes),
        ("05_preset_arrows", sample_preset_arrows),
        ("06_freeform_geometry", sample_freeform_geometry),
        ("07_gradient_pattern_shadow", sample_gradient_pattern_shadow),
        ("08_image_crop_overlay", sample_image_crop_overlay),
        ("09_grouped_icon_text", sample_grouped_icon_text),
        ("10_split_image_group", sample_split_image_group),
        ("11_table", sample_table),
        ("12_bar_chart", sample_bar_chart),
        ("13_line_chart", sample_line_chart),
        ("14_pie_chart", sample_pie_chart),
        ("15_master_layout_theme", sample_master_layout_theme),
    ]

    rows: list[dict[str, Any]] = []
    for index, (name, builder) in enumerate(builders, start=1):
        pptx_path = SAMPLES / f"{name}.pptx"
        builder(pptx_path, image_assets)
        result = generate_pptx_ooxml(pptx_path, f"konva_eval_{index}", render=True)
        asset_by_id = {asset.asset_id: asset for asset in result.assets}
        golden = asset_by_id["slide_render_1"]
        golden_path = GOLDEN / f"{name}.png"
        golden_path.write_bytes(base64.b64decode(golden.content_base64))

        deck, unresolved_assets, fallback_objects = build_deck_payload(
            result.canvas,
            result.blueprint,
            result.assets,
            index,
        )
        payload_path = PAYLOADS / f"{name}.json"
        payload_path.write_text(
            json.dumps({"deck": deck, "slideIndex": 0}, ensure_ascii=False),
            encoding="utf-8",
        )
        rows.append(
            {
                "name": name,
                "pptxPath": relative_path(pptx_path),
                "goldenPath": relative_path(golden_path),
                "payloadPath": relative_path(payload_path),
                "candidatePath": relative_path(CANDIDATE / f"{name}.png"),
                "fallbackObjects": fallback_objects,
                "fullSlideFallbackUsed": any(
                    bool(slide.get("style", {}).get("backgroundImage"))
                    for slide in deck["slides"]
                ),
                "unresolvedAssets": unresolved_assets,
                "warnings": result.warnings,
                "elementCounts": element_counts(result.blueprint["slides"][0]["elements"]),
            }
        )

    manifest = {
        "sampleCount": len(rows),
        "threshold": 0.95,
        "route": "/__deck-render",
        "rows": rows,
    }
    (OUT / "manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    print(json.dumps(manifest, ensure_ascii=False, indent=2))


def build_deck_payload(
    canvas: dict[str, Any],
    blueprint: dict[str, Any],
    assets: list[Any],
    index: int,
) -> tuple[dict[str, Any], list[str], int]:
    asset_urls = {
        f"asset:{asset.asset_id}": data_url(asset.mime_type, asset.content_base64)
        for asset in assets
    }
    unresolved: list[str] = []
    slides = []
    fallback_objects = 0
    for slide_index, slide in enumerate(blueprint["slides"], start=1):
        elements = []
        for element in slide.get("elements", []):
            next_element = json.loads(json.dumps(element))
            props = next_element.get("props")
            if isinstance(props, dict) and isinstance(props.get("src"), str):
                src = props["src"]
                if src.startswith("asset:shape_render_"):
                    fallback_objects += 1
                if src.startswith("asset:"):
                    resolved = asset_urls.get(src)
                    if resolved:
                        props["src"] = resolved
                    else:
                        unresolved.append(src)
            elements.append(next_element)
        slides.append(
            {
                "slideId": f"slide_konva_eval_{index}_{slide_index}",
                "order": slide_index,
                "title": f"Slide {slide_index}",
                "thumbnailUrl": "",
                "style": slide.get("style", {}),
                "speakerNotes": "",
                "elements": elements,
                "keywords": [],
                "animations": [],
                "aiNotes": {"emphasisPoints": [], "sourceEvidence": []},
            }
        )
    return (
        {
            "deckId": f"deck_konva_eval_{index}",
            "projectId": "project_konva_eval",
            "title": "PPTX Konva accuracy fixture",
            "version": 1,
            "metadata": {"language": "ko", "locale": "ko-KR", "sourceType": "import"},
            "canvas": canvas,
            "theme": blueprint["theme"],
            "slides": slides,
        },
        sorted(set(unresolved)),
        fallback_objects,
    )


def data_url(mime_type: str, content_base64: str) -> str:
    return f"data:{mime_type};base64,{content_base64}"


def new_presentation() -> tuple[Presentation, Any]:
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_WIDE)
    presentation.slide_height = Inches(SLIDE_HIGH)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    return presentation, slide


def create_image_assets() -> dict[str, Path]:
    hero = ASSETS / "hero.png"
    split = ASSETS / "split.png"
    image = Image.new("RGB", (520, 320), "#2563EB")
    draw = ImageDraw.Draw(image)
    for x in range(0, 520, 26):
        draw.rectangle(
            [x, 0, x + 25, 320],
            fill="#7C3AED" if (x // 26) % 2 else "#0EA5E9",
        )
    draw.ellipse([130, 70, 390, 270], fill="#F59E0B")
    image.save(hero)

    split_image = Image.new("RGB", (400, 260), "#FFFFFF")
    draw = ImageDraw.Draw(split_image)
    draw.rectangle([0, 0, 200, 130], fill="#EF4444")
    draw.rectangle([200, 0, 400, 130], fill="#10B981")
    draw.rectangle([0, 130, 200, 260], fill="#2563EB")
    draw.rectangle([200, 130, 400, 260], fill="#F59E0B")
    split_image.save(split)
    return {"hero": hero, "split": split}


def sample_text_basic(path: Path, assets: dict[str, Path]) -> None:
    del assets
    prs, slide = new_presentation()
    add_title(slide, "Editable Text Baseline")
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    box.text_frame.text = "Simple paragraph text should stay editable."
    run = box.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(34)
    run.font.color.rgb = RGBColor(17, 24, 39)
    prs.save(path)


def sample_rich_text_cjk(path: Path, assets: dict[str, Path]) -> None:
    del assets
    prs, slide = new_presentation()
    add_title(slide, "Mixed CJK / Latin Runs")
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10.5), Inches(2))
    paragraph = box.text_frame.paragraphs[0]
    for text, size, color, bold in [
        ("한국어 텍스트 ", 32, RGBColor(17, 24, 39), True),
        ("Latin ", 26, RGBColor(37, 99, 235), False),
        ("日本語", 30, RGBColor(124, 58, 237), False),
    ]:
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    prs.save(path)


def sample_vertical_multiline_text(path: Path, assets: dict[str, Path]) -> None:
    del assets
    prs, slide = new_presentation()
    add_title(slide, "Vertical Text")
    box = slide.shapes.add_textbox(Inches(1), Inches(1.7), Inches(1), Inches(4.8))
    box.text_frame.text = "VERTICAL"
    box._element.txBody.bodyPr.set("vert", "vert270")
    body = slide.shapes.add_textbox(Inches(3), Inches(2), Inches(7), Inches(2.5))
    body.text_frame.text = "Line one\nLine two\nLine three"
    prs.save(path)


def sample_basic_shapes(path: Path, assets: dict[str, Path]) -> None:
    del assets
    prs, slide = new_presentation()
    add_title(slide, "Basic Shapes")
    for shape_type, left, color in [
        (MSO_SHAPE.RECTANGLE, 1, RGBColor(37, 99, 235)),
        (MSO_SHAPE.OVAL, 3.2, RGBColor(245, 158, 11)),
        (MSO_SHAPE.ROUNDED_RECTANGLE, 5.4, RGBColor(16, 185, 129)),
    ]:
        shape = slide.shapes.add_shape(
            shape_type, Inches(left), Inches(2.2), Inches(1.7), Inches(1.7)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(17, 24, 39)
    prs.save(path)


def sample_preset_arrows(path: Path, assets: dict[str, Path]) -> None:
    del assets
    prs, slide = new_presentation()
    add_title(slide, "Preset Arrows")
    for shape_type, left in [
        (MSO_SHAPE.RIGHT_ARROW, 1),
        (MSO_SHAPE.CHEVRON, 3.4),
        (MSO_SHAPE.STAR_5_POINT, 5.8),
        (MSO_SHAPE.DONUT, 8.2),
    ]:
        shape = slide.shapes.add_shape(
            shape_type, Inches(left), Inches(2.3), Inches(1.8), Inches(1.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(37, 99, 235)
        shape.line.color.rgb = RGBColor(255, 255, 255)
    prs.save(path)


def sample_freeform_geometry(path: Path, assets: dict[str, Path]) -> None:
    del assets
    prs, slide = new_presentation()
    add_title(slide, "Freeform Geometry")
    builder = slide.shapes.build_freeform(Inches(1), Inches(2))
    builder.move_to(0, 0)
    builder.add_line_segments(
        [(Inches(1.2), Inches(0.2)), (Inches(2.2), Inches(1.4)), (0, Inches(1.8))]
    )
    shape = builder.convert_to_shape()
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(124, 58, 237)
    shape.line.color.rgb = RGBColor(17, 24, 39)
    prs.save(path)


def sample_gradient_pattern_shadow(path: Path, assets: dict[str, Path]) -> None:
    del assets
    prs, slide = new_presentation()
    add_title(slide, "Gradient / Pattern")
    gradient = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(3.5), Inches(2)
    )
    replace_shape_fill_with_gradient(gradient)
    pattern = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(2), Inches(3.5), Inches(2)
    )
    replace_shape_fill_with_pattern(pattern)
    prs.save(path)


def sample_image_crop_overlay(path: Path, assets: dict[str, Path]) -> None:
    prs, slide = new_presentation()
    add_title(slide, "Image Crop + Overlay Text")
    picture = slide.shapes.add_picture(
        str(assets["hero"]), Inches(1), Inches(1.8), Inches(6), Inches(3.7)
    )
    picture.crop_left = 0.12
    picture.crop_top = 0.08
    picture.crop_right = 0.18
    picture.crop_bottom = 0.1
    text = slide.shapes.add_textbox(Inches(1.4), Inches(4.4), Inches(5), Inches(0.8))
    text.text_frame.text = "Text over image"
    text.text_frame.paragraphs[0].runs[0].font.size = Pt(30)
    prs.save(path)


def sample_grouped_icon_text(path: Path, assets: dict[str, Path]) -> None:
    del assets
    prs, slide = new_presentation()
    add_title(slide, "Grouped Icon + Text")
    a = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(2), Inches(1), Inches(1))
    b = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(1.9), Inches(2.1), Inches(2), Inches(0.8)
    )
    c = slide.shapes.add_textbox(Inches(1.2), Inches(3.2), Inches(3), Inches(0.6))
    c.text_frame.text = "Grouped label"
    slide.shapes.add_group_shape([a, b, c])
    prs.save(path)


def sample_split_image_group(path: Path, assets: dict[str, Path]) -> None:
    prs, slide = new_presentation()
    add_title(slide, "Split Image Group")
    pieces = []
    for row in range(2):
        for col in range(2):
            pic = slide.shapes.add_picture(
                str(assets["split"]),
                Inches(1 + col * 2.1),
                Inches(2 + row * 1.4),
                Inches(2),
                Inches(1.3),
            )
            pic.crop_left = col * 0.5
            pic.crop_right = (1 - col) * 0.5
            pic.crop_top = row * 0.5
            pic.crop_bottom = (1 - row) * 0.5
            pieces.append(pic)
    slide.shapes.add_group_shape(pieces)
    prs.save(path)


def sample_table(path: Path, assets: dict[str, Path]) -> None:
    del assets
    prs, slide = new_presentation()
    add_title(slide, "Table")
    table = slide.shapes.add_table(3, 3, Inches(1), Inches(1.8), Inches(7), Inches(3)).table
    for row in range(3):
        for col in range(3):
            table.cell(row, col).text = f"R{row + 1}C{col + 1}"
    prs.save(path)


def sample_bar_chart(path: Path, assets: dict[str, Path]) -> None:
    del assets
    create_chart_sample(path, XL_CHART_TYPE.COLUMN_CLUSTERED, "Bar Chart")


def sample_line_chart(path: Path, assets: dict[str, Path]) -> None:
    del assets
    create_chart_sample(path, XL_CHART_TYPE.LINE_MARKERS, "Line Chart")


def sample_pie_chart(path: Path, assets: dict[str, Path]) -> None:
    del assets
    create_chart_sample(path, XL_CHART_TYPE.PIE, "Pie Chart")


def create_chart_sample(path: Path, chart_type: Any, title: str) -> None:
    prs, slide = new_presentation()
    add_title(slide, title)
    data = ChartData()
    data.categories = ["A", "B", "C", "D"]
    data.add_series("Series 1", (4.3, 2.5, 3.5, 4.5))
    chart = slide.shapes.add_chart(
        chart_type, Inches(1), Inches(1.7), Inches(8), Inches(4.5), data
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    prs.save(path)


def sample_master_layout_theme(path: Path, assets: dict[str, Path]) -> None:
    del assets
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDE)
    prs.slide_height = Inches(SLIDE_HIGH)
    layout = prs.slide_layouts[6]
    layout.shapes._spTree.add_autoshape(
        99, "LayoutBand", "rect", Inches(0), Inches(6.8), Inches(13.33), Inches(0.5)
    )
    shape = layout.shapes[-1]
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(37, 99, 235)
    slide = prs.slides.add_slide(layout)
    add_title(slide, "Master/Layout Decoration")
    prs.save(path)


def add_title(slide: Any, text: str) -> None:
    title = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(10.5), Inches(0.7))
    title.text_frame.text = text
    run = title.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = RGBColor(17, 24, 39)


def replace_shape_fill_with_gradient(shape: object) -> None:
    sp_pr = shape._element.spPr
    remove_fill_nodes(sp_pr)
    sp_pr.insert(
        0,
        parse_xml(
            """
            <a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <a:gsLst>
                <a:gs pos="0"><a:srgbClr val="2563EB"/></a:gs>
                <a:gs pos="100000"><a:srgbClr val="7C3AED"/></a:gs>
              </a:gsLst>
              <a:lin ang="5400000" scaled="1"/>
            </a:gradFill>
            """
        ),
    )


def replace_shape_fill_with_pattern(shape: object) -> None:
    sp_pr = shape._element.spPr
    remove_fill_nodes(sp_pr)
    sp_pr.insert(
        0,
        parse_xml(
            """
            <a:pattFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" prst="pct20">
              <a:fgClr><a:srgbClr val="111827"/></a:fgClr>
              <a:bgClr><a:srgbClr val="F59E0B"/></a:bgClr>
            </a:pattFill>
            """
        ),
    )


def remove_fill_nodes(sp_pr: Any) -> None:
    for child in list(sp_pr):
        if child.tag.rsplit("}", maxsplit=1)[-1] in {
            "solidFill",
            "gradFill",
            "noFill",
            "pattFill",
        }:
            sp_pr.remove(child)


def element_counts(elements: list[dict[str, Any]]) -> dict[str, int]:
    counts: dict[str, int] = {}
    for element in elements:
        key = str(element.get("type", "unknown"))
        counts[key] = counts.get(key, 0) + 1
    return counts


def relative_path(path: Path) -> str:
    return path.resolve().relative_to(ROOT).as_posix()


if __name__ == "__main__":
    main()
